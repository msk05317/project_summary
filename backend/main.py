"""
사업부 진행현황 보고 - 백엔드 메인
- PPT 업로드 → 슬라이드 PNG 변환 → GPT-4o 요약 → 신호등 분류
- 사용자 직접 차트/사진 업로드 및 부서·섹션 매핑
- 프로젝트별 상세 API 제공
"""
import os
import base64
import time
import secrets
import hmac
import json
import re
import shutil
import hashlib
import subprocess
from datetime import datetime, timedelta
from pathlib import Path

from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Request, Response, Cookie, Depends
from typing import Optional
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import JSONResponse, HTMLResponse, RedirectResponse
from dotenv import load_dotenv
from openai import OpenAI
from rag import vector_store as _vs

from project_templates import (
    PROJECT_LABELS,
    aggregate_projects,
    build_project_detail,
)

# ============================================================
# config_loader 기반 분류 enrichment
# - 기존 응답 필드는 절대 변경하지 않는다.
# - division/project 정보는 "추가만" 한다.
# - 분류 실패해도 응답이 깨지지 않게 안전 처리한다.
# ============================================================
import config_loader as _cl


def _safe_text_for_card(card: dict) -> str:
    """카드 텍스트 분류에 사용할 문자열 합치기."""
    parts = [
        card.get("product") or "",
        card.get("headline") or "",
        card.get("project_key") or "",
    ]
    return " ".join([p for p in parts if p]).strip()




# ============================================================
# Dashboard cards 그룹핑 (모델 단위)
# - 같은 project_id + 동일/포함 모델명 → 한 카드
# - 그 안에 이슈(headline) 리스트
# ============================================================
def _normalize_model_name(name: str) -> str:
    """모델명에서 앞 번호 prefix 제거 + 공백 정리"""
    if not isinstance(name, str):
        return ""
    import re as _re
    t = name.strip()
    t = _re.sub(r"^\s*\d+(?:[-.]\d+)*[.)\]]?\s+", "", t)
    t = _re.sub(r"\s+", " ", t)
    return t.strip()

def _model_match_key(name: str, project_id, existing) -> str:
    """
    하이브리드 매칭:
    - 같은 project_id 안에서만 그룹핑
    - 정규화된 이름이 기존 어떤 키의 prefix 거나 그 반대면 같은 모델로 합침
    - existing 형식: {(project_id, key): ...}
    """
    base = _normalize_model_name(name).lower()
    if not base:
        return base
    for (pid, k) in existing:
        if pid != project_id:
            continue
        if not k:
            continue
        a, b = (base, k) if len(base) >= len(k) else (k, base)
        if a == b:
            return k
        if a.startswith(b):
            rest = a[len(b):]
            if rest[:1] in (" ", "(", "[", "{", "-", "_", ".", ",", ":", "/"):
                return k
    return base

_STATUS_SEVERITY = {"RED": 3, "BLUE": 2, "BLACK": 1}

def _worst_status(statuses):
    best = "BLACK"
    best_sev = -1
    for s in statuses:
        sev = _STATUS_SEVERITY.get(s, 0)
        if sev > best_sev:
            best_sev = sev
            best = s
    return best



# ---------- headline AI 요약 (캐시 기반) ----------
import hashlib as _hashlib_hl
import json as _json_hl
from pathlib import Path as _Path_hl

_HEADLINE_CACHE_PATH = _Path_hl(__file__).parent / "headline_cache.json"

def _load_headline_cache() -> dict:
    try:
        return _json_hl.loads(_HEADLINE_CACHE_PATH.read_text(encoding="utf-8"))
    except Exception:
        return {}

def _save_headline_cache(cache: dict) -> None:
    try:
        _HEADLINE_CACHE_PATH.write_text(
            _json_hl.dumps(cache, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
    except Exception:
        pass

def _pick_headline_source(dated_items: list, due_date_min) -> str:
    """
    dated_items: [(due_date, text, raw_str), ...]
    due_date_min 과 같은 날짜의 항목을 우선 선택.
    없으면 첫 항목.
    """
    if not dated_items:
        return ""
    if due_date_min:
        for d, t, _ in dated_items:
            if d == due_date_min:
                return t
    return dated_items[0][1]

def _ai_headline(source_text: str) -> str:
    """
    원문 한 줄을 15자 이내 행동 중심 요약으로 변환.
    캐시 우선. OpenAI 실패 시 원문 앞 15자.
    """
    src = (source_text or "").strip()
    if not src:
        return ""
    # 캐시 키: 원문 해시
    key = _hashlib_hl.sha256(src.encode("utf-8")).hexdigest()[:16]
    cache = _load_headline_cache()
    if key in cache:
        return cache[key]

    # OpenAI 호출
    result = ""
    try:
        if client is not None:
            prompt = (
                "다음 프로젝트 상태 문장에서 가장 중요한 액션 하나만 뽑아 20자 이내 명사형 한 줄로 요약해라. "
                "규칙: 원문을 앞에서부터 그대로 자르지 말고 핵심 행동/상태를 재구성. "
                "조사·어미·부사어 삭제. 여러 항목 나열 금지(가장 시급한 것 하나만). "
                "괄호·날짜·이모지 제거. 숫자와 모델명(W25, CEFEM 등)은 유지. "
                "예시: 'W23 24EA, W24 7개 완료, W25 3개 출하예정 (자재부족)' → 'W25 3개 출하 지연'. "
                "예시: 'CEFEM 2대, VTM 2대 제작중, VTM 1대 자재 지연' → 'VTM 1대 자재 지연'. "
                "요약문만 출력.\n\n"
                f"원문: {src}"
            )
            resp = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.2,
                max_tokens=40,
            )
            result = (resp.choices[0].message.content or "").strip()
            # 따옴표/쌍따옴표 감싸져 오면 제거
            result = result.strip("\"'` ")
            # 개행 제거
            result = result.split("\n")[0].strip()
    except Exception:
        result = ""

    # 폴백: 원문 앞 15자
    if not result:
        result = src[:22]

    # 15자 초과 방어
    if len(result) > 24:
        result = result[:24].rstrip() + "…"

    cache[key] = result
    _save_headline_cache(cache)
    return result

def _normalize_issue_headline(text: str) -> str:
    """
    이슈 헤드라인 dedup 용 정규화 (중간 강도).
    - 앞쪽 번호 prefix 제거
    - 양 끝 공백/구두점 정리
    - 다중 공백 1개로
    - 끝쪽의 보조어 (중/진행/예정/완료/중임/중입니다 등) 제거
    - 한글/영문 모델명 사이 공백 통일
    - 비교용이므로 lowercase
    """
    if not isinstance(text, str):
        return ""
    import re as _re
    t = text.strip()
    # 앞 번호 prefix
    t = _re.sub(r"^\s*\d+(?:[-.]\d+)*[.)\]]?\s+", "", t)
    # 끝쪽 보조어 반복 제거
    suffixes = [
        "중입니다", "입니다", "중임", "되었음", "되었습니다",
        "중", "진행 중", "진행중", "진행", "예정", "완료", "완료됨",
        "진행 예정", "수행 중", "검토 중", "검토중",
    ]
    changed = True
    while changed:
        changed = False
        for suf in suffixes:
            if t.lower().endswith(suf.lower()):
                t = t[: -len(suf)].rstrip(" ,.;:·-")
                changed = True
    # 다중 공백/구두점 정리
    t = _re.sub(r"\s+", " ", t)
    t = t.strip(" ,.·-;:")
    return t.lower()

def _group_dashboard_cards(cards: list) -> list:
    """
    cards 를 모델 단위로 그룹핑.
    """
    groups = {}            # key=(project_id, model_key) → group dict
    order = []             # 그룹 등장 순서 보존
    issue_seen = {}        # 그룹 내 (status, headline) 중복 제거용

    for c in cards:
        product = c.get("product") or ""
        pid = c.get("project_id")  # None 가능 (미분류)
        # 미분류는 project_id 대신 라벨 "__unclassified__" 로 강제
        bucket_pid = pid if pid else "__unclassified__"

        match_key = _model_match_key(product, bucket_pid, groups.keys())
        gkey = (bucket_pid, match_key)

        if gkey not in groups:
            groups[gkey] = {
                "model": _normalize_model_name(product) or product,
                "status": c.get("status") or "BLACK",
                "project_id": pid,
                "project_label": c.get("project_label"),
                "project_badge": c.get("project_badge"),
                "division_id": c.get("division_id"),
                "division_label": c.get("division_label"),
                "report_date": c.get("report_date"),
                "report_family": c.get("report_family"),
                "issues": [],
                "_statuses": [],
            }
            order.append(gkey)
            issue_seen[gkey] = set()
        else:
            # 더 짧은(=상위) 모델명이 나중에 들어오면 그걸로 표시 보정 (예: "챔버 13종" 보다 "챔버" 우선)
            existing_name = groups[gkey]["model"]
            new_name = _normalize_model_name(product) or product
            if len(new_name) < len(existing_name):
                groups[gkey]["model"] = new_name

        # 이슈 dedup (정규화 헤드라인 기반)
        st = c.get("status") or "BLACK"
        head = (c.get("headline") or "").strip()
        # summary_bullets / due_date_min 은 headline 없어도 통과
        sb = c.get("summary_bullets") or []
        if sb and not groups[gkey].get("summary_bullets"):
            groups[gkey]["summary_bullets"] = sb
        ddm = c.get("due_date_min")
        if ddm and not groups[gkey].get("due_date_min"):
            groups[gkey]["due_date_min"] = ddm
        if st and st not in groups[gkey].get("_statuses", []):
            groups[gkey]["_statuses"].append(st)
        if not head:
            continue
        norm = _normalize_issue_headline(head)
        sig = (st, norm) if norm else (st, head.lower())
        if sig in issue_seen[gkey]:
            # 더 길고 풍부한 표현이 들어오면 표시용 headline 만 갱신
            for it in groups[gkey]["issues"]:
                if it.get("status") == st and _normalize_issue_headline(it.get("headline", "")) == norm:
                    if len(head) > len(it.get("headline", "")):
                        it["headline"] = head
                    break
            continue
        issue_seen[gkey].add(sig)
        groups[gkey]["issues"].append({
            "status": st,
            "headline": head,
            "doc_id": c.get("doc_id", ""),
        })
        # summary_bullets / due_date_min 통과
        sb = c.get("summary_bullets") or []
        if sb and not groups[gkey].get("summary_bullets"):
            groups[gkey]["summary_bullets"] = sb
        ddm = c.get("due_date_min")
        if ddm and not groups[gkey].get("due_date_min"):
            groups[gkey]["due_date_min"] = ddm
        groups[gkey]["_statuses"].append(st)

        # report_date 는 가장 최신으로 갱신
        rd_new = c.get("report_date") or ""
        rd_old = groups[gkey].get("report_date") or ""
        if rd_new > rd_old:
            groups[gkey]["report_date"] = rd_new

    result = []
    for gkey in order:
        g = groups[gkey]
        g["status"] = _worst_status(g.pop("_statuses") or [g["status"]])
        result.append(g)

    # 가장 심각한 그룹부터 정렬
    result.sort(key=lambda g: -_STATUS_SEVERITY.get(g["status"], 0))
    return result


def enrich_card(card: dict) -> dict:
    """
    /dashboard 카드에 division/project 분류 정보를 '추가'한다.
    기존 필드는 절대 건드리지 않는다.
    실패 시 새 필드는 None으로만 들어가고, 카드 자체는 정상 반환한다.
    """
    try:
        text = _safe_text_for_card(card)

        # 1순위: 백엔드가 이미 정한 project_key 가 있으면 그걸 신뢰
        hint_pid = card.get("project_key")
        project_id = None
        if hint_pid:
            # project_key 가 config 의 project_id 와 일치하는 경우 우선 사용
            if _cl.get_project(hint_pid):
                project_id = hint_pid

        # 2순위: config_loader 분류기
        if not project_id:
            project_id = _cl.classify_project(text)

        # 3순위: /upload 때 사장님이 지정한 폴백 (자동 분류가 다 실패한 경우)
        if not project_id:
            project_id = card.get("_fallback_project_id")

        division_id = _cl.derive_division_from_project(project_id)
        # 폴백 division 도 받기 (project_id 가 없는 극단적 경우)
        if not division_id:
            division_id = card.get("_fallback_division_id")

        project = _cl.get_project(project_id) if project_id else None
        division = _cl.get_division(division_id) if division_id else None

        # ⚠️ 기존 필드 그대로 두고 새 필드만 추가
        card["division_id"] = division_id
        card["division_label"] = division.get("label") if division else None
        card["project_id"] = project_id
        card["project_label"] = project.get("label") if project else None
        card["project_badge"] = _cl.badge_label_for_project(project_id)
    except Exception as e:
        # 어떤 이유로든 분류 실패해도 응답이 깨지지 않게
        card.setdefault("division_id", None)
        card.setdefault("division_label", None)
        card.setdefault("project_id", None)
        card.setdefault("project_label", None)
        card.setdefault("project_badge", None)
        card["_enrich_error"] = str(e)
    return card


def enrich_project_entry(entry: dict) -> dict:
    """
    /projects 응답 각 항목에 division/project 분류 정보를 '추가'한다.
    """
    try:
        # /projects 항목은 보통 project_key 또는 name/label 을 가짐
        hint_pid = entry.get("project_key") or entry.get("key")
        text = " ".join([
            str(entry.get("label") or ""),
            str(entry.get("name") or ""),
            str(hint_pid or ""),
        ]).strip()

        project_id = None
        if hint_pid and _cl.get_project(hint_pid):
            project_id = hint_pid
        if not project_id:
            project_id = _cl.classify_project(text)

        division_id = _cl.derive_division_from_project(project_id)
        project = _cl.get_project(project_id) if project_id else None
        division = _cl.get_division(division_id) if division_id else None

        entry["division_id"] = division_id
        entry["division_label"] = division.get("label") if division else None
        entry["project_id"] = project_id
        entry["project_label"] = project.get("label") if project else None
        entry["project_badge"] = _cl.badge_label_for_project(project_id)
    except Exception as e:
        entry.setdefault("division_id", None)
        entry.setdefault("division_label", None)
        entry.setdefault("project_id", None)
        entry.setdefault("project_label", None)
        entry.setdefault("project_badge", None)
        entry["_enrich_error"] = str(e)
    return entry


def enrich_project_detail(detail: dict) -> dict:
    """
    /projects/{key} 상세 응답 최상위에 분류 정보를 '추가'한다.
    """
    if not isinstance(detail, dict):
        return detail
    try:
        hint_pid = detail.get("project_key") or detail.get("key")
        text = " ".join([
            str(detail.get("label") or ""),
            str(detail.get("name") or ""),
            str(hint_pid or ""),
        ]).strip()

        project_id = None
        if hint_pid and _cl.get_project(hint_pid):
            project_id = hint_pid
        if not project_id:
            project_id = _cl.classify_project(text)

        division_id = _cl.derive_division_from_project(project_id)
        project = _cl.get_project(project_id) if project_id else None
        division = _cl.get_division(division_id) if division_id else None

        detail["division_id"] = division_id
        detail["division_label"] = division.get("label") if division else None
        detail["project_id"] = project_id
        detail["project_label"] = project.get("label") if project else None
        detail["project_badge"] = _cl.badge_label_for_project(project_id)
    except Exception as e:
        detail.setdefault("division_id", None)
        detail.setdefault("division_label", None)
        detail.setdefault("project_id", None)
        detail.setdefault("project_label", None)
        detail.setdefault("project_badge", None)
        detail["_enrich_error"] = str(e)
    return detail

from chart_extractor import extract_charts_from_pptx


# =========================================================
# 환경 변수
# =========================================================
load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
UPLOAD_PASSWORD = os.getenv("UPLOAD_PASSWORD", "1234")
client = OpenAI(api_key=OPENAI_API_KEY) if OPENAI_API_KEY else None

BASE_DIR = Path(__file__).parent.resolve()

# Railway/프로덕션: /data (영구 Volume)
# 로컬 개발: backend 폴더 안
DATA_DIR = Path(os.getenv("DATA_DIR", str(BASE_DIR)))
DATA_DIR.mkdir(parents=True, exist_ok=True)

UPLOAD_DIR = DATA_DIR / "uploads"
UPLOAD_HASHES_PATH = DATA_DIR / "upload_hashes.json"
SLIDES_DIR = DATA_DIR / "slides"
SLIDE_IMAGES_DIR = DATA_DIR / "slide_images"   # 기존 호환용
CROPPED_DIR = DATA_DIR / "cropped"
LATEST_FILE = DATA_DIR / "reports_latest.json"
HISTORY_FILE = DATA_DIR / "reports_history.json"
# 주간 보고 첨부 자료 (표 JSON / 사진)
NOTE_ASSETS_DIR = DATA_DIR / "note_assets"
NOTE_TABLES_DIR = NOTE_ASSETS_DIR / "tables"
NOTE_PHOTOS_DIR = NOTE_ASSETS_DIR / "photos"

for d in [UPLOAD_DIR, SLIDES_DIR, SLIDE_IMAGES_DIR, CROPPED_DIR, NOTE_ASSETS_DIR, NOTE_TABLES_DIR, NOTE_PHOTOS_DIR]:
    d.mkdir(exist_ok=True)

# DATA_DIR이 git 경로와 다르면(Railway), git에 포함된 시드 자산을 영구 볼륨에 복사
try:
    import shutil
    seed_assets = BASE_DIR / "note_assets"
    if DATA_DIR != BASE_DIR and seed_assets.exists():
        for src in seed_assets.rglob("*"):
            if src.is_file():
                rel = src.relative_to(seed_assets)
                dst = NOTE_ASSETS_DIR / rel
                if not dst.exists():
                    dst.parent.mkdir(parents=True, exist_ok=True)
                    shutil.copy2(src, dst)
        print(f"[seed] note_assets seeded to {NOTE_ASSETS_DIR}")
except Exception as _e:
    print(f"[seed] note_assets seed failed: {_e}")

# notes.json 도 동일하게 시드
try:
    import shutil
    seed_notes = BASE_DIR / "notes.json"
    if DATA_DIR != BASE_DIR and seed_notes.exists() and not (DATA_DIR / "notes.json").exists():
        shutil.copy2(seed_notes, DATA_DIR / "notes.json")
        print(f"[seed] notes.json seeded to {DATA_DIR}")
except Exception as _e:
    print(f"[seed] notes.json seed failed: {_e}")

RETENTION_DAYS = 180


# =========================================================
# FastAPI 앱
# =========================================================

# ─── FCM (Firebase Cloud Messaging) ───
# 대표님 폰으로 상태 변화(RED/ORANGE 신규 전이) 시 푸시 알람 전송
DEVICE_TOKENS_FILE = DATA_DIR / "device_tokens.json"
_fcm_initialized = False
_fcm_lock = None

def _init_fcm():
    """Firebase Admin SDK 초기화. 서비스 계정 JSON은 FIREBASE_SERVICE_ACCOUNT_JSON 환경변수."""
    global _fcm_initialized
    if _fcm_initialized:
        return True
    try:
        import firebase_admin
        from firebase_admin import credentials
        raw = os.getenv("FIREBASE_SERVICE_ACCOUNT_JSON")
        if not raw:
            print("[FCM] FIREBASE_SERVICE_ACCOUNT_JSON 환경변수 없음. FCM 비활성.")
            return False
        cred_dict = json.loads(raw)
        cred = credentials.Certificate(cred_dict)
        try:
            firebase_admin.get_app()
        except ValueError:
            firebase_admin.initialize_app(cred)
        _fcm_initialized = True
        print("[FCM] Firebase Admin 초기화 완료")
        return True
    except Exception as e:
        print(f"[FCM] 초기화 실패: {e}")
        return False


def _load_device_tokens() -> list:
    try:
        with open(DEVICE_TOKENS_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
            return data.get("tokens", []) if isinstance(data, dict) else []
    except Exception:
        return []


def _save_device_tokens(tokens: list) -> None:
    from datetime import datetime
    DEVICE_TOKENS_FILE.parent.mkdir(parents=True, exist_ok=True)
    with open(DEVICE_TOKENS_FILE, "w", encoding="utf-8") as f:
        json.dump({"updated_at": datetime.now().isoformat(), "tokens": tokens}, f, ensure_ascii=False, indent=2)


def _add_device_token(token: str, platform: str = "android", debug: bool = True) -> None:
    from datetime import datetime
    tokens = _load_device_tokens()
    for t in tokens:
        if t.get("token") == token:
            t["last_seen"] = datetime.now().isoformat()
            t["platform"] = platform
            t["debug"] = debug
            _save_device_tokens(tokens)
            return
    tokens.append({
        "token": token,
        "platform": platform,
        "debug": debug,
        "created_at": datetime.now().isoformat(),
        "last_seen": datetime.now().isoformat(),
    })
    _save_device_tokens(tokens)
    print(f"[FCM] 새 토큰 등록: {token[:20]}... (총 {len(tokens)}개)")


def _send_fcm_to_all(title: str, body: str, data: dict = None) -> dict:
    """등록된 모든 토큰에 알람 전송. 오늘 스코프: 디버그 토큰만 대상."""
    if not _init_fcm():
        return {"ok": False, "reason": "fcm_not_initialized"}
    tokens = _load_device_tokens()
    # 안전장치: debug=True 토큰만 (릴리즈 빌드는 서버에 등록 안 되므로 자동 필터)
    # 릴리즈/디버그 모두 전송. debug 라벨은 로깅 참고용.
    target_tokens = [t["token"] for t in tokens if t.get("token")]
    if not target_tokens:
        print("[FCM] 대상 토큰 없음")
        return {"ok": False, "reason": "no_tokens"}

    from firebase_admin import messaging
    sent = 0
    failed = 0
    invalid_tokens = []
    for tok in target_tokens:
        try:
            msg = messaging.Message(
                notification=messaging.Notification(title=title, body=body),
                data={k: str(v) for k, v in (data or {}).items()},
                token=tok,
                android=messaging.AndroidConfig(
                    priority="high",
                    notification=messaging.AndroidNotification(
                        channel_id="briefing_alarm",
                        priority="high",
                    ),
                ),
            )
            messaging.send(msg)
            sent += 1
        except Exception as e:
            failed += 1
            err_str = str(e)
            if "registration-token-not-registered" in err_str or "invalid-argument" in err_str:
                invalid_tokens.append(tok)
            print(f"[FCM] 전송 실패: {e}")

    # 무효 토큰 정리
    if invalid_tokens:
        remaining = [t for t in tokens if t.get("token") not in invalid_tokens]
        _save_device_tokens(remaining)
        print(f"[FCM] 무효 토큰 {len(invalid_tokens)}개 제거")

    print(f"[FCM] 전송 완료: 성공 {sent}, 실패 {failed}")
    return {"ok": True, "sent": sent, "failed": failed}


def _snapshot_notes_status() -> dict:
    """현재 notes.json의 카드별 상태 + 원인 아이템 스냅샷.

    Value 형태:
      {
        "status": "RED",
        "trigger": {"text": "...", "due_date": "2026-08-02", "days_diff": -1}
      }
    """
    from datetime import date as _date

    def _pick_worst_item(card_dict):
        """카드 안 아이템 중 상태가 가장 위험한 것 하나 반환."""
        priority = {"RED": 5, "ORANGE": 4, "BLUE": 3, "GREEN": 2, "BLACK": 1}
        worst = None
        worst_sev = -1
        for sec in card_dict.get("sections", []) or []:
            for it in (sec.get("items", []) or []):
                if not isinstance(it, dict):
                    continue
                due = (it.get("due_date") or "").strip()
                if not due:
                    continue
                s = _due_status_one(due)
                sev = priority.get(s, 0)
                if sev > worst_sev:
                    worst_sev = sev
                    worst = {"item": it, "status": s}
        return worst

    def _collect_hot_items(card_dict):
        """카드 안 RED/ORANGE 아이템 전체 수집 (신규 지연 항목 감지용)."""
        out = []
        for sec in card_dict.get("sections", []) or []:
            for it in (sec.get("items", []) or []):
                if not isinstance(it, dict):
                    continue
                due = (it.get("due_date") or "").strip()
                if not due:
                    continue
                s = _due_status_one(due)
                if s in ("RED", "ORANGE"):
                    out.append({
                        "text": (it.get("text") or it.get("title") or "").strip()[:80],
                        "due_date": due,
                        "status": s,
                    })
        return out

    try:
        notes = _load_notes()
        result = {}
        today = _date.today()
        for div_id, div_data in (notes.get("notes") or {}).items():
            for card in (div_data.get("cards") or []):
                title = card.get("title") or card.get("product") or ""
                if not title:
                    continue
                # 카드에 project_key가 이미 있으면 사용, 없으면 매핑
                project_key = (card.get("project_key") or "").strip()
                if not project_key:
                    try:
                        from project_templates import _match_project_key
                        project_key = _match_project_key(title) or ""
                    except Exception:
                        project_key = ""
                status = _calc_card_status(card) if card.get("sections") else (card.get("status") or "")
                trigger = None
                worst = _pick_worst_item(card) if card.get("sections") else None
                if worst and worst.get("status") == status:
                    it = worst["item"]
                    due = (it.get("due_date") or "").strip()
                    days_diff = None
                    try:
                        y, m, d = due[:10].split("-")
                        days_diff = (_date(int(y), int(m), int(d)) - today).days
                    except Exception:
                        days_diff = None
                    trigger = {
                        "text": (it.get("text") or it.get("title") or "").strip()[:80],
                        "due_date": due,
                        "days_diff": days_diff,
                    }
                result[f"{div_id}/{title}"] = {
                    "status": status,
                    "trigger": trigger,
                    "project_key": project_key,
                    "hot_items": _collect_hot_items(card) if card.get("sections") else [],
                }
        return result
    except Exception as e:
        print(f"[FCM] snapshot 실패: {e}")
        return {}


def _detect_status_transitions(before: dict, after: dict) -> list:
    """이전 → 이후 상태 diff. RED/ORANGE로 새로 전이된 카드만 반환.

    before/after 는 신구 구조 모두 대응 (str or dict).
    """
    def _as_status(v):
        if isinstance(v, dict):
            return v.get("status") or ""
        return v or ""

    def _trigger(v):
        if isinstance(v, dict):
            return v.get("trigger")
        return None

    def _pkey(v):
        if isinstance(v, dict):
            return v.get("project_key") or ""
        return ""

    events = []
    for key, after_val in after.items():
        new_status = _as_status(after_val)
        old_status = _as_status(before.get(key, ""))
        # 상태 전이만 알림. 나빠지는 방향 (GREEN→ORANGE, GREEN→RED, ORANGE→RED) 만 발송
        _worsen = {
            ("", "RED"), ("없음", "RED"), ("GREEN", "RED"), ("BLUE", "RED"), ("ORANGE", "RED"),
            ("", "ORANGE"), ("없음", "ORANGE"), ("GREEN", "ORANGE"), ("BLUE", "ORANGE"),
        }
        if (old_status, new_status) in _worsen:
            div_id, title = key.split("/", 1) if "/" in key else ("", key)
            events.append({
                "division_id": div_id,
                "title": title,
                "old_status": old_status or "없음",
                "new_status": new_status,
                "trigger": _trigger(after_val),
                "project_key": _pkey(after_val),
            })

    # ─── 신규 RED/ORANGE 항목 감지 (카드 색 전이 없어도 발송) ───
    from datetime import date as _date2
    _today2 = _date2.today()
    for key, after_val in after.items():
        if not isinstance(after_val, dict):
            continue
        new_status = after_val.get("status") or ""
        if new_status not in ("RED", "ORANGE"):
            continue
        old_status = _as_status(before.get(key, ""))
        if (old_status, new_status) in _worsen:
            continue  # 색 전이 알람으로 이미 처리됨
        bv = before.get(key)
        before_hot = bv.get("hot_items") if isinstance(bv, dict) else []
        before_set = {(h.get("text") or "", h.get("due_date") or "") for h in (before_hot or []) if isinstance(h, dict)}
        for h in (after_val.get("hot_items") or []):
            if not isinstance(h, dict):
                continue
            sig = (h.get("text") or "", h.get("due_date") or "")
            if sig in before_set:
                continue
            due = h.get("due_date") or ""
            days_diff = None
            try:
                yy, mm, dd = due[:10].split("-")
                days_diff = (_date2(int(yy), int(mm), int(dd)) - _today2).days
            except Exception:
                pass
            div_id, title = key.split("/", 1) if "/" in key else ("", key)
            events.append({
                "division_id": div_id,
                "title": title,
                "old_status": old_status or "없음",
                "new_status": h.get("status") or new_status,
                "trigger": {"text": h.get("text") or "", "due_date": due, "days_diff": days_diff},
                "project_key": _pkey(after_val),
                "kind": "new_item",
            })
    return events


# FCM 중복 알람 쿨다운
FCM_COOLDOWN_FILE = DATA_DIR / "fcm_cooldown.json"
FCM_COOLDOWN_SECONDS = 30 * 60  # 30분


def _load_fcm_cooldown() -> dict:
    try:
        return _read_json(FCM_COOLDOWN_FILE, {}) or {}
    except Exception:
        return {}


def _save_fcm_cooldown(d: dict) -> None:
    try:
        _write_json(FCM_COOLDOWN_FILE, d)
    except Exception as e:
        print(f"[FCM cooldown] save error: {e}")


# FCM 알람 히스토리
FCM_HISTORY_FILE = DATA_DIR / "notification_history.json"
FCM_HISTORY_MAX = 200


def _load_fcm_history() -> list:
    try:
        v = _read_json(FCM_HISTORY_FILE, [])
        return v if isinstance(v, list) else []
    except Exception:
        return []


def _save_fcm_history(hist: list) -> None:
    try:
        # 최신 200건만 유지
        if len(hist) > FCM_HISTORY_MAX:
            hist = hist[-FCM_HISTORY_MAX:]
        _write_json(FCM_HISTORY_FILE, hist)
    except Exception as e:
        print(f"[FCM history] save error: {e}")



# ─── FCM 편집 세션 관리 (알람 억제용) ───
EDIT_SESSION_FILE = Path("/tmp/fcm_edit_session.json")

def _load_edit_session() -> dict:
    """편집 세션 파일 로드. 구조: {doc_id: {"snapshot": {...}, "started_at": timestamp}}"""
    if not EDIT_SESSION_FILE.exists():
        return {}
    try:
        return json.loads(EDIT_SESSION_FILE.read_text(encoding="utf-8"))
    except Exception:
        return {}

def _save_edit_session(doc_id: str, snapshot: dict) -> None:
    """편집 세션 저장 (before snapshot)"""
    import time as _time
    sessions = _load_edit_session()
    sessions[doc_id] = {
        "snapshot": snapshot,
        "started_at": int(_time.time())
    }
    try:
        EDIT_SESSION_FILE.write_text(json.dumps(sessions, ensure_ascii=False), encoding="utf-8")
    except Exception as e:
        print(f"[FCM] edit session save error: {e}")

def _clear_edit_session(doc_id: str) -> dict:
    """편집 세션 클리어 및 저장된 snapshot 반환"""
    sessions = _load_edit_session()
    session_data = sessions.pop(doc_id, None)
    try:
        EDIT_SESSION_FILE.write_text(json.dumps(sessions, ensure_ascii=False), encoding="utf-8")
    except Exception as e:
        print(f"[FCM] edit session clear error: {e}")
    return session_data.get("snapshot") if session_data else None

def _cleanup_expired_sessions() -> None:
    """30분 이상 된 세션 자동 정리"""
    import time as _time
    now = int(_time.time())
    sessions = _load_edit_session()
    expired = [k for k, v in sessions.items() if now - v.get("started_at", 0) > 30 * 60]
    for k in expired:
        sessions.pop(k, None)
        print(f"[FCM] expired session cleaned: {k}")
    if expired:
        try:
            EDIT_SESSION_FILE.write_text(json.dumps(sessions, ensure_ascii=False), encoding="utf-8")
        except Exception:
            pass

def _fire_status_alarms(events: list) -> None:
    """상태 전이 이벤트를 FCM으로 전송. 30분 쿨다운 적용 + 히스토리 저장."""
    if not events:
        return

    import time as _time, uuid as _uuid
    from datetime import datetime as _dt
    now = int(_time.time())
    cooldown = _load_fcm_cooldown()
    history = _load_fcm_history()

    # 오래된 항목 청소 (24시간 이상 지난 것 제거)
    stale_keys = [k for k, v in cooldown.items() if not isinstance(v, (int, float)) or now - int(v) > 24 * 3600]
    for k in stale_keys:
        cooldown.pop(k, None)

    fired = 0
    skipped = 0
    for ev in events:
        div = ev.get("division_id") or ""
        proj_title = ev.get("title") or ""
        new_status = ev.get("new_status") or ""
        _due_key = ((ev.get("trigger") or {}).get("due_date") or "")
        key = f"{div}|{proj_title}|{new_status}|{_due_key}"

        last = cooldown.get(key, 0)
        try:
            last = int(last)
        except Exception:
            last = 0

        elapsed = now - last
        if last and elapsed < FCM_COOLDOWN_SECONDS:
            skipped += 1
            print(f"[FCM cooldown] skip {key} (last fired {elapsed}s ago, threshold={FCM_COOLDOWN_SECONDS}s)")
            continue

        emoji = "🔴" if new_status == "RED" else "🟠"
        # 제목은 프로젝트명만 (사업부명은 body 데이터로만 전달)
        title = f"{emoji} {proj_title}"

        # 본문 조립: 트리거 아이템이 있으면 AI 요약 + D-day 표기
        trig = ev.get("trigger") or {}
        trig_text_raw = (trig.get("text") or "").strip()
        # AI 요약 (캐시 우선, 실패 시 원문 앞 15자)
        try:
            trig_text = _ai_headline(trig_text_raw) if trig_text_raw else ""
        except Exception as _e:
            print(f"[FCM] AI 요약 실패, fallback: {_e}")
            trig_text = trig_text_raw[:20] if trig_text_raw else ""
        if not trig_text:
            trig_text = trig_text_raw[:20] if trig_text_raw else ""
        days_diff = trig.get("days_diff")
        due_date = (trig.get("due_date") or "").strip()

        if isinstance(days_diff, int):
            if days_diff < 0:
                dday_label = f"마감 지남 (D+{abs(days_diff)})"
            elif days_diff == 0:
                dday_label = "오늘 마감 (D-Day)"
            else:
                dday_label = f"마감 임박 (D-{days_diff})"
        elif due_date:
            dday_label = f"마감일 {due_date}"
        else:
            dday_label = f"{ev.get('old_status') or ''} → {new_status}"

        if trig_text:
            body = f"{trig_text} · {dday_label}"
        else:
            body = dday_label
        project_key = ev.get("project_key") or ""
        try:
            _send_fcm_to_all(title, body, data={
                "type": "status_change",
                "division_id": div,
                "title": proj_title,
                "new_status": new_status,
                "project_key": project_key,
            })
            cooldown[key] = now
            fired += 1
            history.append({
                "id": _uuid.uuid4().hex[:12],
                "ts": _dt.now().isoformat(timespec="seconds"),
                "division_id": div,
                "title": proj_title,
                "project_key": project_key,
                "old_status": ev.get("old_status") or "",
                "new_status": new_status,
                "trigger_text": trig_text,
                "due_date": due_date,
                "days_diff": days_diff if isinstance(days_diff, int) else None,
                "dday_label": dday_label,
                "read": False,
            })
        except Exception as e:
            print(f"[FCM] send error for {key}: {e}")

    _save_fcm_cooldown(cooldown)
    _save_fcm_history(history)
    print(f"[FCM] events={len(events)} fired={fired} skipped_by_cooldown={skipped}")


# ─── FCM 엔드포인트 ───

app = FastAPI(title="사업부 진행현황 보고 API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 정적 파일 마운트
app.mount("/slides", StaticFiles(directory=str(SLIDES_DIR)), name="slides")
app.mount("/slide_images", StaticFiles(directory=str(SLIDE_IMAGES_DIR)), name="slide_images")
app.mount("/cropped", StaticFiles(directory=str(CROPPED_DIR)), name="cropped")
app.mount("/note_photos", StaticFiles(directory=str(NOTE_PHOTOS_DIR)), name="note_photos")
_STATIC_DIR = BASE_DIR / "static"
_STATIC_DIR.mkdir(exist_ok=True)
app.mount("/static", StaticFiles(directory=str(_STATIC_DIR)), name="static")

# ============================================================
# 관리자 세션 (8시간)
# ============================================================
ADMIN_SESSION_COOKIE = "admin_auth"
ADMIN_SESSION_MAX_LIFETIME_SEC = int(os.getenv("ADMIN_SESSION_MAX_LIFETIME_SEC", str(24 * 3600)))  # 발급 후 최대 24시간
ADMIN_SESSION_TTL_SEC = 8 * 60 * 60
ADMIN_SESSION_SECRET = os.getenv("ADMIN_SESSION_SECRET", "")
ADMIN_COOKIE_SECURE = os.getenv("ADMIN_COOKIE_SECURE", "false").lower() == "true"
if not ADMIN_SESSION_SECRET:
    ADMIN_SESSION_SECRET = hashlib.sha256(("session::" + UPLOAD_PASSWORD).encode()).hexdigest()

def _b64u(b: bytes) -> str:
    return base64.urlsafe_b64encode(b).rstrip(b"=").decode()

def _b64u_decode(s: str) -> bytes:
    pad = "=" * (-len(s) % 4)
    return base64.urlsafe_b64decode(s + pad)

def _sign_session(exp_ts: int, issued_at: int) -> str:
    msg = f"{issued_at}.{exp_ts}".encode()
    sig = hmac.new(ADMIN_SESSION_SECRET.encode(), msg, hashlib.sha256).digest()
    return f"{issued_at}.{exp_ts}.{_b64u(sig)}"

def _verify_session(token: Optional[str]) -> Optional[int]:
    if not token:
        return None
    try:
        parts = token.split(".")
        if len(parts) != 3:
            return None
        issued_str, exp_str, sig_b64 = parts
        issued_at = int(issued_str)
        exp_ts = int(exp_str)
        msg = f"{issued_at}.{exp_ts}".encode()
        expected = hmac.new(ADMIN_SESSION_SECRET.encode(), msg, hashlib.sha256).digest()
        provided = _b64u_decode(sig_b64)
        if not hmac.compare_digest(expected, provided):
            return None
        now = int(time.time())
        if exp_ts < now:
            return None
        # 발급 후 최대 수명 초과 검사
        if now - issued_at > ADMIN_SESSION_MAX_LIFETIME_SEC:
            return None
        return exp_ts
    except Exception:
        return None

def _issue_session_cookie(response: Response, issued_at: Optional[int] = None):
    now = int(time.time())
    if issued_at is None:
        issued_at = now  # 새 발급
    exp_ts = now + ADMIN_SESSION_TTL_SEC
    # 발급 후 최대 수명 초과 방지: exp_ts를 issued_at + max_lifetime 으로 제한
    hard_limit = issued_at + ADMIN_SESSION_MAX_LIFETIME_SEC
    if exp_ts > hard_limit:
        exp_ts = hard_limit
    token = _sign_session(exp_ts, issued_at)
    response.set_cookie(
        key=ADMIN_SESSION_COOKIE,
        value=token,
        max_age=max(0, exp_ts - now),
        httponly=True,
        samesite="lax",
        secure=ADMIN_COOKIE_SECURE,
        path="/",
    )
    return exp_ts

def get_admin_session(admin_auth: Optional[str] = Cookie(default=None)) -> int:
    exp = _verify_session(admin_auth)
    if not exp:
        raise HTTPException(status_code=401, detail="인증이 필요합니다.")
    return exp

_ADMIN_LOGIN_HTML = """<!doctype html>
<html lang=\"ko\"><head>
<meta charset=\"utf-8\" />
<title>관리자 인증</title>
<meta name=\"viewport\" content=\"width=device-width,initial-scale=1\" />
<style>
  html, body { height:100%; }
  body { font-family: -apple-system, BlinkMacSystemFont, 'Pretendard', sans-serif; background:#f1f5f9; margin:0; min-height:100vh; display:flex; align-items:center; justify-content:center; padding: 16px; box-sizing: border-box; }
  .card { background:#fff; padding:32px; border-radius:14px; box-shadow:0 8px 30px rgba(0,0,0,0.08); width:340px; margin: auto; }
  h2 { margin:0 0 8px; color:#1E3A5F; font-size:20px; }
  p { color:#64748b; font-size:13px; margin:0 0 18px; }
  input[type=password] { width:100%; padding:12px; font-size:14px; border:1px solid #cbd5e1; border-radius:8px; box-sizing:border-box; }
  button { width:100%; padding:12px; margin-top:14px; background:#1E3A5F; color:#fff; border:none; border-radius:8px; font-weight:700; font-size:14px; cursor:pointer; }
  .err { color:#b91c1c; font-size:13px; margin-top:10px; min-height:18px; }
</style></head>
<body>
  <form class="card" id="loginForm">
    <h2>🔒 관리자 인증</h2>
    <p>업로드 비밀번호를 입력하세요. 8시간 유지됩니다.</p>
    <input type=\"password\" id=\"pw\" placeholder=\"비밀번호\" required autofocus />
    <button type=\"submit\">로그인</button>
    <div class=\"err\" id=\"err\"></div>
  </form>
<script>
document.getElementById('loginForm').addEventListener('submit', async (e) => {
  e.preventDefault();
  const pw = document.getElementById('pw').value;
  const err = document.getElementById('err');
  err.textContent = '';
  try {
    const fd = new FormData();
    fd.append('password', pw);
    const r = await fetch('/admin/login', { method: 'POST', body: fd, credentials: 'same-origin' });
    if (r.ok) {
      const next = new URLSearchParams(location.search).get('next') || '/admin/upload';
      location.href = next;
    } else {
      const j = await r.json().catch(()=>({}));
      err.textContent = j.detail || '비밀번호가 올바르지 않습니다.';
    }
  } catch (e) {
    err.textContent = '네트워크 오류';
  }
});
</script>


</body></html>"""

@app.get("/admin/login", response_class=HTMLResponse)
def admin_login_page(
    admin_auth: Optional[str] = Cookie(default=None),
    next: Optional[str] = None,
):
    # next 안전성 체크: 반드시 /admin/ 하위만 허용
    safe_next = "/admin/upload"
    if next and next.startswith("/admin/") and "://" not in next:
        safe_next = next
    if _verify_session(admin_auth):
        return RedirectResponse(url=safe_next, status_code=302)
    # HTML 안에 next 값을 script 로 심어서 로그인 성공 후 리다이렉트에 활용
    html = _ADMIN_LOGIN_HTML.replace(
        "</head>",
        "<script>window.__LOGIN_NEXT__ = " + repr(safe_next) + ";</script></head>",
        1,
    )
    return HTMLResponse(content=html)

@app.post("/admin/login")
def admin_login_submit(
    response: Response,
    password: str = Form(...),
    next: Optional[str] = Form(default=None),
):
    if password != UPLOAD_PASSWORD:
        raise HTTPException(status_code=401, detail="비밀번호가 올바르지 않습니다.")
    exp = _issue_session_cookie(response)
    safe_next = "/admin/upload"
    if next and next.startswith("/admin/") and "://" not in next:
        safe_next = next
    return {"ok": True, "expires_at": exp, "next": safe_next}

@app.post("/admin/logout")
def admin_logout(response: Response):
    response.delete_cookie(ADMIN_SESSION_COOKIE, path="/")
    return {"ok": True}

@app.post("/admin/extend")
def admin_extend(
    response: Response,
    admin_auth: Optional[str] = Cookie(default=None),
    _exp: int = Depends(get_admin_session),
):
    # 기존 토큰의 issued_at 보존 → 최대 수명 24시간 강제
    issued_at = None
    if admin_auth:
        try:
            parts = admin_auth.split(".")
            if len(parts) == 3:
                issued_at = int(parts[0])
        except Exception:
            pass
    now = int(time.time())
    if issued_at is not None and now - issued_at >= ADMIN_SESSION_MAX_LIFETIME_SEC:
        raise HTTPException(status_code=401, detail="세션 최대 수명을 초과했습니다. 다시 로그인해 주세요.")
    exp = _issue_session_cookie(response, issued_at=issued_at)
    return {"ok": True, "expires_at": exp, "max_lifetime_remaining": (issued_at + ADMIN_SESSION_MAX_LIFETIME_SEC - now) if issued_at else ADMIN_SESSION_MAX_LIFETIME_SEC}

@app.get("/admin/session")
def admin_session_info(_exp: int = Depends(get_admin_session)):
    return {"ok": True, "expires_at": _exp, "ttl_sec": ADMIN_SESSION_TTL_SEC}



# =========================================================
# 유틸
# =========================================================
def _read_json(path: Path, default):
    if not path.exists():
        return default
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return default


def _write_json(path: Path, data):
    path.write_text(
        json.dumps(data, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


def _make_doc_id(filename: str) -> str:
    short = hashlib.md5(
        f"{filename}_{datetime.now().isoformat()}".encode("utf-8")
    ).hexdigest()[:12]
    return short


def _convert_pptx_to_png(pptx_path: Path, out_dir: Path):
    """LibreOffice + Poppler 로 PPTX → PDF → PNG 변환"""
    out_dir.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        ["soffice", "--headless", "--convert-to", "pdf",
         "--outdir", str(out_dir), str(pptx_path)],
        check=True, capture_output=True,
    )
    pdf_file = out_dir / (pptx_path.stem + ".pdf")
    if not pdf_file.exists():
        raise RuntimeError("PDF 변환 실패")
    subprocess.run(
        ["pdftoppm", "-png", "-r", "150",
         str(pdf_file), str(out_dir / "slide")],
        check=True, capture_output=True,
    )
    # slide-1.png → slide_01.png 정규화
    for p in sorted(out_dir.glob("slide-*.png")):
        num = int(p.stem.split("-")[-1])
        p.rename(out_dir / f"slide_{num:02d}.png")
    try:
        pdf_file.unlink()
    except Exception:
        pass


def _extract_text_per_slide(pptx_path: Path):
    """슬라이드별 텍스트 추출"""
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    texts = []
    for slide in prs.slides:
        buf = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                buf.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        buf.append(cell.text)
        texts.append("\n".join(buf))
    return texts


def _extract_pptx_text_and_tables(pptx_path: Path):
    """슬라이드별 텍스트(표 제외)와 표를 분리 추출.
    반환: (slide_texts: list[str], slide_tables: list[list[dict]])
    각 표: {"headers": [...], "rows": [[...], ...]}
    """
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    slide_texts = []
    slide_tables = []
    for slide in prs.slides:
        text_buf = []
        tables_in_slide = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text
                if t and t.strip():
                    text_buf.append(t)
            if shape.has_table:
                tbl = shape.table
                rows_data = []
                for row in tbl.rows:
                    rows_data.append([cell.text.strip() for cell in row.cells])
                if rows_data:
                    headers = rows_data[0]
                    body = rows_data[1:] if len(rows_data) > 1 else []
                    tables_in_slide.append({"headers": headers, "rows": body})
        slide_texts.append("\n".join(text_buf))
        slide_tables.append(tables_in_slide)
    return slide_texts, slide_tables


def _classify_status_with_gpt(slide_texts: list[str]) -> dict:
    if not client:
        return {"products": []}

    joined = "\n\n".join([f"=== Slide {i+1} ===\n{t}" for i, t in enumerate(slide_texts)])

    system_prompt = """너는 반도체 사업부의 임원 보고서 분석 어시스턴트다.
PPT 슬라이드 텍스트를 읽고, 등장하는 모든 프로젝트/제품을 JSON으로 정리한다.

[가장 중요] 사장님이 30초 안에 다 파악할 수 있도록 정리한다.
- summary_bullets: 핵심 3줄 (한 줄 25자 이내, 가장 중요한 사실만)
- sections: 상세 내용 (펼쳤을 때 보임)

규칙:
1. 섹션 구조는 PPT 원문에 등장한 그대로 따른다 (양산/개발/EMA/주차별 출하/Dep챔버진행 등).
2. PPT에 없는 섹션은 만들지 말 것.
3. items는 원문을 가능하면 한 줄로 축약 (불필요한 수식어 제거, 핵심 숫자/일정 유지).
4. notes에는 ※ 또는 참고 사항만.
5. status:
   - RED: 출하 미달, PO 미접수, 자재 쇼티지, 일정 지연 등 즉시 확인
   - BLUE: 진행 중, 정상 추진
   - BLACK: 완료/이슈 없음
6. headline: 가장 중요한 한 줄 (15자 이내 권장).
7. summary_bullets: 핵심 3줄 (정확히 3개, 각 25자 이내).

JSON 스키마:
{
  "products": [
    {
      "name": "프로젝트명",
      "category": "분류",
      "status": "RED|BLUE|BLACK",
      "headline": "핵심 한 줄",
      "summary_bullets": ["요약1", "요약2", "요약3"],
      "header_summary": "상단 요약 (선택)",
      "sections": [
        {"title": "섹션 제목", "items": ["불릿1"], "notes": ["※ 참고"]}
      ]
    }
  ]
}"""

    try:
        resp = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"다음 PPT 내용을 정리:\n\n{joined}"},
            ],
            response_format={"type": "json_object"},
            temperature=0.2,
        )
        return json.loads(resp.choices[0].message.content)
    except Exception as e:
        print(f"🔴 GPT 분류 실패: {e}")
        return {"products": []}

def _cleanup_old_files():
    """6개월 지난 파일 자동 삭제"""
    cutoff = datetime.now() - timedelta(days=RETENTION_DAYS)
    for base in [UPLOAD_DIR, SLIDES_DIR, SLIDE_IMAGES_DIR, CROPPED_DIR]:
        if not base.exists():
            continue
        for entry in base.iterdir():
            try:
                mtime = datetime.fromtimestamp(entry.stat().st_mtime)
                if mtime < cutoff:
                    if entry.is_dir():
                        shutil.rmtree(entry, ignore_errors=True)
                    else:
                        entry.unlink()
            except Exception:
                continue


def _update_latest(summary):
    """
    최신 보고서 저장.
    - 같은 doc_id가 있으면 갱신 (중복 방지)
    - 그 외에는 누적
    - report_family는 부가 정보로만 사용 (덮어쓰기 키로 사용 X)
    - 최대 50개까지만 유지 (오래된 것부터 제거)
    """
    latest = _read_json(LATEST_FILE, [])
    new_doc_id = summary.get("doc_id")

    # 같은 doc_id 있으면 제거하고 새로 추가
    latest = [r for r in latest if r.get("doc_id") != new_doc_id]
    latest.append(summary)

    # 업로드 시각 기준 최신순 정렬, 최대 50개
    latest.sort(
        key=lambda r: r.get("uploaded_at") or r.get("report_meta", {}).get("date", ""),
        reverse=True
    )
    latest = latest[:50]

    _write_json(LATEST_FILE, latest)

# =========================================================
# 1. 헬스체크
# =========================================================
@app.get("/health")
def health():
    return {"status": "ok", "time": datetime.now().isoformat()}


# =========================================================
# 2. PPT 업로드
# =========================================================


def _load_upload_hash_index():
    try:
        if not UPLOAD_HASHES_PATH.exists():
            return {}
        return json.loads(UPLOAD_HASHES_PATH.read_text(encoding="utf-8"))
    except Exception:
        return {}

def _save_upload_hash_index(data):
    try:
        UPLOAD_HASHES_PATH.write_text(
            json.dumps(data, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
    except Exception:
        pass

def _find_duplicate_upload_by_hash(file_hash: str):
    idx = _load_upload_hash_index()
    return idx.get(file_hash)

def _record_upload_hash(file_hash: str, saved_name: str, doc_id: str):
    idx = _load_upload_hash_index()
    idx[file_hash] = {
        "filename": saved_name,
        "doc_id": doc_id,
    }
    _save_upload_hash_index(idx)

@app.post("/upload")
async def upload_ppt(
    file: UploadFile = File(...),
    password: Optional[str] = Form(None),
    admin_auth: Optional[str] = Cookie(default=None),
    report_family: str = Form("default"),
    division_id: Optional[str] = Form(None),  # 5-2b: 사업부 사전 매핑 (선택)
    project_id: Optional[str] = Form(None),   # 5-2b: 프로젝트 사전 매핑 (선택)
    allow_duplicate: bool = Form(False),  # 중복 업로드 허용
):
    _has_session = bool(_verify_session(admin_auth))
    if not _has_session and password != UPLOAD_PASSWORD:
        raise HTTPException(status_code=401, detail="비밀번호가 틀립니다.")
    if not file.filename.lower().endswith((".pptx", ".ppt")):
        raise HTTPException(status_code=400, detail=".pptx 또는 .ppt 파일만 업로드 가능합니다.")

    # 파일 내용 읽기 + 해시 계산
    file_bytes = await file.read()
    file_hash = hashlib.sha256(file_bytes).hexdigest()

    # 중복 업로드 검사 (allow_duplicate 이 False 일 때만 차단)
    duplicate_hit = _find_duplicate_upload_by_hash(file_hash)
    if duplicate_hit and not allow_duplicate:
        return JSONResponse(
            status_code=409,
            content={
                "ok": False,
                "code": "duplicate_upload",
                "detail": "이미 업로드된 파일입니다. 같은 파일을 다시 올리려면 '중복 업로드 허용'을 체크하세요.",
                "duplicate_of": duplicate_hit.get("filename"),
            },
        )

    _cleanup_old_files()

    doc_id = _make_doc_id(file.filename)
    saved_pptx_path = UPLOAD_DIR / f"{doc_id}.pptx"
    saved_pptx_path.write_bytes(file_bytes)
    _record_upload_hash(file_hash, saved_pptx_path.name, doc_id)

    # slides/ 와 slide_images/ 둘 다 저장 (호환성)
    slides_out_dir = SLIDES_DIR / doc_id
    try:
        _convert_pptx_to_png(saved_pptx_path, slides_out_dir)
        # slide_images/ 에도 복사
        slide_images_out = SLIDE_IMAGES_DIR / doc_id
        if not slide_images_out.exists():
            shutil.copytree(slides_out_dir, slide_images_out)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"슬라이드 변환 실패: {e}")

    slide_texts = _extract_text_per_slide(saved_pptx_path)
    gpt_result = _classify_status_with_gpt(slide_texts)
    products = gpt_result.get("products", [])

    section_images = {}
    try:
        section_images = extract_charts_from_pptx(
            pptx_path=saved_pptx_path,
            slides_image_dir=slides_out_dir,
            output_dir=CROPPED_DIR,
            doc_id=doc_id,
        )
        print(f"🟢 차트 자동 추출 완료: {list(section_images.keys())}")
    except Exception as e:
        print(f"🔴 차트 자동 추출 실패: {e}")

    slide_images = {
        str(i + 1): f"/slides/{doc_id}/slide_{i+1:02d}.png"
        for i in range(len(slide_texts))
    }


    # ============================================================
    # 5-2b: 사장님이 사전 지정한 division_id/project_id 폴백 적용
    # - 자동 분류 성공 카드: 그대로 유지
    # - 자동 분류 실패 카드만 사장님 의도로 폴백
    # ============================================================
    if division_id or project_id:
        import config_loader as _cl
        try:
            _cl.reload()
        except Exception:
            pass
        for p in products:
            name = p.get("name") or p.get("product", "")
            category = p.get("category", "")
            text = f"{name} {category}".strip()
            try:
                auto_pid = _cl.classify_project(text)
            except Exception:
                auto_pid = None
            if not auto_pid:
                # 자동 분류 실패한 카드만 폴백
                if project_id:
                    p["_fallback_project_id"] = project_id
                if division_id:
                    p["_fallback_division_id"] = division_id

    summary = {
        "report_meta": {
            "report_family": report_family,
            "date": datetime.now().strftime("%Y-%m-%d"),
        },
        "doc_id": doc_id,
        "file_name": file.filename,
        "upload_timestamp": datetime.now().isoformat(),
        "products": products,
        "section_images": section_images,
        "slide_images": slide_images,
    }
    _update_latest(summary)

    # ============================================================
    # KPI 자동 추출 훅 (신규)
    # - 프로파일 등록된 프로젝트만 자동 반영
    # - 실패해도 업로드는 성공
    # ============================================================
    _kpi_extract_result = _auto_extract_kpi_from_pptx(
        saved_pptx_path,
        file.filename,
        doc_id,
        project_id_hint=project_id,
        products=products,
    )

    return {
        "ok": True,
        "doc_id": doc_id,
        "product_count": len(products),
        "slide_count": len(slide_texts),
        "kpi_auto_extract": _kpi_extract_result,
    }


# =========================================================
# 3. 대시보드 / 리포트
# =========================================================
def _due_status_one(due_raw: str) -> str:
    """단일 due_date 문자열 → 상태색."""
    from datetime import date
    due_raw = (due_raw or "").strip()
    if not due_raw:
        return "BLACK"  # 일정 없음 (내용은 있는데 일정 추출 안 됨)
    try:
        y, m, d = due_raw[:10].split("-")
        due = date(int(y), int(m), int(d))
    except Exception:
        return "BLACK"
    diff = (due - date.today()).days
    if diff < 0:
        return "RED"
    if diff <= 3:
        return "ORANGE"
    return "BLUE"


def _calc_card_status(card: dict) -> str:
    """카드 안 items 들의 due_date 들을 보고 가장 위험한 색으로 집계.
    빨강 > 주황 > 파랑 > 초록(매출만) > 검정 순.
    items가 비어있고 매출 데이터도 없으면 → BLACK.
    매출 데이터가 있고 due 항목이 없으면 → GREEN (정상 운영중)."""
    priority = {"RED": 5, "ORANGE": 4, "BLUE": 3, "GREEN": 2, "BLACK": 1}
    best = "BLACK"
    sections = card.get("sections") or []
    has_sales = False
    for sec in sections:
        # 매출 데이터 감지
        if (sec.get("sales_summary") or "").strip():
            has_sales = True
        for it in (sec.get("items") or []):
            if not isinstance(it, dict):
                continue
            s = _due_status_one(it.get("due_date") or "")
            if priority.get(s, 0) > priority.get(best, 0):
                best = s
                if best == "RED":
                    return best
    # due 상태가 BLACK인데 매출은 있으면 GREEN 승격
    if best == "BLACK" and has_sales:
        return "GREEN"
    return best


@app.get("/dashboard")
def dashboard():
    """신호등 카드 - 각 카드에 project_key 포함 → Flutter에서 상세 이동에 사용"""
    from project_templates import _match_project_key  # 키워드 매핑 함수 재사용

    latest = _read_json(LATEST_FILE, [])
    cards = []
    for report in latest:
        report_date = (
            report.get("report_meta", {}).get("date")
            or report.get("report_date")
        )
        report_family = (
            report.get("report_meta", {}).get("report_family")
            or report.get("report_family", "")
        )
        for p in report.get("products", []):
            product_name = p.get("name") or p.get("product", "")
            category = p.get("category", "")
            # 프로젝트 키 매핑 (백엔드 로직 그대로)
            project_key = (
                _match_project_key(product_name)
                or _match_project_key(category)
                or _match_project_key(report_family)
            )
            cards.append({
                "doc_id": report.get("doc_id", ""),
                "product": product_name,
                "status": p.get("status", "BLACK"),
                "headline": p.get("headline", ""),
                "report_date": report_date,
                "report_family": report_family,
                "project_key": project_key,  # ← 추가
            })

    # 🟢 카드별로 division/project 분류 정보 enrichment (기존 필드 변경 없음)
    cards = [enrich_card(c) for c in cards]

    # 🟢 주간 보고 노트 카드도 대시보드에 포함
    # 같은 project_key 가 PPT/노트 양쪽에 있으면 노트 카드를 우선
    try:
        notes_data = _load_notes()
        notes_map = notes_data.get("notes", {}) or {}
        note_cards = []
        for div_id, div_notes in notes_map.items():
            report_date = (div_notes or {}).get("report_date", "") or ""
            for nc in (div_notes or {}).get("cards", []) or []:
                if not isinstance(nc, dict):
                    continue
                title = (nc.get("title") or "").strip()
                if not title:
                    continue
                # 카드 제목 → project_key 매칭
                pkey = (nc.get("project_key") or "").strip()
                if not pkey:
                    try:
                        pkey = _cl.classify_project(title) or ""
                    except Exception:
                        pkey = ""

                # headline / summary_bullets 추출
                headline = ""
                bullets = []
                from datetime import date as _date
                due_items = []        # (date, text, raw_str)
                dated_bullets = []    # 날짜 있는 항목만 (앱 카드 본문용)
                highlight_text = ""
                due_date_min = None
                for sec in (nc.get("sections") or []):
                    for it in (sec.get("items") or []):
                        if not isinstance(it, dict):
                            if isinstance(it, str) and len(bullets) < 3:
                                bullets.append(it)
                            continue
                        t = (it.get("type") or "bullet").lower()
                        txt = (it.get("text") or "").strip()
                        if not txt:
                            continue
                        if t == "highlight" and not highlight_text:
                            highlight_text = txt
                        d = (it.get("due_date") or "").strip()
                        if d:
                            try:
                                y, m, dd = d[:10].split("-")
                                due_d = _date(int(y), int(m), int(dd))
                                due_items.append((due_d, txt, d[:10]))
                                dated_bullets.append(f"{txt} ({d[5:10]})")
                                if due_date_min is None or due_d < due_date_min:
                                    due_date_min = due_d
                            except Exception:
                                pass
                        if t in ("bullet", "group_note", "sub") and len(bullets) < 5:
                            bullets.append(txt)

                                # 카드 본문: 날짜 있는 항목 우선, 없으면 일반 bullets 사용
                # 완전히 빈 카드(bullets/highlight/dated/sales 모두 없음)만 스킵
                # sales_summary 가 있으면 매출 카드로 인정
                sales_summary_any = ""
                for _sec in (nc.get("sections") or []):
                    _ss = (_sec or {}).get("sales_summary")
                    if _ss and str(_ss).strip():
                        sales_summary_any = str(_ss).strip()
                        break

                if not dated_bullets and not bullets and not highlight_text and not sales_summary_any:
                    continue

                if dated_bullets:
                    bullets = dated_bullets[:10]
                else:
                    bullets = bullets[:10]

                # 매출만 있는 카드: 요약을 bullets/headline 으로 사용
                if not bullets and sales_summary_any:
                    bullets = [sales_summary_any]

                # headline: summary_bullets[0] 또는 첫 dated bullet 에서 날짜 제거
                def _shorten(txt: str, limit: int = 20) -> str:
                    txt = (txt or "").strip()
                    if len(txt) <= limit:
                        return txt
                    return txt[:limit].rstrip() + "…"

                # headline: due_date_min 항목 우선 선택 후 AI 15자 요약
                headline_src = _pick_headline_source(due_items, due_date_min)
                headline = _ai_headline(headline_src) if headline_src else ""
                # 매출만 있는 카드: headline 을 매출 요약으로 대체
                if not headline and sales_summary_any:
                    headline = sales_summary_any
                due_date_min_str = due_date_min.isoformat() if due_date_min else None

                computed_status = _calc_card_status(nc)
                note_cards.append({
                    "doc_id": f"note:{div_id}",
                    "product": title,
                    "status": computed_status,
                    "headline": headline,
                    "summary_bullets": bullets,
                    "due_date_min": due_date_min_str,
                    "report_date": report_date,
                    "report_family": "weekly_note",
                    "project_key": pkey or None,
                    "division_id": div_id,
                    "from_note": True,
                })

        # enrich (project_key 가 있으면 division/label 등 자동 채움)
        note_cards = [enrich_card(c) for c in note_cards]

        # 대시보드는 노트 카드만 표시 (PPT 카드는 관리자용 텍스트 추출 도구로만)
        cards = note_cards
    except Exception as _e:
        print(f"노트 대시보드 머지 실패: {_e}")
        cards = []

    severity = {"RED": 5, "ORANGE": 4, "BLUE": 3, "GREEN": 2, "GRAY": 2, "BLACK": 1}
    cards.sort(key=lambda c: -severity.get(c["status"], 0))

    # 🟢 project_key 기준 중복 제거 (같은 프로젝트는 1개만 유지)
    # - severity 높은 순으로 정렬되어 있으므로 첫 번째 카드만 남김
    # - project_key 없으면 product(제목)으로 대체
    seen_keys = set()
    deduped = []
    for c in cards:
        key = c.get("project_key") or c.get("product") or ""
        if not key:
            # 키가 없으면 고유 ID로 처리 (중복 아님)
            key = c.get("doc_id", "") or str(id(c))
        if key in seen_keys:
            continue
        seen_keys.add(key)
        deduped.append(c)
    cards = deduped

    # 🟢 모델 단위 그룹핑 (신규 필드, 옛 cards 는 호환을 위해 그대로 유지)
    grouped_cards = _group_dashboard_cards(cards)

    return {"cards": cards, "grouped_cards": grouped_cards}


def _parse_report_filename(file_name: str, fallback_date: str = "") -> dict:
    """
    파일명에서 프로젝트명 / 날짜 / 주차 추출.

    Examples:
        "260703_파워박스 진행현황.pptx"
            -> {"projects": ["파워박스"], "date": "2026-07-03", "week": 27, "display_title": "파워박스 · W27 주간보고"}
        "260528_내재화프레임, 메이져모듈 진행현황.pptx"
            -> {"projects": ["내재화프레임", "메이져모듈"], "date": "2026-05-28", "week": 22, "display_title": "내재화프레임, 메이져모듈 · W22 주간보고"}
        "Chamber_W26_20260706.pptx"
            -> {"projects": ["Chamber"], "date": "2026-07-06", "week": 26, "display_title": "Chamber · W26 주간보고"}
    """
    import re as _re
    import unicodedata as _ud
    from datetime import datetime as _dt, date as _date

    # macOS 파일명은 NFD(자모 분리) 로 저장되어 있을 수 있음 → NFC 로 정규화
    file_name = _ud.normalize("NFC", file_name or "")
    name = file_name.rsplit(".", 1)[0]  # 확장자 제거

    projects: list = []
    date_str = ""
    week_num: int = 0

    # 패턴 A: YYMMDD_프로젝트명(,프로젝트명) 진행현황
    m = _re.match(r"^(\d{6})[_\s]+(.+)$", name)
    if m:
        yymmdd = m.group(1)
        rest = m.group(2).strip()
        # " 진행현황", "_진행현황", "진행 현황" 등 trailing 제거
        rest = _re.sub(r"[\s_]*진행\s*현황.*$", "", rest).strip()
        rest = rest.rstrip("_ .").strip()
        try:
            dt = _dt.strptime(yymmdd, "%y%m%d").date()
            date_str = dt.isoformat()
            week_num = dt.isocalendar()[1]
        except Exception:
            pass
        # 콤마 분리
        projects = [x.strip() for x in _re.split(r"[,、/]", rest) if x.strip()]

    # 패턴 B: 프로젝트명_W##_YYYYMMDD
    if not projects:
        m2 = _re.match(r"^(.+?)_W(\d{1,2})_(\d{8})$", name)
        if m2:
            projects = [m2.group(1).strip()]
            try:
                week_num = int(m2.group(2))
                dt = _dt.strptime(m2.group(3), "%Y%m%d").date()
                date_str = dt.isoformat()
            except Exception:
                pass

    # fallback: 파일명 전체를 프로젝트명 취급
    if not projects:
        projects = [name]

    # fallback date
    if not date_str and fallback_date:
        try:
            dt = _dt.strptime(fallback_date, "%Y-%m-%d").date()
            date_str = fallback_date
            if not week_num:
                week_num = dt.isocalendar()[1]
        except Exception:
            pass

    display_title = ", ".join(projects)
    if week_num:
        display_title += " · W" + str(week_num) + " 주간보고"

    return {
        "projects": projects,
        "date": date_str,
        "week": week_num,
        "display_title": display_title,
    }


def _normalize_project_name(name) -> str:
    """프로젝트명 정규화: trim + 연속공백을 단일공백으로 + 자주 틀리는 표기 통일."""
    if not name:
        return ""
    s = str(name).strip()
    s = re.sub(r"\s+", " ", s)
    # 자주 틀리는 한글 표기 통일 (dedupe 기준값만; 실제 저장/표시 값은 안 바꿈)
    _alias = {
        "메이져": "메이저",
        "메이져모듈": "메이저모듈",
        "매이저": "메이저",
        "매이져": "메이저",
    }
    for k, v in _alias.items():
        s = s.replace(k, v)
    return s


def _report_sort_key(r: dict):
    """dedupe용 정렬 키. 최신이 클수록 큰 값.
    uploaded_at -> created_at -> parsed.date 순으로 fallback."""
    if not isinstance(r, dict):
        return ""
    for k in ("uploaded_at", "created_at"):
        v = r.get(k)
        if v:
            return str(v)
    parsed = r.get("parsed") or {}
    d = parsed.get("date") or (r.get("report_meta") or {}).get("date") or ""
    return str(d)


def _dedupe_cards_by_project(cards: list) -> list:
    """같은 프로젝트명이 여러 카드에 있으면 최신 것만 남긴다.
    비교 전 이름을 정규화. 순서는 최신순으로 정렬."""
    if not cards:
        return []
    best = {}
    for c in cards:
        if not isinstance(c, dict):
            continue
        parsed = c.get("parsed") or {}
        projs = parsed.get("projects") or []
        proj = _normalize_project_name(projs[0] if projs else "")
        if not proj:
            # 프로젝트명 없는 카드는 doc_id 로 유지
            proj = "__no_project__" + str(c.get("doc_id", ""))
        key = proj
        prev = best.get(key)
        if prev is None or _report_sort_key(c) > _report_sort_key(prev):
            best[key] = c
    # 최신순 정렬
    out = list(best.values())
    out.sort(key=_report_sort_key, reverse=True)
    return out


def _split_report_by_project(r_enriched: dict) -> list:
    """한 리포트에 여러 프로젝트가 들어있으면 프로젝트 단위 카드로 분리.
    project_overrides가 있으면 프로젝트명을 그걸로 대체.
    """
    parsed = (r_enriched or {}).get("parsed") or {}
    projects = parsed.get("projects") or []
    week_num = (r_enriched or {}).get("week_override") or parsed.get("week")
    overrides = (r_enriched or {}).get("project_overrides") or {}
    def _resolve(name):
        return overrides.get(name, name)
    if len(projects) <= 1:
        # 단일 카드도 override 반영
        if projects:
            proj = _resolve(projects[0])
            new_parsed = dict(parsed)
            new_parsed["projects"] = [proj]
            new_parsed["week"] = week_num
            new_title = proj + (" · W" + str(week_num) + " 주간보고" if week_num else "")
            new_parsed["display_title"] = new_title
            r_enriched = dict(r_enriched)
            r_enriched["parsed"] = new_parsed
            r_enriched["display_title"] = new_title
            r_enriched["_split_project"] = projects[0]  # 원본 이름을 키로 유지
        return [r_enriched]
    cards = []
    for proj in projects:
        clone = dict(r_enriched)
        display_proj = _resolve(proj)
        new_parsed = dict(parsed)
        new_parsed["projects"] = [display_proj]
        new_parsed["week"] = week_num
        new_title = display_proj
        if week_num:
            new_title += " · W" + str(week_num) + " 주간보고"
        new_parsed["display_title"] = new_title
        clone["parsed"] = new_parsed
        clone["display_title"] = new_title
        clone["_split_project"] = proj  # 원본 이름을 키로 유지 (override 저장에 사용)
        cards.append(clone)
    return cards


def _classify_report_status(upload_timestamp: str, report_date: str) -> dict:
    """업로드 후 경과 시간과 report_date 기준 상태 분류.

    - 업로드 후 5분 이내: ai_processing (Chamber D-2 예시처럼)
    - report_date가 미래 or 오늘: review_pending (검토 대기)
    - 그 외: published (발행 완료)
    """
    from datetime import datetime as _dt, timedelta as _td
    now = _dt.now()

    status = "published"
    d_day = None

    try:
        if upload_timestamp:
            ts = _dt.fromisoformat(upload_timestamp)
            if now - ts < _td(minutes=5):
                status = "ai_processing"
    except Exception:
        pass

    try:
        if report_date and status != "ai_processing":
            rd = _dt.strptime(report_date, "%Y-%m-%d").date()
            today = now.date()
            delta = (rd - today).days
            d_day = delta
            if delta >= 0:
                status = "review_pending"
            else:
                status = "published"
    except Exception:
        pass

    return {"status": status, "d_day": d_day}


@app.get("/reports")
def list_reports():
    raw = _read_json(LATEST_FILE, [])
    enriched = []
    for r in (raw or []):
        try:
            if (r or {}).get("hidden") is True:
                continue
            meta = (r or {}).get("report_meta") or {}
            fname = (r or {}).get("file_name") or ""
            parsed = _parse_report_filename(fname, meta.get("date", ""))
            manual_projs = (r or {}).get("manual_projects")
            if manual_projs:
                parsed = dict(parsed)
                parsed["projects"] = list(manual_projs)
            mw = (r or {}).get("week_override")
            if mw:
                parsed["week"] = mw
            if parsed.get("projects"):
                _wk = parsed.get("week")
                parsed["display_title"] = ", ".join(parsed["projects"]) + (" · W" + str(_wk) + " 주간보고" if _wk else "")
            classified = _classify_report_status(
                (r or {}).get("upload_timestamp", ""),
                parsed.get("date") or meta.get("date", "")
            )
            r_enriched = dict(r or {})
            r_enriched["parsed"] = parsed
            r_enriched["display_title"] = parsed.get("display_title", "")
            r_enriched["report_status"] = classified.get("status", "published")
            r_enriched["d_day"] = classified.get("d_day")
            enriched.append(r_enriched)
        except Exception:
            enriched.append(r)
    split_cards = []
    for e in enriched:
        try:
            split_cards.extend(_split_report_by_project(e))
        except Exception:
            split_cards.append(e)
    split_cards = _dedupe_cards_by_project(split_cards)
    return {"reports": split_cards}


@app.get("/reports/{doc_id}/{product_name}")
def get_product_detail(doc_id: str, product_name: str):
    """제품 단위 상세 (기존 카드 탭용)"""
    latest = _read_json(LATEST_FILE, [])
    for report in latest:
        if report.get("doc_id") != doc_id:
            continue
        for p in report.get("products", []):
            name = p.get("product") or p.get("name", "")
            if name == product_name:
                return {
                    "doc_id": doc_id,
                    "report_date": (
                        report.get("report_meta", {}).get("date")
                        or report.get("report_date")
                    ),
                    "slide_images": report.get("slide_images", {}),
                    "section_images": report.get("section_images", {}),
                    **p,
                }
    raise HTTPException(status_code=404, detail="제품을 찾을 수 없습니다.")


# =========================================================
# 4. 프로젝트별 보기
# =========================================================
APP_VERSION_FILE = BASE_DIR / "app_version.json" if "BASE_DIR" in globals() else Path("app_version.json")
APP_APK_FILE = Path("app_release.apk")


@app.get("/app/version")
def get_app_version():
    """앱 시작 시 호출. 최신 버전 정보 반환."""
    try:
        with open("app_version.json", "r", encoding="utf-8") as f:
            data = json.load(f)
    except Exception:
        data = {
            "latest_version": "1.0.0",
            "latest_version_code": 1,
            "download_url": "/app/download",
            "release_notes": "",
            "force_update": False,
        }
    print(f"[APP_VERSION] 반환: {data}")
    return data


@app.get("/app/download")
def download_app_apk():
    """최신 APK 다운로드."""
    from fastapi.responses import FileResponse
    apk_path = Path("app_release.apk")
    if not apk_path.exists():
        raise HTTPException(status_code=404, detail="APK not uploaded yet")
    return FileResponse(
        path=str(apk_path),
        media_type="application/vnd.android.package-archive",
        filename="app_release.apk",
    )


NOTES_FILE = DATA_DIR / "notes.json"


def _load_notes() -> dict:
    try:
        with open(NOTES_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {"version": 1, "updated_at": None, "notes": {}}


def _save_notes(data: dict) -> None:
    from datetime import datetime
    data["updated_at"] = datetime.now().isoformat()
    with open(NOTES_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


# ─── 프로젝트 모델 (양산/개발 구조) 헬퍼 ───
MODELS_FILE = DATA_DIR / "models.json"


def _load_models() -> dict:
    try:
        with open(MODELS_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {"version": 1, "updated_at": None, "projects": {}}


def _save_models(data: dict) -> None:
    from datetime import datetime
    data["updated_at"] = datetime.now().isoformat()
    with open(MODELS_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def _normalize_model(raw: dict, existing_ids: set) -> dict | None:
    """dict → 정규화된 모델 dict. 검증 실패 시 None 반환"""
    if not isinstance(raw, dict):
        return None
    mid = str(raw.get("id") or "").strip()
    name = str(raw.get("name") or "").strip()
    if not name:
        return None
    if not mid:
        mid = (name.lower().replace(" ", "-").replace("/", "-")) or f"model-{int(__import__('time').time()*1000)}"
    if mid in existing_ids:
        return None
    group = str(raw.get("group") or "양산").strip()
    if group not in ("양산", "개발"):
        group = "양산"
    try:
        progress = int(raw.get("progress") or 0)
    except (ValueError, TypeError):
        progress = 0
    progress = max(0, min(100, progress))
    status = str(raw.get("status") or "정상").strip()
    if status not in ("정상", "주의", "지연"):
        status = "정상"
    # 판가/재료비 (선택)
    price = 0
    material_cost = 0
    try:
        price = int(raw.get("price") or 0)
    except (ValueError, TypeError):
        pass
    try:
        material_cost = int(raw.get("material_cost") or 0)
    except (ValueError, TypeError):
        pass
    price = max(0, price)
    material_cost = max(0, material_cost)

    return {
        "id": mid,
        "name": name,
        "group": group,
        "progress": progress,
        "status": status,
        "price": price,
        "material_cost": material_cost,
    }


def _get_project_models(project_key: str) -> list:
    """프로젝트 키로 모델 목록 조회 (별칭 매핑 적용)"""
    _alias = {
        "havaplate": "hrva_plate",
        "hrvaplate": "hrva_plate",
        "hrva-plate": "hrva_plate",
    }
    _key = _alias.get(project_key.strip().lower(), project_key.strip())
    data = _load_models()
    proj = (data.get("projects") or {}).get(_key) or {}
    return proj.get("models") or []


# ─── 노트 첨부 자료 (표 JSON / 사진) 헬퍼 ───
import uuid as _uuid_mod
from datetime import datetime as _dt_mod

def _new_asset_id(division_id: str) -> str:
    """사업부별 자산 ID 생성: YYYY-MM-DD_<uuid8>"""
    date_str = _dt_mod.now().strftime("%Y-%m-%d")
    uid = _uuid_mod.uuid4().hex[:8]
    return f"{date_str}_{uid}"

def _table_path(division_id: str, asset_id: str) -> Path:
    div_dir = NOTE_TABLES_DIR / division_id
    div_dir.mkdir(parents=True, exist_ok=True)
    return div_dir / f"{asset_id}.json"

def _photo_path(division_id: str, asset_id: str, ext: str = "jpg") -> Path:
    div_dir = NOTE_PHOTOS_DIR / division_id
    div_dir.mkdir(parents=True, exist_ok=True)
    return div_dir / f"{asset_id}.{ext}"

def _save_note_table(division_id: str, table_data: dict) -> str:
    """표 JSON 저장 → asset_id 반환 (예: 'semiconductor/2026-06-11_a1b2c3d4')"""
    asset_id = _new_asset_id(division_id)
    p = _table_path(division_id, asset_id)
    p.write_text(json.dumps(table_data, ensure_ascii=False, indent=2), encoding="utf-8")
    return f"{division_id}/{asset_id}"

def _save_note_photo(division_id: str, image_bytes: bytes, ext: str = "png") -> str:
    """이미지 bytes → note_photos/에 저장 → photo_ref 반환.
    반환 예: 'semiconductor/2026-07-30_a1b2c3d4.png'
    """
    if not image_bytes:
        raise ValueError("empty image bytes")
    ext = (ext or "png").lower().lstrip(".")
    asset_id = _new_asset_id(division_id)
    dst = _photo_path(division_id, asset_id, ext)
    dst.write_bytes(image_bytes)
    return f"{division_id}/{asset_id}.{ext}"


def _xlsx_file_to_png_bytes(xlsx_path: Path) -> bytes:
    """xlsx → PNG bytes (LibreOffice + pdf2image로 색상/스타일 완벽 반영).
    
    옛 openpyxl+Pillow 직접 렌더링 방식은 표 서식/조건부 서식/테마 색상을
    지원하지 못했음. LibreOffice로 pdf 변환 후 pdf2image로 png 렌더링하면
    원본 엑셀과 시각적으로 100% 동일.
    """
    import subprocess
    import tempfile
    from io import BytesIO
    from pdf2image import convert_from_path
    
    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        
        # 1) LibreOffice: xlsx → pdf
        result = subprocess.run(
            ['libreoffice', '--headless', '--convert-to',
             'pdf:calc_pdf_Export:{"SinglePageSheets":{"type":"boolean","value":"true"}}',
             '--outdir', str(tmp_dir), str(xlsx_path)],
            capture_output=True, timeout=120
        )
        if result.returncode != 0:
            raise RuntimeError(
                f"LibreOffice 변환 실패 (code={result.returncode}): "
                f"{result.stderr.decode('utf-8', errors='ignore')[:500]}"
            )
        
        pdf_files = list(tmp_dir.glob("*.pdf"))
        if not pdf_files:
            raise RuntimeError("LibreOffice가 PDF를 생성하지 못함")
        pdf_path = pdf_files[0]
        
        # 2) pdf → png (첫 페이지, 고해상도)
        images = convert_from_path(str(pdf_path), dpi=150, first_page=1, last_page=1)
        if not images:
            raise RuntimeError("pdf2image로 PNG 변환 실패")
        
        # 3) PNG bytes
        buf = BytesIO()
        images[0].save(buf, 'PNG', optimize=True)
        return buf.getvalue()


def _derive_division_id_from_report(it: dict) -> str:
    """report dict → division_id 유추.
    우선순위: it.division_id → products[i].name derive → 'semiconductor' 폴백.
    _sync_report_to_notes의 로직과 동일하게 유지.
    """
    if not isinstance(it, dict):
        return "semiconductor"
    division_id = (it.get("division_id") or "").strip()
    if division_id:
        return division_id
    products = it.get("products") or []
    for prod in products:
        if isinstance(prod, dict) and prod.get("name"):
            try:
                division_id = _cl.derive_division_from_project(prod["name"]) or ""
            except Exception:
                pass
            if division_id:
                return division_id
    return "semiconductor"


def _load_note_table(table_ref: str):
    """table_ref = 'semiconductor/2026-06-11_a1b2c3d4' → 표 JSON 반환"""
    if not table_ref or "/" not in table_ref:
        return None
    division_id, asset_id = table_ref.split("/", 1)
    p = _table_path(division_id, asset_id)
    if not p.exists():
        return None
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return None

def _delete_note_table(table_ref: str) -> bool:
    if not table_ref or "/" not in table_ref:
        return False
    division_id, asset_id = table_ref.split("/", 1)
    p = _table_path(division_id, asset_id)
    if p.exists():
        p.unlink()
        return True
    return False

def _delete_note_photo(photo_ref: str) -> bool:
    """photo_ref = 'semiconductor/2026-06-11_a1b2c3d4.jpg'"""
    if not photo_ref or "/" not in photo_ref:
        return False
    p = NOTE_PHOTOS_DIR / photo_ref
    if p.exists():
        p.unlink()
        return True
    return False


_NOTE_AI_SYSTEM_PROMPT = """당신은 회사 임원에게 보고되는 주간 보고서를 구조화하는 도우미입니다.

입력으로 들어오는 자유 형식 텍스트를 다음 JSON 스키마로 변환하세요:

{
  "cards": [
    {
      "title": "프로젝트 이름",
      "sections": [
        {
          "title": "섹션 이름",
          "items": [
            {"type": "bullet | highlight | sub | group_note", "text": "항목 내용", "group_id": "g1 (선택)"}
          ]
        }
      ]
    }
  ]
}

규칙:
1. <텍스트>로 감싼 것은 새 카드(프로젝트)
2. 줄 시작에 "1.", "2.", "3." 같은 숫자+점으로 시작하는 줄은 무조건 새 섹션(section)이다. 절대 item으로 만들지 마라.
   - 예: "1. 양산 14종" → section title="양산 14종"
   - 예: "2. 개발 18종" → section title="개발 18종"
   - 예: "3. EMA 33종 (총 33종 중 승인 7종, 개발 26종)" → section title="EMA 33종 (총 33종 중 승인 7종, 개발 26종)"
3. "1)", "2)" 같은 괄호 번호나 들여쓰기는 일반 item(bullet) — section 만들지 마라
4. *로 시작하거나 "현황:" 표시는 type: highlight
4-1. 색상 태그 처리 — 본문에 "[빨간색]" / "[파란색]" / "[주황색]" / "[노란색]" 이 포함된 항목은:
     - 태그 문자열만 제거하고 본문은 그대로 유지
     - 결과 JSON 아이템에 "color": "red" / "blue" / "orange" 필드를 추가
     - 색상 태그는 group_id 부여의 신호가 아니다 (색깔이 같다고 묶지 말 것)
     - 예: "[파란색] W23 출하 24EA" → {"type":"bullet","text":"W23 출하 24EA","color":"blue"}
5. "-", "▸"로 시작하는 줄은 type: bullet
6. "→", "=>", "↳" 또는 명백한 들여쓰기는 type: sub (상위 항목의 부가 설명)
7. "※" 로 시작하는 줄은 type: highlight (참고/주의)
   - "※" 기호는 절대 제거하지 말고 본문 첫 글자로 그대로 유지한다 (예: "※에테르GDX 자재 ...")
   - "※" 로 시작하는 모든 줄은 반드시 출력 JSON 에 포함시켜야 한다 (절대 누락 금지, 요약/병합 금지)
   - 카드 마지막에 있는 "※" 줄들은 가장 가까운(직전) 섹션의 items 끝에 그대로 추가한다
   - "※" 줄을 별도 "기타 특이사항" 섹션으로 옮기지 말고, 원문이 속한 섹션 안에 그대로 둔다
   - "※" 줄에는 group_id 를 절대 부여하지 않는다
8. 카드 제목·섹션 제목에서 선행 번호와 기호는 제거
9. item 텍스트에서 선행 기호(-, *, ▸)는 제거. 본문 안의 →, ↳는 유지
10. "📷 파일명 @@photo_ref=..." 줄은 절대 수정/요약/제거하지 말고 그 자리에 그대로 한 줄 유지
11. 의미는 절대 바꾸지 말고, 임의로 항목을 합치거나 추가하지 마라
12. 빈 항목·중복 항목은 만들지 마라

★ group_id 부여 조건 (중요 — 함부로 묶지 말 것):

다음 신호 중 하나라도 명시적으로 존재할 때만 group_id를 부여한다:
  (a) 줄 끝에 "}" 기호가 있는 경우
  (b) "←", "<—", "<-" 화살표가 줄 시작이나 끝에 있는 경우
  (c) 텍스트에 "포괄적으로 들어가야", "공통으로 적용", "모두 해당" 같은 명시적 표현이 있는 경우

위 신호가 전혀 없으면 group_id를 절대 부여하지 마라. (일반 bullet들을 마음대로 묶지 말 것)

★ 사진 placeholder 규칙 (매우 중요, 절대 어기지 말 것):
  - 본문에 "[PHOTO_KEEP_0]", "[PHOTO_KEEP_1]" 같은 토큰이 포함된 줄이 있으면:
    1) 그 줄 전체를 반드시 별도 item 으로 만들 것 (다른 줄과 병합 금지, 삭제 금지, 요약 금지)
    2) item 의 text 는 원문 그대로 ("📷 파일명 [PHOTO_KEEP_0]") 유지
    3) item 의 type 은 "bullet" 로 둔다
    4) 원문에서 그 줄이 있던 위치 그대로 보존 — 앞뒤 내용 순서 절대 변경 금지, 카드 끝으로 옮기지 말 것
    5) placeholder 줄을 group_id 의 신호로 사용하지 말 것

★ group_id 를 부여하면 안 되는 경우 (오해 방지):
  - 색상 태그([빨간색]/[파란색]/[주황색]/[노란색])만으로는 group_id 부여 금지
  - 같은 섹션에 속한다는 이유로 묶지 말 것
  - 의미가 비슷하다는 이유로 묶지 말 것
  - "※" 로 시작하는 줄은 묶지 말 것

group_id를 부여하는 경우 처리 방법:
  1) 신호가 나타난 줄 바로 위쪽의 연관된 bullet 2-3개에 동일한 group_id ("g1", "g2", ...) 부여
  2) 공통 메모는 별도 항목으로 추가: {"type":"group_note","text":"<메모 본문>","group_id":"g1"}
  3) group_note의 text 에서는 "}", "←" 같은 메타 기호 제거하고 본문만 남긴다

[예시 1 — 번호 소제목]
입력:
<파워박스>
1. 양산 14종
- 에테르 GDX : 연간 1,461대

2. 개발 18종
- 화성 세이버 5종 7월말 승인 타겟

출력:
{
  "cards": [{
    "title": "파워박스",
    "sections": [
      {"title": "양산 14종", "items": [
        {"type":"bullet","text":"에테르 GDX : 연간 1,461대"}
      ]},
      {"title": "개발 18종", "items": [
        {"type":"bullet","text":"화성 세이버 5종 7월말 승인 타겟"}
      ]}
    ]
  }]
}

[예시 2 — } 마커가 있을 때만 group_id]
입력:
<하바플레이트>
*현황: 총141개 모델 중 양산진행 14종
- 신규 장비 40대
- 플라스틱 전용 장비 6대 } 6월말 고객사 승인 예정

출력:
{
  "cards": [{
    "title": "하바플레이트",
    "sections": [
      {"title": "현황", "items": [
        {"type":"highlight","text":"총141개 모델 중 양산진행 14종"},
        {"type":"bullet","text":"신규 장비 40대","group_id":"g1"},
        {"type":"bullet","text":"플라스틱 전용 장비 6대","group_id":"g1"},
        {"type":"group_note","text":"6월말 고객사 승인 예정","group_id":"g1"}
      ]}
    ]
  }]
}

응답은 반드시 위 JSON 형식 한 객체만 출력. 마크다운, 코드블록, 설명 없이 JSON만 출력하세요."""


@app.post("/admin/notes/from_pptx")
def admin_notes_from_pptx(payload: dict, _admin: int = Depends(get_admin_session)):
    """업로드된 PPT를 주간 보고 카드 구조로 자동 변환.
    - 슬라이드 텍스트는 AI로 카드 구조화
    - 표는 자동으로 note_assets/tables/ 에 저장 + table_ref 연결
    """
    doc_id = (payload or {}).get("doc_id", "").strip()
    division_id = (payload or {}).get("division_id", "").strip()
    if not doc_id:
        raise HTTPException(status_code=400, detail="doc_id가 필요합니다")
    if not division_id:
        raise HTTPException(status_code=400, detail="division_id가 필요합니다")

    pptx_path = UPLOAD_DIR / f"{doc_id}.pptx"
    if not pptx_path.exists():
        raise HTTPException(status_code=404, detail=f"PPT 파일을 찾을 수 없음: {doc_id}")

    # 1) 텍스트 + 표 분리 추출
    try:
        slide_texts, slide_tables = _extract_pptx_text_and_tables(pptx_path)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PPT 추출 실패: {e}")

    # 2) AI로 카드 구조 변환 (표 데이터는 텍스트로 함께 전달, AI가 표 참조 자리에 마커 삽입)
    joined_text = "\n\n".join([f"=== Slide {i+1} ===\n{t}" for i, t in enumerate(slide_texts) if t.strip()])
    if not joined_text.strip():
        return {"cards": []}

    user_message = (
        f"[PPT 텍스트 - 슬라이드별]\n{joined_text}\n\n"
        "위 PPT 텍스트를 주간 보고 카드 구조 JSON으로 변환해 주세요. "
        "한 슬라이드는 보통 한 프로젝트입니다. 슬라이드 제목이 프로젝트명입니다."
    )

    try:
        from openai import OpenAI
        import os
        oai = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))
        resp = oai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": _NOTE_AI_SYSTEM_PROMPT},
                {"role": "user", "content": user_message},
            ],
            response_format={"type": "json_object"},
            temperature=0.1,
        )
        result = json.loads(resp.choices[0].message.content)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI 변환 실패: {e}")

    # 3) 표가 있는 슬라이드 → 자동 저장 + 첫 번째 섹션에 table_ref 부착
    cards = result.get("cards", []) or []
    saved_tables = []
    for slide_idx, tables in enumerate(slide_tables):
        if not tables:
            continue
        if slide_idx >= len(cards):
            continue
        card = cards[slide_idx]
        sections = card.get("sections", []) or []
        if not sections:
            continue
        # 첫 번째 표만 첫 섹션에 자동 첨부 (나머지는 별도 섹션으로)
        for t_idx, tbl in enumerate(tables):
            table_obj = {
                "title": card.get("title", "") + (f" 표{t_idx+1}" if t_idx > 0 else ""),
                "headers": tbl["headers"],
                "rows": tbl["rows"],
            }
            try:
                table_ref = _save_note_table(division_id, None, table_obj)
                saved_tables.append(table_ref)
                target_section = sections[0] if t_idx == 0 else sections[min(t_idx, len(sections)-1)]
                target_section["table_ref"] = table_ref
            except Exception as e:
                print(f"표 저장 실패 slide={slide_idx+1} t={t_idx}: {e}")

    return {"cards": cards, "saved_tables": saved_tables}


@app.get("/notes/by_project")
def notes_by_project(project_key: str = ""):
    """프로젝트 키로 가장 최근 노트 카드를 찾아 반환.
    매칭 우선순위:
      1) card.get("project_key") == project_key
      2) card.get("title") 가 projects.json 의 label/aliases 와 매칭
    표는 _load_note_table 으로 inline 포함.
    """
    project_key = (project_key or "").strip()
    if not project_key:
        raise HTTPException(status_code=400, detail="project_key 필요")

    try:
        proj = _cl.get_project(project_key)
    except Exception:
        proj = None
    project_label = (proj or {}).get("label", "")
    project_aliases = set((proj or {}).get("aliases", []) or [])
    project_division = (proj or {}).get("division_id", "")
    if project_label:
        project_aliases.add(project_label)
    project_aliases.add(project_key)

    data = _load_notes()
    notes_map = data.get("notes", {}) or {}

    # 사업부 우선 검색, 없으면 전체 사업부 스캔
    div_ids = [project_division] if project_division else list(notes_map.keys())
    if project_division and project_division not in notes_map:
        div_ids = list(notes_map.keys())

    best = None
    best_date = ""
    best_division = ""
    for did in div_ids:
        ndata = notes_map.get(did) or {}
        cards = ndata.get("cards", []) or []
        report_date = ndata.get("report_date", "") or ""
        for c in cards:
            ck = (c.get("project_key") or "").strip()
            ctitle = (c.get("title") or "").strip()
            matched = False
            if ck and ck == project_key:
                matched = True
            elif ctitle and (ctitle == project_label or ctitle in project_aliases):
                matched = True
            if not matched:
                continue
            # 가장 최근 날짜 선택
            if report_date > best_date:
                best = c
                best_date = report_date
                best_division = did

    if not best:
        return {"card": None}

    # table_ref inline 로드
    sections = best.get("sections", []) or []
    for sec in sections:
        # 섹션 레벨 table_ref (기존 호환)
        sec_ref = sec.get("table_ref")
        if sec_ref:
            try:
                sec["table_data"] = _load_note_table(sec_ref)
            except Exception:
                sec["table_data"] = None
        # 아이템 흐름 안의 table/photo
        items = sec.get("items", []) or []
        for it in items:
            if not isinstance(it, dict):
                continue
            t_ref = it.get("table_ref")
            if t_ref:
                try:
                    it["table_data"] = _load_note_table(t_ref)
                except Exception:
                    it["table_data"] = None

    return {
        "card": best,
        "report_date": best_date,
        "division_id": best_division,
    }






# ─── 노트 엑셀(표) 마커 전/후처리 (ai_parse 용) ───
_EXCEL_MARKER_RE = re.compile(
    "\U0001F4CA\s*([^\n\r]+?)\s*@@table_ref=([^\s\n\r]+)",
    re.IGNORECASE,
)

def _extract_excel_markers(text: str):
    """텍스트에서 엑셀 마커 추출.
    반환: (clean_text, [{filename, table_ref, anchor}])
    """
    lines = text.split("\n")
    out_lines = []
    tables = []
    last_nonempty = ""
    for ln in lines:
        m = _EXCEL_MARKER_RE.search(ln)
        if m:
            fname = (m.group(1) or "").strip()
            ref = (m.group(2) or "").strip()
            if fname and ref:
                tables.append({
                    "filename": fname,
                    "table_ref": ref,
                    "anchor": last_nonempty,
                })
            continue
        out_lines.append(ln)
        if ln.strip():
            last_nonempty = ln.strip()
    return ("\n".join(out_lines), tables)


def _inject_tables_into_cards(cards, tables):
    """정리된 카드 구조에 type:'table' item 을 앵커 기준으로 주입."""
    if not tables or not isinstance(cards, list):
        return cards

    def _norm(s):
        return re.sub(r"\s+", " ", str(s or "").strip().lower())

    used = set()

    for tb in tables:
        anchor_n = _norm(tb.get("anchor"))
        ref = tb.get("table_ref")
        fname = tb.get("filename") or ""
        if not ref or ref in used:
            continue

        # table_data 도 함께 로드해서 inline 포함
        table_data = _load_note_table(ref) if ref else None

        injected = False
        if anchor_n:
            for card in cards:
                if injected:
                    break
                for sec in (card.get("sections") or []):
                    if injected:
                        break
                    items = sec.get("items") or []
                    for idx, it in enumerate(items):
                        it_text = _norm(it.get("text") if isinstance(it, dict) else it)
                        if not it_text:
                            continue
                        if anchor_n == it_text or anchor_n in it_text or it_text in anchor_n:
                            new_item = {
                                "type": "table",
                                "text": fname,
                                "table_ref": ref,
                            }
                            if table_data:
                                new_item["table_data"] = table_data
                            items.insert(idx + 1, new_item)
                            sec["items"] = items
                            injected = True
                            used.add(ref)
                            break

        if not injected and cards:
            card0 = cards[0]
            secs = card0.setdefault("sections", [])
            if not secs:
                secs.append({"title": "첨부", "items": []})
            last_sec = secs[-1]
            new_item = {
                "type": "table",
                "text": fname,
                "table_ref": ref,
            }
            if table_data:
                new_item["table_data"] = table_data
            last_sec.setdefault("items", []).append(new_item)
            used.add(ref)

    return cards


# ─── 노트 사진 마커 전/후처리 (ai_parse 용) ───
import re as _re_photo
_PHOTO_MARKER_RE = _re_photo.compile(
    r"\u{1F4F7}\s*([^\n\r]+?)\s*@@photo_ref=([^\s\n\r]+)".replace(r"\u{1F4F7}", "\U0001F4F7"),
    _re_photo.IGNORECASE,
)

def _extract_photo_markers(text: str):
    """텍스트에서 사진 마커 추출 (next_anchor 기반 위치 보존).
    반환: (clean_text, [{filename, photo_ref, anchor, next_anchor, placeholder}])
    anchor      = 마커 직전의 비어있지 않은 줄
    next_anchor = 마커 다음에 나오는 비어있지 않은 줄 (마커가 아닌 줄)
    """
    lines = text.split("\n")

    marker_positions = []
    for i, ln in enumerate(lines):
        m = _PHOTO_MARKER_RE.search(ln)
        if m:
            fname = (m.group(1) or "").strip()
            ref = (m.group(2) or "").strip()
            if fname and ref:
                marker_positions.append((i, fname, ref))

    photos = []
    out_lines = []
    last_nonempty = ""
    marker_idx_set = {pos[0] for pos in marker_positions}

    def find_next_anchor(start):
        for j in range(start + 1, len(lines)):
            if j in marker_idx_set:
                continue
            v = lines[j].strip()
            if v:
                return v
        return ""

    pos_map = {pos[0]: pos for pos in marker_positions}
    for i, ln in enumerate(lines):
        if i in pos_map:
            _, fname, ref = pos_map[i]
            idx = len(photos)
            placeholder = f"[PHOTO_KEEP_{idx}]"
            placeholder_line = f"\U0001F4F7 {fname} {placeholder}"
            nxt = find_next_anchor(i)
            photos.append({
                "filename": fname,
                "photo_ref": ref,
                "anchor": last_nonempty,
                "next_anchor": nxt,
                "placeholder": placeholder,
            })
            out_lines.append(placeholder_line)
            last_nonempty = placeholder_line
            continue
        out_lines.append(ln)
        if ln.strip():
            last_nonempty = ln.strip()
    return ("\n".join(out_lines), photos)


_DONE_KEYWORDS = ("완료", "완료됨", "출고완료", "승인완료", "끝")
_PENDING_KEYWORDS = ("예정", "예상", "타겟", "목표", "진행중", "준비중", "까지", "예약")


def _split_for_date_scope(txt: str):
    """문장을 ',', '/', '|', '·' 기준으로 조각내서 (시작idx, 끝idx, 조각텍스트) 리스트 반환."""
    chunks = []
    start = 0
    for m in re.finditer(r"[,/|·]", txt):
        end = m.start()
        chunks.append((start, end, txt[start:end]))
        start = m.end()
    chunks.append((start, len(txt), txt[start:]))
    return chunks


def _is_done_scope(txt: str, idx: int) -> bool:
    """idx 위치의 날짜가 '완료된 날짜'인지 판정.
    한국어 어순: 날짜 → 키워드 (예: "5/29 승인완료", "W25 출하예정")
    규칙:
      1) idx 뒤에서 가장 가까이 나타나는 키워드(완료/예정)를 찾는다.
         - 완료 키워드면: 완료 (True, 날짜 버림)
         - 예정 키워드면: 살림 (False)
      2) idx 뒤에 키워드 없으면 idx 앞 12자 윈도우에 완료 키워드 있으면 완료.
      3) 그 외엔 살림.
    """
    if not txt:
        return False

    # 1) idx 뒤에서 가장 가까운 완료/예정 키워드 위치
    nearest_pos = -1
    nearest_kind = None  # 'done' or 'pending'
    for k in _DONE_KEYWORDS:
        i = txt.find(k, idx)
        if i != -1 and (nearest_pos == -1 or i < nearest_pos):
            nearest_pos = i
            nearest_kind = "done"
    for k in _PENDING_KEYWORDS:
        i = txt.find(k, idx)
        if i != -1 and (nearest_pos == -1 or i < nearest_pos):
            nearest_pos = i
            nearest_kind = "pending"

    # 1-b) done 키워드 뒤에 곧바로 "예정" 이 나오면 pending 으로 재분류
    # 예: "완료 예정", "완료예정", "출하예정 완료" 등
    if nearest_kind == "done" and nearest_pos != -1:
        # done 키워드 뒤 6자 이내에 예정/예상 이 있는지
        window_end = min(len(txt), nearest_pos + 12)
        after = txt[nearest_pos:window_end]
        for pk in _PENDING_KEYWORDS:
            if pk in after:
                nearest_kind = "pending"
                break

    # 1-b) done 키워드 뒤에 곧바로 "예정" 이 나오면 pending 으로 재분류
    # 예: "완료 예정", "완료예정", "출하예정 완료" 등
    if nearest_kind == "done" and nearest_pos != -1:
        # done 키워드 뒤 6자 이내에 예정/예상 이 있는지
        window_end = min(len(txt), nearest_pos + 12)
        after = txt[nearest_pos:window_end]
        for pk in _PENDING_KEYWORDS:
            if pk in after:
                nearest_kind = "pending"
                break

    # 같은 줄에서 너무 멀면 무시 (다른 절일 가능성)
    if nearest_pos != -1 and (nearest_pos - idx) <= 25:
        return nearest_kind == "done"

    # 2) idx 앞 12자 윈도우에 완료 키워드만 있으면 완료
    lo = max(0, idx - 12)
    front = txt[lo:idx]
    if any(k in front for k in _DONE_KEYWORDS) and not any(k in front for k in _PENDING_KEYWORDS):
        return True

    return False


def _extract_due_date_from_text(txt: str) -> str:
    """본문 텍스트에서 due_date(YYYY-MM-DD) 추출. 못 찾으면 빈 문자열.
    
    화살표 인식형 파서 (마커 없이 자동 판단):
      1) 텍스트를 화살표(-->, ->, →, ➜, ⇒, ~>)로 split → 단계별로 나눔.
      2) 마지막 세그먼트만 대상으로 날짜 후보 수집 (앞 단계는 완료로 간주).
      3) 마지막 세그먼트에 날짜 없으면 전체 텍스트에서 fallback.
      4) 후보 중 최적 선택:
         - 오늘 이후(포함) 날짜 있으면 → 그중 가장 가까운 미래
         - 없으면 → 지난 날짜 중 가장 최근 (지연 D+n 표시용)
    
    지원 패턴:
      - [YYYY-MM-DD] / YYYY-MM-DD
      - YY.M.D / YY-M-D
      - M월말 (해당 월 마지막 평일)
      - M월 D일
      - (M/D), M/D
      - (M/D-M/D), (M/D~M/D)  → 종료일 채택
      - W23~W53 (해당 ISO 주차 금요일)
    """
    from datetime import date, timedelta
    import calendar
    if not txt:
        return ""
    YEAR = 2026
    today = date.today()

    # ─── Step 1: 화살표로 split ───
    arrow_pattern = r"-->|->|→|➜|⇒|~>|»"
    segments = re.split(arrow_pattern, txt)
    last_seg = segments[-1] if segments else txt

    def _collect_dates(s: str):
        """문자열 s 에서 모든 날짜 후보를 date 리스트로 반환. done_scope 검사 없음(마커 없이 자동판단)."""
        out = []

        # 범위 (M/D-M/D) / (M/D~M/D) → 종료일
        for m in re.finditer(r"\(?\s*(\d{1,2})/(\d{1,2})\s*[-~]\s*(\d{1,2})/(\d{1,2})\s*\)?", s):
            try:
                mo = int(m.group(3)); dd = int(m.group(4))
                if 1 <= mo <= 12 and 1 <= dd <= 31:
                    out.append(date(YEAR, mo, dd))
            except Exception:
                pass

        # YYYY-MM-DD
        for m in re.finditer(r"(20\d{2})[-./](\d{1,2})[-./](\d{1,2})", s):
            try:
                out.append(date(int(m.group(1)), int(m.group(2)), int(m.group(3))))
            except Exception:
                pass

        # YY.M.D / YY-M-D  (2자리 연도)
        for m in re.finditer(r"(\d{2})[-./](\d{1,2})[-./](\d{1,2})", s):
            try:
                out.append(date(2000 + int(m.group(1)), int(m.group(2)), int(m.group(3))))
            except Exception:
                pass

        # X월말
        for m in re.finditer(r"(\d{1,2})\s*월\s*말", s):
            try:
                mo = int(m.group(1))
                last_day = calendar.monthrange(YEAR, mo)[1]
                d = date(YEAR, mo, last_day)
                if d.weekday() == 6:
                    d = d - timedelta(days=2)
                out.append(d)
            except Exception:
                pass

        # M월 D일
        for m in re.finditer(r"(\d{1,2})\s*월\s*(\d{1,2})\s*일", s):
            try:
                out.append(date(YEAR, int(m.group(1)), int(m.group(2))))
            except Exception:
                pass

        # (M/D) 또는 M/D  (분수 제외 — 경계 있을 때만)
        for m in re.finditer(r"(?:^|[\s\(\[~])(\d{1,2})/(\d{1,2})(?=[\s\)\]\,\.~]|$)", s):
            try:
                mo = int(m.group(1)); dd = int(m.group(2))
                if 1 <= mo <= 12 and 1 <= dd <= 31:
                    out.append(date(YEAR, mo, dd))
            except Exception:
                pass

        # W23 ~ W53
        for m in re.finditer(r"W\s*(\d{1,2})", s, flags=re.IGNORECASE):
            try:
                wk = int(m.group(1))
                if 1 <= wk <= 53:
                    out.append(date.fromisocalendar(YEAR, wk, 5))
            except Exception:
                pass

        return out

    # ─── Step 2: 마지막 세그먼트에서 먼저 시도 ───
    candidates = _collect_dates(last_seg)

    # ─── Step 3: 없으면 전체에서 fallback ───
    if not candidates and len(segments) == 1:
        # split 안 됐으면 이미 전체 = last_seg, 중복 fallback 불필요
        pass
    elif not candidates:
        candidates = _collect_dates(txt)

    if not candidates:
        return ""

    # ─── Step 4: 최적 선택 ───
    future = [d for d in candidates if d >= today]
    if future:
        return min(future).isoformat()  # 가장 임박한 미래
    # 전부 과거 → 가장 최근 과거 (지연 표시용)
    return max(candidates).isoformat()



def _normalize_note_item(it: dict, section_title: str = "") -> dict:
    """AI 결과 아이템 후처리: 색상 태그, ※ 항목, 화살표 중복 제거.
    
    ★ 2026-07-27~ 스키마 확장:
      - item_id: md5(section_title + normalized_text)[:8] 자동 생성
      - due_date_auto: 파서가 뽑은 자동값 (매번 재계산)
      - due_date_override: 사용자 지정값 (있으면 유지)
      - due_date: override or auto (앱 호환용 최종값)
    """
    if not isinstance(it, dict):
        return it

    txt = (it.get("text") or "").strip()
    typ = (it.get("type") or "bullet").strip().lower()
    color = (it.get("color") or "").strip().lower()

    # [빨간색]/[파란색]/[주황색]/[노란색] -> color 필드로 승격, 본문에서는 제거
    if "[빨간색]" in txt:
        color = "red"
        txt = txt.replace("[빨간색]", "").strip()
    if "[파란색]" in txt:
        color = "blue"
        txt = txt.replace("[파란색]", "").strip()
    if "[주황색]" in txt:
        color = "orange"
        txt = txt.replace("[주황색]", "").strip()
    if "[노란색]" in txt:
        color = "orange"
        txt = txt.replace("[노란색]", "").strip()

    # ※ 로 시작하는 줄은 절대 group_note/특이사항으로 빼지 말고 일반 bullet 유지
    if txt.startswith("※"):
        if typ == "group_note":
            typ = "bullet"
        it.pop("group_id", None)

    # sub: 선행 화살표가 여러 개거나 중복이면 1개로 정규화
    if typ == "sub":
        txt = re.sub(r"^(?:\s*(?:→|↳|=>)\s*)+", "→ ", txt).strip()

    # bullet/highlight 도 "→ → ..." 처럼 화살표 중복 시 1개로 축약
    txt = re.sub(r"^(→\s*){2,}", "→ ", txt)

    # ─── item_id 생성/갱신 (section_title + normalized_text 기반) ───
    normalized = re.sub(r"\s+", " ", txt).strip()
    if normalized:
        new_id = hashlib.md5(
            f"{section_title}||{normalized}".encode("utf-8")
        ).hexdigest()[:8]
        it["item_id"] = new_id

    # ─── due_date_auto 재계산 (매번, 텍스트는 원문 유지) ───
    _, auto_iso = _extract_due_date(txt)
    # 사용자가 자동 날짜를 숨긴 경우(auto_due_hidden=True)는 재계산돼도 유지
    _prev_auto = (it.get("due_date_auto") or "").strip()
    _new_auto = auto_iso or ""
    it["due_date_auto"] = _new_auto

    # ─── due_date_override 처리 ───
    override_raw = (it.get("due_date_override") or "").strip()
    override = override_raw if override_raw else ""

    # ─── 최종 due_date = override or auto (앱 호환) ───
    effective = override or (auto_iso or "")
    if effective:
        it["due_date"] = effective
    else:
        # 아무 값 없으면 due_date 필드 제거 (하위 호환)
        it.pop("due_date", None)

    # override 명시 저장 (빈 값이면 삭제)
    if override:
        it["due_date_override"] = override
    else:
        it.pop("due_date_override", None)

    it["type"] = typ
    it["text"] = txt
    if color:
        it["color"] = color

    return it


def _normalize_note_cards(parsed: dict) -> dict:
    """parsed["cards"] 전체에 _normalize_note_item 적용 + 기타 특이사항 섹션 병합."""
    if not isinstance(parsed, dict):
        return parsed
    cards = parsed.get("cards") or []
    _ETC_TITLES = {"기타 특이사항", "기타특이사항", "특이사항", "기타"}
    for card in cards:
        if not isinstance(card, dict):
            continue

        # 1) 각 섹션 아이템 정규화
        for sec in (card.get("sections") or []):
            if not isinstance(sec, dict):
                continue
            new_items = []
            for it in (sec.get("items") or []):
                if isinstance(it, dict):
                    new_items.append(_normalize_note_item(it, sec.get('title', '')))
                else:
                    new_items.append(it)
            sec["items"] = new_items

        # 2) "기타 특이사항" 류 섹션은 직전 섹션에 흡수 후 삭제
        sections = card.get("sections") or []
        merged_sections = []
        for sec in sections:
            if not isinstance(sec, dict):
                merged_sections.append(sec)
                continue
            title = (sec.get("title") or "").strip()
            if title in _ETC_TITLES and merged_sections:
                # 직전 섹션 마지막에 합치기
                prev = merged_sections[-1]
                if isinstance(prev, dict):
                    prev_items = prev.get("items") or []
                    prev_items.extend(sec.get("items") or [])
                    prev["items"] = prev_items
                    continue
            merged_sections.append(sec)
        card["sections"] = merged_sections

    return parsed


def _inject_photos_into_cards(cards, photos):
    """정리된 카드 구조에 type:'photo' item 을 주입.
    1순위: placeholder 텍스트 ('__PHOTO_PLACEHOLDER_N__') 가진 item 을 photo item 으로 직접 치환 → 원문 위치 보존
    2순위(fallback): anchor 텍스트 매칭 (구버전 호환)
    """
    if not photos or not isinstance(cards, list):
        return cards

    def _norm(s):
        v = re.sub(r"^\s*\d+[.)\]]?\s*", "", str(s or "").strip())
        return re.sub(r"\s+", " ", v.lower())

    used = set()

    # ── 1순위: placeholder 직접 치환 ──
    for ph in photos:
        placeholder = ph.get("placeholder") or ""
        ref = ph.get("photo_ref")
        fname = ph.get("filename") or ""
        if not ref or not placeholder or ref in used:
            continue

        replaced = False
        for card in cards:
            if replaced:
                break
            for sec in (card.get("sections") or []):
                if replaced:
                    break
                items = sec.get("items") or []
                for idx, it in enumerate(items):
                    if not isinstance(it, dict):
                        continue
                    txt = str(it.get("text") or "")
                    if placeholder in txt:
                        # 해당 item 을 photo item 으로 치환
                        items[idx] = {
                            "type": "photo",
                            "text": fname,
                            "photo_ref": ref,
                        }
                        sec["items"] = items
                        replaced = True
                        used.add(ref)
                        break

        if replaced:
            continue

        # ── 1.5순위: next_anchor 매칭 (마커 다음 줄 앞에 삽입) ──
        next_anchor_n = _norm(ph.get("next_anchor"))
        if next_anchor_n:
            inserted_by_next = False
            # (a) item 텍스트 매칭 → 그 앞에 삽입
            for card in cards:
                if inserted_by_next:
                    break
                for sec in (card.get("sections") or []):
                    if inserted_by_next:
                        break
                    items = sec.get("items") or []
                    for idx, it in enumerate(items):
                        it_text = _norm(it.get("text") if isinstance(it, dict) else it)
                        if not it_text:
                            continue
                        if next_anchor_n == it_text or next_anchor_n in it_text or it_text in next_anchor_n:
                            items.insert(idx, {
                                "type": "photo",
                                "text": fname,
                                "photo_ref": ref,
                            })
                            sec["items"] = items
                            inserted_by_next = True
                            used.add(ref)
                            break
            # (b) section title 매칭 → 그 섹션의 첫 item 으로 삽입
            if not inserted_by_next:
                for card in cards:
                    if inserted_by_next:
                        break
                    for sec in (card.get("sections") or []):
                        sec_title = _norm(sec.get("title"))
                        if not sec_title:
                            continue
                        if next_anchor_n == sec_title or next_anchor_n in sec_title or sec_title in next_anchor_n:
                            items = sec.get("items") or []
                            items.insert(0, {
                                "type": "photo",
                                "text": fname,
                                "photo_ref": ref,
                            })
                            sec["items"] = items
                            inserted_by_next = True
                            used.add(ref)
                            break
            if inserted_by_next:
                continue

        # ── 2순위(fallback): anchor 텍스트 매칭 ──
        anchor_n = _norm(ph.get("anchor"))
        if anchor_n and anchor_n.startswith("__photo_placeholder_"):
            anchor_n = ""  # placeholder 자체는 anchor 로 쓰지 않음
        injected = False
        if anchor_n:
            # item 텍스트 매칭
            for card in cards:
                if injected:
                    break
                for sec in (card.get("sections") or []):
                    if injected:
                        break
                    items = sec.get("items") or []
                    for idx, it in enumerate(items):
                        it_text = _norm(it.get("text") if isinstance(it, dict) else it)
                        if not it_text:
                            continue
                        if anchor_n == it_text or anchor_n in it_text or it_text in anchor_n:
                            insert_at = idx + 1
                            matched_gid = (it.get("group_id") if isinstance(it, dict) else "") or ""
                            if matched_gid:
                                k = idx + 1
                                while k < len(items):
                                    nxt = items[k]
                                    if not isinstance(nxt, dict):
                                        break
                                    if (nxt.get("group_id") or "") != matched_gid:
                                        break
                                    insert_at = k + 1
                                    k += 1
                            else:
                                k = idx + 1
                                while k < len(items):
                                    nxt = items[k]
                                    if not isinstance(nxt, dict):
                                        break
                                    if (nxt.get("type") or "") != "group_note":
                                        break
                                    insert_at = k + 1
                                    k += 1
                            items.insert(insert_at, {
                                "type": "photo",
                                "text": fname,
                                "photo_ref": ref,
                            })
                            sec["items"] = items
                            injected = True
                            used.add(ref)
                            break

            # section title 매칭
            if not injected:
                for card in cards:
                    if injected:
                        break
                    for sec in (card.get("sections") or []):
                        sec_title = _norm(sec.get("title"))
                        if not sec_title:
                            continue
                        if anchor_n == sec_title or anchor_n in sec_title or sec_title in anchor_n:
                            items = sec.get("items") or []
                            items.insert(0, {
                                "type": "photo",
                                "text": fname,
                                "photo_ref": ref,
                            })
                            sec["items"] = items
                            injected = True
                            used.add(ref)
                            break

            # card title 매칭
            if not injected:
                for card in cards:
                    card_title = _norm(card.get("title"))
                    if not card_title:
                        continue
                    if anchor_n == card_title or anchor_n in card_title or card_title in anchor_n:
                        secs = card.setdefault("sections", [])
                        if not secs:
                            secs.append({"title": "첨부", "items": []})
                        secs[-1].setdefault("items", []).append({
                            "type": "photo",
                            "text": fname,
                            "photo_ref": ref,
                        })
                        injected = True
                        used.add(ref)
                        break

        # 앵커 못 찾으면 첫 카드 마지막 섹션 끝
        if not injected and cards:
            card0 = cards[0]
            secs = card0.setdefault("sections", [])
            if not secs:
                secs.append({"title": "첨부", "items": []})
            last_sec = secs[-1]
            last_sec.setdefault("items", []).append({
                "type": "photo",
                "text": fname,
                "photo_ref": ref,
            })
            used.add(ref)

    # ── 후처리: 사용되지 않은 placeholder item 제거 (AI 가 위치 보존 안 한 경우 대비) ──
    placeholder_re = re.compile(r"\[PHOTO_KEEP_\d+\]")
    for card in cards:
        for sec in (card.get("sections") or []):
            items = sec.get("items") or []
            sec["items"] = [
                it for it in items
                if not (isinstance(it, dict) and placeholder_re.match(str(it.get("text") or "").strip()))
            ]

    return cards


@app.post("/admin/notes/ai_parse")
def admin_notes_ai_parse(payload: dict, _admin: int = Depends(get_admin_session)):
    """자유 형식 텍스트를 AI로 구조화 JSON으로 변환 (증분 병합 지원)."""
    text = (payload or {}).get("text", "").strip()
    division_id = (payload or {}).get("division_id", "").strip()
    if not text:
        raise HTTPException(status_code=400, detail="text가 비어있습니다")

    # 사진 마커 추출 (AI 가 보지 않도록 제거 후, 응답에 주입)
    text, _extracted_photos = _extract_photo_markers(text)
    text, _extracted_tables = _extract_excel_markers(text)

    user_message = f"다음 텍스트를 정리:\n\n{text}"

    try:
        from openai import OpenAI
        import os
        client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": _NOTE_AI_SYSTEM_PROMPT},
                {"role": "user", "content": user_message},
            ],
            response_format={"type": "json_object"},
            temperature=0.1,
        )
        content = response.choices[0].message.content
        parsed = json.loads(content)
        if isinstance(parsed, dict) and isinstance(parsed.get("cards"), list):
            parsed["cards"] = _inject_photos_into_cards(parsed["cards"], _extracted_photos)
            parsed = _normalize_note_cards(parsed)
            # parsed["cards"] = _inject_tables_into_cards(parsed["cards"], _extracted_tables)  # 엑셀은 photo로 처리
        return parsed
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI 파싱 실패: {str(e)}")


@app.post("/admin/notes")
def admin_save_note(payload: dict, _admin: int = Depends(get_admin_session)):
    """사업부별 노트 저장."""
    division_id = (payload or {}).get("division_id", "").strip()
    report_date = (payload or {}).get("report_date", "").strip()
    raw_text = (payload or {}).get("raw_text", "").strip()
    cards = (payload or {}).get("cards", [])

    if not division_id:
        raise HTTPException(status_code=400, detail="division_id 필수")
    if not isinstance(cards, list):
        raise HTTPException(status_code=400, detail="cards는 배열이어야 합니다")

    # 저장 전 카드/아이템 정규화
    # - text 안의 due_date 패턴 추출
    # - item 구조 표준화
    # - 향후 AI 정리/수기 저장/직접 API 호출 모두 동일 규칙 적용
    try:
        normalized = _normalize_note_cards({"cards": cards}) or {}
        cards = normalized.get("cards") or cards
    except Exception as _e:
        print(f"admin_save_note normalize 실패: {_e}")

    from datetime import datetime
    data = _load_notes()
    notes_map = data.setdefault("notes", {})
    existing = notes_map.get(division_id) or {}
    existing_cards = existing.get("cards") or []

    # 카드 제목 기준으로 merge (같은 title은 새 내용으로 교체, 새 title은 추가)
    def _norm_title(t):
        return re.sub(r"\s+", "", (t or "").strip()).lower()

    by_title = {}
    order = []
    for c in existing_cards:
        if not isinstance(c, dict):
            continue
        tkey = _norm_title(c.get("title"))
        if not tkey or tkey in by_title:
            continue
        by_title[tkey] = c
        order.append(tkey)

    for c in cards:
        if not isinstance(c, dict):
            continue
        tkey = _norm_title(c.get("title"))
        if not tkey:
            continue
        if tkey in by_title:
            by_title[tkey] = c  # 교체
        else:
            by_title[tkey] = c
            order.append(tkey)

    merged_cards = [by_title[k] for k in order]

    notes_map[division_id] = {
        "report_date": report_date or existing.get("report_date", ""),
        "updated_at": datetime.now().isoformat(),
        "raw_text": raw_text if isinstance(raw_text, str) else existing.get("raw_text", ""),
        "cards": merged_cards,
    }
    _save_notes(data)
    return {"ok": True, "division_id": division_id, "card_count": len(merged_cards)}


# ─── 알림 히스토리 (앱 종 아이콘용) ───
@app.get("/notifications")
def get_notifications(limit: int = 100, unread_only: bool = False):
    """최근 알람 목록. 최신순."""
    hist = _load_fcm_history()
    if unread_only:
        hist = [h for h in hist if not h.get("read")]
    # 최신순
    hist = list(reversed(hist))[:max(1, min(limit, 200))]
    unread_count = sum(1 for h in _load_fcm_history() if not h.get("read"))
    return {"items": hist, "unread_count": unread_count}


@app.post("/notifications/mark_read")
def mark_notifications_read(payload: dict = None):
    """id 지정 시 해당 항목만, 없으면 전체 읽음 처리."""
    payload = payload or {}
    ids = payload.get("ids")
    hist = _load_fcm_history()
    changed = 0
    for h in hist:
        if ids is None or h.get("id") in ids:
            if not h.get("read"):
                h["read"] = True
                changed += 1
    _save_fcm_history(hist)
    return {"ok": True, "changed": changed}


@app.get("/notes")
def get_notes(division_id: str = ""):
    """앱이 호출하는 공개 API. division_id 지정 시 해당 사업부만 반환."""
    data = _load_notes()
    notes = data.get("notes", {})
    if division_id:
        item = notes.get(division_id)
        if not item:
            return {"division_id": division_id, "report_date": "", "cards": [], "updated_at": None}
        return {"division_id": division_id, **item}
    return {"notes": notes}


# ─────────────────────────────────────────────
# item 단위 due_date override (사용자 수동 지정)
#
# 규칙:
#   - due_date_auto: 파서 자동값 (매번 재계산)
#   - due_date_override: 사용자가 지정한 값 (있으면 유지)
#   - due_date: override or auto (최종 표시값)
#
# POST → override 지정, DELETE → override 제거 (자동값 복귀)
# ─────────────────────────────────────────────
def _find_item_by_id(notes_data: dict, division_id: str, card_title: str, item_id: str):
    """(item_ref, card_ref, section_ref) 반환. 못 찾으면 (None, None, None)."""
    div = (notes_data.get("notes") or {}).get(division_id)
    if not isinstance(div, dict):
        return None, None, None
    for card in div.get("cards", []) or []:
        if not isinstance(card, dict):
            continue
        if (card.get("title") or "").strip() != (card_title or "").strip():
            continue
        for sec in card.get("sections", []) or []:
            if not isinstance(sec, dict):
                continue
            for it in sec.get("items", []) or []:
                if isinstance(it, dict) and it.get("item_id") == item_id:
                    return it, card, sec
    return None, None, None


def _apply_effective_due_date(it: dict) -> None:
    """it 안의 override/auto 를 보고 최종 due_date 필드를 갱신."""
    override = (it.get("due_date_override") or "").strip()
    auto_raw = (it.get("due_date_auto") or "").strip()
    # auto_due_hidden 이 True 면 자동값을 무시
    auto = "" if it.get("auto_due_hidden") else auto_raw
    effective = override or auto
    if effective:
        it["due_date"] = effective
    else:
        it.pop("due_date", None)
    # 빈 override 는 필드 자체 제거
    if not override:
        it.pop("due_date_override", None)


@app.post("/admin/notes/item/due_override")
def admin_notes_item_due_override(
    payload: dict,
    _admin: int = Depends(get_admin_session),
):
    """특정 item 의 due_date_override 지정.
    
    Request:
      { division_id, card_title, item_id, due_date: "YYYY-MM-DD" }
    """
    from datetime import date as _date

    division_id = (payload or {}).get("division_id", "").strip()
    card_title = (payload or {}).get("card_title", "").strip()
    item_id = (payload or {}).get("item_id", "").strip()
    due_date_raw = (payload or {}).get("due_date", "").strip()

    if not division_id or not card_title or not item_id:
        raise HTTPException(status_code=400, detail="division_id, card_title, item_id 필수")
    if not due_date_raw:
        raise HTTPException(status_code=400, detail="due_date 필수 (지우려면 DELETE 사용)")
    
    # 형식 검증 (YYYY-MM-DD)
    try:
        parsed = _date.fromisoformat(due_date_raw)
        due_iso = parsed.isoformat()
    except Exception:
        raise HTTPException(status_code=400, detail="due_date 형식 오류 (YYYY-MM-DD 이어야 함)")

    notes_data = _load_notes()
    item, card, sec = _find_item_by_id(notes_data, division_id, card_title, item_id)
    if item is None:
        raise HTTPException(status_code=404, detail=f"item_id={item_id} 을(를) 찾을 수 없음")

    # FCM alarm hook: 저장 전 스냅샷
    _fcm_before = _snapshot_notes_status()

    item["due_date_override"] = due_iso
    _apply_effective_due_date(item)
    _save_notes(notes_data)

    # FCM alarm hook: 저장 후 diff → RED/ORANGE 신규 전이 알람
    try:
        _fcm_after = _snapshot_notes_status()
        _events = _detect_status_transitions(_fcm_before, _fcm_after)
        if _events:
            print(f"[FCM] due_override POST 상태 전이: {_events}")
            _fire_status_alarms(_events)
    except Exception as _e:
        print(f"[FCM due_override POST hook] {_e}")

    return {
        "ok": True,
        "item_id": item_id,
        "due_date": item.get("due_date"),
        "due_date_auto": item.get("due_date_auto"),
        "due_date_override": item.get("due_date_override"),
        "auto_due_hidden": bool(item.get("auto_due_hidden")),
    }


@app.delete("/admin/notes/item/due_override")
def admin_notes_item_due_override_reset(
    payload: dict,
    _admin: int = Depends(get_admin_session),
):
    """item 의 due_date_override 제거 → 자동값(due_date_auto) 복귀.
    
    Request:
      { division_id, card_title, item_id }
    """
    division_id = (payload or {}).get("division_id", "").strip()
    card_title = (payload or {}).get("card_title", "").strip()
    item_id = (payload or {}).get("item_id", "").strip()

    if not division_id or not card_title or not item_id:
        raise HTTPException(status_code=400, detail="division_id, card_title, item_id 필수")

    notes_data = _load_notes()
    item, card, sec = _find_item_by_id(notes_data, division_id, card_title, item_id)
    if item is None:
        raise HTTPException(status_code=404, detail=f"item_id={item_id} 을(를) 찾을 수 없음")

    # FCM alarm hook: 저장 전 스냅샷
    _fcm_before = _snapshot_notes_status()

    item.pop("due_date_override", None)
    _apply_effective_due_date(item)
    _save_notes(notes_data)

    # FCM alarm hook: DELETE 이후에도 상태 재판정 → 신규 RED/ORANGE 알람
    try:
        _fcm_after = _snapshot_notes_status()
        _events = _detect_status_transitions(_fcm_before, _fcm_after)
        if _events:
            print(f"[FCM] due_override DELETE 상태 전이: {_events}")
            _fire_status_alarms(_events)
    except Exception as _e:
        print(f"[FCM due_override DELETE hook] {_e}")

    return {
        "ok": True,
        "item_id": item_id,
        "due_date": item.get("due_date"),
        "due_date_auto": item.get("due_date_auto"),
        "due_date_override": item.get("due_date_override"),
        "auto_due_hidden": bool(item.get("auto_due_hidden")),
    }


@app.post("/admin/notes/item/hide_auto_due")
def admin_notes_item_hide_auto_due(
    payload: dict,
    _admin: int = Depends(get_admin_session),
):
    """특정 item 의 자동 파싱 날짜를 숨김 처리. 수동값이 있으면 그것도 함께 제거.
    
    Request:
      { division_id, card_title, item_id }
    """
    division_id = (payload or {}).get("division_id", "").strip()
    card_title = (payload or {}).get("card_title", "").strip()
    item_id = (payload or {}).get("item_id", "").strip()

    if not division_id or not card_title or not item_id:
        raise HTTPException(status_code=400, detail="division_id, card_title, item_id 필수")

    notes_data = _load_notes()
    item, card, sec = _find_item_by_id(notes_data, division_id, card_title, item_id)
    if item is None:
        raise HTTPException(status_code=404, detail=f"item_id={item_id} 을(를) 찾을 수 없음")

    # FCM alarm hook: 저장 전 스냅샷
    _fcm_before = _snapshot_notes_status()

    # 자동값 숨김 + 수동값도 함께 제거
    item["auto_due_hidden"] = True
    item.pop("due_date_override", None)
    _apply_effective_due_date(item)
    _save_notes(notes_data)

    # FCM alarm hook
    try:
        _fcm_after = _snapshot_notes_status()
        _events = _detect_status_transitions(_fcm_before, _fcm_after)
        if _events:
            print(f"[FCM] hide_auto_due POST 상태 전이: {_events}")
            _fire_status_alarms(_events)
    except Exception as _e:
        print(f"[FCM hide_auto_due POST hook] {_e}")

    return {
        "ok": True,
        "item_id": item_id,
        "due_date": item.get("due_date"),
        "due_date_auto": item.get("due_date_auto"),
        "due_date_override": item.get("due_date_override"),
        "auto_due_hidden": True,
    }


@app.delete("/admin/notes/item/hide_auto_due")
def admin_notes_item_hide_auto_due_reset(
    payload: dict,
    _admin: int = Depends(get_admin_session),
):
    """item 의 auto_due_hidden 해제 → 자동값(due_date_auto) 복원.
    
    Request:
      { division_id, card_title, item_id }
    """
    division_id = (payload or {}).get("division_id", "").strip()
    card_title = (payload or {}).get("card_title", "").strip()
    item_id = (payload or {}).get("item_id", "").strip()

    if not division_id or not card_title or not item_id:
        raise HTTPException(status_code=400, detail="division_id, card_title, item_id 필수")

    notes_data = _load_notes()
    item, card, sec = _find_item_by_id(notes_data, division_id, card_title, item_id)
    if item is None:
        raise HTTPException(status_code=404, detail=f"item_id={item_id} 을(를) 찾을 수 없음")

    # FCM alarm hook: 저장 전 스냅샷
    _fcm_before = _snapshot_notes_status()

    item.pop("auto_due_hidden", None)
    _apply_effective_due_date(item)
    _save_notes(notes_data)

    # FCM alarm hook
    try:
        _fcm_after = _snapshot_notes_status()
        _events = _detect_status_transitions(_fcm_before, _fcm_after)
        if _events:
            print(f"[FCM] hide_auto_due DELETE 상태 전이: {_events}")
            _fire_status_alarms(_events)
    except Exception as _e:
        print(f"[FCM hide_auto_due DELETE hook] {_e}")

    return {
        "ok": True,
        "item_id": item_id,
        "due_date": item.get("due_date"),
        "due_date_auto": item.get("due_date_auto"),
        "due_date_override": item.get("due_date_override"),
        "auto_due_hidden": False,
    }


@app.delete("/admin/notes/{division_id}")
def admin_delete_note(division_id: str, _admin: int = Depends(get_admin_session)):
    """노트 삭제."""
    data = _load_notes()
    if division_id in data.get("notes", {}):
        del data["notes"][division_id]
        _save_notes(data)
        return {"ok": True}
    raise HTTPException(status_code=404, detail="해당 사업부 노트 없음")



# ─── 노트 표/사진 첨부 API ───

# ─── 노트 표/사진 첨부 API ───
@app.post("/admin/notes/table")
def admin_save_note_table(payload: dict, _admin: int = Depends(get_admin_session)):
    """표 생성/수정. payload = {division_id, table_ref?, table: {title, headers, rows}}"""
    division_id = (payload or {}).get("division_id", "").strip()
    table = (payload or {}).get("table")
    existing_ref = (payload or {}).get("table_ref", "").strip()

    if not division_id:
        raise HTTPException(status_code=400, detail="division_id 필수")
    if not isinstance(table, dict):
        raise HTTPException(status_code=400, detail="table 필수")

    # 기존 표 수정인 경우 같은 파일에 덮어쓰기
    if existing_ref and "/" in existing_ref:
        div_id, asset_id = existing_ref.split("/", 1)
        p = _table_path(div_id, asset_id)
        p.write_text(json.dumps(table, ensure_ascii=False, indent=2), encoding="utf-8")
        return {"table_ref": existing_ref, "updated": True}

    # 신규 생성
    table_ref = _save_note_table(division_id, table)
    return {"table_ref": table_ref, "updated": False}


@app.get("/notes/table/{division_id}/{asset_id}")
def get_note_table(division_id: str, asset_id: str):
    """표 JSON 조회 (모바일/관리자 공용)"""
    table_ref = f"{division_id}/{asset_id}"
    data = _load_note_table(table_ref)
    if data is None:
        raise HTTPException(status_code=404, detail="표를 찾을 수 없음")
    return {"table_ref": table_ref, "table": data}


@app.delete("/admin/notes/table/{division_id}/{asset_id}")
def admin_delete_note_table(division_id: str, asset_id: str, _admin: int = Depends(get_admin_session)):
    table_ref = f"{division_id}/{asset_id}"
    ok = _delete_note_table(table_ref)
    if not ok:
        raise HTTPException(status_code=404, detail="표가 없음")
    return {"deleted": True, "table_ref": table_ref}


@app.post("/admin/notes/photo")
async def admin_upload_note_photo(
    division_id: str = Form(...),
    file: UploadFile = File(...),
    _admin: int = Depends(get_admin_session),
):
    """사진 업로드 → photo_ref(예: 'semiconductor/2026-06-11_abcd1234.jpg') 반환"""
    division_id = division_id.strip()
    if not division_id:
        raise HTTPException(status_code=400, detail="division_id 필수")

    # 확장자 결정
    orig_name = (file.filename or "").lower()
    ext = "jpg"
    for cand in ["jpg", "jpeg", "png", "gif", "webp"]:
        if orig_name.endswith("." + cand):
            ext = cand if cand != "jpeg" else "jpg"
            break

    asset_id = _new_asset_id(division_id)
    p = _photo_path(division_id, asset_id, ext)
    content = await file.read()
    p.write_bytes(content)
    photo_ref = f"{division_id}/{asset_id}.{ext}"
    return {"photo_ref": photo_ref, "url": f"/note_photos/{photo_ref}"}




# ─── 엑셀 업로드 (드래그&드롭) ───









def _excel_sheet_to_preview_data_url(ws):
    from io import BytesIO
    import base64
    from PIL import Image, ImageDraw, ImageFont, ImageChops

    NL = chr(10)

    def _effective_sheet_bounds():
        coords = []

        def _has_value(v):
            return v is not None and str(v).strip() != ""

        def _has_border(cell):
            try:
                b = cell.border
                return any([
                    b.left and b.left.style,
                    b.right and b.right.style,
                    b.top and b.top.style,
                    b.bottom and b.bottom.style,
                ])
            except Exception:
                return False

        def _has_fill(cell):
            try:
                f = cell.fill
                if f and f.fgColor and f.fgColor.rgb:
                    rgb = str(f.fgColor.rgb)
                    # 투명/흰색 제외
                    if rgb not in ("00000000", "FFFFFFFF", "FFFFFF", "None"):
                        return True
            except Exception:
                pass
            return False

        # 값 OR 테두리 OR 배경색 있는 셀 모두 포함
        for row in ws.iter_rows():
            for cell in row:
                if _has_value(cell.value) or _has_border(cell) or _has_fill(cell):
                    coords.append((cell.row, cell.column))

        # 병합 셀은 anchor 값이 있으면 병합 범위 전체 포함
        for rng in ws.merged_cells.ranges:
            min_col2, min_row2, max_col2, max_row2 = rng.bounds
            anchor = ws.cell(min_row2, min_col2)
            if _has_value(anchor.value) or _has_border(anchor) or _has_fill(anchor):
                coords.append((min_row2, min_col2))
                coords.append((max_row2, max_col2))

        if not coords:
            return 1, 1, 1, 1

        min_row2 = min(r for r, c in coords)
        max_row2 = max(r for r, c in coords)
        min_col2 = min(c for r, c in coords)
        max_col2 = max(c for r, c in coords)
        return min_row2, max_row2, min_col2, max_col2

    min_row, max_row, min_col, max_col = _effective_sheet_bounds()

    anchor_map = {}
    skip = set()
    for rng in ws.merged_cells.ranges:
        mc, mr, max_col2, max_row2 = rng.bounds
        anchor_map[(mr, mc)] = (max_row2, max_col2)
        for rr in range(mr, max_row2 + 1):
            for cc in range(mc, max_col2 + 1):
                if (rr, cc) != (mr, mc):
                    skip.add((rr, cc))

    # 한글 폰트 후보 (Mac → Linux/Railway 순)
    _font_candidates = [
        "/System/Library/Fonts/AppleSDGothicNeo.ttc",
        "/System/Library/Fonts/Supplemental/AppleGothic.ttf",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJKkr-Regular.otf",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
        "/usr/share/fonts/truetype/nanum/NanumBarunGothic.ttf",
    ]
    _font_bold_candidates = [
        "/System/Library/Fonts/AppleSDGothicNeo.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJKkr-Bold.otf",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Bold.ttc",
        "/usr/share/fonts/truetype/nanum/NanumGothicBold.ttf",
        "/usr/share/fonts/truetype/nanum/NanumBarunGothicBold.ttf",
    ]

    def _load_font(candidates, size):
        import os
        for fp in candidates:
            if os.path.exists(fp):
                try:
                    return ImageFont.truetype(fp, size)
                except Exception:
                    continue
        return ImageFont.load_default()

    font = _load_font(_font_candidates, 22)
    font_bold = _load_font(_font_bold_candidates, 22)

    def text_size(txt, f):
        try:
            bbox = f.getbbox(txt)
            return bbox[2] - bbox[0], bbox[3] - bbox[1]
        except Exception:
            return (len(txt) * 11, 22)

    def format_value(cell):
        v = cell.value
        if v is None:
            return ""
        nf = (cell.number_format or "").lower()
        if isinstance(v, (int, float)) and ("#,##0" in nf or "0" in nf):
            try:
                iv = int(v)
                if iv == 0:
                    return "-"
                return f"{iv:,}" if iv >= 0 else f"-{abs(iv):,}"
            except Exception:
                return str(v)
        if isinstance(v, str):
            try:
                if v.strip() and v.strip().lstrip("-").isdigit():
                    iv = int(v)
                    if iv == 0:
                        return "-"
                    return f"{iv:,}" if iv >= 0 else f"-{abs(iv):,}"
            except Exception:
                pass
        return str(v)

    cell_text = {}
    for r in range(min_row, max_row + 1):
        for c in range(min_col, max_col + 1):
            cell_text[(r, c)] = format_value(ws.cell(r, c))
    
    col_widths = [0] * (max_col + 2)
    for c in range(min_col, max_col + 1):
        widest = 72
        for r in range(min_row, max_row + 1):
            if (r, c) in skip:
                continue
            t = cell_text.get((r, c), "")
            if not t:
                continue
            use_font = font_bold if (ws.cell(r, c).font and ws.cell(r, c).font.b) else font
            for line in t.split(NL):
                w, _ = text_size(line, use_font)
                if w + 24 > widest:
                    widest = min(w + 24, 320)
        col_widths[c] = widest

    row_heights = [0] * (max_row + 2)
    for r in range(min_row, max_row + 1):
        lines_max = 1
        for c in range(min_col, max_col + 1):
            t = cell_text.get((r, c), "")
            lines = max(1, t.count(NL) + 1)
            if lines > lines_max:
                lines_max = lines
        row_heights[r] = max(40, 14 + 26 * lines_max)

    # 표 범위만 그리도록 인덱스 맵 통일
    cols = list(range(min_col, max_col + 1))
    rows = list(range(min_row, max_row + 1))
    col_idx = {c: i for i, c in enumerate(cols)}
    row_idx = {r: i for i, r in enumerate(rows)}

    W = sum(col_widths[c] for c in cols)
    H = sum(row_heights[r] for r in rows)

    img = Image.new("RGB", (W, H), "white")
    draw = ImageDraw.Draw(img)

    xs = [0]
    for c in cols:
        xs.append(xs[-1] + col_widths[c])
    ys = [0]
    for r in rows:
        ys.append(ys[-1] + row_heights[r])

    def argb_to_hex(argb):
        if not argb:
            return None
        v = str(argb)
        if v in ("00000000", "None"):
            return None
        if len(v) == 8:
            v = v[2:]
        if len(v) == 6:
            try:
                int(v, 16)
                return "#" + v.lower()
            except Exception:
                return None
        return None

    def get_fill(cell):
        try:
            fg = cell.fill.fgColor
            if fg is not None and getattr(fg, "type", None) == "rgb":
                return argb_to_hex(fg.rgb)
        except Exception:
            pass
        return None

    def get_font_color(cell):
        try:
            c = cell.font.color
            if c is not None and getattr(c, "type", None) == "rgb":
                return argb_to_hex(c.rgb)
        except Exception:
            pass
        return None

    def is_dark(hex_color):
        try:
            v = hex_color.lstrip("#")
            r = int(v[0:2], 16)
            g = int(v[2:4], 16)
            b = int(v[4:6], 16)
            return (0.299 * r + 0.587 * g + 0.114 * b) < 160
        except Exception:
            return False

    for r in range(min_row, max_row + 1):
        for c in range(min_col, max_col + 1):
            if (r, c) in skip:
                continue
            x1, y1 = xs[col_idx[c]], ys[row_idx[r]]
            rr, cc = anchor_map.get((r, c), (r, c))
            x2, y2 = xs[col_idx[cc] + 1], ys[row_idx[rr] + 1]
            cell = ws.cell(r, c)
            bg = get_fill(cell) or "#ffffff"
            fc = get_font_color(cell)
            if not fc:
                fc = "#ffffff" if is_dark(bg) else "#111111"
            draw.rectangle([x1, y1, x2 - 1, y2 - 1], fill=bg, outline="#cbd5e1", width=1)
            val = cell_text.get((r, c), "")
            if val:
                lines = val.split(NL)
                line_h = 26
                total_h = line_h * len(lines)
                ty = y1 + max(4, ((y2 - y1) - total_h) // 2)
                use_font = font_bold if (cell.font and cell.font.b) else font
                cell_w = x2 - x1
                for li, line in enumerate(lines):
                    w, _ = text_size(line, use_font)
                    tx = x1 + max(4, (cell_w - w) // 2)
                    draw.text((tx, ty + li * line_h), line, fill=fc, font=use_font)

    # 마지막 흰 여백 소폭 trim
    try:
        bg = Image.new(img.mode, img.size, "white")
        diff = ImageChops.difference(img, bg)
        bbox = diff.getbbox()
        if bbox:
            l, t, r, b = bbox
            # 공백 최소화: 거의 딱 맞게 자르되 2px만 여유
            pad_x = 2
            pad_y = 2
            img = img.crop((
                max(0, l - pad_x),
                max(0, t - pad_y),
                min(img.width, r + pad_x),
                min(img.height, b + pad_y),
            ))
    except Exception:
        pass

    bio = BytesIO()
    img.save(bio, format="PNG")
    return "data:image/png;base64," + base64.b64encode(bio.getvalue()).decode("ascii")




@app.post("/admin/notes/excel")
async def admin_upload_note_excel(
    division_id: str = Form(...),
    file: UploadFile = File(...),
    _admin: int = Depends(get_admin_session),
):
    """엑셀(.xlsx/.xlsm) 업로드 → 첫 시트를 PNG 이미지로 변환 → 사진처럼 저장 → photo_ref 반환."""
    division_id = (division_id or "").strip()
    if not division_id:
        raise HTTPException(status_code=400, detail="division_id 필수")

    orig_name = file.filename or "upload.xlsx"
    lower = orig_name.lower()
    if not (lower.endswith(".xlsx") or lower.endswith(".xlsm")):
        raise HTTPException(status_code=400, detail="xlsx/xlsm 파일만 허용")

    raw = await file.read()
    if not raw:
        raise HTTPException(status_code=400, detail="빈 파일")

    try:
        import openpyxl
        from io import BytesIO
        import base64
        wb = openpyxl.load_workbook(BytesIO(raw), data_only=True)
        ws = wb.worksheets[0]
        data_url = _excel_sheet_to_preview_data_url(ws)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"엑셀 변환 실패: {e}")

    # data:image/png;base64,... → 바이트 디코드
    try:
        b64 = data_url.split(",", 1)[1]
        png_bytes = base64.b64decode(b64)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PNG 디코딩 실패: {e}")

    asset_id = "xls_" + _new_asset_id(division_id)
    out_path = _photo_path(division_id, asset_id, "png")
    out_path.write_bytes(png_bytes)
    photo_ref = f"{division_id}/{asset_id}.png"

    return {
        "ok": True,
        "filename": orig_name,
        "photo_ref": photo_ref,
        "url": f"/note_photos/{photo_ref}",
    }

@app.delete("/admin/notes/photo/{division_id}/{filename}")
def admin_delete_note_photo(division_id: str, filename: str, _admin: int = Depends(get_admin_session)):
    photo_ref = f"{division_id}/{filename}"
    ok = _delete_note_photo(photo_ref)
    if not ok:
        raise HTTPException(status_code=404, detail="사진이 없음")
    return {"deleted": True, "photo_ref": photo_ref}


@app.get("/divisions")
def list_divisions():
    """공개 사업부 목록 + 각 사업부의 visible 프로젝트 리스트.
    매핑 관리 대시보드 / 모바일 사업부 화면에서 사용."""
    divs = _cl.get_divisions(visible_only=True)
    
    # 모델/현황/주차별 계획 데이터 로드 (프로젝트별 데이터 유무)
    models_data = _load_models()
    models_by_project = {}
    for proj_key, proj_data in models_data.get("projects", {}).items():
        has_models = len(proj_data.get("models", [])) > 0
        has_note = bool((proj_data.get("status_note") or "").strip())
        has_plan = bool(proj_data.get("weekly_plan", {}).get("photo_ref"))
        if has_models or has_note or has_plan:
            models_by_project[proj_key] = len(proj_data.get("models", []))
    
    result = []
    for d in divs:
        div_id = d.get("id")
        # 해당 사업부의 visible 프로젝트만 (order 순)
        try:
            ps = _cl.get_projects(div_id, visible_only=True)
        except Exception:
            ps = []
        ps_sorted = sorted(ps, key=lambda x: x.get("order", 999))
        result.append({
            "id": div_id,
            "label": d.get("label"),
            "order": d.get("order", 999),
            "badge_short_label": d.get("badge_short_label"),
            "projects": [
                {
                    "id": p.get("id"),
                    "label": p.get("label"),
                    "order": p.get("order", 999),
                    "has_models": p.get("id") in models_by_project,
                    "model_count": models_by_project.get(p.get("id"), 0),
                }
                for p in ps_sorted
            ],
        })
    result.sort(key=lambda x: x.get("order", 999))
    return {"divisions": result}


@app.get("/divisions/updates")
def get_divisions_updates():
    """각 사업부의 최신 업데이트 타임스탬프 반환.
    카드의 project_key 가 없을 수 있으므로 텍스트 기반 분류기를 사용한다.
    """
    try:
        with open(LATEST_FILE, "r", encoding="utf-8") as f:
            reports = json.load(f)
    except Exception:
        reports = []

    div_latest: dict[str, str] = {}

    for item in reports:
        ts = item.get("upload_timestamp") or ""
        products = item.get("products") or []
        for p in products:
            project_id = p.get("project_key") or p.get("project") or None
            if not project_id:
                parts = [
                    p.get("name") or "",
                    p.get("headline") or "",
                    " ".join(p.get("summary_bullets") or []),
                ]
                text = " ".join(s for s in parts if s)
                if text.strip():
                    try:
                        project_id = _cl.classify_project(text)
                    except Exception:
                        project_id = None

            div_id = _cl.derive_division_from_project(project_id) if project_id else None
            if not div_id:
                div_id = "unclassified"

            card_ts = p.get("updated_at") or ts
            if not card_ts:
                continue
            cur = div_latest.get(div_id, "")
            if card_ts > cur:
                div_latest[div_id] = card_ts

    all_divs = _cl.get_divisions(visible_only=True)
    result = []
    for d in all_divs:
        did = d.get("id")
        result.append({
            "division_id": did,
            "label": d.get("label"),
            "latest_updated_at": div_latest.get(did, ""),
        })

    return {"divisions": result}

@app.get("/projects")
def list_projects():
    """프로젝트 버튼 목록 + 모델 데이터 유무"""
    latest = _read_json(LATEST_FILE, [])
    grouped = aggregate_projects(latest)
    
    # 모델 데이터 로드
    models_data = _load_models()
    models_by_project = {}
    for proj_key, proj_data in models_data.get("projects", {}).items():
        models_list = proj_data.get("models", [])
        if models_list:
            models_by_project[proj_key] = len(models_list)
    
    projects = [
        {
            "key": k,
            "label": v["label"],
            "status": v["status"],
            "report_date": v.get("report_date"),
            "has_models": k in models_by_project,  # 모델 데이터 유무 추가
            "model_count": models_by_project.get(k, 0),  # 모델 개수 추가
        }
        for k, v in grouped.items()
    ]

    # 보고서에 없지만 모델이 있는 프로젝트도 목록에 추가
    for proj_key, model_count in models_by_project.items():
        if proj_key not in grouped:
            projects.append({
                "key": proj_key,
                "label": proj_key.upper(),
                "status": "BLACK",
                "report_date": None,
                "has_models": True,
                "model_count": model_count,
            })
    severity = {"RED": 3, "BLUE": 2, "BLACK": 1}

    # 🟢 프로젝트 목록 enrichment (기존 필드 변경 없음)
    projects = [enrich_project_entry(p) for p in projects]

    projects.sort(key=lambda p: -severity.get(p["status"], 0))
    return {"projects": projects}


@app.get("/projects/{project_key}/models")
def get_project_models(project_key: str):
    """프로젝트 모델 목록 + 그룹별 요약 (앱 8번 화면용)

    응답:
    {
      "project_key": "cup",
      "has_models": true,
      "total": 4,
      "groups": {"양산": {"count": 4, "models": [...]}, "개발": {...}},
      "models": [...]
    }
    """
    models = _get_project_models(project_key)
    groups: dict = {}
    for m in models:
        g = (m.get("group") or "기타").strip() or "기타"
        groups.setdefault(g, {"count": 0, "models": []})
        groups[g]["count"] += 1
        groups[g]["models"].append(m)
    return {
        "project_key": project_key,
        "has_models": len(models) > 0,
        "total": len(models),
        "groups": groups,
        "models": models,
    }


@app.get("/projects/{project_key}")
def get_project_detail(project_key: str):
    """프로젝트 상세"""
    # 별칭 매핑: 앱/구 데이터에서 오는 키를 실제 프로젝트 키로 변환
    _alias_map = {
        "havaplate": "hrva_plate",
        "hrvaplate": "hrva_plate",
        "hrva-plate": "hrva_plate",
    }
    project_key = _alias_map.get(project_key.strip().lower(), project_key.strip())
    latest = _read_json(LATEST_FILE, [])
    grouped = aggregate_projects(latest)
    detail = build_project_detail(project_key, grouped)
    if not detail:
        # PPT 없어도 노트만 있으면 노트 카드로 응답
        try:
            notes_data = _read_json(NOTES_FILE, {})
            for div_id, div_obj in (notes_data.get("notes") or {}).items():
                for card in (div_obj.get("cards") or []):
                    if (card.get("title") or "").strip() == project_key.strip():
                        detail = {
                            "project_key": project_key,
                            "label": project_key,
                            "name": project_key,
                            "status": _calc_card_status(card),
                            "sections": card.get("sections") or [],
                            "note_only": True,
                            "report_date": div_obj.get("report_date"),
                            "updated_at": div_obj.get("updated_at"),
                        }
                        break
                if detail:
                    break
        except Exception as _e:
            pass
        if not detail:
            raise HTTPException(status_code=404, detail="프로젝트를 찾을 수 없습니다.")
    
    # === notes.json 매칭 sections 우선 사용 + str items → obj items 변환 ===
    try:
        _notes_data = _read_json(NOTES_FILE, {})
        _project_label = (detail.get("label") or detail.get("project_label") or "").strip()
        _notes_card = None
        for _div_id, _div_obj in (_notes_data.get("notes") or {}).items():
            for _card in (_div_obj.get("cards") or []):
                if (_card.get("title") or "").strip() == _project_label:
                    _notes_card = _card
                    break
            if _notes_card:
                break

        if _notes_card and _notes_card.get("sections"):
            detail["sections"] = _notes_card["sections"]
        else:
            for _sec in detail.get("sections") or []:
                _raw_items = _sec.get("items", []) or []
                _raw_notes = _sec.get("notes", []) or []
                _new_items = []
                for _txt in _raw_items:
                    if isinstance(_txt, dict):
                        _new_items.append(_txt)
                        continue
                    _s = str(_txt).strip()
                    if _s:
                        _new_items.append({"type": "bullet", "text": _s})
                for _txt in _raw_notes:
                    if isinstance(_txt, dict):
                        _new_items.append(_txt)
                        continue
                    _s = str(_txt).strip()
                    if not _s:
                        continue
                    if _s.startswith("※"):
                        _new_items.append({"type": "highlight", "text": _s.lstrip("※ ").strip()})
                    else:
                        _new_items.append({"type": "sub", "text": _s})
                _sec["items"] = _new_items
                _sec.pop("notes", None)
    except Exception as _e:
        pass

    # 🟢 프로젝트 상세 enrichment (기존 필드 변경 없음)
    detail = enrich_project_detail(detail)

    # KPI 카드 자동 첨부 (project_key 별 화이트리스트)
    try:
        if project_key == "major_module":
            detail["kpi_card"] = _build_major_module_kpi_card()
            detail["issue_lines"] = _build_major_module_issue_lines()
    except Exception as _e:
        # KPI 계산 실패해도 기본 상세는 정상 반환되어야 함
        pass

    return detail


# =========================================================
# 5. 사용자 직접 차트/사진 업로드
# =========================================================
# ============================================================
# Admin Config API (5단계: admin 페이지 동적 dropdown용)
# - 운영 도구 페이지에서 사업부/프로젝트/섹션을 동적으로 로드하기 위한 API
# - 모두 GET, 인증 불필요 (admin 페이지 내부에서만 호출됨)
# ============================================================

@app.get("/admin/config/divisions")
def admin_config_divisions(_exp: int = Depends(get_admin_session)):
    """
    admin 페이지의 사업부 dropdown 용.
    config_loader.get_divisions() 결과를 그대로 노출.
    """
    try:
        items = _cl.get_divisions(visible_only=True)
        return {
            "divisions": [
                {
                    "id": d.get("id"),
                    "label": d.get("label"),
                    "mode": d.get("mode"),
                    "order": d.get("order"),
                    "badge_short_label": d.get("badge_short_label"),
                }
                for d in items
            ]
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"divisions load error: {e}")


# === Learned Aliases (사용자 승인 오타 -> 프로젝트 매핑) ===
LEARNED_ALIASES_FILE = Path("learned_aliases.json")

def _load_learned_aliases() -> dict:
    try:
        with open(LEARNED_ALIASES_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {"version": 1, "updated_at": None, "aliases": {}}

def _save_learned_aliases(data: dict) -> None:
    from datetime import datetime
    data["updated_at"] = datetime.now().isoformat()
    with open(LEARNED_ALIASES_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

def _find_project_by_learned_alias(query: str) -> Optional[dict]:
    q = (query or "").strip().lower()
    if not q:
        return None
    data = _load_learned_aliases()
    for project_id, alias_list in (data.get("aliases") or {}).items():
        for a in alias_list:
            if str(a).strip().lower() == q:
                try:
                    all_projects = _cl.get_projects(visible_only=True)
                    for p in all_projects:
                        if p.get("id") == project_id:
                            return {
                                "id": p.get("id"),
                                "label": (p.get("label") or "").strip(),
                                "division_id": p.get("division_id"),
                                "matched_on": a,
                                "source": "learned_alias",
                            }
                except Exception:
                    pass
    return None


@app.post("/admin/projects/resolve")
def admin_projects_resolve(payload: dict, _admin: int = Depends(get_admin_session)):
    """
    fuzzy로 못 잡은 오타를 AI로 판정.
    payload: {"query": "참바", "division_id": "semiconductor" (optional)}
    """
    if client is None:
        return {"resolved": None, "reason": "OPENAI_API_KEY not configured"}

    query = str((payload or {}).get("query") or "").strip()
    division_id = (payload or {}).get("division_id") or None
    if not query:
        return {"resolved": None, "reason": "empty query"}

    # 먼저 learned aliases 체크 (AI 호출 안 하고)
    learned = _find_project_by_learned_alias(query)
    if learned:
        return {"resolved": learned, "source": "learned_alias", "confidence": 1.0}

    try:
        all_projects = _cl.get_projects(division_id=division_id, visible_only=True)
    except Exception:
        all_projects = []

    if not all_projects:
        return {"resolved": None, "reason": "no candidates"}

    # 후보 목록 문자열 생성
    candidates_desc = []
    for p in all_projects:
        label = (p.get("label") or "").strip()
        aliases = p.get("aliases") or []
        candidates_desc.append(f"- id={p.get('id')}, label={label}, aliases={aliases}")
    candidates_text = "\n".join(candidates_desc)

    system_msg = "너는 한국어/영어 프로젝트명 오타 판정 어시스턴트다. 사용자가 입력한 이름이 후보 목록 중 어느 프로젝트를 의미하는지 판단한다. 반드시 JSON만 응답한다."

    user_msg = f"""사용자 입력: "{query}"

후보 프로젝트 목록:
{candidates_text}

응답 형식 (JSON only, no markdown):
{{
  "matched_id": "chamber" 또는 null,
  "confidence": 0.0~1.0,
  "reason": "짧은 설명"
}}

규칙:
- 명백히 오타/약어/유사어면 matched_id 반환
- 완전히 다른 이름이면 matched_id: null
- confidence 0.6 미만이면 matched_id: null
- 후보에 없는 새로운 프로젝트로 보이면 matched_id: null
"""

    try:
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": system_msg},
                {"role": "user", "content": user_msg},
            ],
            temperature=0.0,
            response_format={"type": "json_object"},
        )
        raw = resp.choices[0].message.content or "{}"
        result = json.loads(raw)
    except Exception as e:
        return {"resolved": None, "reason": f"ai error: {e}"}

    matched_id = result.get("matched_id")
    confidence = float(result.get("confidence") or 0.0)
    reason = result.get("reason") or ""

    if not matched_id or confidence < 0.6:
        return {"resolved": None, "reason": reason, "confidence": confidence}

    # matched_id로 프로젝트 정보 찾기
    for p in all_projects:
        if p.get("id") == matched_id:
            return {
                "resolved": {
                    "id": p.get("id"),
                    "label": (p.get("label") or "").strip(),
                    "division_id": p.get("division_id"),
                    "matched_on": query,
                    "source": "ai",
                },
                "confidence": confidence,
                "reason": reason,
            }

    return {"resolved": None, "reason": "matched_id not found in candidates"}


@app.post("/admin/projects/alias")
def admin_projects_add_alias(payload: dict, _admin: int = Depends(get_admin_session)):
    """
    사용자 승인 alias를 learned_aliases.json에 저장.
    payload: {"project_id": "chamber", "alias": "참바"}
    """
    project_id = str((payload or {}).get("project_id") or "").strip()
    alias = str((payload or {}).get("alias") or "").strip()

    if not project_id or not alias:
        raise HTTPException(status_code=400, detail="project_id and alias required")

    # 프로젝트 존재 확인
    try:
        all_projects = _cl.get_projects(visible_only=True)
        found = any(p.get("id") == project_id for p in all_projects)
        if not found:
            raise HTTPException(status_code=404, detail=f"project_id not found: {project_id}")
    except HTTPException:
        raise
    except Exception:
        pass

    data = _load_learned_aliases()
    aliases = data.setdefault("aliases", {})
    lst = aliases.setdefault(project_id, [])

    # 중복 방지 (case insensitive)
    alias_lower = alias.lower()
    if not any(str(a).lower() == alias_lower for a in lst):
        lst.append(alias)
        _save_learned_aliases(data)
        return {"ok": True, "added": True, "project_id": project_id, "alias": alias, "total_aliases": len(lst)}
    else:
        return {"ok": True, "added": False, "reason": "alias already exists"}


# === 한/영 자판 변환 유틸 ===
_EN_TO_KO_JAMO = {
    "q":"ㅂ","w":"ㅈ","e":"ㄷ","r":"ㄱ","t":"ㅅ","y":"ㅛ","u":"ㅕ","i":"ㅑ","o":"ㅐ","p":"ㅔ",
    "a":"ㅁ","s":"ㄴ","d":"ㅇ","f":"ㄹ","g":"ㅎ","h":"ㅗ","j":"ㅓ","k":"ㅏ","l":"ㅣ",
    "z":"ㅋ","x":"ㅌ","c":"ㅊ","v":"ㅍ","b":"ㅠ","n":"ㅜ","m":"ㅡ",
    "Q":"ㅃ","W":"ㅉ","E":"ㄸ","R":"ㄲ","T":"ㅆ","O":"ㅒ","P":"ㅖ",
}
_KO_JAMO_TO_EN = {v: k for k, v in _EN_TO_KO_JAMO.items()}

_CHOSEONG = list("ㄱㄲㄴㄷㄸㄹㅁㅂㅃㅅㅆㅇㅈㅉㅊㅋㅌㅍㅎ")
_JUNGSEONG = list("ㅏㅐㅑㅒㅓㅔㅕㅖㅗㅘㅙㅚㅛㅜㅝㅞㅟㅠㅡㅢㅣ")
_JONGSEONG = [""] + list("ㄱㄲㄳㄴㄵㄶㄷㄹㄺㄻㄼㄽㄾㄿㅀㅁㅂㅄㅅㅆㅇㅈㅊㅋㅌㅍㅎ")

_JUNG_COMBINE = {
    ("ㅗ","ㅏ"):"ㅘ",("ㅗ","ㅐ"):"ㅙ",("ㅗ","ㅣ"):"ㅚ",
    ("ㅜ","ㅓ"):"ㅝ",("ㅜ","ㅔ"):"ㅞ",("ㅜ","ㅣ"):"ㅟ",
    ("ㅡ","ㅣ"):"ㅢ",
}
_JONG_COMBINE = {
    ("ㄱ","ㅅ"):"ㄳ",("ㄴ","ㅈ"):"ㄵ",("ㄴ","ㅎ"):"ㄶ",
    ("ㄹ","ㄱ"):"ㄺ",("ㄹ","ㅁ"):"ㄻ",("ㄹ","ㅂ"):"ㄼ",
    ("ㄹ","ㅅ"):"ㄽ",("ㄹ","ㅌ"):"ㄾ",("ㄹ","ㅍ"):"ㄿ",
    ("ㄹ","ㅎ"):"ㅀ",("ㅂ","ㅅ"):"ㅄ",
}

def _en_to_ko(text: str) -> str:
    if not text:
        return text
    jamos = []
    for ch in text:
        if ch in _EN_TO_KO_JAMO:
            jamos.append(_EN_TO_KO_JAMO[ch])
        else:
            jamos.append(ch)
    result = []
    i = 0
    n = len(jamos)
    while i < n:
        j = jamos[i]
        if j not in _CHOSEONG and j not in _JUNGSEONG:
            result.append(j)
            i += 1
            continue
        if j not in _CHOSEONG:
            result.append(j)
            i += 1
            continue
        cho = _CHOSEONG.index(j)
        if i + 1 >= n or jamos[i+1] not in _JUNGSEONG:
            result.append(j)
            i += 1
            continue
        jung_char = jamos[i+1]
        i += 2
        if i < n and jamos[i] in _JUNGSEONG:
            combined = _JUNG_COMBINE.get((jung_char, jamos[i]))
            if combined:
                jung_char = combined
                i += 1
        jung = _JUNGSEONG.index(jung_char)
        jong = 0
        if i < n and jamos[i] in _CHOSEONG:
            if i + 1 >= n or jamos[i+1] not in _JUNGSEONG:
                jong_char = jamos[i]
                if jong_char in _JONGSEONG:
                    if i + 1 < n and jamos[i+1] in _CHOSEONG:
                        if i + 2 >= n or jamos[i+2] not in _JUNGSEONG:
                            combined = _JONG_COMBINE.get((jong_char, jamos[i+1]))
                            if combined and combined in _JONGSEONG:
                                jong = _JONGSEONG.index(combined)
                                i += 2
                            else:
                                jong = _JONGSEONG.index(jong_char)
                                i += 1
                        else:
                            jong = _JONGSEONG.index(jong_char)
                            i += 1
                    else:
                        jong = _JONGSEONG.index(jong_char)
                        i += 1
        code = 0xAC00 + (cho * 21 + jung) * 28 + jong
        result.append(chr(code))
    return "".join(result)

def _ko_to_en(text: str) -> str:
    if not text:
        return text
    result = []
    for ch in text:
        code = ord(ch)
        if 0xAC00 <= code <= 0xD7A3:
            base = code - 0xAC00
            cho_i = base // (21 * 28)
            jung_i = (base % (21 * 28)) // 28
            jong_i = base % 28
            cho = _CHOSEONG[cho_i]
            jung = _JUNGSEONG[jung_i]
            jong = _JONGSEONG[jong_i]
            for jamo in [cho, jung, jong]:
                if not jamo:
                    continue
                if jamo in ["ㅘ","ㅙ","ㅚ","ㅝ","ㅞ","ㅟ","ㅢ"]:
                    for k, v in _JUNG_COMBINE.items():
                        if v == jamo:
                            for j2 in k:
                                if j2 in _KO_JAMO_TO_EN:
                                    result.append(_KO_JAMO_TO_EN[j2])
                            break
                elif jamo in ["ㄳ","ㄵ","ㄶ","ㄺ","ㄻ","ㄼ","ㄽ","ㄾ","ㄿ","ㅀ","ㅄ"]:
                    for k, v in _JONG_COMBINE.items():
                        if v == jamo:
                            for j2 in k:
                                if j2 in _KO_JAMO_TO_EN:
                                    result.append(_KO_JAMO_TO_EN[j2])
                            break
                elif jamo in _KO_JAMO_TO_EN:
                    result.append(_KO_JAMO_TO_EN[jamo])
        elif ch in _KO_JAMO_TO_EN:
            result.append(_KO_JAMO_TO_EN[ch])
        else:
            result.append(ch)
    return "".join(result)


@app.get("/admin/projects/suggest")
def admin_projects_suggest(query: str = "", limit: int = 5):
    """
    프로젝트명 fuzzy match.
    - learned_aliases.json 우선 체크
    - 한글 자모 단위 유사도 (챕버 <-> 챔버)
    - 한/영 자판 변환 (xhffhs -> 톨론, coaqj -> 챔버)
    """
    from difflib import SequenceMatcher
    import unicodedata

    def to_jamo(s: str) -> str:
        return unicodedata.normalize("NFKD", s or "")

    def get_choseong(s: str) -> str:
        # 한글 완성형에서 초성만 추출
        result = []
        for ch in s:
            code = ord(ch)
            if 0xAC00 <= code <= 0xD7A3:
                base = code - 0xAC00
                cho_i = base // (21 * 28)
                result.append(_CHOSEONG[cho_i])
            else:
                result.append(ch)
        return "".join(result)

    def sim_ratio(a: str, b: str) -> float:
        char_sim = SequenceMatcher(None, a, b).ratio()
        jamo_a, jamo_b = to_jamo(a), to_jamo(b)
        jamo_sim = SequenceMatcher(None, jamo_a, jamo_b).ratio() if jamo_a and jamo_b else 0.0
        base_sim = max(char_sim, jamo_sim)
        # 초성 일치 보너스 (한글 오타 특성)
        cho_a = get_choseong(a)
        cho_b = get_choseong(b)
        if cho_a and cho_b and len(cho_a) == len(cho_b):
            cho_sim = SequenceMatcher(None, cho_a, cho_b).ratio()
            if cho_sim >= 0.99:  # 초성 완전 일치
                base_sim = min(1.0, base_sim + 0.15)
            elif cho_sim >= 0.7:  # 초성 대부분 일치
                base_sim = min(1.0, base_sim + 0.08)
        return base_sim

    q_raw = (query or "").strip()
    q = q_raw.lower()
    if not q:
        return {"suggestions": [], "exact_match": None}

    learned = _find_project_by_learned_alias(q)
    if learned:
        return {
            "suggestions": [{
                "id": learned["id"],
                "label": learned["label"],
                "division_id": learned["division_id"],
                "similarity": 1.0,
                "matched_on": learned["matched_on"],
                "exact": True,
                "source": "learned_alias",
            }],
            "exact_match": learned,
        }

    q_variants = [q]
    has_en = any(("a" <= c <= "z") or ("A" <= c <= "Z") for c in q_raw)
    has_ko = any((0xAC00 <= ord(c) <= 0xD7A3) or (0x3131 <= ord(c) <= 0x318E) for c in q_raw)
    if has_en and not has_ko:
        converted = _en_to_ko(q_raw).lower()
        if converted and converted != q:
            q_variants.append(converted)
    elif has_ko and not has_en:
        converted = _ko_to_en(q_raw).lower()
        if converted and converted != q:
            q_variants.append(converted)

    try:
        all_projects = _cl.get_projects(visible_only=True)
    except Exception:
        all_projects = []

    scored = []
    exact = None

    for p in all_projects:
        candidates = []
        label = (p.get("label") or "").strip()
        if label:
            candidates.append(label)
        for a in (p.get("aliases") or []):
            if a: candidates.append(str(a))
        for k in (p.get("keywords") or []):
            if k: candidates.append(str(k))

        best_sim = 0.0
        best_matched = ""
        found_exact_here = False
        for c in candidates:
            cl = c.lower().strip()
            for qv in q_variants:
                if cl == qv:
                    exact = {
                        "id": p.get("id"),
                        "label": label,
                        "division_id": p.get("division_id"),
                        "matched_on": c,
                    }
                    best_sim = 1.0
                    best_matched = c
                    found_exact_here = True
                    break
                sim = sim_ratio(qv, cl)
                # 부분 포함 보정: alias가 label보다 짧으면 페널티
                if qv in cl or cl in qv:
                    # cl이 매우 짧으면(2자 이하) 부분 매칭 보너스 낮춤
                    if len(cl) <= 2 and c != label:
                        sim = max(sim, 0.5)
                    else:
                        sim = max(sim, 0.75)
                if sim > best_sim:
                    best_sim = sim
                    best_matched = c
            if found_exact_here:
                break

        if best_sim >= 0.55:
            scored.append({
                "id": p.get("id"),
                "label": label,
                "division_id": p.get("division_id"),
                "similarity": round(best_sim, 3),
                "matched_on": best_matched,
                "exact": best_sim >= 0.999,
            })

    scored.sort(key=lambda x: x["similarity"], reverse=True)
    return {
        "suggestions": scored[:limit],
        "exact_match": exact,
    }


@app.get("/admin/config/projects")
def admin_config_projects(division_id: str | None = None):
    """
    admin 페이지의 프로젝트 dropdown 용.
    division_id 가 주어지면 해당 사업부의 프로젝트만 반환.
    """
    try:
        items = _cl.get_projects(division_id=division_id, visible_only=True)
        return {
            "division_id": division_id,
            "projects": [
                {
                    "id": p.get("id"),
                    "label": p.get("label"),
                    "badge_label": p.get("badge_label"),
                    "group": p.get("group"),
                    "order": p.get("order"),
                }
                for p in items
            ]
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"projects load error: {e}")


@app.get("/admin/config/sections")
def admin_config_sections(project_key: str, _exp: int = Depends(get_admin_session)):
    """
    admin 이미지 업로드 탭의 섹션 dropdown 용.
    현재 분석된 PPT 결과에서 해당 프로젝트의 실제 GPT 섹션 제목들을 가져옴.
    이것이 핵심 — 더 이상 자유 입력 안 함, GPT 가 만든 실제 섹션만 옵션으로 노출.
    """
    try:
        latest = _read_json(LATEST_FILE, [])
        grouped = aggregate_projects(latest)
        detail = build_project_detail(project_key, grouped)
        if not detail:
            return {"project_key": project_key, "sections": []}

        sections = detail.get("sections", [])
        return {
            "project_key": project_key,
            "project_label": detail.get("label"),
            "sections": [
                {
                    "title": s.get("title"),
                    "image_count": len(s.get("image_urls", []) or []),
                }
                for s in sections
            ]
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"sections load error: {e}")


@app.get("/admin/config/stats")
def admin_config_stats(_exp: int = Depends(get_admin_session)):
    """
    매핑 관리 탭용 통계.
    - 현재 사업부/프로젝트 수
    - 자동 분류된 카드/실패 카드 수
    - 실패 카드의 product/category 샘플
    """
    try:
        divs = _cl.get_divisions(visible_only=True)
        all_projects = _cl.get_projects(visible_only=True)

        latest = _read_json(LATEST_FILE, [])
        total_cards = 0
        unclassified_cards = []
        for report in latest:
            for p in report.get("products", []):
                total_cards += 1
                name = p.get("name") or p.get("product", "")
                category = p.get("category", "")
                text = f"{name} {category}".strip()
                pid = _cl.classify_project(text)
                if not pid:
                    unclassified_cards.append({
                        "name": name,
                        "category": category,
                        "report_family": report.get("report_meta", {}).get("report_family", ""),
                    })

        return {
            "divisions_count": len(divs),
            "projects_count": len(all_projects),
            "total_cards": total_cards,
            "unclassified_count": len(unclassified_cards),
            "unclassified_samples": unclassified_cards[:20],
            "projects_by_division": {
                d.get("id"): [
                    {"id": p.get("id"), "label": p.get("label")}
                    for p in _cl.get_projects(division_id=d.get("id"), visible_only=True)
                ]
                for d in divs
            }
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"stats load error: {e}")

# =========================================================
# 6. 어드민 페이지 (브라우저용)
# =========================================================


@app.get("/admin/v2", response_class=HTMLResponse)
def admin_v2_page(admin_auth: Optional[str] = Cookie(default=None)):
    if not _verify_session(admin_auth):
        return RedirectResponse(url="/admin/login?next=/admin/v2", status_code=302)
    """OneView Admin v2 — 시안 기반 새 관리자 페이지."""
    return HTMLResponse(_ADMIN_V2_HTML)


_ADMIN_V2_HTML = """<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="utf-8">
<title>OneView Admin</title>
<meta name="viewport" content="width=device-width, initial-scale=1">
<style>
  * { box-sizing: border-box; }
  html, body { margin: 0; padding: 0; height: 100%; font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", "Pretendard", "Apple SD Gothic Neo", sans-serif; color: #1F2937; background: #F5F7FA; }
  a { color: inherit; text-decoration: none; }

  .app { display: flex; min-height: 100vh; }

  /* 좌측 사이드바 */
  .sidebar { width: 240px; flex-shrink: 0; background: #0B2B5C; color: #E5E7EB; padding: 22px 16px; display: flex; flex-direction: column; gap: 18px; }
  .sidebar .brand { display: flex; align-items: baseline; gap: 8px; padding: 0 6px; }
  .sidebar .brand .name { font-size: 22px; font-weight: 800; color: #fff; letter-spacing: -0.3px; }
  .sidebar .brand .sub { font-size: 12px; color: #93C5FD; font-weight: 600; }

  .biz-select { background: #12356F; border-radius: 10px; padding: 10px 12px; display: flex; flex-direction: column; gap: 4px; }
  .biz-select .label { font-size: 12px; color: #93C5FD; font-weight: 600; letter-spacing: 0.5px; }
  .biz-select select { background: transparent; color: #fff; border: none; font-size: 15px; font-weight: 600; padding: 2px 0; outline: none; cursor: pointer; }
  .biz-select select option { color: #111; }

  .nav-title { font-size: 12px; color: #93C5FD; letter-spacing: 1px; padding: 0 6px; font-weight: 700; }
  .nav { display: flex; flex-direction: column; gap: 4px; }
  .nav-item { display: flex; align-items: center; gap: 12px; padding: 12px 14px; border-radius: 10px; font-size: 15px; font-weight: 500; color: #CBD5E1; cursor: pointer; transition: background .15s, color .15s; }
  .nav-item:hover { background: rgba(255,255,255,0.06); color: #fff; }
  .nav-item.active { background: #fff; color: #0B2B5C; font-weight: 700; }
  .nav-item .icon { width: 22px; text-align: center; font-size: 16px; }

  .sidebar .spacer { flex: 1; }
  .sidebar .footer { font-size: 12px; color: #93C5FD; padding: 0 6px; opacity: 0.7; }

  /* 메인 */
  .main { flex: 1; display: flex; flex-direction: column; min-width: 0; }

  /* 상단 헤더 */
  .topbar { background: #fff; border-bottom: 1px solid #E5E7EB; padding: 16px 28px; display: flex; align-items: center; gap: 16px; }
  .breadcrumb { font-size: 13px; color: #6B7280; }
  .breadcrumb .sep { margin: 0 6px; color: #9CA3AF; }
  .breadcrumb .cur { color: #111827; font-weight: 600; }
  .topbar .grow { flex: 1; }
  .search { position: relative; }
  .search input { padding: 8px 12px 8px 34px; border: 1px solid #E5E7EB; border-radius: 999px; background: #F9FAFB; font-size: 13px; width: 220px; outline: none; }
  .search input:focus { border-color: #3B82F6; background: #fff; }
  .search::before { content: "🔍"; position: absolute; left: 12px; top: 50%; transform: translateY(-50%); font-size: 12px; opacity: 0.6; }
  .btn-primary { background: #0F2C59; color: #fff; border: none; border-radius: 8px; padding: 9px 16px; font-size: 13px; font-weight: 600; cursor: pointer; }
  .btn-primary:hover { background: #12356F; }

  /* 본문 영역 */
  .content { padding: 28px; flex: 1; }
  .page-title { font-size: 24px; font-weight: 800; color: #111827; margin: 0 0 4px 0; }
  .page-sub { font-size: 14px; color: #6B7280; margin: 0 0 20px 0; }

  .placeholder { background: #fff; border-radius: 14px; border: 1px solid #E5E7EB; padding: 60px 28px; text-align: center; color: #6B7280; }
  .placeholder .big { font-size: 40px; margin-bottom: 12px; opacity: 0.5; }
  .placeholder .title { font-size: 16px; font-weight: 700; color: #374151; margin-bottom: 6px; }
  .placeholder .sub { font-size: 13px; color: #9CA3AF; }
</style>
</head>
<body>
<div class="app">
  <!-- 사이드바 -->
  <aside class="sidebar">
    <div class="brand">
      <span class="name">OneView</span>
      <span class="sub">Admin</span>
    </div>

    <div class="biz-select">
      <span class="label">사업부</span>
      <select id="v2-division-select">
        <option value="">로딩 중...</option>
      </select>
    </div>

    <div class="nav-title">MENU</div>
    <nav class="nav">
      <div class="nav-item" data-page="home"><span class="icon">🏠</span><span>홈 대시보드</span></div>
      <div class="nav-item active" data-page="models"><span class="icon">🧩</span><span>모델 관리</span></div>
    </nav>

    <div class="spacer"></div>
    <div class="footer">v2 skeleton</div>
  </aside>

  <!-- 메인 -->
  <main class="main">
    <!-- 상단바 -->
    <div class="topbar">
      <div class="breadcrumb">
        <span id="v2-crumb-biz">반도체사업부</span>
        <span class="sep">›</span>
        <span class="cur" id="v2-crumb-page">모델 관리</span>
      </div>
      <div class="grow"></div>
      <div class="search"><input type="text" placeholder="검색"></div>
      
    </div>

    <!-- 페이지 컨텐츠 -->
    <div class="content" id="v2-content">
      <div id="ov-report-root"></div>
    </div>
  </main>
</div>

<script>
window._normalizeAiNumberedHtml = function(html){
  if (!html) return html;
  var s = String(html);
  // <br> 기준으로 라인 분리 (대소문자 무시)
  var lines = s.split(/<br\s*\/?>/i);
  // 각 라인의 &nbsp; 정규화
  lines = lines.map(function(l){ return l.replace(/&nbsp;/g, ' '); });
  // stripTag 함수: 태그 벗겨서 텍스트만 봤을 때 번호로 시작하는지 확인
  function stripTag(x){ return x.replace(/<[^>]+>/g, ''); }
  // 각 라인이 새 항목 시작인지 판별
  var items = [];
  var current = null;
  for (var i = 0; i < lines.length; i++){
    var raw = lines[i];
    var text = stripTag(raw).replace(/\u00A0/g, ' ').trim();
    if (!text) {
      // 빈 라인은 현재 항목에 <br> 로 유지
      if (current !== null) current += '<br>';
      continue;
    }
    // "1)" "2." 등 새 항목 시작
    var m = text.match(/^(\d+)[)\.]\s+(.*)$/);
    if (m) {
      // 새 항목 시작 - raw에서 "숫자)" 또는 "숫자." 부분 제거
      var stripped = raw.replace(/^(\s|&nbsp;|\u00A0)*\d+[)\.]\s*/, '');
      if (current !== null) items.push(current);
      current = stripped;
    } else {
      // 이어붙임
      if (current === null) {
        current = raw;
      } else {
        current += '<br>' + raw;
      }
    }
  }
  if (current !== null) items.push(current);
  // 항목이 2개 이상일 때만 <ol> 로 감싸기
  if (items.length < 2) return html;
  var out = '<ol class="ov-num ov-num-paren" data-start="1">';
  for (var j = 0; j < items.length; j++){
    out += '<li>' + items[j] + '</li>';
  }
  out += '</ol>';
  return out;
};

window._attachAutoListBehavior = function(el){
  // === AI 리스트 헬퍼 (depth/style/tab 관리) ===
  if (!window._aiListHelpersReady) {
    window._aiListHelpersReady = true;

    window._circledNumber = function(n){
      var map = {1:'①',2:'②',3:'③',4:'④',5:'⑤',6:'⑥',7:'⑦',8:'⑧',9:'⑨',10:'⑩',11:'⑪',12:'⑫',13:'⑬',14:'⑭',15:'⑮',16:'⑯',17:'⑰',18:'⑱',19:'⑲',20:'⑳'};
      return map[n] || (n + '.');
    };

    window._aiStyleForDepth = function(depth){
      // Word식 반복 순환: dot → paren → circled → dot → paren → circled ...
      if (depth <= 0) return 'dot';
      var styles = ['dot', 'paren', 'circled'];
      return styles[depth % 3];
    };

    window._aiGetClosestLi = function(node, root){
      var cur = node;
      while (cur && cur !== root) {
        if (cur.nodeType === 1 && cur.tagName === 'LI') return cur;
        cur = cur.parentNode;
      }
      return null;
    };

    window._aiGetCurrentLi = function(editor){
      var sel = window.getSelection();
      if (!sel || !sel.rangeCount) return null;
      return window._aiGetClosestLi(sel.getRangeAt(0).startContainer, editor);
    };

    window._aiGetCurrentOl = function(editor){
      var li = window._aiGetCurrentLi(editor);
      if (!li) return null;
      return (li.parentNode && li.parentNode.tagName === 'OL') ? li.parentNode : null;
    };

    window._aiGetListDepth = function(ol, editor){
      var depth = 0;
      var cur = ol.parentNode;
      while (cur && cur !== editor) {
        if (cur.nodeType === 1 && cur.tagName === 'OL') {
          var ovr = cur.dataset.depthOverride;
          if (ovr != null && ovr !== '') {
            var od = parseInt(ovr, 10);
            if (!isNaN(od) && od >= 0) {
              // 조상 ol에 depthOverride가 있으면 그걸 base로 사용
              // 자식 depth = 조상의 effective depth + (그 사이 ol 개수) + 1
              return od + depth + 1;
            }
          }
          depth++;
        }
        cur = cur.parentNode;
      }
      return depth;
    };

    window._aiRefreshList = function(ol){
      if (!ol) return;
      var start = parseInt(ol.dataset.start || '1', 10);
      if (!start || start < 1) start = 1;
      if (ol.classList.contains('ov-num-circled')) {
        var lis = Array.prototype.filter.call(ol.children, function(x){ return x.tagName === 'LI'; });
        lis.forEach(function(li, idx){
          li.setAttribute('data-marker', window._circledNumber(start + idx));
        });
      } else {
        ol.style.counterReset = 'aiitem ' + (start - 1);
      }
    };

    window._aiApplyStart = function(ol, start){
      start = parseInt(start, 10);
      if (!start || start < 1) start = 1;
      ol.dataset.start = String(start);
      window._aiRefreshList(ol);
    };

    window._aiSetOlStyleByDepth = function(ol, depth){
      ol.classList.remove('ov-num-dot', 'ov-num-paren', 'ov-num-circled');
      var effectiveDepth = depth;
      var override = ol.dataset.depthOverride;
      if (override != null && override !== '') {
        var od = parseInt(override, 10);
        if (!isNaN(od) && od >= 0) effectiveDepth = od;
      }
      var style = window._aiStyleForDepth(effectiveDepth);
      ol.classList.add('ov-num', 'ov-num-' + style);
      // depth-override 있으면 margin-left CSS도 반영
      if (override != null && override !== '') {
        ol.style.marginLeft = (effectiveDepth * 24) + 'px';
      } else {
        ol.style.marginLeft = '';
      }
      if (!ol.dataset.start) ol.dataset.start = '1';
      window._aiRefreshList(ol);
    };

    // 지정 li가 ol의 첫 번째면 원본 ol 반환, 아니면 그 li부터 분할한 새 ol 반환
    window._aiSplitOlAtLi = function(ol, li){
      var lis = Array.prototype.filter.call(ol.children, function(x){ return x.tagName === 'LI'; });
      var idx = lis.indexOf(li);
      if (idx <= 0) return ol; // 첫번째 li거나 li 못찾음 → 분할 불필요
      var newOl = document.createElement('ol');
      newOl.className = ol.className;
      newOl.dataset.start = '1';
      newOl.dataset.manualSplit = '1';
      ol.dataset.manualSplit = '1';
      for (var i = idx; i < lis.length; i++) {
        newOl.appendChild(lis[i]);
      }
      if (ol.nextSibling) ol.parentNode.insertBefore(newOl, ol.nextSibling);
      else ol.parentNode.appendChild(newOl);
      return newOl;
    };

    // 편집기 안에서 지정 li보다 앞에 있는 마지막 ol 찾기 (li의 자기 ol 제외)
    window._aiFindPrevOl = function(editor, targetOl){
      var ols = Array.prototype.slice.call(editor.querySelectorAll('ol.ov-num'));
      var myIdx = ols.indexOf(targetOl);
      if (myIdx <= 0) return null;
      return ols[myIdx - 1];
    };

    // li를 리스트에서 빼서 일반 텍스트 div로 변환, 뒤 li들은 새 ol로 분리
    window._aiUnlistCurrentLi = function(editor, targetLi){
      if (!targetLi || targetLi.tagName !== 'LI') return;
      var srcOl = targetLi.parentNode;
      if (!srcOl || srcOl.tagName !== 'OL') return;

      var lis = Array.prototype.filter.call(srcOl.children, function(x){ return x.tagName === 'LI'; });
      var idx = lis.indexOf(targetLi);
      if (idx < 0) return;

      // 새 div 만들어 li 내용 옮김
      var block = document.createElement('div');
      block.innerHTML = targetLi.innerHTML || '<br>';

      // 뒤에 남은 li들을 새 ol로 분리
      var afterOl = null;
      if (idx < lis.length - 1) {
        afterOl = document.createElement('ol');
        afterOl.className = srcOl.className;
        afterOl.dataset.start = '1';
        afterOl.dataset.manualSplit = '1';
        srcOl.dataset.manualSplit = '1';
        for (var i = idx + 1; i < lis.length; i++) {
          afterOl.appendChild(lis[i]);
        }
      }

      // 원본 li 제거
      targetLi.remove();

      // 원본 ol 뒤에 block, afterOl 순서로 삽입
      var insertAfter = srcOl.nextSibling;
      var parent = srcOl.parentNode;
      parent.insertBefore(block, insertAfter);
      if (afterOl) parent.insertBefore(afterOl, insertAfter);

      // 원본 ol이 비었으면 제거
      if (!srcOl.querySelector('li')) srcOl.remove();

      // 스타일/번호 재계산
      if (window._aiRefreshAllLists) window._aiRefreshAllLists(editor);

      // 커서를 block 안으로
      var r = document.createRange();
      r.selectNodeContents(block);
      r.collapse(true);
      var s = window.getSelection();
      s.removeAllRanges();
      s.addRange(r);
    };

    // 일반 텍스트 block의 indent를 조정
    window._aiFindTextBlock = function(node, editor){
      var cur = node;
      while (cur && cur !== editor) {
        if (cur.nodeType === 1) {
          var tag = cur.tagName;
          // li 안이면 텍스트 블록 아님 (리스트 처리)
          if (tag === 'LI') return null;
          if (tag === 'OL' || tag === 'UL') return null;
          // 편집박스 직계 텍스트 블록
          if (cur.parentNode === editor && (tag === 'DIV' || tag === 'P')) return cur;
        }
        cur = cur.parentNode;
      }
      return null;
    };

    window._aiGetTextBlockAtCursor = function(editor){
      var sel = window.getSelection();
      if (!sel || !sel.rangeCount) return null;
      var node = sel.getRangeAt(0).startContainer;
      // 텍스트 노드면 부모부터 탐색
      if (node.nodeType === 3) node = node.parentNode;
      // 편집박스 직계 자식 찾기
      var cur = node;
      while (cur && cur.parentNode !== editor) {
        if (cur === editor) return null;
        cur = cur.parentNode;
      }
      if (!cur) return null;
      if (cur.tagName === 'OL' || cur.tagName === 'UL') return null;
      return cur;
    };

    window._aiIndentTextBlock = function(editor){
      // data-indent 대신 텍스트 앞 공백 2칸 삽입 (저장 안정성 확보)
      document.execCommand('insertText', false, '  ');
      return true;
    };

    window._aiOutdentTextBlock = function(editor){
      // 커서 위치에서 앞 공백 2칸 제거
      var sel = window.getSelection();
      if (!sel.rangeCount) return false;
      var node = sel.anchorNode;
      if (node.nodeType === 3) {  // 텍스트 노드
        var text = node.textContent;
        var offset = sel.anchorOffset;
        // 커서 앞의 공백 2칸 제거
        var before = text.substring(0, offset);
        var after = text.substring(offset);
        if (before.endsWith('  ')) {
          node.textContent = before.slice(0, -2) + after;
          sel.setPosition(node, offset - 2);
        } else if (before.endsWith(' ')) {
          node.textContent = before.slice(0, -1) + after;
          sel.setPosition(node, offset - 1);
        }
      }
      return true;
    };

    // 커서 컨테이너에서 편집박스 직계 텍스트 블록 또는 li 찾기
    window._aiFindNodeForRange = function(node, editor){
      var cur = node;
      if (cur && cur.nodeType === 3) cur = cur.parentNode;
      while (cur && cur !== editor) {
        if (cur.nodeType === 1) {
          if (cur.tagName === 'LI') return { type: 'li', node: cur };
          if (cur.parentNode === editor && (cur.tagName === 'DIV' || cur.tagName === 'P')) {
            return { type: 'block', node: cur };
          }
        }
        cur = cur.parentNode;
      }
      return null;
    };

    // 현재 selection이 걸친 모든 li와 텍스트 블록 수집
    window._aiCollectSelectedNodes = function(editor){
      var sel = window.getSelection();
      if (!sel || !sel.rangeCount) return { lis: [], blocks: [], single: true };
      var range = sel.getRangeAt(0);
      if (range.collapsed) {
        // 단일 커서: 한 노드만 반환
        var one = window._aiFindNodeForRange(range.startContainer, editor);
        return {
          lis: one && one.type === 'li' ? [one.node] : [],
          blocks: one && one.type === 'block' ? [one.node] : [],
          single: true
        };
      }

      var lis = [];
      var blocks = [];
      var seen = new Set();

      // range 시작과 끝의 노드부터 감지
      var startInfo = window._aiFindNodeForRange(range.startContainer, editor);
      var endInfo = window._aiFindNodeForRange(range.endContainer, editor);

      function addNode(info){
        if (!info) return;
        if (seen.has(info.node)) return;
        seen.add(info.node);
        if (info.type === 'li') lis.push(info.node);
        else blocks.push(info.node);
      }

      // 편집박스 안의 모든 li, 편집박스 직계 텍스트 블록 순회하며 range와 교차하는지 확인
      var allLis = Array.prototype.slice.call(editor.querySelectorAll('li'));
      var allBlocks = Array.prototype.filter.call(editor.children, function(x){
        return x.nodeType === 1 && (x.tagName === 'DIV' || x.tagName === 'P');
      });

      function intersects(el){
        try {
          var r = document.createRange();
          r.selectNodeContents(el);
          // range 시작이 el 끝 이후이거나 range 끝이 el 시작 이전이면 교차 없음
          if (range.compareBoundaryPoints(Range.START_TO_END, r) < 0) return false;
          if (range.compareBoundaryPoints(Range.END_TO_START, r) > 0) return false;
          return true;
        } catch(e) { return false; }
      }

      allLis.forEach(function(li){ if (intersects(li)) addNode({ type: 'li', node: li }); });
      allBlocks.forEach(function(b){ if (intersects(b)) addNode({ type: 'block', node: b }); });

      // 시작/끝 노드도 안전하게 추가
      addNode(startInfo);
      addNode(endInfo);

      return { lis: lis, blocks: blocks, single: false };
    };

    window._aiSplitTextByNewlines = function(editor){
      if (!editor) return;
      if (editor.__aiLineSplit) return;
      editor.__aiLineSplit = true;
      var LF = String.fromCharCode(10);
      var children = Array.prototype.slice.call(editor.children);
      children.forEach(function(child){
        var skipTags = ['OL', 'UL', 'TABLE', 'BLOCKQUOTE', 'PRE'];
        if (skipTags.indexOf(child.tagName) >= 0) return;
        if (child.tagName !== 'DIV' && child.tagName !== 'P') return;
        var text = child.textContent || '';
        if (text.indexOf(LF) < 0) return;
        var html = child.innerHTML;
        var parts = html.split(LF);
        if (parts.length <= 1) return;
        var origIndent = child.dataset.indent || '';
        var origStyle = child.getAttribute('style') || '';
        var frag = document.createDocumentFragment();
        parts.forEach(function(part){
          var newDiv = document.createElement('div');
          newDiv.innerHTML = part;
          if (!newDiv.textContent.trim() && !newDiv.querySelector('br')) {
            newDiv.innerHTML = '<br>';
          }
          if (origIndent) newDiv.dataset.indent = origIndent;
          if (origStyle) newDiv.setAttribute('style', origStyle);
          frag.appendChild(newDiv);
        });
        child.parentNode.insertBefore(frag, child);
        child.remove();
      });
    };

    window._aiRefreshAllLists = function(root){
      // legacy 클래스(ov-list-paren) 흡수
      root.querySelectorAll('ol.ov-list-paren').forEach(function(ol){
        ol.classList.remove('ov-list-paren');
        if (!ol.classList.contains('ov-num')) ol.classList.add('ov-num');
        if (!ol.dataset.start) ol.dataset.start = '1';
      });
      root.querySelectorAll('ol.ov-num').forEach(function(ol){
        var depth = window._aiGetListDepth(ol, root);
        window._aiSetOlStyleByDepth(ol, depth);
      });
    };

    window._aiEnsureChildOl = function(parentLi, style){
      var child = null;
      for (var i = 0; i < parentLi.children.length; i++) {
        if (parentLi.children[i].tagName === 'OL') { child = parentLi.children[i]; break; }
      }
      if (!child) {
        child = document.createElement('ol');
        child.className = 'ov-num ov-num-' + style;
        child.dataset.start = '1';
        parentLi.appendChild(child);
      }
      return child;
    };

    window._aiIndentCurrentLi = function(editor){
      var li = window._aiGetCurrentLi(editor);
      if (!li) return false;
      var parentOl = li.parentNode;
      var prevLi = li.previousElementSibling;
      if (prevLi && prevLi.tagName === 'LI') {
        // 형제 li 있음 → 자식 ol로 이동
        var nextDepth = window._aiGetListDepth(parentOl, editor) + 1;
        var childOl = window._aiEnsureChildOl(prevLi, window._aiStyleForDepth(nextDepth));
        childOl.appendChild(li);
        if (!parentOl.querySelector('li')) parentOl.remove();
      } else {
        // 형제 li 없음 → 자기 ol depth override 증가
        var siblings = Array.prototype.filter.call(parentOl.children, function(x){ return x.tagName === 'LI'; });
        var targetOl = parentOl;
        if (siblings.length > 1) {
          var isolatedOl = document.createElement('ol');
          isolatedOl.className = parentOl.className;
          isolatedOl.dataset.start = '1';
          isolatedOl.dataset.manualSplit = '1';
          parentOl.dataset.manualSplit = '1';
          isolatedOl.appendChild(li);
          if (parentOl.nextSibling) parentOl.parentNode.insertBefore(isolatedOl, parentOl.nextSibling);
          else parentOl.parentNode.appendChild(isolatedOl);
          targetOl = isolatedOl;
        }
        var currentDepth = window._aiGetListDepth(targetOl, editor);
        var override = targetOl.dataset.depthOverride;
        var currentEffective = (override != null && override !== '') ? parseInt(override, 10) : currentDepth;
        var newDepth = currentEffective + 1;
        if (newDepth <= 5) targetOl.dataset.depthOverride = String(newDepth);
      }
      window._aiRefreshAllLists(editor);
      var r = document.createRange();
      r.selectNodeContents(li);
      r.collapse(false);
      var s = window.getSelection();
      s.removeAllRanges();
      s.addRange(r);
      return true;
    };

    window._aiOutdentCurrentLi = function(editor){
      var li = window._aiGetCurrentLi(editor);
      if (!li) return false;
      var parentOl = li.parentNode;
      if (!parentOl) return false;

      // depth-override가 있으면 그것부터 감소
      var override = parentOl.dataset.depthOverride;
      if (override != null && override !== '') {
        var currentOverride = parseInt(override, 10);
        var currentDepth = window._aiGetListDepth(parentOl, editor);
        if (currentOverride > currentDepth) {
          // override 감소
          var newOverride = currentOverride - 1;
          if (newOverride <= currentDepth) {
            delete parentOl.dataset.depthOverride;
          } else {
            parentOl.dataset.depthOverride = String(newOverride);
          }
          window._aiRefreshAllLists(editor);
          return true;
        }
      }

      // 일반 outdent (nesting 기반)
      var hostLi = parentOl.parentNode;
      if (!hostLi || hostLi.tagName !== 'LI') return false;
      var grandOl = hostLi.parentNode;
      if (!grandOl || grandOl.tagName !== 'OL') return false;
      if (hostLi.nextSibling) grandOl.insertBefore(li, hostLi.nextSibling);
      else grandOl.appendChild(li);
      if (!parentOl.querySelector('li')) parentOl.remove();
      window._aiRefreshAllLists(editor);
      var r = document.createRange();
      r.selectNodeContents(li);
      r.collapse(false);
      var s = window.getSelection();
      s.removeAllRanges();
      s.addRange(r);
      return true;
    };

    window._aiContinueOl = function(editor, targetLi){
      if (!targetLi || targetLi.tagName !== 'LI') return;
      var srcOl = targetLi.parentNode;
      if (!srcOl || srcOl.tagName !== 'OL') return;
      var workOl = window._aiSplitOlAtLi(srcOl, targetLi);
      var prev = window._aiFindPrevOl(editor, workOl);
      if (!prev) { alert('앞에 이어붙일 리스트가 없어'); if (window._aiRefreshAllLists) window._aiRefreshAllLists(editor); return; }
      var prevStart = parseInt(prev.dataset.start || '1', 10) || 1;
      var prevCount = Array.prototype.filter.call(prev.children, function(x){ return x.tagName === 'LI'; }).length;
      window._aiApplyStart(workOl, prevStart + prevCount);
      if (window._aiRefreshAllLists) window._aiRefreshAllLists(editor);
    };

    window._aiRestartOl = function(editor, targetLi){
      if (!targetLi || targetLi.tagName !== 'LI') return;
      var srcOl = targetLi.parentNode;
      if (!srcOl || srcOl.tagName !== 'OL') return;
      var workOl = window._aiSplitOlAtLi(srcOl, targetLi);
      window._aiApplyStart(workOl, 1);
      if (window._aiRefreshAllLists) window._aiRefreshAllLists(editor);
    };

    window._aiSetStartOl = function(editor, targetLi){
      if (!targetLi || targetLi.tagName !== 'LI') return;
      var srcOl = targetLi.parentNode;
      if (!srcOl || srcOl.tagName !== 'OL') return;
      // 현재 li의 실제 번호 (data-start + index)
      var lis = Array.prototype.filter.call(srcOl.children, function(x){ return x.tagName === 'LI'; });
      var idx = lis.indexOf(targetLi);
      var srcStart = parseInt(srcOl.dataset.start || '1', 10) || 1;
      var current = srcStart + idx;
      var v = prompt('이 항목의 번호', String(current));
      if (v == null) return;
      var n = parseInt(String(v).trim(), 10);
      if (!n || n < 1) { alert('1 이상의 숫자를 입력해'); return; }
      var workOl = window._aiSplitOlAtLi(srcOl, targetLi);
      window._aiApplyStart(workOl, n);
      if (window._aiRefreshAllLists) window._aiRefreshAllLists(editor);
    };
  }

  if (!el) return;

  // observer는 attached 플래그와 독립적으로 항상 부착 시도
  if (!el.__aiListObserver) {
    function _mergeAdjacentAiListsGlobal(root){
      if (!root) return;
      var ols = root.querySelectorAll('ol.ov-num');
      for (var i = 0; i < ols.length; i++){
        var ol = ols[i];
        var prev = ol.previousElementSibling;
        while (prev && prev.tagName === 'OL' && prev.classList.contains('ov-num') && prev.dataset.manualSplit !== '1' && ol.dataset.manualSplit !== '1') {
          while (ol.firstChild) prev.appendChild(ol.firstChild);
          ol.remove();
          ol = prev;
          prev = ol.previousElementSibling;
        }
        var next = ol.nextElementSibling;
        while (next && next.tagName === 'OL' && next.classList.contains('ov-num') && next.dataset.manualSplit !== '1' && ol.dataset.manualSplit !== '1') {
          while (next.firstChild) ol.appendChild(next.firstChild);
          next.remove();
          next = ol.nextElementSibling;
        }
        ol.style.counterReset = '';
      }
    }
    var _mergeTimer = null;
    el.__aiListObserver = new MutationObserver(function(){
      if (_mergeTimer) return;
      _mergeTimer = setTimeout(function(){
        _mergeTimer = null;
        _mergeAdjacentAiListsGlobal(el);
      }, 50);
    });
    el.__aiListObserver.observe(el, { childList: true, subtree: true });
  }

  if (el.__autoListAttached) return;
  el.__autoListAttached = true;

  var NL = String.fromCharCode(10);

  function _getCurrentLineText(){
    var sel = window.getSelection();
    if (!sel.rangeCount) return { line: '', range: null };
    var range = sel.getRangeAt(0);
    var node = range.startContainer;
    if (node.nodeType !== 3) return { line: '', range: range };
    var text = node.textContent || '';
    var leftText = text.substring(0, range.startOffset);
    var lastBr = leftText.lastIndexOf(NL);
    var lineStart = lastBr === -1 ? 0 : lastBr + 1;
    var line = leftText.substring(lineStart);
    return { line: line, range: range, node: node, lineStart: lineStart };
  }

  function _detectExistingListLine(line){
    var m = line.match(/^(\s*)(\d+)([\)\.])\s+(.*)$/);
    if (m) return { type: 'num', indent: m[1], num: parseInt(m[2],10), sep: m[3], content: m[4] };
    m = line.match(/^(\s*)([-*\u00B7])\s+(.*)$/);
    if (m) return { type: 'bullet', indent: m[1], bullet: m[2], content: m[3] };
    return null;
  }

  el.addEventListener('keydown', function(e){
    if (e.key !== 'Enter' || e.shiftKey) return;
    var info = _getCurrentLineText();
    var existing = _detectExistingListLine(info.line);
    if (!existing) return;

    if (!existing.content.trim()) {
      e.preventDefault();
      var sel = window.getSelection();
      if (sel.rangeCount && info.node) {
        var range = document.createRange();
        range.setStart(info.node, info.lineStart);
        range.setEnd(info.node, sel.getRangeAt(0).startOffset);
        range.deleteContents();
      }
      document.execCommand('insertLineBreak');
      return;
    }

    e.preventDefault();
    var nextPrefix;
    if (existing.type === 'num') {
      nextPrefix = existing.indent + (existing.num + 1) + existing.sep + ' ';
    } else {
      nextPrefix = existing.indent + existing.bullet + ' ';
    }
    document.execCommand('insertLineBreak');
    document.execCommand('insertText', false, nextPrefix);
  });
};
</script>

<script>
(function(){
  const PAGE_LABEL = {
    home: '홈 대시보드',
    report: '보고',
    production: '생산',
    inbound: '입고',
    outbound: '출하',
  };

  document.querySelectorAll('.nav-item').forEach(function(el){
    el.addEventListener('click', function(){
      document.querySelectorAll('.nav-item').forEach(function(x){ x.classList.remove('active'); });
      el.classList.add('active');
      const pg = el.getAttribute('data-page');
      document.getElementById('v2-crumb-page').textContent = PAGE_LABEL[pg] || pg;
      renderPage(pg);
    });
  });

  function renderPage(pg){
    const c = document.getElementById('v2-content');
    if (pg === 'report') {
      c.innerHTML = renderReportPageHTML();
      loadAdminV2Reports();
      bindAdminV2Upload();
    } else {
      const icons = { home:'🏠', production:'🏭', inbound:'📥', outbound:'📤' };
      c.innerHTML = '<div class="placeholder"><div class="big">'+(icons[pg]||'✦')+'</div><div class="title">'+(PAGE_LABEL[pg]||pg)+'</div><div class="sub">개발 중 — 곧 오픈됩니다</div></div>';
    }
  }

  // ============================================================
  // Admin v2 · Report Page (KPI + Upload + List)
  // ============================================================
  window.renderReportPageHTML = function(){
    return `
      <style>
        .ov-page { padding: 0; }
        .ov-kpi-grid {
          display: grid;
          grid-template-columns: repeat(4, minmax(0, 1fr));
          gap: 16px;
          margin-bottom: 20px;
        }
        .ov-kpi-card {
          background: #fff;
          border: 1px solid #E6EBF2;
          border-radius: 20px;
          padding: 18px;
          box-shadow: 0 6px 18px rgba(15,44,89,0.04);
        }
        .ov-kpi-label { font-size: 13px; color: #6E7785; font-weight: 600; margin-bottom: 10px; }
        .ov-kpi-value { font-size: 28px; color: #0F2C59; font-weight: 800; line-height: 1.1; margin-bottom: 6px; }
        .ov-kpi-sub { font-size: 13px; color: #7A8595; font-weight: 600; }

        .ov-upload-card {
          background: #fff;
          border: 1px solid #E6EBF2;
          border-radius: 24px;
          padding: 24px;
          margin-bottom: 20px;
          box-shadow: 0 6px 18px rgba(15,44,89,0.04);
        }
        .ov-upload-drop {
          border: 2px dashed #BFD2EA;
          background: #F7FBFF;
          border-radius: 22px;
          padding: 34px 20px;
          text-align: center;
        }
        .ov-upload-icon { font-size: 34px; margin-bottom: 10px; }
        .ov-upload-title { font-size: 22px; color: #173B72; font-weight: 800; margin-bottom: 8px; }
        .ov-upload-desc { font-size: 14px; color: #6F7B8C; margin-bottom: 16px; }
        .ov-upload-meta {
          display: inline-block; margin-top: 14px;
          background: #EEF4FB; color: #2E5B94;
          padding: 8px 12px; border-radius: 999px;
          font-size: 13px; font-weight: 700;
        }

        .ov-list-card {
          background: #fff;
          border: 1px solid #E6EBF2;
          border-radius: 24px;
          padding: 20px;
          box-shadow: 0 6px 18px rgba(15,44,89,0.04);
        }
        .ov-list-head { display:flex; justify-content:space-between; align-items:center; margin-bottom: 16px; }
        .ov-new-btn {
          background: #0F2C59; color: #fff; border: none; border-radius: 12px;
          padding: 10px 16px; font-size: 14px; font-weight: 800; cursor: pointer;
        }
        .ov-new-btn:hover { background: #173B72; }
        .ov-modal-mask {
          position: fixed; inset: 0; background: rgba(15,44,89,0.35);
          display: none; align-items: center; justify-content: center; z-index: 9999;
        }
        .ov-modal-mask.open { display: flex; }
        .ov-modal {
          background: #fff; border-radius: 20px; padding: 28px 26px; width: 420px; max-width: 92vw;
          box-shadow: 0 20px 60px rgba(15,44,89,0.25);
        }
        .ov-modal h3 { margin: 0 0 16px 0; font-size: 20px; color: #0F2C59; font-weight: 800; }
        .ov-modal label { display: block; font-size: 13px; color: #4A5568; font-weight: 700; margin: 10px 0 6px; }
        .ov-modal input {
          width: 100%; box-sizing: border-box; padding: 10px 12px;
          border: 1px solid #D0D9E5; border-radius: 10px; font-size: 14px;
        }
        .ov-modal input:focus { outline: none; border-color: #4A6FA5; }
        .ov-modal-actions { display: flex; justify-content: flex-end; gap: 8px; margin-top: 20px; }
        .ov-modal-cancel {
          background: #F4F6FA; color: #4A5568; border: 1px solid #E6EBF2;
          border-radius: 10px; padding: 10px 16px; font-size: 14px; font-weight: 700; cursor: pointer;
        }
        .ov-modal-confirm {
          background: #0F2C59; color: #fff; border: none;
          border-radius: 10px; padding: 10px 18px; font-size: 14px; font-weight: 800; cursor: pointer;
        }
        .ov-manual-badge {
          display: inline-block; background: #FFF4EC; color: #B4380F;
          padding: 2px 8px; border-radius: 999px; font-size: 11px; font-weight: 700; margin-left: 8px;
        }
        .ov-list-title { font-size: 20px; color: #173B72; font-weight: 800; }
        .ov-list-sub { font-size: 13px; color: #748092; font-weight: 600; }

        .ov-report-list { display: flex; flex-direction: column; gap: 14px; }
        .ov-report-item.is-hidden { opacity: 0.45; background: #F4F6FA; }
        .ov-report-item.is-hidden .ov-report-project::after {
          content: " · 숨김"; color: #B4380F; font-size: 12px; font-weight: 700;
        }
        .ov-hide-btn {
          background: #F4F6FA; color: #4A5568;
          border: 1px solid #E6EBF2; border-radius: 999px;
          padding: 6px 12px; font-size: 12px; font-weight: 700;
          cursor: pointer; margin-right: 8px;
        }
        .ov-hide-btn:hover { background: #E6EBF2; }
        .ov-hide-btn.is-active { background: #FFF4EC; color: #B4380F; border-color: #F5C9AE; }
        .ov-del-btn {
          background: #FEE2E2; color: #DC2626;
          border: 1px solid #FCA5A5; border-radius: 999px;
          padding: 6px 12px; font-size: 12px; font-weight: 700;
          cursor: pointer;
        }
        .ov-del-btn:hover { background: #FCA5A5; color: #ffffff; border-color: #DC2626; }
        .ov-side-top-actions { display: flex; gap: 8px; align-items: center; }
        .ov-report-item {
          display: grid;
          grid-template-columns: 6px minmax(0, 1fr) 180px;
          gap: 18px;
          align-items: center;
          border: 1px solid #E8EDF4;
          border-radius: 22px;
          overflow: hidden;
          background: #FFFFFF;
          min-height: 112px;
        }
        .ov-report-bar {
          align-self: stretch;
          min-height: 100%;
          width: 6px;
        }
        .ov-report-bar.orange { background: #F59E0B; }
        .ov-report-bar.red { background: #EF4444; }
        .ov-report-bar.green { background: #10B981; }
        .ov-report-main {
          padding: 18px 0;
          min-width: 0;
        }
        .ov-report-project {
          font-size: 20px;
          line-height: 1.2;
          color: #12325F;
          font-weight: 800;
          margin-bottom: 8px;
          letter-spacing: -0.2px;
        }
        .ov-report-file {
          font-size: 14px;
          line-height: 1.35;
          color: #7A8595;
          margin-bottom: 8px;
          white-space: nowrap;
          overflow: hidden;
          text-overflow: ellipsis;
        }
        .ov-report-preview {
          font-size: 14px;
          line-height: 1.4;
          color: #667085;
          font-weight: 500;
        }
        .ov-report-side {
          display: flex;
          flex-direction: column;
          justify-content: center;
          align-items: flex-end;
          gap: 14px;
          padding: 18px 18px 18px 0;
        }
        .ov-badge {
          display:inline-flex;
          align-items:center;
          justify-content:center;
          border-radius: 999px;
          padding: 9px 16px;
          font-size: 13px;
          font-weight: 800;
          min-width: 108px;
        }
        .ov-badge.blue { background: #E7F0FB; color: #23538C; }
        .ov-badge.amber { background: #FFF3D9; color: #9A6700; }
        .ov-badge.green { background: #E7F8F0; color: #117A52; }
        .ov-open-btn {
          border: 0;
          border-radius: 12px;
          background: #163A70;
          color: #fff;
          padding: 11px 18px;
          font-size: 15px;
          font-weight: 800;
          cursor: pointer;
          width: 118px;
          height: 44px;
        }
        .ov-load-more { display:flex; justify-content:center; margin-top: 16px; }
        .ov-load-more button {
          border: 1px solid #DCE4EF; background: #fff; color: #35527C;
          border-radius: 999px; padding: 12px 18px;
          font-size: 14px; font-weight: 700; cursor: pointer;
        }
        @media (max-width: 1200px) { .ov-kpi-grid { grid-template-columns: repeat(2, minmax(0, 1fr)); } }
        @media (max-width: 760px) {
          .ov-kpi-grid { grid-template-columns: 1fr; }
          .ov-report-item { grid-template-columns: 6px 1fr; }
          .ov-report-side { grid-column: 2; align-items: flex-start; padding: 0 16px 16px 0; }
        }
      </style>

      <div class="ov-page">
        <div class="ov-kpi-grid">
          <div class="ov-kpi-card">
            <div class="ov-kpi-label">이번 주</div>
            <div class="ov-kpi-value" id="ov-kpi-week">-</div>
            <div class="ov-kpi-sub">업로드된 보고 수</div>
          </div>
          <div class="ov-kpi-card">
            <div class="ov-kpi-label">AI 처리중</div>
            <div class="ov-kpi-value" id="ov-kpi-ai">-</div>
            <div class="ov-kpi-sub">현재 처리 대기/진행</div>
          </div>
          <div class="ov-kpi-card">
            <div class="ov-kpi-label">검토 대기</div>
            <div class="ov-kpi-value" id="ov-kpi-review">-</div>
            <div class="ov-kpi-sub">관리자 확인 필요</div>
          </div>
          <div class="ov-kpi-card">
            <div class="ov-kpi-label">발행 완료</div>
            <div class="ov-kpi-value" id="ov-kpi-published">-</div>
            <div class="ov-kpi-sub">앱 반영 가능 상태</div>
          </div>
        </div>

        <div class="ov-upload-card">
          <div class="ov-upload-drop" id="ov-upload-drop">
            <div class="ov-upload-icon">📤</div>
            <div class="ov-upload-title">보고 PPT를 여기로 끌어다 놓으세요</div>
            <div class="ov-upload-desc">또는 버튼으로 파일 선택 · .pptx / .ppt · 최대 50MB</div>
            <input type="file" id="ov-upload-input" accept=".pptx,.ppt" style="display:none;" />
            <button class="btn-primary" type="button" id="ov-upload-btn">파일 선택</button>
            <div class="ov-upload-meta">✨ 표 · 목록 · 이미지 자동 파싱</div>
            <div id="ov-upload-status" style="margin-top:14px;font-size:13px;font-weight:600;"></div>
          </div>
        </div>

        <div class="ov-list-card">
          <div class="ov-list-head">
            <div>
              <div class="ov-list-title">최근 보고</div>
              <div class="ov-list-sub">가장 최근 등록된 보고를 확인하세요</div>
            </div>
            <button class="ov-new-btn" type="button" id="ov-new-project-btn">+ 새 프로젝트</button>
          </div>
          <div class="ov-report-list" id="ov-report-list">
            <div style="color:#9CA3AF;font-size:14px;padding:20px;text-align:center;">불러오는 중...</div>
          </div>
          <div class="ov-load-more">
            <button type="button" id="ov-load-more-btn">보고 더 보기</button>
          </div>
        </div>
        <div class="ov-modal-mask" id="ov-newproj-mask">
          <div class="ov-modal" style="max-width:640px;width:90vw;">
            <h3>+ 새 프로젝트</h3>
            <div style="margin-top:8px;font-size:13px;color:#374151;font-weight:600;">등록된 프로젝트 <span id="ov-np-current-div" style="color:#6B7280;font-weight:400;font-size:12px;"></span></div>
            <div style="font-size:12px;color:#9CA3AF;margin-top:2px;">클릭하면 바로 생성됩니다</div>
            <div id="ov-np-project-buttons" style="display:flex;flex-wrap:wrap;gap:8px;margin-top:8px;min-height:44px;padding:8px;background:#F9FAFB;border-radius:8px;"></div>
            <div style="margin-top:16px;display:flex;align-items:center;gap:8px;">
              <div style="flex:1;height:1px;background:#E5E7EB;"></div>
              <span style="color:#9CA3AF;font-size:12px;">또는 새 이름으로</span>
              <div style="flex:1;height:1px;background:#E5E7EB;"></div>
            </div>
            <label for="ov-np-name" style="margin-top:8px;">프로젝트명</label>
            <input type="text" id="ov-np-name" name="np-name" placeholder="예: Chamber" autocomplete="off" />
            <div id="ov-np-suggest" style="margin-top:8px;min-height:0;"></div>
            <div style="font-size:12px;color:#7A8595;margin-top:8px;">주차는 현재 주차로 자동 설정됩니다.</div>
            <div class="ov-modal-actions">
              <button class="ov-modal-cancel" type="button" id="ov-np-cancel">취소</button>
              <button class="ov-modal-confirm" type="button" id="ov-np-confirm">생성</button>
            </div>
          </div>
        </div>
        <div class="ov-modal-mask" id="ov-delconf-mask">
          <div class="ov-modal" style="max-width:480px;width:90vw;">
            <h3 style="color:#DC2626;">프로젝트 삭제</h3>
            <div id="ov-delconf-body" style="margin-top:12px;font-size:14px;color:#374151;line-height:1.6;"></div>
            <div style="margin-top:12px;padding:10px 12px;background:#FEF2F2;border:1px solid #FCA5A5;border-radius:6px;font-size:13px;color:#991B1B;">⚠️ 되돌릴 수 없습니다. 관련 카드와 PPT 파일 모두 제거됩니다.</div>
            <div id="ov-delconf-status" style="margin-top:8px;font-size:13px;min-height:18px;"></div>
            <div class="ov-modal-actions">
              <button class="ov-modal-cancel" type="button" id="ov-delconf-cancel">취소</button>
              <button class="ov-modal-confirm" type="button" id="ov-delconf-confirm" style="background:#DC2626;color:#fff;">삭제</button>
            </div>
          </div>
        </div>
      </div>
    `;
  };

  window.openNewProjectModal = async function(){
    var mask = document.getElementById('ov-newproj-mask');
    var nameEl = document.getElementById('ov-np-name');
    var sidebarSel = document.getElementById('v2-division-select');
    var divisionId = sidebarSel ? sidebarSel.value : 'semiconductor';
    if (!mask) return;
    if (nameEl) nameEl.value = '';
    if (window._npClearSuggest) window._npClearSuggest();
    mask.classList.add('open');
    var divLabel = '';
    if (sidebarSel) {
      var opt = sidebarSel.options[sidebarSel.selectedIndex];
      if (opt) divLabel = opt.textContent || '';
    }
    var curDivEl = document.getElementById('ov-np-current-div');
    if (curDivEl) curDivEl.textContent = divLabel ? '(' + divLabel + ')' : '';
    await window._npLoadProjectButtons(divisionId);
    setTimeout(function(){ if (nameEl) nameEl.focus(); }, 50);
  };

  window._npLoadProjectButtons = async function(divisionId){
    var box = document.getElementById('ov-np-project-buttons');
    if (!box) return;
    while (box.firstChild) box.removeChild(box.firstChild);
    if (!divisionId) {
      var hint = document.createElement('div');
      hint.style.color = '#9CA3AF';
      hint.style.fontSize = '13px';
      hint.textContent = '사업부를 선택하세요';
      box.appendChild(hint);
      return;
    }
    var loading = document.createElement('div');
    loading.style.color = '#9CA3AF';
    loading.style.fontSize = '13px';
    loading.textContent = '불러오는 중...';
    box.appendChild(loading);
    try {
      var r = await fetch('/admin/config/projects?division_id=' + encodeURIComponent(divisionId), { credentials: 'same-origin' });
      var d = await r.json();
      var items = (d && d.projects) || [];
      while (box.firstChild) box.removeChild(box.firstChild);
      if (!items.length) {
        var empty = document.createElement('div');
        empty.style.color = '#9CA3AF';
        empty.style.fontSize = '13px';
        empty.textContent = '등록된 프로젝트가 없습니다';
        box.appendChild(empty);
        return;
      }
      items.forEach(function(p){
        var label = p.label || p.id;
        var btn = document.createElement('button');
        btn.type = 'button';
        btn.textContent = label;
        btn.style.padding = '8px 14px';
        btn.style.border = '1px solid #D1D5DB';
        btn.style.borderRadius = '6px';
        btn.style.background = '#ffffff';
        btn.style.fontSize = '13px';
        btn.style.cursor = 'pointer';
        btn.addEventListener('mouseenter', function(){ btn.style.background = '#F3F4F6'; btn.style.borderColor = '#2563EB'; });
        btn.addEventListener('mouseleave', function(){ btn.style.background = '#ffffff'; btn.style.borderColor = '#D1D5DB'; });
        btn.addEventListener('click', function(){
          var nameEl = document.getElementById('ov-np-name');
          if (nameEl) nameEl.value = label;
          window.submitNewProject();
        });
        box.appendChild(btn);
      });
    } catch(e) {
      while (box.firstChild) box.removeChild(box.firstChild);
      var err = document.createElement('div');
      err.style.color = '#EF4444';
      err.style.fontSize = '13px';
      err.textContent = '로드 실패';
      box.appendChild(err);
    }
  };

  window._pendingDeleteDocId = null;
  window.openDeleteConfirm = function(docId, title){
    var mask = document.getElementById('ov-delconf-mask');
    var body = document.getElementById('ov-delconf-body');
    var status = document.getElementById('ov-delconf-status');
    if (!mask) return;
    window._pendingDeleteDocId = docId;
    if (body) {
      while (body.firstChild) body.removeChild(body.firstChild);
      var p1 = document.createElement('div');
      p1.textContent = '다음 항목을 삭제하시겠습니까?';
      var p2 = document.createElement('div');
      p2.style.marginTop = '6px';
      p2.style.fontWeight = '600';
      p2.style.color = '#1F2937';
      p2.textContent = title || docId;
      body.appendChild(p1);
      body.appendChild(p2);
    }
    if (status) { status.textContent = ''; status.style.color = ''; }
    mask.classList.add('open');
  };
  window.closeDeleteConfirm = function(){
    var mask = document.getElementById('ov-delconf-mask');
    if (mask) mask.classList.remove('open');
    window._pendingDeleteDocId = null;
  };
  window.executeDelete = async function(){
    var docId = window._pendingDeleteDocId;
    var status = document.getElementById('ov-delconf-status');
    var confirmBtn = document.getElementById('ov-delconf-confirm');
    if (!docId) return;
    if (status) { status.textContent = '삭제 중...'; status.style.color = '#6B7280'; }
    if (confirmBtn) { confirmBtn.disabled = true; confirmBtn.style.opacity = '0.6'; }
    try {
      var res = await fetch('/admin/reports/' + encodeURIComponent(docId), {
        method: 'DELETE',
        credentials: 'same-origin',
      });
      if (!res.ok) {
        var detail = 'HTTP ' + res.status;
        try {
          var err = await res.json();
          if (err && err.detail) detail = err.detail;
        } catch(_) {}
        if (status) { status.textContent = '❌ 삭제 실패: ' + detail; status.style.color = '#DC2626'; }
        return;
      }
      if (status) { status.textContent = '✅ 삭제 완료'; status.style.color = '#059669'; }
      setTimeout(function(){
        window.closeDeleteConfirm();
        if (window.loadAdminV2Reports) window.loadAdminV2Reports();
      }, 500);
    } catch(e) {
      if (status) { status.textContent = '❌ 네트워크 오류: ' + (e.message || e); status.style.color = '#DC2626'; }
    } finally {
      if (confirmBtn) { confirmBtn.disabled = false; confirmBtn.style.opacity = '1'; }
    }
  };

    var _npSuggestTimer = null;
  window._npClearSuggest = function(){
    var box = document.getElementById('ov-np-suggest');
    if (!box) return;
    while (box.firstChild) box.removeChild(box.firstChild);
  };
  window._npBuildSuggestBox = function(kind, msg, btnLabel, onAccept, onReject){
    var box = document.getElementById('ov-np-suggest');
    if (!box) return;
    while (box.firstChild) box.removeChild(box.firstChild);
    var wrap = document.createElement('div');
    wrap.style.padding = '10px 12px';
    wrap.style.borderRadius = '6px';
    wrap.style.fontSize = '13px';
    wrap.style.display = 'flex';
    wrap.style.flexWrap = 'wrap';
    wrap.style.alignItems = 'center';
    wrap.style.gap = '8px';
    if (kind === 'warn') {
      wrap.style.background = '#FEF3C7';
      wrap.style.border = '1px solid #F59E0B';
    } else {
      wrap.style.background = '#DBEAFE';
      wrap.style.border = '1px solid #3B82F6';
    }
    var txt = document.createElement('span');
    txt.textContent = msg;
    txt.style.flex = '1';
    wrap.appendChild(txt);
    var acceptBtn = document.createElement('button');
    acceptBtn.type = 'button';
    acceptBtn.textContent = btnLabel;
    acceptBtn.style.padding = '4px 10px';
    acceptBtn.style.border = 'none';
    acceptBtn.style.borderRadius = '4px';
    acceptBtn.style.fontSize = '12px';
    acceptBtn.style.cursor = 'pointer';
    acceptBtn.style.color = '#ffffff';
    acceptBtn.style.background = (kind === 'warn') ? '#F59E0B' : '#3B82F6';
    acceptBtn.addEventListener('click', onAccept);
    wrap.appendChild(acceptBtn);
    if (onReject) {
      var rejectBtn = document.createElement('button');
      rejectBtn.type = 'button';
      rejectBtn.textContent = '아니오';
      rejectBtn.style.padding = '4px 10px';
      rejectBtn.style.border = '1px solid #D1D5DB';
      rejectBtn.style.borderRadius = '4px';
      rejectBtn.style.fontSize = '12px';
      rejectBtn.style.cursor = 'pointer';
      rejectBtn.style.background = '#ffffff';
      rejectBtn.style.color = '#374151';
      rejectBtn.addEventListener('click', onReject);
      wrap.appendChild(rejectBtn);
    }
    box.appendChild(wrap);
  };
  window._npBuildLoadingBox = function(msg){
    var box = document.getElementById('ov-np-suggest');
    if (!box) return;
    while (box.firstChild) box.removeChild(box.firstChild);
    var wrap = document.createElement('div');
    wrap.style.padding = '10px 12px';
    wrap.style.borderRadius = '6px';
    wrap.style.fontSize = '13px';
    wrap.style.background = '#F3F4F6';
    wrap.style.border = '1px solid #D1D5DB';
    wrap.style.color = '#6B7280';
    wrap.textContent = msg || 'AI로 확인 중...';
    box.appendChild(wrap);
  };

  window._npBuildAiSuggestBox = function(resolved, confidence, query, sidebarDivId){
    var box = document.getElementById('ov-np-suggest');
    if (!box) return;
    while (box.firstChild) box.removeChild(box.firstChild);
    var wrap = document.createElement('div');
    wrap.style.padding = '10px 12px';
    wrap.style.borderRadius = '6px';
    wrap.style.fontSize = '13px';
    wrap.style.display = 'flex';
    wrap.style.flexWrap = 'wrap';
    wrap.style.alignItems = 'center';
    wrap.style.gap = '8px';
    wrap.style.background = '#EDE9FE';
    wrap.style.border = '1px solid #8B5CF6';
    var badge = document.createElement('span');
    badge.textContent = '🤖 AI';
    badge.style.fontSize = '11px';
    badge.style.padding = '2px 6px';
    badge.style.background = '#8B5CF6';
    badge.style.color = '#ffffff';
    badge.style.borderRadius = '4px';
    badge.style.fontWeight = '700';
    wrap.appendChild(badge);
    var txt = document.createElement('span');
    txt.textContent = '혹시 ' + resolved.label + ' 프로젝트인가요? (신뢰도 ' + Math.round(confidence * 100) + '%)';
    txt.style.flex = '1';
    wrap.appendChild(txt);
    var acceptBtn = document.createElement('button');
    acceptBtn.type = 'button';
    acceptBtn.textContent = '네, ' + resolved.label;
    acceptBtn.style.padding = '4px 10px';
    acceptBtn.style.border = 'none';
    acceptBtn.style.borderRadius = '4px';
    acceptBtn.style.fontSize = '12px';
    acceptBtn.style.cursor = 'pointer';
    acceptBtn.style.color = '#ffffff';
    acceptBtn.style.background = '#8B5CF6';
    acceptBtn.addEventListener('click', async function(){
      // alias 저장 (백그라운드, 실패해도 진행)
      try {
        await fetch('/admin/projects/alias', {
          method: 'POST',
          headers: { 'Content-Type': 'application/json' },
          credentials: 'same-origin',
          body: JSON.stringify({ project_id: resolved.id, alias: query })
        });
      } catch(_) {}
      var nameEl = document.getElementById('ov-np-name');
      if (nameEl) nameEl.value = resolved.label;
      window.submitNewProject();
    });
    wrap.appendChild(acceptBtn);
    var rejectBtn = document.createElement('button');
    rejectBtn.type = 'button';
    rejectBtn.textContent = '아니오';
    rejectBtn.style.padding = '4px 10px';
    rejectBtn.style.border = '1px solid #D1D5DB';
    rejectBtn.style.borderRadius = '4px';
    rejectBtn.style.fontSize = '12px';
    rejectBtn.style.cursor = 'pointer';
    rejectBtn.style.background = '#ffffff';
    rejectBtn.style.color = '#374151';
    rejectBtn.addEventListener('click', function(){
      window._npClearSuggest();
    });
    wrap.appendChild(rejectBtn);
    box.appendChild(wrap);
  };

  window._onNewProjectNameInput = function(){
    var nameEl = document.getElementById('ov-np-name');
    if (!nameEl) return;
    var q = (nameEl.value || '').trim();
    if (_npSuggestTimer) clearTimeout(_npSuggestTimer);
    if (!q) { window._npClearSuggest(); return; }
    _npSuggestTimer = setTimeout(async function(){
      try {
        var r = await fetch('/admin/projects/suggest?query=' + encodeURIComponent(q) + '&limit=3', { credentials: 'same-origin' });
        var d = await r.json();
        var exact = d.exact_match;
        var sug = d.suggestions || [];
        if (exact) {
          window._npBuildSuggestBox(
            'warn',
            '이미 등록된 프로젝트입니다: ' + exact.label,
            '이 프로젝트로 진행',
            function(){
              nameEl.value = exact.label;
              window.submitNewProject();
            },
            null
          );
          return;
        }
        if (sug.length > 0 && sug[0].similarity >= 0.6) {
          var top = sug[0];
          window._npBuildSuggestBox(
            'info',
            '혹시 ' + top.label + ' 프로젝트인가요?',
            '네, ' + top.label,
            function(){
              nameEl.value = top.label;
              window.submitNewProject();
            },
            function(){
              window._npClearSuggest();
            }
          );
          return;
        }
        // fuzzy 실패 -> AI resolver 호출 (최소 2자)
        if (q.length < 2) {
          window._npClearSuggest();
          return;
        }
        var sidebarSel = document.getElementById('v2-division-select');
        var divisionId = sidebarSel ? sidebarSel.value : '';
        window._npBuildLoadingBox('🤖 AI로 확인 중...');
        try {
          var ar = await fetch('/admin/projects/resolve', {
            method: 'POST',
            headers: { 'Content-Type': 'application/json' },
            credentials: 'same-origin',
            body: JSON.stringify({ query: q, division_id: divisionId })
          });
          var ad = await ar.json();
          if (ad && ad.resolved && ad.confidence >= 0.6) {
            window._npBuildAiSuggestBox(ad.resolved, ad.confidence, q, divisionId);
          } else {
            window._npClearSuggest();
          }
        } catch(_) {
          window._npClearSuggest();
        }
      } catch(e) {
        window._npClearSuggest();
      }
    }, 500);
  };
  window.closeNewProjectModal = function(){
    const mask = document.getElementById('ov-newproj-mask');
    if (mask) mask.classList.remove('open');
  };
  window.submitNewProject = async function(){
    const name = (document.getElementById('ov-np-name').value || '').trim();
    if (!name) { alert('프로젝트명을 입력해주세요'); return; }
    const divisionSel = document.getElementById('v2-division-select');
    const divisionId = ((divisionSel && divisionSel.value) ? divisionSel.value : 'semiconductor').trim() || 'semiconductor';
    // ★ 중복 검사 — 같은 division에 같은 프로젝트명이 이미 있으면 차단
    try {
      const _sel = document.getElementById('v2-division-select');
      const _curDiv = (_sel && _sel.value ? _sel.value : '').trim();
      const _all = Array.isArray(window._v2AllReports) ? window._v2AllReports : [];
      const _dup = _all.find(function(r){
        const _rd = (r.division_id
                || (r.products && r.products[0] && r.products[0].division_id)
                || '').trim();
        const _sameDiv = _curDiv === 'semiconductor'
          ? (!_rd || _rd === 'semiconductor')
          : (_rd === _curDiv);
        if (!_sameDiv) return false;
        const _pname = ((r.products && r.products[0] && r.products[0].name) || '').trim();
        return _pname === name;
      });
      if (_dup) {
        alert('⚠️ "' + name + '" 프로젝트가 이미 존재합니다. 기존 보고서를 열어서 편집해주세요.');
        window.closeNewProjectModal && window.closeNewProjectModal();
        return;
      }
    } catch(_e) {
      console.warn('중복 검사 실패, 계속 진행:', _e);
    }
    // 현재 ISO 주차 자동 계산
    function getISOWeek(d){
      const dt = new Date(Date.UTC(d.getFullYear(), d.getMonth(), d.getDate()));
      const dayNum = dt.getUTCDay() || 7;
      dt.setUTCDate(dt.getUTCDate() + 4 - dayNum);
      const yearStart = new Date(Date.UTC(dt.getUTCFullYear(),0,1));
      return Math.ceil((((dt - yearStart) / 86400000) + 1)/7);
    }
    const week = getISOWeek(new Date());
    try {
      const res = await fetch('/admin/reports', {
        method: 'POST',
        headers: { 'Content-Type': 'application/json' },
        body: JSON.stringify({ project_name: name, week: week, headline: '', division_id: divisionId })
      });
      if (!res.ok) { alert('생성 실패'); return; }
      window.closeNewProjectModal();
      if (window.loadAdminV2Reports) window.loadAdminV2Reports();
    } catch(e) {
      alert('오류: ' + e.message);
    }
  };

// ─────────────────────────────────────────────
// admin/v2 사업부 필터 헬퍼
// ─────────────────────────────────────────────
window._v2FilterByDivision = function(all){
    const sel = document.getElementById('v2-division-select');
    const currentDiv = (sel && sel.value ? sel.value : '').trim();
    const arr = Array.isArray(all) ? all : [];
    if (!currentDiv) return arr;
    return arr.filter(function(r){
      const rd = (r.division_id
              || (r.products && r.products[0] && r.products[0].division_id)
              || (r.products && r.products[0] && r.products[0]._fallback_division_id)
              || '').trim();
      if (currentDiv === 'semiconductor') {
        return !rd || rd === 'semiconductor';
      }
      return rd === currentDiv;
    });
};

window.renderAdminV2ByDivision = function(){
    if (typeof window.loadAdminV2Reports === 'function') {
      window.loadAdminV2Reports();
    }
};

  window.loadAdminV2Reports = async function(){
    try {
      const res = await fetch('/admin/reports/all');
      const data = await res.json();
      window._v2AllReports = Array.isArray(data.reports) ? data.reports : [];
    let reports = window._v2FilterByDivision(window._v2AllReports);
      // 최신순 정렬 (upload_timestamp desc)
      reports.sort(function(a, b){
        return (b.upload_timestamp || '').localeCompare(a.upload_timestamp || '');
      });

      const listEl = document.getElementById('ov-report-list');
      if (!listEl) return;

      const __showAll = window._v2ShowAllReports === true;
      const recent = __showAll ? reports : reports.slice(0, 3);

      // 더 보기 버튼 처리
      const __moreBtn = document.getElementById('ov-load-more-btn');
      if (__moreBtn) {
        if (reports.length <= 3) {
          __moreBtn.style.display = 'none';
        } else {
          __moreBtn.style.display = '';
          __moreBtn.textContent = __showAll ? '접기' : ('보고 더 보기 (' + (reports.length - 3) + '개)');
          __moreBtn.onclick = function(){
            window._v2ShowAllReports = !window._v2ShowAllReports;
            window.loadAdminV2Reports();
          };
        }
      }
      if (recent.length === 0) {
        listEl.innerHTML = '<div style="color:#9CA3AF;font-size:14px;padding:20px;text-align:center;">등록된 보고가 없습니다.</div>';
      } else {
        listEl.innerHTML = recent.map(function(r){
          const title = r.display_title || r.file_name || '제목 없음';
          const fileName = r.file_name || '파일 없음';
          const size = r.file_size || '';
          const firstProduct = (r.products && r.products[0]) || {};
          const preview = firstProduct.headline || '';
          const status = (r.report_status || 'published').toString();
          const dDay = r.d_day;

          let badgeClass = 'green';
          let badgeText = '발행 완료';
          let barClass = 'green';
          let dDayText = '';
          let dDayColor = '#7A8595';

          if (status === 'ai_processing') {
            badgeClass = 'blue'; badgeText = 'AI 처리 중'; barClass = 'orange';
          } else if (status === 'review_pending') {
            badgeClass = 'amber'; badgeText = '검토 대기'; barClass = 'red';
          }

          if (dDay === null || dDay === undefined) {
            dDayText = '';
          } else if (dDay === 0) {
            dDayText = 'D-day';
            dDayColor = '#C1272D';
          } else if (dDay > 0) {
            dDayText = 'D-' + dDay;
            dDayColor = '#C1272D';
          } else {
            dDayText = '완료';
            dDayColor = '#7A8595';
          }

          // 수기 프로젝트는 파일 없음 → 상태배지/D-day/AI 검토중 preview 다 숨김
          const isManual = !!r.is_manual;
          const hideStatus = isManual;
          const showPreview = !isManual && preview && preview !== 'AI 검토중';

          const sizeText = size ? ' · ' + size : '';
          const isHidden = !!r.hidden;
          const hideLabel = isHidden ? '숨김 해제' : '숨기기';
          const itemClass = 'ov-report-item' + (isHidden ? ' is-hidden' : '');
          return ''
            + '<div class="' + itemClass + '">'
            +   '<div class="ov-report-bar ' + barClass + '"></div>'
            +   '<div class="ov-report-main">'
            +     '<div class="ov-report-project">' + title + '</div>'
            +     '<div class="ov-report-file">📎 ' + fileName + sizeText + '</div>'
            +     (showPreview ? '<div class="ov-report-preview">' + preview + '</div>' : '')
            +   '</div>'
            +   '<div class="ov-report-side">'
            +     (hideStatus ? '' : ('<div style="display:flex;align-items:center;gap:14px;">'
            +       '<div class="ov-badge ' + badgeClass + '">' + badgeText + '</div>'
            +       (dDayText ? '<div style="font-size:15px;font-weight:800;color:' + dDayColor + ';min-width:46px;text-align:right;">' + dDayText + '</div>' : '')
            +     '</div>'))
            +     '<div class="ov-side-top-actions">'
            +       '<button class="ov-hide-btn ov-hide-report' + (isHidden ? ' is-active' : '') + '" type="button" data-doc="' + (r.doc_id || '') + '">' + hideLabel + '</button>'
            +       '<button class="ov-del-btn ov-del-report" type="button" data-doc="' + (r.doc_id || '') + '" data-title="' + ((r.display_title || r.doc_id || '')).replace(/"/g, '&quot;') + '">삭제</button>'
            +     '</div>'
            +     '<button class="ov-open-btn ov-open-report" type="button" data-doc="' + (r.doc_id || '') + '" data-product="' + (((r.products && r.products[0] && r.products[0].name) || '')).replace(/"/g, '') + '" data-split-project="' + ((r._split_project || '')).replace(/"/g, '') + '">열기</button>'
            +   '</div>'
            + '</div>';
        }).join('');
      }

      const published = reports.filter(function(r){ return (r.report_status || 'published') === 'published'; }).length;
      const review = reports.filter(function(r){ return (r.report_status || '') === 'review_pending'; }).length;
      const ai = reports.filter(function(r){ return (r.report_status || '') === 'ai_processing'; }).length;

      // 이번 주 = 오늘 기준 지난 7일 이내 업로드
      const now = new Date();
      const sevenDaysAgo = new Date(now.getTime() - 7 * 24 * 60 * 60 * 1000);
      const thisWeek = reports.filter(function(r){
        try { return new Date(r.upload_timestamp) >= sevenDaysAgo; } catch(_){ return false; }
      }).length;

      const wkEl = document.getElementById('ov-kpi-week');
      const aiEl = document.getElementById('ov-kpi-ai');
      const rvEl = document.getElementById('ov-kpi-review');
      const pbEl = document.getElementById('ov-kpi-published');
      if (wkEl) wkEl.textContent = String(thisWeek);
      if (aiEl) aiEl.textContent = String(ai);
      if (rvEl) rvEl.textContent = String(review);
      if (pbEl) pbEl.textContent = String(published);

      // + 새 프로젝트 버튼 (재바인딩 안전)
      const newBtn = document.getElementById('ov-new-project-btn');
      if (newBtn && !newBtn._bound) {
        newBtn._bound = true;
        newBtn.addEventListener('click', window.openNewProjectModal);
      }
      const cancelBtn = document.getElementById('ov-np-cancel');
      if (cancelBtn && !cancelBtn._bound) {
        cancelBtn._bound = true;
        cancelBtn.addEventListener('click', window.closeNewProjectModal);
      }
      const confirmBtn = document.getElementById('ov-np-confirm');
      if (confirmBtn && !confirmBtn._bound) {
        confirmBtn._bound = true;
        confirmBtn.addEventListener('click', window.submitNewProject);
      }
      const maskEl = document.getElementById('ov-newproj-mask');
      if (maskEl && !maskEl._bound) {
        maskEl._bound = true;
        maskEl.addEventListener('click', function(e){
          if (e.target === maskEl) window.closeNewProjectModal();
        });
      }
      const delCancel = document.getElementById('ov-delconf-cancel');
      if (delCancel && !delCancel._bound) {
        delCancel._bound = true;
        delCancel.addEventListener('click', window.closeDeleteConfirm);
      }
      const delConfirm = document.getElementById('ov-delconf-confirm');
      if (delConfirm && !delConfirm._bound) {
        delConfirm._bound = true;
        delConfirm.addEventListener('click', window.executeDelete);
      }
      const delMask = document.getElementById('ov-delconf-mask');
      if (delMask && !delMask._bound) {
        delMask._bound = true;
        delMask.addEventListener('click', function(e){
          if (e.target === delMask) window.closeDeleteConfirm();
        });
      }
      // 모달 사업부 드롭다운은 제거됨 - 사이드바 v2-division-select 값 자동 참조
      const nameInputBind = document.getElementById('ov-np-name');
      if (nameInputBind && !nameInputBind._suggestBound) {
        nameInputBind._suggestBound = true;
        nameInputBind.addEventListener('input', window._onNewProjectNameInput);
      }

      // 열기 버튼 이벤트 바인딩
      document.querySelectorAll('.ov-open-report').forEach(function(btn){
        btn.addEventListener('click', function(){
          const doc = btn.getAttribute('data-doc') || '';
          const prod = btn.getAttribute('data-product') || '';
          const splitProject = btn.getAttribute('data-split-project') || '';
          if (window.openReportEdit) window.openReportEdit(doc, prod, splitProject);
        });
      });
      document.querySelectorAll('.ov-del-report').forEach(function(btn){
        btn.addEventListener('click', function(){
          var doc = btn.getAttribute('data-doc') || '';
          var title = btn.getAttribute('data-title') || '';
          if (doc) window.openDeleteConfirm(doc, title);
        });
      });
      document.querySelectorAll('.ov-hide-report').forEach(function(btn){
        btn.addEventListener('click', async function(){
          const doc = btn.getAttribute('data-doc') || '';
          if (!doc) return;
          btn.disabled = true;
          try {
            const res = await fetch('/admin/reports/' + encodeURIComponent(doc) + '/hide', { method: 'POST' });
            if (!res.ok) { alert('숨김 처리 실패'); btn.disabled = false; return; }
            if (window.loadAdminV2Reports) window.loadAdminV2Reports();
          } catch(e) {
            alert('숨김 처리 오류: ' + e.message);
            btn.disabled = false;
          }
        });
      });
    } catch (e) {
      console.error('loadAdminV2Reports failed', e);
    }
  };

  // ============================================================
  // Admin v2 · 보고 편집 (raw_text)
  // ============================================================
  window._currentEditContext = { divisionId: '', cardTitle: '' };

  window.renderEditPageHTML = function(ctx){
    ctx = ctx || {};
    const projectLabel = ctx.projectLabel || '보고';
    const statusLabel = ctx.statusLabel || '';
    const statusClass = ctx.statusClass || 'green';
    const subLine = ctx.subLine || '';
    return ''
      + '<style>'
      + '.ov-edit-topbar{display:flex;align-items:center;justify-content:space-between;background:#fff;border:1px solid #E6EBF2;border-radius:20px;padding:14px 18px;margin-bottom:20px;box-shadow:0 6px 18px rgba(15,44,89,0.04);}'
      + '.ov-back-btn{border:1px solid #DCE4EF;background:#fff;color:#35527C;border-radius:12px;padding:9px 14px;font-size:14px;font-weight:700;cursor:pointer;}'
      + '.ov-edit-title{font-size:16px;color:#111827;font-weight:800;}'
      + '.ov-save-btn{border:0;background:#0F2C59;color:#fff;border-radius:12px;padding:10px 18px;font-size:14px;font-weight:800;cursor:pointer;}'
      + '.ov-header-card{background:#fff;border:1px solid #E6EBF2;border-radius:22px;padding:22px 24px;margin-bottom:18px;box-shadow:0 6px 18px rgba(15,44,89,0.04);}'
      + '.ov-header-title{display:flex;align-items:center;gap:12px;font-size:26px;font-weight:800;color:#0F2C59;letter-spacing:-0.3px;margin-bottom:8px;}'
      + '.ov-header-sub{font-size:14px;color:#6E7785;font-weight:600;}'
      + '.ov-tabs{display:inline-flex;gap:6px;background:#F1F5FB;border-radius:14px;padding:4px;margin-bottom:20px;}'
      + '.ov-tab{border:0;background:transparent;color:#6E7785;padding:8px 16px;font-size:13px;font-weight:700;border-radius:10px;cursor:pointer;}'
      + '.ov-tab.active{background:#0F2C59;color:#fff;}'
      + '.ov-tab-wrap{display:inline-flex;align-items:center;gap:4px;}'
      + '.ov-tab-actions{display:inline-flex;align-items:center;gap:4px;opacity:0;pointer-events:none;transition:opacity .15s ease;}'
      + '.ov-tab-wrap:hover .ov-tab-actions{opacity:1;pointer-events:auto;}'
      + '.ov-tab-icon{width:24px;height:24px;border:0;border-radius:7px;background:#E8EEF7;color:#5B7BB0;font-size:12px;cursor:pointer;display:inline-flex;align-items:center;justify-content:center;padding:0;}'
      + '.ov-tab-icon:hover{background:#D8E5F7;}'
      + '.ov-tab-icon.ov-tab-del:hover{background:#FEE7E7;color:#B8302E;}'
      + '.ov-section{background:#fff;border:1px solid #E6EBF2;border-radius:22px;padding:20px;margin-bottom:18px;box-shadow:0 6px 18px rgba(15,44,89,0.04);}'
      + '.ov-section-title{font-size:18px;font-weight:800;color:#0F2C59;margin-bottom:14px;display:flex;align-items:center;gap:10px;}'
      + '.ov-section-num{display:inline-flex;align-items:center;justify-content:center;width:26px;height:26px;border-radius:50%;background:#0F2C59;color:#fff;font-size:13px;font-weight:800;}'
      + '.ov-placeholder{border:2px dashed #D9E3F1;background:#F8FBFF;border-radius:16px;padding:22px;text-align:center;color:#8593A6;font-size:14px;font-weight:600;}'
      + '.ov-kpi-slot-wrap{margin-top:16px;}'
      + '.ov-kpi-slot-label{font-size:13px;color:#6E7785;font-weight:700;margin-bottom:10px;padding-left:4px;}'
      + '.ov-kpi-slot{background:#F8FBFF;border:1px solid #E6EBF2;border-radius:18px;padding:18px;}'
      + '.ov-badge-inline{display:inline-flex;align-items:center;border-radius:999px;padding:6px 12px;font-size:12px;font-weight:800;}'
      + '.ov-badge-inline.green{background:#E7F8F0;color:#117A52;}'
      + '.ov-badge-inline.amber{background:#FFF3D9;color:#9A6700;}'
      + '.ov-badge-inline.blue{background:#E7F0FB;color:#23538C;}'
      + '.ov-badge-inline.red{background:#FEE7E7;color:#B8302E;}'
      + '.ov-kpi-grid-months{display:grid;grid-template-columns:repeat(3, minmax(0,1fr));gap:12px;margin-bottom:12px;}'
      + '.ov-kpi-grid-weeks{display:grid;grid-template-columns:repeat(5, minmax(0,1fr));gap:10px;}'
      + '.ov-kpi-cell{background:#fff;border:1px solid #E6EBF2;border-radius:14px;padding:14px;}'
      + '.ov-kpi-cell.actual{border-color:#B9E1CE;}'
      + '.ov-kpi-cell.plan{border-color:#D9E3F1;}'
      + '.ov-kpi-cell-label{font-size:12px;color:#6E7785;font-weight:700;margin-bottom:4px;}'
      + '.ov-kpi-cell-value{font-size:20px;color:#0F2C59;font-weight:800;line-height:1.1;}'
      + '.ov-kpi-cell-unit{font-size:12px;color:#7A8595;font-weight:600;margin-left:4px;}'
      + '.ov-kpi-cell-sub{font-size:11px;color:#8593A6;font-weight:600;margin-top:4px;}'
      + '.ov-kpi-cell-badge{font-size:10px;font-weight:800;padding:2px 6px;border-radius:6px;margin-left:6px;vertical-align:middle;}'
      + '.ov-kpi-cell-badge.actual{background:#E7F8F0;color:#117A52;}'
      + '.ov-kpi-cell-badge.plan{background:#EAF0FB;color:#2E5B94;}'
      + '.ov-issue-list{display:flex;flex-direction:column;gap:10px;}'
      + '.ov-issue-row{display:flex;justify-content:space-between;align-items:flex-start;gap:12px;padding:12px 14px;background:#F8FBFF;border:1px solid #E6EBF2;border-radius:14px;}'
      + '.ov-issue-text{font-size:14px;color:#111827;font-weight:700;line-height:1.45;flex:1;}'
      + '.ov-issue-dday{font-size:12px;font-weight:800;color:#B8302E;white-space:nowrap;}'
      + '.ov-list-dash{padding-left:20px;} .ov-list-dash li{position:relative;list-style:none;} .ov-list-dash li::before{content:"－";position:absolute;left:-16px;color:#6E7785;}'
      + '.ov-list-triangle{padding-left:22px;} .ov-list-triangle li{position:relative;list-style:none;} .ov-list-triangle li::before{content:"▶";position:absolute;left:-18px;color:#2E5B94;font-size:0.85em;top:0.15em;}'
      + '.ov-list-check{padding-left:22px;} .ov-list-check li{position:relative;list-style:none;} .ov-list-check li::before{content:"✓";position:absolute;left:-18px;color:#117A52;font-weight:800;}'
      + '.ov-list-star{padding-left:22px;} .ov-list-star li{position:relative;list-style:none;} .ov-list-star li::before{content:"★";position:absolute;left:-20px;color:#D9A400;}'
      + '.ov-list-arrow{padding-left:24px;} .ov-list-arrow li{position:relative;list-style:none;} .ov-list-arrow li::before{content:"➜";position:absolute;left:-20px;color:#2E5B94;}'
      + '.ov-list-diamond{padding-left:22px;} .ov-list-diamond li{position:relative;list-style:none;} .ov-list-diamond li::before{content:"◆";position:absolute;left:-18px;color:#8B5CF6;}'
      + '.ov-list-dot{padding-left:18px;} .ov-list-dot li{position:relative;list-style:none;} .ov-list-dot li::before{content:"·";position:absolute;left:-12px;color:#6E7785;font-weight:800;}'
      + '.ov-list-circled{padding-left:26px;counter-reset:ov-circled;} .ov-list-circled li{position:relative;list-style:none;counter-increment:ov-circled;} .ov-list-circled li::before{content:counter(ov-circled,decimal);position:absolute;left:-22px;display:inline-flex;align-items:center;justify-content:center;width:18px;height:18px;border:1.5px solid #12325F;border-radius:50%;font-size:10px;font-weight:800;color:#111827;top:0.2em;}'
      + '.ov-list-paren{padding-left:22px;counter-reset:ov-paren;} .ov-list-paren li{position:relative;list-style:none;counter-increment:ov-paren;} .ov-list-paren li::before{content:counter(ov-paren) ")";position:absolute;left:-18px;color:#111827;font-weight:700;}'
      + '</style>'
      + '<div class="ov-edit-topbar" style="display:flex;align-items:center;justify-content:space-between;gap:12px;">'
      + '  <button class="ov-back-btn" type="button" id="ov-back-btn" style="flex-shrink:0;">← 목록으로</button>'
      + '  <div class="ov-edit-title" id="ov-edit-title" style="flex:1;text-align:center;font-size:16px;font-weight:800;color:#111827;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;">' + projectLabel + '</div>'
      + '  <div style="flex-shrink:0;width:100px;"></div>'
      + '</div>'
      + (ctx.isManual
          ? '<div id="ov-sections-container"></div>'
          : ('<div class="ov-tabs" id="ov-ppt-tabs">'
            + '  <button class="ov-tab active" type="button" data-ppt-jump="0">① 현황</button>'
            + '  <button class="ov-tab" type="button" data-ppt-jump="1">② 주차별 출하실적</button>'
            + '  <button class="ov-tab" type="button" data-ppt-jump="2">③ 주요내용</button>'
            + '</div>'
            + '<div class="ov-section">'
            + '  <div class="ov-section-title"><span class="ov-section-num">1</span>현황</div>'
            + '  <div id="ov-current-status" class="ov-placeholder">현황 요약 (headline) 이 여기에 표시됩니다.</div>'
            + '</div>'
            + '<div class="ov-section">'
            + '  <div class="ov-section-title"><span class="ov-section-num">2</span>주차별 출하실적</div>'
            + '  <div class="ov-placeholder" style="margin-bottom:14px;">📥 출하실적 엑셀 업로드 · 개발 중</div>'
            + '  <div class="ov-placeholder" style="margin-bottom:14px;">📊 주차별 출하실적 표 · 개발 중</div>'
            + '  <div class="ov-kpi-slot-wrap">'
            + '    <div class="ov-kpi-slot-label">↓ 표에서 파생 · KPI 카드 (자동 계산)</div>'
            + '    <div class="ov-kpi-slot" id="ov-kpi-slot">'
            + '      <div style="color:#8593A6;font-size:13px;text-align:center;padding:12px;">KPI 데이터 로드 중...</div>'
            + '    </div>'
            + '  </div>'
            + '</div>'
            + '<div class="ov-section">'
            + '  <div class="ov-section-title"><span class="ov-section-num">3</span>주요내용</div>'
            + '  <div id="ov-issues-slot">'
            + '    <div class="ov-placeholder">주요내용 블록 · 개발 중</div>'
            + '  </div>'
            + '</div>'));
  };

  // KPI 카드 렌더 헬퍼
  window.renderKpiCardInline = function(card, issues){
    if (!card) {
      return '<div style="color:#8593A6;font-size:13px;text-align:center;padding:12px;">이 프로젝트는 KPI 데이터가 없습니다.</div>';
    }
    const months = Array.isArray(card.months) ? card.months : [];
    const weeks = Array.isArray(card.weeks) ? card.weeks : [];
    const unit = card.unit_label || '만불';

    let html = '';

    if (months.length > 0) {
      html += '<div class="ov-kpi-grid-months">';
      months.forEach(function(m){
        const typeClass = m.type === 'actual' ? 'actual' : 'plan';
        const badgeText = m.type === 'actual' ? '실적' : '계획';
        const monthKey = m.month_key || m.key || '';
        html += ''
          + '<div class="ov-kpi-cell ' + typeClass + '">'
          + '  <div class="ov-kpi-cell-label">' + (m.month || '') + '<span class="ov-kpi-cell-badge ' + typeClass + '">' + badgeText + '</span></div>'
          + '  <div class="ov-kpi-cell-value">' + (m.total != null ? Number(m.total).toFixed(2) : '-') + '<span class="ov-kpi-cell-unit">' + unit + '</span></div>'
          + '  <div class="ov-kpi-cell-sub">EFEM ' + (m.efem != null ? Number(m.efem).toFixed(2) : '-') + ' · VTM ' + (m.vtm != null ? Number(m.vtm).toFixed(2) : '-') + '</div>'
          + '</div>';
      });
      html += '</div>';
    }

    if (weeks.length > 0) {
      html += '<div class="ov-kpi-grid-weeks">';
      weeks.forEach(function(w){
        const typeClass = w.type === 'actual' ? 'actual' : 'plan';
        const badgeText = w.type === 'actual' ? '실적' : '계획';
        const weekKey = w.week_key || w.key || w.week || '';
        html += ''
          + '<div class="ov-kpi-cell ' + typeClass + '">'
          + '  <div class="ov-kpi-cell-label">' + (w.week || '') + '<span class="ov-kpi-cell-badge ' + typeClass + '">' + badgeText + '</span></div>'
          + '  <div class="ov-kpi-cell-value">' + (w.total != null ? Number(w.total).toFixed(2) : '-') + '<span class="ov-kpi-cell-unit">' + unit + '</span></div>'
          + '  <div class="ov-kpi-cell-sub">EFEM ' + (w.efem != null ? Number(w.efem).toFixed(2) : '-') + ' · VTM ' + (w.vtm != null ? Number(w.vtm).toFixed(2) : '-') + '</div>'
          + '</div>';
      });
      html += '</div>';
    }

    if (Array.isArray(issues) && issues.length > 0) {
      html += '<div style="margin-top:14px;padding-top:14px;border-top:1px dashed #D9E3F1;">';
      html += '<div style="font-size:12px;color:#6E7785;font-weight:700;margin-bottom:8px;">이슈 라인 (참고)</div>';
      html += '<div class="ov-issue-list">';
      issues.forEach(function(ii){
        const text = ii.text || '';
        const dueDate = ii.due_date || '';
        let ddayText = '';
        if (dueDate && ii.show_dday !== false) {
          try {
            const d = new Date(dueDate);
            const today = new Date();
            const diff = Math.round((d - today) / (1000*60*60*24));
            const mmdd = (d.getMonth()+1) + '/' + d.getDate();
            if (diff === 0) ddayText = 'D-day · ' + mmdd;
            else if (diff > 0) ddayText = 'D-' + diff + ' · ' + mmdd;
            else ddayText = 'D+' + Math.abs(diff) + ' 지남 · ' + mmdd;
          } catch(_){}
        }
        html += ''
          + '<div class="ov-issue-row">'
          + '  <div class="ov-issue-text">' + text + '</div>'
          + (ddayText ? '<div class="ov-issue-dday">' + ddayText + '</div>' : '')
          + '</div>';
      });
      html += '</div></div>';
    }

    return html;
  };

  window.openReportEdit = async function(docId, productName, splitProject){
    const c = document.getElementById('v2-content-inner') || document.getElementById('v2-content');
    if (!c) return;

    // 1) 보고 메타 조회 (admin은 숨김 포함 전체를 봐야 하므로 /admin/reports/all 사용)
    let target = null;
    try {
      const res = await fetch('/admin/reports/all');
      const data = await res.json();
      const reports = Array.isArray(data.reports) ? data.reports : [];
      target = reports.find(function(r){
        return r.doc_id === docId && ((r._split_project || '') === (splitProject || ''));
      }) || reports.find(function(r){
        return r.doc_id === docId;
      }) || reports.find(function(r){
        return ((r._split_project || '') === (splitProject || ''));
      }) || reports.find(function(r){
        const names = []
          .concat(r.manual_projects || [])
          .concat(((r.parsed || {}).projects) || [])
          .concat((r.products || []).map(function(p){ return p && p.name ? p.name : ''; }));
        return names.some(function(n){ return String(n || '').trim() === String(productName || '').trim(); });
      }) || null;
      if (!target) {
        console.warn('열기: 해당 보고를 찾을 수 없음', docId, splitProject);
        alert('해당 보고를 찾을 수 없습니다. 목록을 새로고침합니다.');
        if (window.backToReportList) window.backToReportList();
        return;
      }
    } catch(e){
      console.error(e);
      alert('보고 조회 실패: ' + e.message);
      return;
    }

    // 2) 헤더 컨텍스트 구성
    const parsed = (target && target.parsed) || {};
    const projects = parsed.projects || [];
    const _base = projects.length ? projects.join(', ') : (target && target.file_name) || '보고';
    // week_override > parsed.week > 현재 주차 (자동 계산)
    let _week = (target && target.week_override) || parsed.week || '';
    if (!_week) {
      const _now = new Date();
      const _start = new Date(_now.getFullYear(), 0, 1);
      const _days = Math.floor((_now - _start) / (24 * 60 * 60 * 1000));
      _week = Math.ceil((_days + _start.getDay() + 1) / 7);
    }
    const projectLabel = _base
      ? (_week ? _base + ' · W' + _week + ' 주간보고' : _base)
      : ((target && target.display_title) || '보고');
    const week = _week || '';
    const date = parsed.date || (target && (target.report_meta || {}).date) || '';
    const firstProduct = (target && target.products && target.products[0]) || {};
    const headline = firstProduct.headline || '';

    let statusLabel = '발행 완료';
    let statusClass = 'green';
    if (target && target.report_status === 'ai_processing') { statusLabel = 'AI 처리 중'; statusClass = 'amber'; }
    else if (target && target.report_status === 'review_pending') { statusLabel = '검토 대기'; statusClass = 'red'; }

    const subParts = [];
    if (date) subParts.push('보고일자 ' + date);
    if (week) subParts.push('W' + week);
    if (headline) subParts.push(headline);
    const subLine = subParts.join(' · ');

    // 3) 스켈레톤 렌더
    const isManual = !!(target && target.is_manual);
    c.innerHTML = window.renderEditPageHTML({
      fileName: (target && target.file_name) || (isManual ? '(수기 프로젝트)' : '편집'),
      projectLabel: projectLabel,
      statusLabel: isManual ? '' : statusLabel,
      statusClass: statusClass,
      subLine: isManual ? '' : subLine,
      isManual: isManual
    });

    // 수기 프로젝트: sections CRUD 지원
    let sectionsState = [];
    let sectionsDirty = false;
    const SECTION_ALIASES = {
      '현황': ['현황','상황','진행현황','진행상황'],
      '주차별 출하실적': ['주차별 출하실적','출하','출하실적','주차실적','실적','shipment'],
      '주차별 계획': ['주차별 계획','계획','주차계획','주간계획','plan'],
      '주요내용': ['주요내용','주요','핵심내용','핵심','main'],
      '이슈/리스크': ['이슈/리스크','이슈','리스크','문제','특이사항','issue','risk'],
      '요청사항': ['요청사항','요청','지원요청','request']
    };
    function _normSectionName(s){
      return String(s || '').trim().toLowerCase().replace(/\s+/g,' ');
    }
    function _findAliasCanonical(input){
      const n = _normSectionName(input);
      if (!n) return null;
      for (const canonical in SECTION_ALIASES) {
        const list = SECTION_ALIASES[canonical];
        for (let i=0;i<list.length;i++){
          if (_normSectionName(list[i]) === n) return canonical;
        }
      }
      return null;
    }
    function _genSecId(){
      return 'sec_' + Math.random().toString(36).slice(2,10);
    }
    function _markSectionsDirty(){
      sectionsDirty = true;
      markDirty();
    }

    if (isManual) {
      sectionsState = JSON.parse(JSON.stringify(firstProduct.sections || []));

      // ─── 편집 컨텍스트 설정 (일정 관리 패널용) ───
      try {
        // divisionId 우선순위: target.division_id → 사이드바 선택값
        let divIdFromTarget = (target && (target.division_id || target.divisionId)) || '';
        if (!divIdFromTarget) {
          const sidebarSel = document.getElementById('v2-division-select')
                          || document.getElementById('sidebarDivisionSelect')
                          || document.querySelector('select[name="division"]')
                          || document.querySelector('[data-role="division-select"]');
          if (sidebarSel) divIdFromTarget = sidebarSel.value || '';
        }
        // 마지막 fallback: URL/hash에서 division 추출
        if (!divIdFromTarget) {
          const m = location.hash.match(/division[=/]([\w-]+)/);
          if (m) divIdFromTarget = m[1];
        }

        // cardTitle 우선순위: manual_projects[0] → firstProduct.name → projectLabel 앞부분
        let cardTitleClean = '';
        const mp = (target && target.manual_projects) || [];
        if (mp.length > 0) cardTitleClean = String(mp[0]).trim();
        if (!cardTitleClean && firstProduct && firstProduct.name) {
          cardTitleClean = String(firstProduct.name).trim();
        }
        if (!cardTitleClean && projectLabel) {
          // "챔버 · W30 주간보고" 같은 형태에서 "챔버" 부분만
          cardTitleClean = String(projectLabel).split('·')[0].trim();
        }

        window._currentEditContext = {
          divisionId: divIdFromTarget,
          cardTitle: cardTitleClean
        };
        console.log('[일정 패널] 컨텍스트:', window._currentEditContext);
      } catch(e) { console.error('컨텍스트 설정 실패', e); }

      // ─── 일정 관리 패널용 items 캐시 로드 (현재 카드만 notes에서 로드) ───
      window._notesItemsCache = {};
      (async function _loadNotesForSchedule(){
        try {
          let divId = (window._currentEditContext || {}).divisionId || '';
          const cardTitle = (window._currentEditContext || {}).cardTitle || '';
          if (!divId) {
            const sel = document.getElementById('v2-division-select');
            if (sel) divId = sel.value || '';
          }
          if (!divId || !cardTitle) return;
          const r = await fetch('/notes?division_id=' + encodeURIComponent(divId), { credentials: 'same-origin' });
          if (!r.ok) return;
          const j = await r.json();
          const cards = (j.cards) || ((j.notes || {}).cards) || [];
          const card = cards.find(function(c){ return (c.title || '').trim() === cardTitle.trim(); });
          if (!card) return;
          const map = {};
          (card.sections || []).forEach(function(sec){
            map[sec.title || ''] = Array.isArray(sec.items) ? sec.items : [];
          });
          window._notesItemsCache = map;
          if (window._renderManualSections) window._renderManualSections();
        } catch(e) { console.error('notes items 로드 실패', e); }
      })();
      const sectionsRoot = document.getElementById('ov-sections-container');

      // === Tab/Shift+Tab/Backspace/Ctrl+Z Document-level Delegation ===
      // 재렌더링돼도 리스너 안 죽게 document에 1회만 등록.
      // 기존 editBox 개별 리스너(5349,5370,5511,5661,5964,5980)는
      // stopImmediatePropagation()으로 무력화됨.
      if (!document.__aiTabDelegationInstalled) {
        document.__aiTabDelegationInstalled = true;

        document.addEventListener('keydown', function(ev){
          var target = ev.target;
          if (!target || !target.closest) return;

          var aiBox = target.closest('.ov-ai-diff-edit');
          var origBox = target.closest('.ov-block-text-body');
          var editBox = aiBox || origBox;
          if (!editBox) return;

          var isAi = !!aiBox;
          var isMod = ev.metaKey || ev.ctrlKey;

          // undo API 분기
          var pushUndo = isAi ? editBox._aiPushUndo : editBox._origPushUndo;
          var undoFn = isAi ? editBox._aiUndo : editBox._origUndo;
          var redoFn = isAi ? editBox._aiRedo : editBox._origRedo;
          var getCur = isAi ? editBox._aiGetCursorOffset : editBox._origGetCursorOffset;
          var setCur = isAi ? editBox._aiSetCursorOffset : editBox._origSetCursorOffset;

          // ─────────────────────────────────
          // Ctrl+Z / Cmd+Z / Ctrl+Shift+Z / Ctrl+Y
          // ─────────────────────────────────
          if (isMod && (ev.key === 'z' || ev.key === 'Z')) {
            if (!undoFn || !redoFn) return;
            ev.preventDefault();
            if (ev.stopImmediatePropagation) ev.stopImmediatePropagation();
            if (ev.shiftKey) redoFn();
            else undoFn();
            return;
          }
          if (isMod && (ev.key === 'y' || ev.key === 'Y')) {
            if (!redoFn) return;
            ev.preventDefault();
            if (ev.stopImmediatePropagation) ev.stopImmediatePropagation();
            redoFn();
            return;
          }

          // ─────────────────────────────────
          // Tab / Shift+Tab
          // ─────────────────────────────────
          if (ev.key === 'Tab') {
            if (!window._aiCollectSelectedNodes) return;
            // 자동 마이그레이션 + selection 복원
            // 개행 문자로 뭉친 텍스트를 개별 div로 분리하고 커서를 새 div로 재설정
            if (window._aiSplitTextByNewlines && !editBox.__aiLineSplit) {
              var _savedSel = window.getSelection();
              var _savedAnchor = _savedSel && _savedSel.anchorNode;
              var _savedOffset = _savedSel && _savedSel.anchorOffset;
              var _savedText = null;
              if (_savedAnchor && _savedAnchor.nodeType === 3) {
                _savedText = _savedAnchor.nodeValue;
              }
              window._aiSplitTextByNewlines(editBox);
              if (_savedText) {
                try {
                  var walker = document.createTreeWalker(editBox, NodeFilter.SHOW_TEXT, null, false);
                  var found = null;
                  while (walker.nextNode()) {
                    var n = walker.currentNode;
                    if (n.nodeValue === _savedText ||
                        (_savedText.length > 5 && n.nodeValue.indexOf(_savedText.substring(0, 20)) >= 0)) {
                      found = n;
                      break;
                    }
                  }
                  if (found) {
                    var range = document.createRange();
                    var newOffset = Math.min(_savedOffset || 0, found.nodeValue.length);
                    range.setStart(found, newOffset);
                    range.collapse(true);
                    var s = window.getSelection();
                    s.removeAllRanges();
                    s.addRange(range);
                  }
                } catch(e) {}
              }
            }
            var collected = window._aiCollectSelectedNodes(editBox);
            if (collected.lis.length === 0 && collected.blocks.length === 0) return;

            ev.preventDefault();
            if (ev.stopImmediatePropagation) ev.stopImmediatePropagation();
            if (pushUndo) pushUndo();

            var savedOffset = getCur ? getCur() : -1;
            var lis = collected.lis.slice();
            var blocks = collected.blocks.slice();

            if (ev.shiftKey) {
              // Shift+Tab: outdent / unlist (아래부터 처리)
              var sortedLisOut = lis.slice().sort(function(a, b){
                var pos = a.compareDocumentPosition(b);
                if (pos & Node.DOCUMENT_POSITION_FOLLOWING) return 1;
                if (pos & Node.DOCUMENT_POSITION_PRECEDING) return -1;
                return 0;
              });
              sortedLisOut.forEach(function(li){
                if (!li.parentNode) return;
                var parentOl = li.parentNode;
                var ovr = parentOl.dataset.depthOverride;
                if (ovr != null && ovr !== '') {
                  var currentOvr = parseInt(ovr, 10);
                  var currentDep = window._aiGetListDepth(parentOl, editBox);
                  if (currentOvr > currentDep) {
                    var newOvr = currentOvr - 1;
                    if (newOvr <= currentDep) delete parentOl.dataset.depthOverride;
                    else parentOl.dataset.depthOverride = String(newOvr);
                    return;
                  }
                }
                var hostLi = parentOl ? parentOl.parentNode : null;
                var isTopLevel = !hostLi || hostLi.tagName !== 'LI';
                if (isTopLevel) {
                  if (window._aiUnlistCurrentLi) window._aiUnlistCurrentLi(editBox, li);
                } else {
                  var grandOl = hostLi.parentNode;
                  if (grandOl && grandOl.tagName === 'OL') {
                    if (hostLi.nextSibling) grandOl.insertBefore(li, hostLi.nextSibling);
                    else grandOl.appendChild(li);
                    if (!parentOl.querySelector('li')) parentOl.remove();
                  }
                }
              });
              blocks.forEach(function(block){
                var cur = parseInt(block.dataset.indent || '0', 10) || 0;
                if (cur > 0) {
                  block.setAttribute('data-indent', String(cur - 1));
                  if (String(cur - 1) === '0') block.removeAttribute('data-indent');
                } else if (block.dataset.indent) {
                  delete block.dataset.indent;
                }
              });
            } else {
              // Tab: indent (옵션 3 - selection 모든 li depth 증가)
              var sortedLis = lis.slice().sort(function(a, b){
                var pos = a.compareDocumentPosition(b);
                if (pos & Node.DOCUMENT_POSITION_FOLLOWING) return -1;
                if (pos & Node.DOCUMENT_POSITION_PRECEDING) return 1;
                return 0;
              });
              sortedLis.forEach(function(li, i){
                if (!li.parentNode) return;
                var parentOl = li.parentNode;
                var prevLi = li.previousElementSibling;
                var hasPrevLi = prevLi && prevLi.tagName === 'LI' && sortedLis.indexOf(prevLi) < 0;

                if (hasPrevLi) {
                  var childOl = window._aiEnsureChildOl(prevLi, window._aiStyleForDepth(window._aiGetListDepth(parentOl, editBox) + 1));
                  childOl.appendChild(li);
                  if (!parentOl.querySelector('li')) parentOl.remove();
                } else {
                  // 이전 li가 없어도 자리는 유지: 원래 ol을 li 위치에서 분할
                  // [앞 li들이 있던 원래 ol] [분리된 li의 새 ol] [뒤 li들의 새 ol]
                  var siblings = Array.prototype.filter.call(parentOl.children, function(x){ return x.tagName === 'LI'; });
                  var targetOl = parentOl;
                  if (siblings.length > 1) {
                    // li 뒤에 남는 li들 수집
                    var trailingLis = [];
                    var cur = li.nextElementSibling;
                    while (cur) {
                      var next = cur.nextElementSibling;
                      if (cur.tagName === 'LI') trailingLis.push(cur);
                      cur = next;
                    }

                    // 새 ol 생성 (분리된 li용)
                    var isolatedOl = document.createElement('ol');
                    isolatedOl.className = parentOl.className;
                    isolatedOl.dataset.start = '1';
                    isolatedOl.dataset.manualSplit = '1';
                    parentOl.dataset.manualSplit = '1';

                    // 뒤 li가 있으면 별도 ol로 감쌈 (parentOl 뒤에)
                    var trailingOl = null;
                    if (trailingLis.length > 0) {
                      trailingOl = document.createElement('ol');
                      trailingOl.className = parentOl.className;
                      trailingOl.dataset.manualSplit = '1';
                      // parentOl의 start 값 계산: 원래 start + 지금까지 남은 li 개수
                      var origStart = parseInt(parentOl.dataset.start || parentOl.getAttribute('start') || '1', 10) || 1;
                      var remainingCount = siblings.indexOf(li); // li 앞의 li 개수
                      trailingOl.dataset.start = String(origStart + remainingCount + 1);
                      trailingLis.forEach(function(tli){ trailingOl.appendChild(tli); });
                    }

                    // 순서: parentOl(앞 li들) → isolatedOl(분리 li) → trailingOl(뒤 li들)
                    // li를 isolatedOl로 이동
                    isolatedOl.appendChild(li);

                    // isolatedOl을 parentOl 바로 뒤에 삽입
                    if (parentOl.nextSibling) parentOl.parentNode.insertBefore(isolatedOl, parentOl.nextSibling);
                    else parentOl.parentNode.appendChild(isolatedOl);

                    // trailingOl을 isolatedOl 바로 뒤에 삽입
                    if (trailingOl) {
                      if (isolatedOl.nextSibling) isolatedOl.parentNode.insertBefore(trailingOl, isolatedOl.nextSibling);
                      else isolatedOl.parentNode.appendChild(trailingOl);
                    }

                    // parentOl이 비었으면 제거 (li가 첫 번째였을 때)
                    if (!parentOl.querySelector('li')) parentOl.remove();

                    targetOl = isolatedOl;
                  }
                  var currentDepth = window._aiGetListDepth(targetOl, editBox);
                  var override = targetOl.dataset.depthOverride;
                  var currentEffective = (override != null && override !== '') ? parseInt(override, 10) : currentDepth;
                  var newDepth = currentEffective + 1;
                  if (newDepth <= 5) targetOl.dataset.depthOverride = String(newDepth);
                }
              });
              blocks.forEach(function(block){
                var cur = parseInt(block.dataset.indent || '0', 10) || 0;
                if (cur < 5) block.setAttribute('data-indent', String(cur + 1));
              });
            }

            if (window._aiRefreshAllLists) window._aiRefreshAllLists(editBox);
            if (savedOffset >= 0 && setCur) setCur(savedOffset);

            // body 동기화 + dirty
            try {
              var si2 = parseInt(editBox.getAttribute('data-sec-idx'), 10);
              var bi2 = parseInt(editBox.getAttribute('data-blk-idx'), 10);
              if (sectionsState[si2] && sectionsState[si2].blocks && sectionsState[si2].blocks[bi2]) {
                // data-indent 속성을 실제 HTML 속성으로 동기화 (innerHTML에 포함되도록)
                editBox.querySelectorAll('[data-indent]').forEach(function(el) {
                  if (el.dataset.indent && el.dataset.indent !== '0') {
                    el.setAttribute('data-indent', el.dataset.indent);
                  }
                });
                // HTML → items 변환 (구조화된 저장)
                var items = window._htmlToItems(editBox.innerHTML);
                sectionsState[si2].blocks[bi2].items = items;  // HTML 대신 items 저장
                sectionsState[si2].blocks[bi2].body = editBox.innerHTML;  // 기존 호환용
                _markSectionsDirty();
              }
            } catch(e) {}
            return;
          }

          // ─────────────────────────────────
          // Backspace at line start (outdent/unlist)
          // ─────────────────────────────────
          if (ev.key === 'Backspace') {
            if (!window._aiIsAtLineStart) return;
            var check = window._aiIsAtLineStart(editBox);
            if (!check || !check.atStart) return;
            var info = check.info;
            if (!info || info.type !== 'li') return;

            var li = info.node;
            var parentOl = li.parentNode;
            var hostLi = parentOl ? parentOl.parentNode : null;
            var isTopLevel = !hostLi || hostLi.tagName !== 'LI';

            ev.preventDefault();
            if (ev.stopImmediatePropagation) ev.stopImmediatePropagation();
            if (pushUndo) pushUndo();

            if (isTopLevel) {
              if (window._aiUnlistCurrentLi) window._aiUnlistCurrentLi(editBox, li);
            } else {
              var grandOl = hostLi.parentNode;
              if (grandOl && grandOl.tagName === 'OL') {
                if (hostLi.nextSibling) grandOl.insertBefore(li, hostLi.nextSibling);
                else grandOl.appendChild(li);
                if (!parentOl.querySelector('li')) parentOl.remove();
                if (window._aiRefreshAllLists) window._aiRefreshAllLists(editBox);
                var r = document.createRange();
                r.selectNodeContents(li);
                r.collapse(true);
                var s = window.getSelection();
                if (s) { s.removeAllRanges(); s.addRange(r); }
              }
            }

            try {
              var si3 = parseInt(editBox.getAttribute('data-sec-idx'), 10);
              var bi3 = parseInt(editBox.getAttribute('data-blk-idx'), 10);
              if (sectionsState[si3] && sectionsState[si3].blocks && sectionsState[si3].blocks[bi3]) {
                // data-indent 속성을 HTML에 확실히 반영
                editBox.querySelectorAll('[data-indent]').forEach(function(el) {
                  if (el.dataset.indent && el.dataset.indent !== '0') {
                    el.setAttribute('data-indent', el.dataset.indent);
                  }
                });
                                sectionsState[si3].blocks[bi3].body = editBox.innerHTML;
                _markSectionsDirty();
              }
            } catch(e) {}
            return;
          }
        }, true); // ← capture phase
      }
      // === /Document-level Delegation ===

      // 섹션별 편집 모드 상태 (기본 표시 모드)
      if (!window._sectionEditMode) window._sectionEditMode = {};

      async function _reloadNotesCache(){
        try {
          let divId = (window._currentEditContext || {}).divisionId || '';
          const cardTitle = (window._currentEditContext || {}).cardTitle || '';
          if (!divId) {
            const sel = document.getElementById('v2-division-select');
            if (sel) divId = sel.value || '';
          }
          if (!divId || !cardTitle) return;
          const r = await fetch('/notes?division_id=' + encodeURIComponent(divId), { credentials: 'same-origin' });
          if (!r.ok) return;
          const j = await r.json();
          const cards = (j.cards) || ((j.notes || {}).cards) || [];
          const card = cards.find(function(c){ return (c.title || '').trim() === cardTitle.trim(); });
          const map = {};
          if (card) {
            (card.sections || []).forEach(function(sec){
              map[sec.title || ''] = sec.items || [];
            });
          }
          window._notesItemsCache = map;
        } catch(e) { console.error('notes cache reload 실패', e); }
      }

            window._renderManualSections = function(){
        if (!sectionsRoot) return;
        let html = '';
        // 상단 탭바 (시각적 표시 + 스크롤)
        if (sectionsState.length) {
          html += '<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:14px;gap:12px;flex-wrap:wrap;">'
               +   '<div class="ov-tabs" style="margin-bottom:0;flex:1;overflow-x:auto;">'
               +     sectionsState.map(function(sec, idx){
                       const t = (sec.title || '(제목 없음)');
                       const nums = ['①','②','③','④','⑤','⑥','⑦','⑧','⑨','⑩'];
                       const numStr = idx < nums.length ? nums[idx] : (idx+1) + '.';
                       return '<div class="ov-tab-wrap">'
                       + '<button type="button" class="ov-tab ov-sec-draggable ' + (idx === 0 ? 'active' : '') + '" draggable="true" data-sec-jump="' + idx + '" data-sec-drag-idx="' + idx + '" title="드래그하여 순서 변경">' + numStr + ' ' + t + '</button>'
                       + '<span class="ov-tab-actions">'
                       +   '<button type="button" class="ov-tab-icon ov-sec-tab-rename-btn" data-sec-idx="' + idx + '" title="섹션명 수정">✏️</button>'
                       +   '<button type="button" class="ov-tab-icon ov-tab-del ov-sec-del-btn" data-sec-idx="' + idx + '" title="섹션 삭제">🗑️</button>'
                       + '</span>'
                       + '</div>';
                     }).join('')
               +   '</div>'
               +   '<button type="button" id="ov-sec-add-btn" style="background:#0F2C59;color:#fff;border:0;border-radius:10px;padding:8px 16px;font-size:13px;font-weight:700;cursor:pointer;white-space:nowrap;">+ 섹션 추가</button>'
               + '</div>';
        } else {
          html += '<div style="display:flex;justify-content:flex-end;margin-bottom:14px;">'
               +   '<button type="button" id="ov-sec-add-btn" style="background:#0F2C59;color:#fff;border:0;border-radius:10px;padding:8px 16px;font-size:13px;font-weight:700;cursor:pointer;">+ 섹션 추가</button>'
               + '</div>';
        }
        if (!sectionsState.length) {
          html += '<div class="ov-section"><div class="ov-placeholder">섹션이 없습니다. 위의 「+ 섹션 추가」로 추가하세요.</div></div>';
        } else {
          html += sectionsState.map(function(sec, idx){
            const title = sec.title || '(제목 없음)';
            const secId = sec.id || '';
            const blocks = sec.blocks || [];
            let body = '';
            if (!blocks.length) {
              body = '<div class="ov-placeholder">비어 있는 섹션입니다. (Phase 4: 텍스트/파일 블록 추가 예정)</div>';
            } else {
              const editing = !!window._sectionEditMode[idx];
              body = blocks.map(function(b, bIdx){
                if (b && b.kind === 'text') {
                  const txt = (b.body || '');
                  if (editing) {
                    const editorId = 'txt-' + idx + '-' + bIdx;
                    const toolbar = ''
                      + '<div class="ov-rt-toolbar" style="display:flex;flex-wrap:wrap;gap:4px;padding:6px 8px;background:#EEF4FB;border:1px solid #D9E3F1;border-bottom:0;border-radius:10px 10px 0 0;font-size:12px;" data-editor="' + editorId + '">'
                      +   '<select data-cmd="formatBlock" title="크기" style="border:1px solid #D9E3F1;border-radius:6px;padding:2px 4px;background:#fff;font-size:12px;cursor:pointer;">'
                      +     '<option value="div">본문</option>'
                      +     '<option value="h2">제목</option>'
                      +     '<option value="h3">소제목</option>'
                      +   '</select>'
                      +   '<button type="button" data-cmd="bold" title="굵게" style="width:26px;height:26px;border:1px solid #D9E3F1;background:#fff;border-radius:6px;cursor:pointer;font-weight:800;">B</button>'
                      +   '<button type="button" data-cmd="italic" title="기울임" style="width:26px;height:26px;border:1px solid #D9E3F1;background:#fff;border-radius:6px;cursor:pointer;font-style:italic;">I</button>'
                      +   '<button type="button" data-cmd="underline" title="밑줄" style="width:26px;height:26px;border:1px solid #D9E3F1;background:#fff;border-radius:6px;cursor:pointer;text-decoration:underline;">U</button>'
                      +   '<button type="button" data-cmd="strikeThrough" title="취소선" style="width:26px;height:26px;border:1px solid #D9E3F1;background:#fff;border-radius:6px;cursor:pointer;text-decoration:line-through;">S</button>'
                      +   '<button type="button" class="word-color-btn" data-color-cmd="foreColor" data-default-color="#000000" title="글자색" style="display:inline-flex;align-items:center;justify-content:center;gap:4px;border:1px solid #D9E3F1;background:#fff;border-radius:6px;padding:0 6px;height:26px;cursor:pointer;">'
                      +   '<span style="font-weight:800;font-size:13px;line-height:1;color:#111827;">A</span>'
                      +   '<span class="word-color-preview" style="display:inline-block;width:12px;height:3px;border-radius:999px;background:#000000;"></span>'
                      +   '</button>'
                      +   '<button type="button" class="word-color-btn" data-color-cmd="hiliteColor" data-default-color="#FFF3B0" title="배경색" style="display:inline-flex;align-items:center;justify-content:center;gap:4px;border:1px solid #D9E3F1;background:#fff;border-radius:6px;padding:0 6px;height:26px;cursor:pointer;">'
                      +   '<span style="font-weight:700;font-size:12px;line-height:1;color:#111827;">형광</span>'
                      +   '<span class="word-color-preview" style="display:inline-block;width:12px;height:8px;border-radius:3px;background:#FFF3B0;border:1px solid #E5E7EB;"></span>'
                      +   '</button>'
                      +   '<button type="button" data-cmd="justifyLeft" title="왼쪽 정렬" style="width:28px;height:26px;border:1px solid #D9E3F1;background:#fff;border-radius:6px;cursor:pointer;padding:0;display:inline-flex;align-items:center;justify-content:center;"><svg width=\"14\" height=\"14\" viewBox=\"0 0 24 24\" fill=\"none\" stroke=\"#4b5563\" stroke-width=\"2\" stroke-linecap=\"round\" stroke-linejoin=\"round\"><line x1=\"17\" y1=\"10\" x2=\"3\" y2=\"10\"></line><line x1=\"21\" y1=\"6\" x2=\"3\" y2=\"6\"></line><line x1=\"21\" y1=\"14\" x2=\"3\" y2=\"14\"></line><line x1=\"17\" y1=\"18\" x2=\"3\" y2=\"18\"></line></svg></button>'
                      +   '<button type="button" data-cmd="justifyCenter" title="가운데 정렬" style="width:28px;height:26px;border:1px solid #D9E3F1;background:#fff;border-radius:6px;cursor:pointer;padding:0;display:inline-flex;align-items:center;justify-content:center;"><svg width=\"14\" height=\"14\" viewBox=\"0 0 24 24\" fill=\"none\" stroke=\"#4b5563\" stroke-width=\"2\" stroke-linecap=\"round\" stroke-linejoin=\"round\"><line x1=\"18\" y1=\"10\" x2=\"6\" y2=\"10\"></line><line x1=\"21\" y1=\"6\" x2=\"3\" y2=\"6\"></line><line x1=\"21\" y1=\"14\" x2=\"3\" y2=\"14\"></line><line x1=\"18\" y1=\"18\" x2=\"6\" y2=\"18\"></line></svg></button>'
                      +   '<button type="button" data-cmd="justifyRight" title="오른쪽 정렬" style="width:28px;height:26px;border:1px solid #D9E3F1;background:#fff;border-radius:6px;cursor:pointer;padding:0;display:inline-flex;align-items:center;justify-content:center;"><svg width=\"14\" height=\"14\" viewBox=\"0 0 24 24\" fill=\"none\" stroke=\"#4b5563\" stroke-width=\"2\" stroke-linecap=\"round\" stroke-linejoin=\"round\"><line x1=\"21\" y1=\"10\" x2=\"7\" y2=\"10\"></line><line x1=\"21\" y1=\"6\" x2=\"3\" y2=\"6\"></line><line x1=\"21\" y1=\"14\" x2=\"3\" y2=\"14\"></line><line x1=\"21\" y1=\"18\" x2=\"7\" y2=\"18\"></line></svg></button>'
                      +   '<button type="button" data-cmd="outdent" title="내어쓰기" style="width:28px;height:26px;border:1px solid #D9E3F1;background:#fff;border-radius:6px;cursor:pointer;padding:0;display:inline-flex;align-items:center;justify-content:center;"><svg width=\"14\" height=\"14\" viewBox=\"0 0 24 24\" fill=\"none\" stroke=\"#4b5563\" stroke-width=\"2\" stroke-linecap=\"round\" stroke-linejoin=\"round\"><polyline points=\"7 8 3 12 7 16\"></polyline><line x1=\"21\" y1=\"12\" x2=\"11\" y2=\"12\"></line><line x1=\"21\" y1=\"6\" x2=\"11\" y2=\"6\"></line><line x1=\"21\" y1=\"18\" x2=\"11\" y2=\"18\"></line></svg></button>'
                      +   '<button type="button" data-cmd="indent" title="들여쓰기" style="width:28px;height:26px;border:1px solid #D9E3F1;background:#fff;border-radius:6px;cursor:pointer;padding:0;display:inline-flex;align-items:center;justify-content:center;"><svg width=\"14\" height=\"14\" viewBox=\"0 0 24 24\" fill=\"none\" stroke=\"#4b5563\" stroke-width=\"2\" stroke-linecap=\"round\" stroke-linejoin=\"round\"><polyline points=\"3 8 7 12 3 16\"></polyline><line x1=\"21\" y1=\"12\" x2=\"11\" y2=\"12\"></line><line x1=\"21\" y1=\"6\" x2=\"11\" y2=\"6\"></line><line x1=\"21\" y1=\"18\" x2=\"11\" y2=\"18\"></line></svg></button>'
                      +   '<button type="button" data-cmd="insertUnorderedList" title="글머리 기호" style="width:28px;height:26px;border:1px solid #D9E3F1;background:#fff;border-radius:6px;cursor:pointer;padding:0;display:inline-flex;align-items:center;justify-content:center;"><svg width=\"14\" height=\"14\" viewBox=\"0 0 24 24\" fill=\"none\" stroke=\"#4b5563\" stroke-width=\"2\" stroke-linecap=\"round\" stroke-linejoin=\"round\"><line x1=\"8\" y1=\"6\" x2=\"21\" y2=\"6\"></line><line x1=\"8\" y1=\"12\" x2=\"21\" y2=\"12\"></line><line x1=\"8\" y1=\"18\" x2=\"21\" y2=\"18\"></line><circle cx=\"4\" cy=\"6\" r=\"1.5\" fill=\"#4b5563\"></circle><circle cx=\"4\" cy=\"12\" r=\"1.5\" fill=\"#4b5563\"></circle><circle cx=\"4\" cy=\"18\" r=\"1.5\" fill=\"#4b5563\"></circle></svg></button>'
                      +   '<select data-list-style="ul" title="불릿 종류" style="border:1px solid #D9E3F1;border-radius:6px;padding:2px 4px;background:#fff;font-size:12px;cursor:pointer;">'
                      +     '<option value="disc">●</option>'
                      +     '<option value="circle">○</option>'
                      +     '<option value="square">■</option>'
                      +     '<option value="ov-dash">－</option>'
                      +     '<option value="ov-triangle">▶</option>'
                      +     '<option value="ov-check">✓</option>'
                      +     '<option value="ov-star">★</option>'
                      +     '<option value="ov-arrow">➜</option>'
                      +     '<option value="ov-diamond">◆</option>'
                      +     '<option value="ov-dot">·</option>'
                      +   '</select>'
                      +   '<button type="button" data-cmd="insertOrderedList" title="번호 매기기" style="width:28px;height:26px;border:1px solid #D9E3F1;background:#fff;border-radius:6px;cursor:pointer;padding:0;display:inline-flex;align-items:center;justify-content:center;"><svg width=\"14\" height=\"14\" viewBox=\"0 0 24 24\" fill=\"none\" stroke=\"#4b5563\" stroke-width=\"2\" stroke-linecap=\"round\" stroke-linejoin=\"round\"><line x1=\"10\" y1=\"6\" x2=\"21\" y2=\"6\"></line><line x1=\"10\" y1=\"12\" x2=\"21\" y2=\"12\"></line><line x1=\"10\" y1=\"18\" x2=\"21\" y2=\"18\"></line><text x=\"2\" y=\"8\" font-size=\"6\" fill=\"#4b5563\" stroke=\"none\" font-family=\"Arial\" font-weight=\"700\">1</text><text x=\"2\" y=\"14\" font-size=\"6\" fill=\"#4b5563\" stroke=\"none\" font-family=\"Arial\" font-weight=\"700\">2</text><text x=\"2\" y=\"20\" font-size=\"6\" fill=\"#4b5563\" stroke=\"none\" font-family=\"Arial\" font-weight=\"700\">3</text></svg></button>'
                      +   '<select data-list-style="ol" title="번호 종류 (오른쪽 클릭으로 이어/재시작)" style="border:1px solid #D9E3F1;border-radius:6px;padding:2px 4px;background:#fff;font-size:12px;cursor:pointer;">'
                      +     '<option value="decimal">1.</option>'
                      +     '<option value="ov-circled">①</option>'
                      +     '<option value="ov-paren">1)</option>'
                      +     '<option value="lower-alpha">a.</option>'
                      +     '<option value="upper-alpha">A.</option>'
                      +     '<option value="lower-roman">i.</option>'
                      +     '<option value="upper-roman">I.</option>'
                      +   '</select>'
                                            + '</div>';
                    return '<div class="ov-block ov-block-text" data-sec-idx="' + idx + '" data-blk-idx="' + bIdx + '" style="position:relative;margin-bottom:10px;background:#F8FBFF;border:1px solid #E6EBF2;border-radius:12px;">'
                      +   toolbar
                      +   '<div contenteditable="true" spellcheck="false" class="ov-block-text-body" data-sec-idx="' + idx + '" data-blk-idx="' + bIdx + '" id="' + editorId + '" style="outline:none;font-size:14px;color:#111827;line-height:1.6;white-space:pre-wrap;min-height:60px;padding:12px 14px;border-top:0;">' + txt + '</div>'
                      +   '<div style="position:absolute;top:6px;right:6px;display:flex;gap:2px;z-index:2;">'
                      +     '<button type="button" class="ov-block-up" data-sec-idx="' + idx + '" data-blk-idx="' + bIdx + '" ' + (bIdx === 0 ? 'disabled' : '') + ' style="background:transparent;color:' + (bIdx === 0 ? '#CBD5E1' : '#5B7BB0') + ';border:0;font-size:14px;cursor:' + (bIdx === 0 ? 'not-allowed' : 'pointer') + ';padding:2px 6px;" title="위로">▲</button>'
                      +     '<button type="button" class="ov-block-down" data-sec-idx="' + idx + '" data-blk-idx="' + bIdx + '" ' + (bIdx === blocks.length - 1 ? 'disabled' : '') + ' style="background:transparent;color:' + (bIdx === blocks.length - 1 ? '#CBD5E1' : '#5B7BB0') + ';border:0;font-size:14px;cursor:' + (bIdx === blocks.length - 1 ? 'not-allowed' : 'pointer') + ';padding:2px 6px;" title="아래로">▼</button>'
                      +     '<button type="button" class="ov-block-ai" data-sec-idx="' + idx + '" data-blk-idx="' + bIdx + '" style="background:#F3EEFB;color:#7C3AED;border:1px solid #DDD6FE;font-size:11px;cursor:pointer;padding:3px 8px;font-weight:700;border-radius:6px;margin-right:2px;" title="AI 정리">🤖 AI 정리</button>'
                      +     '<button type="button" class="ov-block-del" data-sec-idx="' + idx + '" data-blk-idx="' + bIdx + '" style="background:transparent;color:#B8302E;border:0;font-size:16px;cursor:pointer;padding:2px 6px;" title="삭제">✕</button>'
                      +   '</div>'
                      + '</div>';
                  }
                  // 표시 모드: 그냥 예쁘게 + 일정 관리 패널
                  const secTitle = sec.title || '';
                  const itemsForSec = (window._notesItemsCache || {})[secTitle] || [];
                  const visibleItems = itemsForSec.filter(function(it){
                    return it && it.text && it.text.trim() && it.type !== 'photo';
                  });
                  let panelHtml = '';
                  if (visibleItems.length > 0) {
                    const rows = visibleItems.map(function(it){
                      const autoRaw = it.due_date_auto || '';
                      const override = it.due_date_override || '';
                      const autoHidden = !!it.auto_due_hidden;
                      // autoHidden 이면 auto 를 표시상 없는 값으로 처리
                      const auto = autoHidden ? '' : autoRaw;
                      const isOverride = !!override;
                      const effective = override || auto;
                      const mmdd = effective ? (effective.slice(5,7) + '/' + effective.slice(8,10)) : '';
                      // 상태 판별:
                      //   - autoHidden=true: 상태 4 (자동값 숨겨짐) → + 날짜 지정 + ↺(자동 복원)
                      //   - effective 있음: 상태 2/3 → chip 표시
                      //   - 그 외: 상태 1 → + 날짜 지정
                      const chipLabel = effective
                        ? (isOverride ? ('수동 ' + mmdd) : ('자동 ' + mmdd))
                        : '+ 날짜 지정';
                      const chipStyle = effective
                        ? (isOverride
                            ? 'background:#DBEAFE;color:#1E40AF;border:1px solid #93C5FD;'
                            : 'background:#F1F3F5;color:#6B7280;border:1px solid #E5E7EB;')
                        : 'background:transparent;color:#9CA3AF;border:1px dashed #D1D5DB;';
                      // ↺ 버튼: 수동값 있을 때 (상태 3) or 자동 숨겨짐 (상태 4)
                      let resetBtn = '';
                      if (isOverride) {
                        resetBtn = '<button type="button" class="ov-due-reset" data-item-id="' + (it.item_id || '') + '" data-auto="' + auto + '" data-mode="override" title="자동값(' + (auto ? auto.slice(5).replace('-','/') : '') + ')으로 되돌리기" style="background:transparent;border:0;color:#1E40AF;cursor:pointer;padding:0 4px;font-size:14px;line-height:1;margin-left:2px;">↺</button>';
                      } else if (autoHidden) {
                        resetBtn = '<button type="button" class="ov-due-reset" data-item-id="' + (it.item_id || '') + '" data-mode="hidden" title="자동값 복원" style="background:transparent;border:0;color:#6B7280;cursor:pointer;padding:0 4px;font-size:14px;line-height:1;margin-left:2px;">↺</button>';
                      }
                      const textEsc = (it.text || '').replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;');
                      return '<div class="ov-schedule-row" style="display:flex;align-items:center;gap:10px;padding:6px 10px;border-bottom:1px solid #F0F4F9;">'
                        +   '<div style="flex:1;font-size:13px;color:#111827;line-height:1.5;overflow:hidden;text-overflow:ellipsis;white-space:pre-wrap;">' + textEsc + '</div>'
                        +   '<button type="button" class="ov-due-chip" data-item-id="' + (it.item_id || '') + '" data-auto="' + auto + '" data-override="' + override + '" data-section-title="' + secTitle.replace(/"/g,'&quot;') + '" title="자동: ' + (auto || '없음') + (override ? (' / 수동: ' + override) : '') + ' | 클릭하여 날짜 변경" style="' + chipStyle + 'border-radius:12px;padding:3px 10px;font-size:11px;font-weight:700;cursor:pointer;white-space:nowrap;flex-shrink:0;">' + chipLabel + '</button>'
                        +   resetBtn
                        + '</div>';
                    }).join('');
                    panelHtml = '<div class="ov-schedule-panel" data-sec-idx="' + idx + '" style="margin-top:8px;padding:8px 12px 4px;background:#F8FBFF;border:1px solid #E6EBF2;border-radius:10px;">'
                      +   '<div style="font-size:12px;font-weight:700;color:#5B7BB0;margin-bottom:6px;">📅 이 섹션의 일정</div>'
                      +   rows
                      + '</div>';
                  }
                  return '<div class="ov-rt-view" style="padding:6px 4px;font-size:14px;color:#111827;line-height:1.65;white-space:pre-wrap;">' + txt + '</div>' + panelHtml;
                }
                if (b && b.kind === 'file') {
                  const fname = b.file_name || '(파일)';
                  const url = b.url || '';
                  const isImage = /\.(png|jpg|jpeg|gif|webp|bmp)$/i.test(fname);
                  const isExcel = /\.(xlsx|xls)$/i.test(fname);
                  const previewHtml = (isImage && url)
                    ? '<div style="margin:8px 0;text-align:center;"><img src="' + url + '" alt="' + fname + '" style="max-width:100%;max-height:320px;border-radius:8px;border:1px solid #E6EBF2;"></div>'
                    : (isExcel && url)
                      ? '<div class="ov-xlsx-preview" data-doc="' + docId + '" data-file="' + fname + '" style="margin:8px 0;padding:10px;background:#fff;border:1px dashed #D9E3F1;border-radius:8px;font-size:12px;color:#8593A6;">엑셀 미리보기 로드 중...</div>'
                      : '';
                  if (editing) {
                    return '<div class="ov-block ov-block-file" data-sec-idx="' + idx + '" data-blk-idx="' + bIdx + '" style="position:relative;padding:10px 14px;margin-bottom:10px;background:#F8FBFF;border:1px solid #E6EBF2;border-radius:12px;font-size:13px;color:#111827;">'
                      +   '<div style="display:flex;align-items:center;gap:8px;padding-right:26px;">'
                      +     '<span>📎</span>'
                      +     (url ? '<a href="' + url + '" target="_blank" style="color:#2E5B94;text-decoration:none;font-weight:700;">' + fname + '</a>' : '<span style="font-weight:700;">' + fname + '</span>')
                      +   '</div>'
                      +   previewHtml
                      +   '<div style="position:absolute;top:6px;right:6px;display:flex;gap:2px;z-index:2;">'
                      +     '<button type="button" class="ov-block-up" data-sec-idx="' + idx + '" data-blk-idx="' + bIdx + '" ' + (bIdx === 0 ? 'disabled' : '') + ' style="background:transparent;color:' + (bIdx === 0 ? '#CBD5E1' : '#5B7BB0') + ';border:0;font-size:14px;cursor:' + (bIdx === 0 ? 'not-allowed' : 'pointer') + ';padding:2px 6px;" title="위로">▲</button>'
                      +     '<button type="button" class="ov-block-down" data-sec-idx="' + idx + '" data-blk-idx="' + bIdx + '" ' + (bIdx === blocks.length - 1 ? 'disabled' : '') + ' style="background:transparent;color:' + (bIdx === blocks.length - 1 ? '#CBD5E1' : '#5B7BB0') + ';border:0;font-size:14px;cursor:' + (bIdx === blocks.length - 1 ? 'not-allowed' : 'pointer') + ';padding:2px 6px;" title="아래로">▼</button>'
                      +     '<button type="button" class="ov-block-del" data-sec-idx="' + idx + '" data-blk-idx="' + bIdx + '" style="background:transparent;color:#B8302E;border:0;font-size:16px;cursor:pointer;padding:2px 6px;" title="삭제">✕</button>'
                      +   '</div>'
                      + '</div>';
                  }
                  // 표시 모드
                  return '<div style="padding:8px 12px;margin-bottom:6px;background:#F8FBFF;border:1px solid #E6EBF2;border-radius:10px;font-size:13px;color:#111827;">'
                    +   '<div style="display:flex;align-items:center;gap:8px;">'
                    +     '<span>📎</span>'
                    +     (url ? '<a href="' + url + '" target="_blank" style="color:#2E5B94;text-decoration:none;font-weight:700;">' + fname + '</a>' : '<span style="font-weight:700;">' + fname + '</span>')
                    +   '</div>'
                    +   previewHtml
                    + '</div>';
                }
                return '';
              }).join('');
              // 편집 모드일 때만 블록 추가 버튼 표시
              if (editing) {
                body += '<div style="display:flex;gap:8px;margin-top:10px;">'
                  +      '<button type="button" class="ov-block-add-text" data-sec-idx="' + idx + '" style="background:#EEF4FB;color:#2E5B94;border:0;border-radius:8px;padding:6px 12px;font-size:12px;font-weight:700;cursor:pointer;">+ 텍스트</button>'
                  +      '<label class="ov-block-add-file-label" style="background:#EEF4FB;color:#2E5B94;border-radius:8px;padding:6px 12px;font-size:12px;font-weight:700;cursor:pointer;">+ 파일<input type="file" data-sec-idx="' + idx + '" class="ov-block-add-file" style="display:none;"></label>'
                  +    '</div>';
              }
            }
            // 빈 섹션이면서 편집 모드라면 블록 추가 버튼 노출
            if (!blocks.length && !!window._sectionEditMode[idx]) {
              body = '<div class="ov-placeholder">비어 있는 섹션입니다. 아래 「+ 텍스트」 「+ 파일」 로 추가하세요.</div>'
                +    '<div style="display:flex;gap:8px;margin-top:10px;">'
                +      '<button type="button" class="ov-block-add-text" data-sec-idx="' + idx + '" style="background:#EEF4FB;color:#2E5B94;border:0;border-radius:8px;padding:6px 12px;font-size:12px;font-weight:700;cursor:pointer;">+ 텍스트</button>'
                +      '<label class="ov-block-add-file-label" style="background:#EEF4FB;color:#2E5B94;border-radius:8px;padding:6px 12px;font-size:12px;font-weight:700;cursor:pointer;">+ 파일<input type="file" data-sec-idx="' + idx + '" class="ov-block-add-file" style="display:none;"></label>'
                +    '</div>';
            }
            const isEditing = !!window._sectionEditMode[idx];
            const headBtns = isEditing
              ? ('<button type="button" class="ov-sec-done-btn" data-sec-idx="' + idx + '" style="background:#0F2C59;color:#fff;border:0;border-radius:8px;padding:6px 12px;font-size:12px;font-weight:700;cursor:pointer;">완료</button>'
                + '<button type="button" class="ov-sec-rename-btn" data-sec-idx="' + idx + '" style="background:#EEF4FB;color:#2E5B94;border:0;border-radius:8px;padding:6px 10px;font-size:12px;font-weight:700;cursor:pointer;">이름 수정</button>'
                + '<button type="button" class="ov-sec-del-btn" data-sec-idx="' + idx + '" style="background:#FEE7E7;color:#B8302E;border:0;border-radius:8px;padding:6px 10px;font-size:12px;font-weight:700;cursor:pointer;">삭제</button>')
              : ('<button type="button" class="ov-sec-edit-btn" data-sec-idx="' + idx + '" style="background:#EEF4FB;color:#2E5B94;border:0;border-radius:8px;padding:6px 12px;font-size:12px;font-weight:700;cursor:pointer;">✏ 수정</button>');

            // 매출 계산 박스 (주차별 계획 섹션 + 편집 모드 전용)
            var salesBoxHtml = '';
            var _secT = (sec.title || '').trim();
            var _isWeeklyPlan = (_secT === '주차별 계획' || _secT === '주차별계획' || _secT === '주차계획' || _secT === '주간계획');
            if (_isWeeklyPlan && isEditing) {
              var _sd = sec.sales_data || {};
              var _prices = _sd.prices || {};
              var _weeks = Array.isArray(_sd.weeks) ? _sd.weeks : [];
              var _ss = String(sec.sales_summary || '(완료 버튼을 누르면 자동 계산됩니다)').replace(/</g,'&lt;').replace(/>/g,'&gt;');
              var _sc = String(sec.sales_computed_at || '');
              var _open = (Object.keys(_prices).length || _weeks.length) ? '▼' : '▶';
              var _display = (Object.keys(_prices).length || _weeks.length) ? 'block' : 'none';
              
              // 판가 행 렌더
              var _pricesHtml = '';
              Object.keys(_prices).forEach(function(mname){
                _pricesHtml += ''
                  + '<div class="ov-sales-price-row" data-model="' + mname + '" style="display:flex;align-items:center;gap:6px;margin-bottom:4px;">'
                  +   '<input type="text" class="ov-sales-price-name" value="' + mname + '" style="width:100px;padding:4px 8px;border:1px solid #C5D0E0;border-radius:4px;font-size:12px;" placeholder="모델명">'
                  +   '<input type="number" class="ov-sales-price-val" value="' + (_prices[mname] || 0) + '" step="0.1" style="width:80px;padding:4px 8px;border:1px solid #C5D0E0;border-radius:4px;font-size:12px;text-align:right;">'
                  +   '<span style="font-size:11px;color:#7C8594;">만불</span>'
                  +   '<button type="button" class="ov-sales-price-del" style="margin-left:auto;background:#FEE7E7;color:#B8302E;border:0;border-radius:4px;padding:3px 8px;font-size:11px;cursor:pointer;">삭제</button>'
                  + '</div>';
              });
              if (!Object.keys(_prices).length) {
                _pricesHtml = '<div style="color:#7C8594;font-size:11px;padding:4px 0;">모델을 추가하세요.</div>';
              }
              
              // 주차별 렌더
              var _weeksHtml = '';
              _weeks.forEach(function(w, wi){
                var _wnum = w.week || '';
                var _wmodels = w.models || {};
                var _rowsHtml = '';
                Object.keys(_prices).forEach(function(mname){
                  var _mv = _wmodels[mname] || {plan: 0, actual: 0};
                  _rowsHtml += ''
                    + '<div style="display:flex;align-items:center;gap:6px;margin:3px 0 3px 20px;">'
                    +   '<span style="width:80px;font-size:12px;color:#111827;">' + mname + '</span>'
                    +   '<span style="font-size:11px;color:#7C8594;">계</span>'
                    +   '<input type="number" class="ov-sales-plan" data-week="' + _wnum + '" data-model="' + mname + '" value="' + (_mv.plan || 0) + '" style="width:56px;padding:3px 6px;border:1px solid #C5D0E0;border-radius:4px;font-size:12px;text-align:right;">'
                    +   '<span style="font-size:11px;color:#7C8594;">실</span>'
                    +   '<input type="number" class="ov-sales-actual" data-week="' + _wnum + '" data-model="' + mname + '" value="' + (_mv.actual || 0) + '" style="width:56px;padding:3px 6px;border:1px solid #C5D0E0;border-radius:4px;font-size:12px;text-align:right;">'
                    + '</div>';
                });
                if (!Object.keys(_prices).length) {
                  _rowsHtml = '<div style="color:#7C8594;font-size:11px;padding:4px 0 4px 20px;">먼저 모델을 추가하세요.</div>';
                }
                _weeksHtml += ''
                  + '<div class="ov-sales-week-box" data-week-idx="' + wi + '" style="padding:8px;margin-bottom:6px;background:#fff;border:1px solid #E5EAF0;border-radius:6px;">'
                  +   '<div style="display:flex;align-items:center;gap:6px;margin-bottom:4px;">'
                  +     '<span style="font-size:12px;font-weight:700;color:#111827;">주차 W</span>'
                  +     '<input type="number" class="ov-sales-week-num" value="' + _wnum + '" style="width:60px;padding:4px 8px;border:1px solid #C5D0E0;border-radius:4px;font-size:12px;text-align:center;">'
                  +     '<button type="button" class="ov-sales-week-del" style="margin-left:auto;background:#FEE7E7;color:#B8302E;border:0;border-radius:4px;padding:3px 8px;font-size:11px;cursor:pointer;">삭제</button>'
                  +   '</div>'
                  +   _rowsHtml
                  + '</div>';
              });
              if (!_weeks.length) {
                _weeksHtml = '<div style="color:#7C8594;font-size:11px;padding:4px 0;">주차를 추가하세요.</div>';
              }
              
              salesBoxHtml = ''
                + '<div class="ov-sales-box" data-sec-idx="' + idx + '" style="margin-top:12px;padding:10px;background:#F7F9FC;border:1px dashed #C5D0E0;border-radius:8px;">'
                +   '<div style="display:flex;align-items:center;justify-content:space-between;gap:8px;">'
                +     '<button type="button" class="ov-sales-toggle" data-sec-idx="' + idx + '" style="background:none;border:0;color:#111827;font-weight:700;font-size:13px;cursor:pointer;padding:4px 0;">'
                +       _open + ' 💰 매출 계산'
                +     '</button>'
                +     '<label class="ov-sales-vis-label" data-sec-idx="' + idx + '" style="display:inline-flex;align-items:center;gap:4px;font-size:11px;color:#7C8594;cursor:pointer;user-select:none;">'
                +       '<input type="checkbox" class="ov-sales-vis" data-sec-idx="' + idx + '" ' + ((sec.sales_visible === false) ? '' : 'checked') + ' style="cursor:pointer;">'
                +       '<span>앱에 매출 표시</span>'
                +     '</label>'
                +   '</div>'
                +   '<div class="ov-sales-body" style="display:' + _display + ';margin-top:8px;">'
                +     '<div class="ov-sales-info" data-sec-idx="' + idx + '" style="font-size:11px;color:#7C8594;margin-bottom:8px;padding:6px 8px;background:#EEF3FB;border-radius:6px;">엑셀 파일 인식 중...</div>'
                +     '<div style="font-size:11px;font-weight:700;color:#111827;margin-bottom:6px;">💵 판가 입력 (만불)</div>'
                +     '<div class="ov-sales-prices" data-sec-idx="' + idx + '" style="padding:4px 0;"><div style="color:#7C8594;font-size:11px;">엑셀을 먼저 첨부해 주세요.</div></div>'
                +     '<div style="margin-top:10px;padding:8px 10px;background:#EEF3FB;border-radius:6px;font-size:13px;color:#111827;">'
                +       '<span style="opacity:0.6;">계산 결과:</span>'
                +       '<div style="margin-top:4px;font-weight:700;">' + _ss + '</div>'
                +       (_sc ? '<div style="margin-top:4px;font-size:11px;color:#7C8594;">계산 시각: ' + _sc + '</div>' : '')
                +     '</div>'
                +     '<div style="margin-top:6px;font-size:11px;color:#7C8594;">완료 버튼을 누르면 자동 저장되고 매출이 계산됩니다.</div>'
                +   '</div>'
                + '</div>';
            }
            
                        return '<div class="ov-section" data-sec-id="' + secId + '" data-sec-idx="' + idx + '">'
              +   '<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:14px;">'
              +     '<div class="ov-section-title" style="margin:0;">'
              +       '<span class="ov-section-num">' + (idx+1) + '</span>'
              +       '<span class="ov-sec-title-text" data-sec-idx="' + idx + '" style="cursor:text;">' + title + '</span>'
              +     '</div>'
              +     '<div style="display:flex;gap:6px;">' + headBtns + '</div>'
              +   '</div>'
              +   body
              +   salesBoxHtml
              + '</div>';
          }).join('');
        }
        sectionsRoot.innerHTML = html;

        // ─── 일정 칩 이벤트 바인딩 ───
        sectionsRoot.querySelectorAll('.ov-due-chip').forEach(function(chip){
          chip.addEventListener('click', function(ev){
            ev.stopPropagation();
            const itemId = chip.getAttribute('data-item-id');
            const auto = chip.getAttribute('data-auto') || '';
            const override = chip.getAttribute('data-override') || '';
            if (!itemId) { alert('item_id 없음 — 저장 후 다시 시도'); return; }

            // 기존 popover 제거
            document.querySelectorAll('.date-popover-custom').forEach(function(n){ n.remove(); });

            // 커스텀 popover 생성
            const pop = document.createElement('div');
            pop.className = 'date-popover-custom';
            pop.style.cssText = 'position:fixed;background:#fff;border:1px solid #d1d5db;border-radius:10px;box-shadow:0 6px 20px rgba(0,0,0,0.15);padding:14px;z-index:9999;display:flex;flex-direction:column;gap:12px;min-width:260px;';
            const rect = chip.getBoundingClientRect();
            // 화면 밖으로 나가지 않게 위치 조정
            const popW = 280;
            let popLeft = rect.left;
            if (popLeft + popW > window.innerWidth - 10) {
              popLeft = window.innerWidth - popW - 10;
            }
            pop.style.left = popLeft + 'px';
            pop.style.top = (rect.bottom + 6) + 'px';

            const inp = document.createElement('input');
            inp.type = 'date';
            inp.value = override || auto || '';
            inp.style.cssText = 'padding:6px 10px;border:1px solid #d1d5db;border-radius:6px;font-size:13px;width:100%;';

            const btnRow = document.createElement('div');
            btnRow.style.cssText = 'display:flex;gap:8px;justify-content:flex-end;';

            const btnSave = document.createElement('button');
            btnSave.type = 'button';
            btnSave.textContent = '저장';
            btnSave.style.cssText = 'padding:6px 14px;background:#2563eb;color:#fff;border:0;border-radius:6px;font-size:13px;font-weight:700;cursor:pointer;';

            const btnDelete = document.createElement('button');
            btnDelete.type = 'button';
            btnDelete.textContent = '삭제';
            btnDelete.style.cssText = 'padding:6px 14px;background:#fff;color:#b91c1c;border:1px solid #d1d5db;border-radius:6px;font-size:13px;font-weight:700;cursor:pointer;';

            btnRow.appendChild(btnDelete);
            btnRow.appendChild(btnSave);
            pop.appendChild(inp);
            pop.appendChild(btnRow);
            document.body.appendChild(pop);
            setTimeout(function(){ inp.focus(); try { inp.showPicker && inp.showPicker(); } catch(e) {} }, 100);

            // 외부 클릭 시 닫기
            function closePop(){ if (pop.parentNode) pop.remove(); document.removeEventListener('mousedown', outsideClick, true); }
            function outsideClick(e){ if (!pop.contains(e.target) && e.target !== chip) closePop(); }
            setTimeout(function(){ document.addEventListener('mousedown', outsideClick, true); }, 100);

            btnSave.addEventListener('click', async function(){
              const val = inp.value;
              if (!val) { alert('날짜를 선택하세요'); return; }
              try {
                const divId = (window._currentEditContext || {}).divisionId || '';
                const cardTitle = (window._currentEditContext || {}).cardTitle || '';
                const r = await fetch('/admin/notes/item/due_override', {
                  method: 'POST',
                  headers: {'Content-Type':'application/json'},
                  credentials: 'same-origin',
                  body: JSON.stringify({
                    division_id: divId,
                    card_title: cardTitle,
                    item_id: itemId,
                    due_date: val
                  })
                });
                const j = await r.json();
                if (!r.ok) { alert('저장 실패: ' + (j.detail || r.status)); return; }
                alert('일정 저장 완료: ' + val);
                closePop();
                await _reloadNotesCache();
                if (window._renderManualSections) window._renderManualSections();
              } catch(e) { alert('오류: ' + e.message); }
            });

            btnDelete.addEventListener('click', async function(){
              closePop();
              if (!confirm('이 일정을 삭제하시겠습니까? (자동 파싱된 날짜도 함께 숨겨집니다)')) return;
              try {
                const divId = (window._currentEditContext || {}).divisionId || '';
                const cardTitle = (window._currentEditContext || {}).cardTitle || '';
                const r = await fetch('/admin/notes/item/hide_auto_due', {
                  method: 'POST',
                  headers: {'Content-Type':'application/json'},
                  credentials: 'same-origin',
                  body: JSON.stringify({
                    division_id: divId,
                    card_title: cardTitle,
                    item_id: itemId
                  })
                });
                const j = await r.json();
                if (!r.ok) { alert('삭제 실패: ' + (j.detail || r.status)); return; }
                alert('일정이 삭제되었습니다');
                closePop();
                await _reloadNotesCache();
                if (window._renderManualSections) window._renderManualSections();
              } catch(e) { alert('오류: ' + e.message); }
            });
          });
        });

        sectionsRoot.querySelectorAll('.ov-due-reset').forEach(function(btn){
          btn.addEventListener('click', async function(ev){
            ev.stopPropagation();
            const mode = btn.getAttribute('data-mode') || 'override';
            const _itemId = btn.getAttribute('data-item-id');
            if (!_itemId) return;
            
            // 상태 4: 자동값 복원
            if (mode === 'hidden') {
              if (!confirm('숨긴 자동 일정을 복원하시겠습니까?')) return;
              try {
                const divId = (window._currentEditContext || {}).divisionId || '';
                const cardTitle = (window._currentEditContext || {}).cardTitle || '';
                const r = await fetch('/admin/notes/item/hide_auto_due', {
                  method: 'DELETE',
                  headers: {'Content-Type':'application/json'},
                  credentials: 'same-origin',
                  body: JSON.stringify({
                    division_id: divId,
                    card_title: cardTitle,
                    item_id: _itemId
                  })
                });
                const j = await r.json();
                if (!r.ok) { alert('복원 실패: ' + (j.detail || r.status)); return; }
                alert('자동 일정이 복원되었습니다');
                await _reloadNotesCache();
                if (window._renderManualSections) window._renderManualSections();
              } catch(e) { alert('오류: ' + e.message); }
              return;
            }
            
            // 상태 3: 기존 수동값 리셋 (아래는 원래 로직)
            const itemId = btn.getAttribute('data-item-id');
            if (!itemId) return;
            const autoVal = btn.getAttribute('data-auto') || '';
            const autoLabel = autoVal ? autoVal.slice(5).replace('-','/') : '자동값';
            if (!confirm('수동 지정을 취소하고 자동값(' + autoLabel + ')으로 되돌릴까요?')) return;
            try {
              const divId = (window._currentEditContext || {}).divisionId || '';
              const cardTitle = (window._currentEditContext || {}).cardTitle || '';
              const r = await fetch('/admin/notes/item/due_override', {
                method: 'DELETE',
                headers: {'Content-Type':'application/json'},
                credentials: 'same-origin',
                body: JSON.stringify({
                  division_id: divId,
                  card_title: cardTitle,
                  item_id: itemId
                })
              });
              const j = await r.json();
              if (!r.ok) { alert('리셋 실패: ' + (j.detail || r.status)); return; }
              alert('일정이 삭제되었습니다 (자동값으로 복귀)');
              await _reloadNotesCache();
              if (window._renderManualSections) window._renderManualSections();
            } catch(e) { alert('오류: ' + e.message); }
          });
        });

        // 이벤트 바인딩
        const addBtn = document.getElementById('ov-sec-add-btn');
        if (addBtn) addBtn.addEventListener('click', function(){ window._addSectionFlow(); });
        // 탭 클릭 → 해당 섹션으로 스크롤
        // ─── 섹션 탭 드래그 앤 드롭 (순서 변경) ───
        (function(){
          let _dragSrcIdx = null;
          sectionsRoot.querySelectorAll('.ov-sec-draggable').forEach(function(btn){
            btn.addEventListener('dragstart', function(e){
              _dragSrcIdx = parseInt(btn.getAttribute('data-sec-drag-idx'), 10);
              try { e.dataTransfer.effectAllowed = 'move'; e.dataTransfer.setData('text/plain', String(_dragSrcIdx)); } catch(_){}
              btn.style.opacity = '0.4';
            });
            btn.addEventListener('dragend', function(){
              btn.style.opacity = '';
              sectionsRoot.querySelectorAll('.ov-sec-draggable').forEach(function(b){
                b.style.borderLeft = '';
                b.style.borderRight = '';
              });
            });
            btn.addEventListener('dragover', function(e){
              e.preventDefault();
              try { e.dataTransfer.dropEffect = 'move'; } catch(_){}
              const tgtIdx = parseInt(btn.getAttribute('data-sec-drag-idx'), 10);
              if (_dragSrcIdx === null || _dragSrcIdx === tgtIdx) return;
              sectionsRoot.querySelectorAll('.ov-sec-draggable').forEach(function(b){
                b.style.borderLeft = '';
                b.style.borderRight = '';
              });
              if (tgtIdx > _dragSrcIdx) {
                btn.style.borderRight = '3px solid #0F2C59';
              } else {
                btn.style.borderLeft = '3px solid #0F2C59';
              }
            });
            btn.addEventListener('dragleave', function(){
              btn.style.borderLeft = '';
              btn.style.borderRight = '';
            });
            btn.addEventListener('drop', function(e){
              e.preventDefault();
              btn.style.borderLeft = '';
              btn.style.borderRight = '';
              const tgtIdx = parseInt(btn.getAttribute('data-sec-drag-idx'), 10);
              if (_dragSrcIdx === null || _dragSrcIdx === tgtIdx) { _dragSrcIdx = null; return; }
              if (!Array.isArray(sectionsState)) { _dragSrcIdx = null; return; }
              const moved = sectionsState.splice(_dragSrcIdx, 1)[0];
              sectionsState.splice(tgtIdx, 0, moved);
              _dragSrcIdx = null;
              try { _markSectionsDirty && _markSectionsDirty(); } catch(_){}
              if (window._renderManualSections) window._renderManualSections();
            });
          });
        })();

        // ─── 섹션 탭 드래그 앤 드롭 ───
        (function(){
          let _dragSrcIdx = null;
          sectionsRoot.querySelectorAll('.ov-sec-draggable').forEach(function(btn){
            btn.addEventListener('dragstart', function(e){
              _dragSrcIdx = parseInt(btn.getAttribute('data-sec-drag-idx'), 10);
              try {
                e.dataTransfer.effectAllowed = 'move';
                e.dataTransfer.setData('text/plain', String(_dragSrcIdx));
              } catch(_){}
              btn.style.opacity = '0.4';
            });
            btn.addEventListener('dragend', function(){
              btn.style.opacity = '';
              sectionsRoot.querySelectorAll('.ov-sec-draggable').forEach(function(b){
                b.style.borderLeft = '';
                b.style.borderRight = '';
              });
            });
            btn.addEventListener('dragover', function(e){
              e.preventDefault();
              try { e.dataTransfer.dropEffect = 'move'; } catch(_){}
              const tgtIdx = parseInt(btn.getAttribute('data-sec-drag-idx'), 10);
              if (_dragSrcIdx === null || _dragSrcIdx === tgtIdx) return;
              sectionsRoot.querySelectorAll('.ov-sec-draggable').forEach(function(b){
                b.style.borderLeft = '';
                b.style.borderRight = '';
              });
              if (tgtIdx > _dragSrcIdx) btn.style.borderRight = '3px solid #0F2C59';
              else btn.style.borderLeft = '3px solid #0F2C59';
            });
            btn.addEventListener('dragleave', function(){
              btn.style.borderLeft = '';
              btn.style.borderRight = '';
            });
            btn.addEventListener('drop', function(e){
              e.preventDefault();
              btn.style.borderLeft = '';
              btn.style.borderRight = '';
              const tgtIdx = parseInt(btn.getAttribute('data-sec-drag-idx'), 10);
              if (_dragSrcIdx === null || _dragSrcIdx === tgtIdx) { _dragSrcIdx = null; return; }
              if (!Array.isArray(sectionsState)) { _dragSrcIdx = null; return; }
              const moved = sectionsState.splice(_dragSrcIdx, 1)[0];
              sectionsState.splice(tgtIdx, 0, moved);
              _dragSrcIdx = null;
              try { _markSectionsDirty && _markSectionsDirty(); } catch(_){}
              if (window._renderManualSections) window._renderManualSections();
            });
          });
        })();

        sectionsRoot.querySelectorAll('[data-sec-jump]').forEach(function(tab){
          tab.addEventListener('click', function(){
            const i = parseInt(tab.getAttribute('data-sec-jump'), 10);
            const targetEl = sectionsRoot.querySelector('[data-sec-idx="' + i + '"]');
            if (targetEl) {
              targetEl.scrollIntoView({behavior: 'smooth', block: 'start'});
              // active 상태 갱신
              sectionsRoot.querySelectorAll('[data-sec-jump]').forEach(function(t){ t.classList.remove('active'); });
              tab.classList.add('active');
            }
          });
        });
        // 수정/완료 버튼 바인딩
        sectionsRoot.querySelectorAll('.ov-sec-edit-btn').forEach(function(btn){
          btn.addEventListener('click', function(){
            const i = parseInt(btn.getAttribute('data-sec-idx'), 10);
            window._sectionEditMode[i] = true;
            window._renderManualSections();
          });
        });
        // 매출 계산 박스 토글 + 입력 → sectionsState 동기화
        // ===== 매출 계산 박스 (폼 UI) =====
        (function bindSalesVisibility(){
          sectionsRoot.querySelectorAll('.ov-sales-vis').forEach(function(cb){
            cb.addEventListener('change', function(){
              var sIdx = parseInt(cb.getAttribute('data-sec-idx'), 10);
              if (isNaN(sIdx)) return;
              if (!sectionsState[sIdx]) return;
              sectionsState[sIdx].sales_visible = !!cb.checked;
              if (typeof _markSectionsDirty === 'function') _markSectionsDirty();
            });
          });
        })();

        // ============================================================
        // [SALES v3] xlsx 자동 파싱 + 판가만 입력
        // ============================================================
        function _findSectionXlsx(secIdx){
          // sectionsState[secIdx] 안에서 xlsx 파일 참조 찾기
          // reports_latest.json 구조: section.blocks[].kind='file', file_name, url
          // notes.json fallback: section.items[].photo_ref
          var sec = sectionsState[secIdx];
          if (!sec) return null;
          
          function extractFilename(raw){
            if (!raw) return null;
            var s = String(raw);
            var lo = s.toLowerCase();
            if (lo.indexOf('.xlsx') < 0 && lo.indexOf('.xlsm') < 0) return null;
            // /admin/manual-files/{doc_id}/{filename} 패턴
            var m = s.match(/\/admin\/manual-files\/[^\/]+\/([^\/?#]+)/);
            if (m) return decodeURIComponent(m[1]);
            // 그냥 파일명일 수도
            return s.split('/').pop().split('?')[0];
          }
          
          // 1) blocks[] 탐색 (편집 화면 기본 구조)
          var blocks = sec.blocks || [];
          for (var i = 0; i < blocks.length; i++) {
            var b = blocks[i] || {};
            if (b.kind === 'file' || b.kind === 'xlsx') {
              var fn = extractFilename(b.url) || extractFilename(b.file_url) || extractFilename(b.file_name) || extractFilename(b.href);
              if (fn) return fn;
            }
            // kind 상관없이 URL 필드 훑기
            var candidates = [b.url, b.file_url, b.file_name, b.href, b.text, b.photo_ref];
            for (var k = 0; k < candidates.length; k++) {
              var fn2 = extractFilename(candidates[k]);
              if (fn2) return fn2;
            }
          }
          
          // 2) items[] 탐색 (notes.json fallback)
          var items = sec.items || [];
          for (var j = 0; j < items.length; j++) {
            var it = items[j] || {};
            var refs = [it.url, it.file_url, it.file_name, it.href, it.text, it.photo_ref];
            for (var m2 = 0; m2 < refs.length; m2++) {
              var fn3 = extractFilename(refs[m2]);
              if (fn3) return fn3;
            }
          }
          
          return null;
        }
        
        function _renderSalesPrices(secIdx, parsed, existingPrices){
          var box = sectionsRoot.querySelector('.ov-sales-box[data-sec-idx="' + secIdx + '"]');
          if (!box) return;
          var infoEl = box.querySelector('.ov-sales-info');
          var pricesEl = box.querySelector('.ov-sales-prices');
          if (!pricesEl) return;
          
          var models = (parsed && parsed.models) || [];
          var weeks = (parsed && parsed.weeks) || [];
          
          if (!models.length) {
            if (infoEl) infoEl.textContent = '⚠️ 엑셀에서 모델을 인식하지 못했습니다.';
            pricesEl.innerHTML = '<div style="color:#7C8594;font-size:11px;">엑셀 형식을 확인해 주세요.</div>';
            return;
          }
          
          var weekLabels = weeks.map(function(w){ return 'W' + w.week; }).join(', ');
          if (infoEl) {
            infoEl.innerHTML = '📊 인식된 모델: <b>' + models.join(', ') + '</b><br>📅 인식된 주차: ' + weekLabels;
          }
          
          var prices = existingPrices || {};
          var html = models.map(function(m){
            var v = (prices[m] !== undefined && prices[m] !== null) ? prices[m] : '';
            return '<div class="ov-sales-price-row" data-model="' + m + '" style="display:flex;align-items:center;gap:8px;margin-bottom:6px;">'
              +   '<div style="min-width:120px;font-size:12px;font-weight:600;color:#111827;">' + m + '</div>'
              +   '<input type="number" step="0.1" class="ov-sales-price-input" data-sec-idx="' + secIdx + '" data-model="' + m + '" value="' + v + '" placeholder="판가" style="flex:1;padding:4px 8px;border:1px solid #C5D0E0;border-radius:4px;font-size:12px;" />'
              +   '<span style="font-size:11px;color:#7C8594;">만불</span>'
              + '</div>';
          }).join('');
          pricesEl.innerHTML = html;
          
          // 입력 이벤트 바인딩
          pricesEl.querySelectorAll('.ov-sales-price-input').forEach(function(inp){
            if (inp._bound) return;
            inp._bound = true;
            inp.addEventListener('input', function(){
              var idx = parseInt(inp.dataset.secIdx, 10);
              var name = inp.dataset.model;
              if (!sectionsState[idx]) return;
              sectionsState[idx].sales_prices = sectionsState[idx].sales_prices || {};
              sectionsState[idx].sales_prices[name] = parseFloat(inp.value || '0') || 0;
              if (typeof _markSectionsDirty === 'function') _markSectionsDirty();
            });
          });
        }
        
        async function _loadSalesForSection(secIdx){
          var box = sectionsRoot.querySelector('.ov-sales-box[data-sec-idx="' + secIdx + '"]');
          if (!box) return;
          var infoEl = box.querySelector('.ov-sales-info');
          var pricesEl = box.querySelector('.ov-sales-prices');
          
          var xlsxName = _findSectionXlsx(secIdx);
          if (!xlsxName) {
            if (infoEl) infoEl.textContent = '📎 엑셀 파일이 아직 첨부되지 않았습니다.';
            if (pricesEl) pricesEl.innerHTML = '<div style="color:#7C8594;font-size:11px;">주차별 계획 섹션에 엑셀을 첨부하세요.</div>';
            return;
          }
          
          if (infoEl) infoEl.textContent = '⏳ 엑셀 파싱 중: ' + xlsxName;
          
          try {
            var res = await fetch('/admin/reports/' + encodeURIComponent(docId) + '/section-parse-xlsx', {
              method: 'POST',
              headers: {'Content-Type': 'application/json'},
              body: JSON.stringify({filename: xlsxName})
            });
            if (!res.ok) {
              var t = await res.text();
              throw new Error('HTTP ' + res.status + ': ' + t);
            }
            var data = await res.json();
            var existingPrices = (sectionsState[secIdx] && sectionsState[secIdx].sales_prices) || {};
            _renderSalesPrices(secIdx, data.parsed, existingPrices);
          } catch(e) {
            if (infoEl) infoEl.textContent = '❌ 파싱 실패: ' + e.message;
            if (pricesEl) pricesEl.innerHTML = '';
          }
        }
        
        // toggle 이벤트 (열릴 때 xlsx 자동 파싱)
        sectionsRoot.querySelectorAll('.ov-sales-toggle').forEach(function(btn){
          if (btn._bound) return;
          btn._bound = true;
          btn.addEventListener('click', function(){
            var body = btn.parentElement.querySelector('.ov-sales-body');
            if (!body) return;
            var isHidden = body.style.display === 'none' || !body.style.display;
            body.style.display = isHidden ? 'block' : 'none';
            btn.textContent = (isHidden ? '▼' : '▶') + ' 💰 매출 계산';
            if (isHidden) {
              var secIdx = parseInt(btn.dataset.secIdx, 10);
              _loadSalesForSection(secIdx);
            }
          });
        });
        
        // 이미 펼쳐진 박스가 있으면 자동 로드
        sectionsRoot.querySelectorAll('.ov-sales-box').forEach(function(box){
          var body = box.querySelector('.ov-sales-body');
          if (body && body.style.display === 'block') {
            var secIdx = parseInt(box.dataset.secIdx, 10);
            _loadSalesForSection(secIdx);
          }
        });
        
                        // ─── 편집 세션 알람 지연 발송 헬퍼 ───
        window._manualEditAlarmSession = window._manualEditAlarmSession || {
          active: false,
          sent: false,
          docId: null
        };

        window._markManualEditAlarmSessionActive = function() {
          try {
            if (typeof isManual !== 'undefined' && isManual && typeof docId !== 'undefined' && docId) {
              window._manualEditAlarmSession.active = true;
              window._manualEditAlarmSession.sent = false;
              window._manualEditAlarmSession.docId = docId;
            }
          } catch (e) {
            console.warn('edit session mark failed:', e);
          }
        }

        window._flushManualEditAlarmSession = function() {
          try {
            var s = window._manualEditAlarmSession;
            if (!s || !s.active || s.sent || !s.docId) return;
            s.sent = true;
            var url = '/admin/reports/' + encodeURIComponent(s.docId) + '/edit_done';
            if (navigator.sendBeacon) {
              navigator.sendBeacon(url, '');
            } else {
              fetch(url, { method: 'POST', keepalive: true, credentials: 'same-origin' }).catch(function(){});
            }
          } catch (e) {
            console.warn('edit_done flush failed:', e);
          }
        }

        if (!window._manualEditAlarmSessionBound) {
          window._manualEditAlarmSessionBound = true;

          // 탭 닫기 / 새로고침 / 다른 페이지 이동
          window.addEventListener('pagehide', function() {
            window._flushManualEditAlarmSession();
          });

          // 사파리/일부 브라우저 보완
          document.addEventListener('visibilitychange', function() {
            if (document.visibilityState === 'hidden') {
              window._flushManualEditAlarmSession();
            }
          });

          // SPA 내부 네비게이션: 캡처 단계에서 모든 클릭 감지
          // 보고서 카드 / 네비 아이템 / 뒤로가기 등 클릭 시 세션 flush
          document.addEventListener('click', function(ev) {
            try {
              var s = window._manualEditAlarmSession;
              if (!s || !s.active || s.sent || !s.docId) return;
              var t = ev.target;
              if (!t || !t.closest) return;
              // 네비게이션성 클릭 대상: 사이드바, 보고서 카드 열기, 목록 버튼 등
              var nav = t.closest('.nav-item, .ov-open-report, .ov-report-item, a[href], button[data-nav]');
              // 단, 현재 보고서 내부의 편집 관련 클릭은 제외
              var inner = t.closest('.ov-sec-done-btn, .ov-block-add-text, .ov-block-add-file, .ov-rt-toolbar, [contenteditable]');
              if (nav && !inner) {
                window._flushManualEditAlarmSession();
              }
            } catch (e) {}
          }, true);
        }

        sectionsRoot.querySelectorAll('.ov-sec-done-btn').forEach(function(btn){
          btn.addEventListener('click', async function(){
            const i = parseInt(btn.getAttribute('data-sec-idx'), 10);
            // 1) 서버 저장 (수기 프로젝트만)
            if (typeof isManual !== 'undefined' && isManual && typeof docId !== 'undefined') {
              const origLabel = btn.textContent;
              btn.disabled = true;
              btn.textContent = '💾 저장 중...';
              try {
                const payload = { products: [{ sections: sectionsState }], edit_mode: true };
                const res = await fetch('/admin/reports/' + encodeURIComponent(docId), {
                  method: 'PUT',
                  headers: { 'Content-Type': 'application/json' },
                  body: JSON.stringify(payload)
                });
                if (!res.ok) {
                  alert('저장 실패: HTTP ' + res.status);
                  btn.textContent = origLabel;
                  btn.disabled = false;
                  return;
                }
                btn.textContent = '✓ 저장됨';
                window._markManualEditAlarmSessionActive();
                // ── P3: 저장 후 캐시 리로드 (auto값이 바뀌었을 수 있음) ──
                try {
                  if (typeof _reloadNotesCache === 'function') {
                    await _reloadNotesCache();
                  }
                } catch(e) { console.warn('cache reload skipped:', e); }
                setTimeout(function(){
                  window._sectionEditMode[i] = false;
                  window._renderManualSections();
                }, 400);
              } catch(e) {
                alert('저장 오류: ' + e.message);
                btn.textContent = origLabel;
                btn.disabled = false;
                return;
              }
            } else {
              // 비수기 프로젝트: 편집 모드만 종료
              window._sectionEditMode[i] = false;
              window._renderManualSections();
            }
          });
        });
        // 블록 이벤트 바인딩
        sectionsRoot.querySelectorAll('.ov-block-add-text').forEach(function(btn){
          btn.addEventListener('click', function(){
            const i = parseInt(btn.getAttribute('data-sec-idx'), 10);
            if (!sectionsState[i]) return;
            if (!Array.isArray(sectionsState[i].blocks)) sectionsState[i].blocks = [];
            sectionsState[i].blocks.push({ kind: 'text', body: '' });
            _markSectionsDirty();
            window._renderManualSections();
          });
        });
        sectionsRoot.querySelectorAll('.ov-block-add-file').forEach(function(input){
          input.addEventListener('change', async function(){
            const i = parseInt(input.getAttribute('data-sec-idx'), 10);
            const f = input.files && input.files[0];
            if (!f) return;
            const fd = new FormData();
            fd.append('file', f);
            try {
              const res = await fetch('/admin/reports/' + encodeURIComponent(docId) + '/section-file', {
                method: 'POST', body: fd
              });
              if (!res.ok) { alert('파일 업로드 실패'); return; }
              const j = await res.json();
              if (!sectionsState[i]) return;
              if (!Array.isArray(sectionsState[i].blocks)) sectionsState[i].blocks = [];
              sectionsState[i].blocks.push({ kind: 'file', file_name: j.file_name, url: j.url, size: j.size });
              _markSectionsDirty();
              window._renderManualSections();
            } catch(e){ alert('업로드 오류: ' + e.message); }
          });
        });
        sectionsRoot.querySelectorAll('.ov-block-up').forEach(function(btn){
          btn.addEventListener('click', function(){
            if (btn.hasAttribute('disabled')) return;
            const si = parseInt(btn.getAttribute('data-sec-idx'), 10);
            const bi = parseInt(btn.getAttribute('data-blk-idx'), 10);
            const sec = sectionsState[si];
            if (!sec || !Array.isArray(sec.blocks)) return;
            if (bi <= 0) return;
            const tmp = sec.blocks[bi - 1];
            sec.blocks[bi - 1] = sec.blocks[bi];
            sec.blocks[bi] = tmp;
            _markSectionsDirty();
            window._renderManualSections();
          });
        });
        sectionsRoot.querySelectorAll('.ov-block-down').forEach(function(btn){
          btn.addEventListener('click', function(){
            if (btn.hasAttribute('disabled')) return;
            const si = parseInt(btn.getAttribute('data-sec-idx'), 10);
            const bi = parseInt(btn.getAttribute('data-blk-idx'), 10);
            const sec = sectionsState[si];
            if (!sec || !Array.isArray(sec.blocks)) return;
            if (bi >= sec.blocks.length - 1) return;
            const tmp = sec.blocks[bi + 1];
            sec.blocks[bi + 1] = sec.blocks[bi];
            sec.blocks[bi] = tmp;
            _markSectionsDirty();
            window._renderManualSections();
          });
        });
        sectionsRoot.querySelectorAll('.ov-block-ai').forEach(function(btn){
          btn.addEventListener('click', async function(){
            const si = parseInt(btn.getAttribute('data-sec-idx'), 10);
            const bi = parseInt(btn.getAttribute('data-blk-idx'), 10);
            const sec = sectionsState[si];
            const blk = sec && Array.isArray(sec.blocks) ? sec.blocks[bi] : null;
            if (!sec || !blk) return;
            const _kind = blk.kind || blk.type;
            if (_kind !== 'text') return;

            const original = String(blk.body || '').trim();
            if (!original) {
              alert('비어 있는 텍스트는 다듬을 수 없습니다.');
              return;
            }

            const oldText = btn.textContent;
            btn.textContent = '…';
            btn.disabled = true;

            try {
              const r = await fetch('/admin/reports/' + encodeURIComponent(docId) + '/polish-text', {
                method: 'POST',
                credentials: 'same-origin',
                headers: { 'Content-Type': 'application/json' },
                body: JSON.stringify({
                  text: original,
                  section_title: String(sec.title || '')
                })
              });

              let j = {};
              try { j = await r.json(); } catch(_) {}

              if (!r.ok) {
                throw new Error((j && j.detail) ? String(j.detail) : ('HTTP ' + r.status));
              }

              const polished = String((j && j.polished) || '').trim();
              if (!polished) {
                throw new Error('빈 응답');
              }

              const diffHtmlRaw = String((j && j.diff_html) || '');
              const diffHtml = (window._normalizeAiNumberedHtml ? window._normalizeAiNumberedHtml(diffHtmlRaw) : diffHtmlRaw);

              // 블록에 diff overlay 삽입 + 적용/취소 버튼
              const blockEl = btn.closest('.ov-block');
              const target = blockEl ? blockEl.querySelector('.ov-block-text-body') : null;
              if (!target) {
                if (!confirm('결과를 적용할까요?')) return;
                blk.body = polished;
                _markSectionsDirty();
                window._renderManualSections();
                return;
              }

              const originalHtml = target.innerHTML;
              target.setAttribute('contenteditable', 'false');
              target.innerHTML = ''
                + '<div style="font-size:11px;font-weight:700;color:#7C3AED;margin-bottom:8px;display:flex;align-items:center;gap:6px;">'
                +   '<span>🤖 AI 정리 결과 (빨강=삭제, 초록=추가, 클릭해서 직접 수정 가능)</span>'
                + '</div>'
                + '<style>'
                +   '.ov-ai-diff-edit ol.ov-num{margin:0;padding-left:1.8em;list-style:none;}'
                +   '.ov-ai-diff-edit ol.ov-num > li{position:relative;margin:0 0 6px 0;line-height:1.8;}'
                +   '.ov-ai-diff-edit ol.ov-num-dot{counter-reset:aiitem 0;}'
                +   '.ov-ai-diff-edit ol.ov-num-dot > li{counter-increment:aiitem;}'
                +   '.ov-ai-diff-edit ol.ov-num-dot > li::before{content:counter(aiitem) ". ";position:absolute;left:-1.8em;top:0;font-weight:700;color:#111827;}'
                +   '.ov-ai-diff-edit ol.ov-num-paren{counter-reset:aiitem 0;}'
                +   '.ov-ai-diff-edit ol.ov-num-paren > li{counter-increment:aiitem;}'
                +   '.ov-ai-diff-edit ol.ov-num-paren > li::before{content:counter(aiitem) ") ";position:absolute;left:-1.8em;top:0;font-weight:700;color:#111827;}'
                +   '.ov-ai-diff-edit ol.ov-num-circled > li::before{content:attr(data-marker) " ";position:absolute;left:-1.8em;top:0;font-weight:700;color:#111827;}'
                +   '.ov-ai-diff-edit [data-indent="1"]{margin-left:24px;}'
                +   '.ov-ai-diff-edit [data-indent="2"]{margin-left:48px;}'
                +   '.ov-ai-diff-edit [data-indent="3"]{margin-left:72px;}'
                +   '.ov-ai-diff-edit [data-indent="4"]{margin-left:96px;}'
                +   '.ov-ai-diff-edit [data-indent="5"]{margin-left:120px;}'
                +   '.ov-ai-diff-edit ol.ov-num ol.ov-num{margin-top:6px;}'
                +   '.ov-ai-diff-edit ol.ov-list-paren{counter-reset:none !important;padding-left:1.8em !important;}'
                +   '.ov-ai-diff-edit ol.ov-list-paren > li{counter-increment:unset !important;}'
                +   '.ov-ai-diff-edit ol.ov-list-paren > li::before{content:none !important;}'
                + '</style>'
                + '<div class="ov-ai-diff-edit" contenteditable="true" spellcheck="false" style="outline:none;padding:12px 14px;background:#FFFFFF;border:1px solid #DDD6FE;border-radius:8px;line-height:1.9;min-height:80px;">'
                + diffHtml
                + '</div>'
                + '<div style="font-size:11px;color:#8593A6;margin-top:6px;">💡 팁: 빨강(삭제 예정)을 지우고 초록(추가)만 남기려면 <b>[적용]</b>을 누르세요. 원본 유지하려면 삭제선 텍스트를 그대로 두면 됩니다.</div>';

              const editBox = target.querySelector('.ov-ai-diff-edit');
              // === AI edit box: Tab / Shift+Tab + 우클릭 메뉴 ===
              if (editBox && !editBox.__aiAttachContextMenu) {
                editBox.__aiAttachContextMenu = true;

                // === undo/redo 스택 ===
                editBox._aiUndoStack = [];
                editBox._aiRedoStack = [];
                editBox._aiUndoDebounce = null;
                editBox._aiUndoMax = 50;

                editBox._aiGetCursorOffset = function(){
                  var sel = window.getSelection();
                  if (!sel.rangeCount) return -1;
                  var range = sel.getRangeAt(0);
                  if (!editBox.contains(range.startContainer)) return -1;
                  var pre = range.cloneRange();
                  pre.selectNodeContents(editBox);
                  pre.setEnd(range.startContainer, range.startOffset);
                  return pre.toString().length;
                };

                editBox._aiSetCursorOffset = function(offset){
                  if (offset < 0) return;
                  var walker = document.createTreeWalker(editBox, NodeFilter.SHOW_TEXT, null, false);
                  var node, count = 0;
                  while ((node = walker.nextNode())) {
                    var len = node.textContent.length;
                    if (count + len >= offset) {
                      var range = document.createRange();
                      range.setStart(node, Math.max(0, offset - count));
                      range.collapse(true);
                      var sel = window.getSelection();
                      sel.removeAllRanges();
                      sel.addRange(range);
                      return;
                    }
                    count += len;
                  }
                };

                editBox._aiPushUndo = function(){
                  var snap = { html: editBox.innerHTML, offset: editBox._aiGetCursorOffset() };
                  var top = editBox._aiUndoStack[editBox._aiUndoStack.length - 1];
                  if (top && top.html === snap.html) return;
                  editBox._aiUndoStack.push(snap);
                  if (editBox._aiUndoStack.length > editBox._aiUndoMax) editBox._aiUndoStack.shift();
                  editBox._aiRedoStack = [];
                };

                editBox._aiPushUndoDebounced = function(){
                  if (editBox._aiUndoDebounce) clearTimeout(editBox._aiUndoDebounce);
                  editBox._aiUndoDebounce = setTimeout(function(){
                    editBox._aiPushUndo();
                    editBox._aiUndoDebounce = null;
                  }, 500);
                };

                editBox._aiUndo = function(){
                  if (editBox._aiUndoStack.length === 0) return;
                  var current = { html: editBox.innerHTML, offset: editBox._aiGetCursorOffset() };
                  var prev = editBox._aiUndoStack.pop();
                  editBox._aiRedoStack.push(current);
                  editBox.innerHTML = prev.html;
                  if (window._aiRefreshAllLists) window._aiRefreshAllLists(editBox);
                  editBox._aiSetCursorOffset(prev.offset);
                };

                editBox._aiRedo = function(){
                  if (editBox._aiRedoStack.length === 0) return;
                  var current = { html: editBox.innerHTML, offset: editBox._aiGetCursorOffset() };
                  var next = editBox._aiRedoStack.pop();
                  editBox._aiUndoStack.push(current);
                  editBox.innerHTML = next.html;
                  if (window._aiRefreshAllLists) window._aiRefreshAllLists(editBox);
                  editBox._aiSetCursorOffset(next.offset);
                };

                // 초기 스냅샷
                setTimeout(function(){ editBox._aiPushUndo(); }, 10);

                // 타이핑 시 debounced 스냅샷
                editBox.addEventListener('input', function(){
                  editBox._aiPushUndoDebounced();
                });

                // Ctrl+Z / Cmd+Z 인터셉트
                editBox.addEventListener('keydown', function(ev){
                  var isMod = ev.metaKey || ev.ctrlKey;
                  if (!isMod) return;
                  if (ev.key === 'z' || ev.key === 'Z') {
                    ev.preventDefault();
                    if (ev.stopImmediatePropagation) ev.stopImmediatePropagation();
                    if (ev.shiftKey) editBox._aiRedo();
                    else editBox._aiUndo();
                  } else if (ev.key === 'y' || ev.key === 'Y') {
                    ev.preventDefault();
                    if (ev.stopImmediatePropagation) ev.stopImmediatePropagation();
                    editBox._aiRedo();
                  }
                }, true);

                // 렌더 직후 모든 ol에 depth 기반 style 적용
                setTimeout(function(){
                  if (window._aiRefreshAllLists) window._aiRefreshAllLists(editBox);
                }, 0);

                // Tab / Shift+Tab keydown
                editBox.addEventListener('keydown', function(ev){
                  if (ev.key !== 'Tab') return;
                  if (!window._aiCollectSelectedNodes) return;
                  var collected = window._aiCollectSelectedNodes(editBox);
                  if (collected.lis.length === 0 && collected.blocks.length === 0) return;

                  ev.preventDefault();
                  if (ev.stopImmediatePropagation) ev.stopImmediatePropagation();
                  if (editBox._aiPushUndo) editBox._aiPushUndo();

                  // 커서 위치 저장 (텍스트 오프셋)
                  var savedOffset = editBox._aiGetCursorOffset ? editBox._aiGetCursorOffset() : -1;

                  // li 처리: outdent는 상위부터, indent는 하위부터 (안전한 순회)
                  var lis = collected.lis.slice();
                  var blocks = collected.blocks.slice();

                  if (ev.shiftKey) {
                    // Shift+Tab: outdent / unlist
                    // 역순 정렬 (아래부터 처리해야 위 li 위치가 안 밀림)
                    var sortedLisOut = lis.slice().sort(function(a, b){
                      var pos = a.compareDocumentPosition(b);
                      if (pos & Node.DOCUMENT_POSITION_FOLLOWING) return 1;
                      if (pos & Node.DOCUMENT_POSITION_PRECEDING) return -1;
                      return 0;
                    });
                    sortedLisOut.forEach(function(li){
                      if (!li.parentNode) return;
                      var parentOl = li.parentNode;

                      // depth-override 감소 먼저
                      var ovr = parentOl.dataset.depthOverride;
                      if (ovr != null && ovr !== '') {
                        var currentOvr = parseInt(ovr, 10);
                        var currentDep = window._aiGetListDepth(parentOl, editBox);
                        if (currentOvr > currentDep) {
                          var newOvr = currentOvr - 1;
                          if (newOvr <= currentDep) delete parentOl.dataset.depthOverride;
                          else parentOl.dataset.depthOverride = String(newOvr);
                          return;
                        }
                      }

                      var hostLi = parentOl ? parentOl.parentNode : null;
                      var isTopLevel = !hostLi || hostLi.tagName !== 'LI';
                      if (isTopLevel) {
                        window._aiUnlistCurrentLi(editBox, li);
                      } else {
                        var grandOl = hostLi.parentNode;
                        if (grandOl && grandOl.tagName === 'OL') {
                          if (hostLi.nextSibling) grandOl.insertBefore(li, hostLi.nextSibling);
                          else grandOl.appendChild(li);
                          if (!parentOl.querySelector('li')) parentOl.remove();
                        }
                      }
                    });
                    blocks.forEach(function(block){
                      var cur = parseInt(block.dataset.indent || '0', 10) || 0;
                      if (cur > 0) {
                        block.setAttribute('data-indent', String(cur - 1));
                        if (String(cur - 1) === '0') block.removeAttribute('data-indent');
                      } else if (block.dataset.indent) {
                        delete block.dataset.indent;
                      }
                    });
                  } else {
                    // Tab: indent (옵션 3: selection 모든 li의 depth 무조건 증가)
                    var sortedLis = lis.slice().sort(function(a, b){
                      var pos = a.compareDocumentPosition(b);
                      if (pos & Node.DOCUMENT_POSITION_FOLLOWING) return -1;
                      if (pos & Node.DOCUMENT_POSITION_PRECEDING) return 1;
                      return 0;
                    });
                    sortedLis.forEach(function(li, i){
                      if (!li.parentNode) return;
                      var parentOl = li.parentNode;
                      var prevLi = li.previousElementSibling;
                      var hasPrevLi = prevLi && prevLi.tagName === 'LI' && sortedLis.indexOf(prevLi) < 0;

                      if (hasPrevLi) {
                        // 앞에 selection 아닌 형제 li 있음 → 그 자식 ol로 이동
                        var childOl = window._aiEnsureChildOl(prevLi, window._aiStyleForDepth(window._aiGetListDepth(parentOl, editBox) + 1));
                        childOl.appendChild(li);
                        if (!parentOl.querySelector('li')) parentOl.remove();
                      } else {
                        // 앞에 형제 li 없음 → 자기 ol 자체의 depth 증가
                        // 자기 ol에 다른 li 있으면 분리
                        var siblings = Array.prototype.filter.call(parentOl.children, function(x){ return x.tagName === 'LI'; });
                        var targetOl = parentOl;
                        if (siblings.length > 1) {
                          // 이 li만 새 ol로 분리
                          var isolatedOl = document.createElement('ol');
                          isolatedOl.className = parentOl.className;
                          isolatedOl.dataset.start = '1';
                          isolatedOl.dataset.manualSplit = '1';
                          parentOl.dataset.manualSplit = '1';
                          isolatedOl.appendChild(li);
                          // parentOl 뒤에 삽입 (li가 원래 첫이면 앞에, 아니면 뒤에)
                          if (parentOl.nextSibling) parentOl.parentNode.insertBefore(isolatedOl, parentOl.nextSibling);
                          else parentOl.parentNode.appendChild(isolatedOl);
                          targetOl = isolatedOl;
                        }
                        // targetOl의 depth override 증가
                        var currentDepth = window._aiGetListDepth(targetOl, editBox);
                        var override = targetOl.dataset.depthOverride;
                        var currentEffective = (override != null && override !== '') ? parseInt(override, 10) : currentDepth;
                        var newDepth = currentEffective + 1;
                        if (newDepth <= 5) {
                          targetOl.dataset.depthOverride = String(newDepth);
                        }
                      }
                    });
                    blocks.forEach(function(block){
                      var cur = parseInt(block.dataset.indent || '0', 10) || 0;
                      if (cur < 5) block.setAttribute('data-indent', String(cur + 1));
                    });
                  }

                  if (window._aiRefreshAllLists) window._aiRefreshAllLists(editBox);
                  if (savedOffset >= 0 && editBox._aiSetCursorOffset) editBox._aiSetCursorOffset(savedOffset);
                }, true);

                // Backspace outdent: 라인 맨 앞에서 Backspace 누르면 indent 해제
                window._aiIsAtLineStart = function(editor){
                  var sel = window.getSelection();
                  if (!sel || !sel.rangeCount) return false;
                  var range = sel.getRangeAt(0);
                  if (!range.collapsed) return false;

                  // 시작 컨테이너로부터 부모 li/block 찾기
                  var info = window._aiFindNodeForRange(range.startContainer, editor);
                  if (!info) return false;

                  // 노드 시작 지점부터 커서까지 텍스트가 모두 비어있는지 확인
                  var testRange = document.createRange();
                  testRange.selectNodeContents(info.node);
                  testRange.setEnd(range.startContainer, range.startOffset);
                  var text = testRange.toString();
                  return { atStart: text.length === 0, info: info };
                };

                editBox.addEventListener('keydown', function(ev){
                  if (ev.key !== 'Backspace') return;
                  var check = window._aiIsAtLineStart(editBox);
                  if (!check || !check.atStart) return;

                  var info = check.info;
                  if (info.type === 'li') {
                    var li = info.node;
                    var parentOl = li.parentNode;
                    var hostLi = parentOl ? parentOl.parentNode : null;
                    var isTopLevel = !hostLi || hostLi.tagName !== 'LI';

                    ev.preventDefault();
                    if (ev.stopImmediatePropagation) ev.stopImmediatePropagation();
                    if (editBox._aiPushUndo) editBox._aiPushUndo();

                    if (isTopLevel) {
                      window._aiUnlistCurrentLi(editBox, li);
                    } else {
                      // outdent
                      var grandOl = hostLi.parentNode;
                      if (grandOl && grandOl.tagName === 'OL') {
                        if (hostLi.nextSibling) grandOl.insertBefore(li, hostLi.nextSibling);
                        else grandOl.appendChild(li);
                        if (!parentOl.querySelector('li')) parentOl.remove();
                        if (window._aiRefreshAllLists) window._aiRefreshAllLists(editBox);
                        var r = document.createRange();
                        r.selectNodeContents(li);
                        r.collapse(true);
                        var s = window.getSelection();
                        s.removeAllRanges();
                        s.addRange(r);
                      }
                    }
                    return;
                  }

                  if (info.type === 'block') {
                    var block = info.node;
                    var curIndent = parseInt(block.dataset.indent || '0', 10) || 0;
                    if (curIndent > 0) {
                      ev.preventDefault();
                      if (ev.stopImmediatePropagation) ev.stopImmediatePropagation();
                      if (editBox._aiPushUndo) editBox._aiPushUndo();
                      block.dataset.indent = String(curIndent - 1);
                      if (String(cur - 1) === '0') { block.removeAttribute('data-indent'); delete block.dataset.indent; } else { block.setAttribute('data-indent', String(cur - 1)); block.dataset.indent = String(cur - 1); }
                    }
                    // indent 0이면 기본 동작 (앞 라인과 병합)
                  }
                }, true);

                // 우클릭 메뉴
                editBox.addEventListener('contextmenu', function(ev){
                  var node = ev.target;
                  var inLi = null;
                  var inOl = null;
                  while (node && node !== editBox) {
                    if (node.nodeType === 1) {
                      if (!inLi && node.tagName === 'LI') inLi = node;
                      if (node.tagName === 'OL') { inOl = node; break; }
                    }
                    node = node.parentNode;
                  }
                  if (!inOl) return;

                  // li가 아직 확정 안 됐거나 이 ol의 직계가 아니면 재감지
                  if (!inLi || inLi.parentNode !== inOl) {
                    var lis = Array.prototype.filter.call(inOl.children, function(x){ return x.tagName === 'LI'; });

                    // 1) 좌표 기반: 클릭한 Y와 각 li rect를 비교
                    var candidate = null;
                    var cy = ev.clientY;
                    for (var i = 0; i < lis.length; i++) {
                      var r = lis[i].getBoundingClientRect();
                      if (cy >= r.top && cy <= r.bottom) { candidate = lis[i]; break; }
                    }
                    // 2) 좌표가 li 사이 gap이면 가장 가까운 li
                    if (!candidate && lis.length) {
                      var minDist = Infinity;
                      for (var j = 0; j < lis.length; j++) {
                        var rr = lis[j].getBoundingClientRect();
                        var mid = (rr.top + rr.bottom) / 2;
                        var d = Math.abs(cy - mid);
                        if (d < minDist) { minDist = d; candidate = lis[j]; }
                      }
                    }
                    // 3) selection fallback
                    if (!candidate) {
                      var sel = window.getSelection();
                      if (sel && sel.rangeCount) {
                        var cur = sel.getRangeAt(0).startContainer;
                        while (cur && cur !== inOl) {
                          if (cur.nodeType === 1 && cur.tagName === 'LI' && cur.parentNode === inOl) { candidate = cur; break; }
                          cur = cur.parentNode;
                        }
                      }
                    }
                    // 4) 그래도 없으면 첫 li
                    if (!candidate) candidate = lis[0] || null;

                    inLi = candidate;
                  }
                  ev.preventDefault();
                  if (ev.stopPropagation) ev.stopPropagation();
                  if (ev.stopImmediatePropagation) ev.stopImmediatePropagation();

                  var legacy = document.getElementById('ov-ol-ctxmenu');
                  if (legacy) legacy.remove();

                  var prev = document.getElementById('ov-ai-ol-ctxmenu');
                  if (prev) prev.remove();
                  var menu = document.createElement('div');
                  menu.id = 'ov-ai-ol-ctxmenu';
                  menu.style.cssText = 'position:fixed;z-index:9999;background:#fff;border:1px solid #D9E3F1;border-radius:8px;box-shadow:0 4px 12px rgba(0,0,0,0.12);padding:4px;min-width:180px;font-size:13px;';
                  menu.style.left = ev.clientX + 'px';
                  menu.style.top = ev.clientY + 'px';
                  menu.innerHTML = ''
                    + '<div class="ov-ai-ctx-item" data-action="continue" style="padding:8px 12px;cursor:pointer;border-radius:4px;">이전 번호에서 이어쓰기</div>'
                    + '<div class="ov-ai-ctx-item" data-action="restart" style="padding:8px 12px;cursor:pointer;border-radius:4px;">번호 다시 시작</div>'
                    + '<div class="ov-ai-ctx-item" data-action="set-start" style="padding:8px 12px;cursor:pointer;border-radius:4px;">번호 설정...</div>'
                    + '<div style="height:1px;background:#E5EAF2;margin:4px 0;"></div>'
                    + '<div class="ov-ai-ctx-item" data-action="unlist" style="padding:8px 12px;cursor:pointer;border-radius:4px;">리스트에서 빼기</div>';
                  document.body.appendChild(menu);
                  menu.querySelectorAll('.ov-ai-ctx-item').forEach(function(item){
                    item.addEventListener('mouseenter', function(){ item.style.background = '#EEF4FB'; });
                    item.addEventListener('mouseleave', function(){ item.style.background = 'transparent'; });
                    item.addEventListener('click', function(){
                      var action = item.getAttribute('data-action');
                      if (editBox._aiPushUndo) editBox._aiPushUndo();
                      if (action === 'continue') window._aiContinueOl(editBox, inLi);
                      else if (action === 'restart') window._aiRestartOl(editBox, inLi);
                      else if (action === 'set-start') window._aiSetStartOl(editBox, inLi);
                      else if (action === 'unlist') window._aiUnlistCurrentLi(editBox, inLi);
                      menu.remove();
                    });
                  });
                  setTimeout(function(){
                    document.addEventListener('click', function _closeOnce(){
                      var m = document.getElementById('ov-ai-ol-ctxmenu');
                      if (m) m.remove();
                      document.removeEventListener('click', _closeOnce);
                    }, { once: true });
                  }, 0);
                });
              }
              if (editBox) {
                if (window._aiSplitTextByNewlines) window._aiSplitTextByNewlines(editBox);
                window._attachAutoListBehavior(editBox);
              }
              if (editBox) {
                // 자동 리스트 변환 방지
                editBox.addEventListener('keydown', function(e){
                  if (e.key === 'Enter') {
                    var sel = window.getSelection();
                    var node = sel && sel.rangeCount ? sel.getRangeAt(0).startContainer : null;
                    var cur = node;
                    var insideAiList = false;
                    while (cur && cur !== editBox) {
                      if (cur.nodeType === 1 && cur.tagName === 'LI' && cur.parentNode && cur.parentNode.classList && cur.parentNode.classList.contains('ov-num')) {
                        insideAiList = true;
                        break;
                      }
                      cur = cur.parentNode;
                    }
                    if (insideAiList) {
                      e.preventDefault();
                      if (e.stopImmediatePropagation) e.stopImmediatePropagation();

                      var li = cur;
                      var ol = li && li.parentNode;
                      if (!li || !ol || !sel || !sel.rangeCount) return;

                      function _htmlOfRange(r){
                        var wrap = document.createElement('div');
                        wrap.appendChild(r.cloneContents());
                        return wrap.innerHTML;
                      }

                      function _cleanStartBr(html){
                        return String(html || '').replace(/^(\s|&nbsp;|<br\s*\/?>)+/ig, '');
                      }

                      function _cleanEndBr(html){
                        return String(html || '').replace(/(<br\s*\/?>|\s|&nbsp;)+$/ig, '');
                      }

                      var rangeNow = sel.getRangeAt(0);
                      var liText = (li.textContent || '').replace(/\u00A0/g, ' ').trim();

                      // 빈 항목에서 Enter → 그 자리에서 번호만 제거하고 plain 빈 줄로 전환
                      if (!liText) {
                        var parent = ol.parentNode;
                        var lis = Array.from(ol.children).filter(function(x){ return x.tagName === 'LI'; });
                        var idx = lis.indexOf(li);
                        var afterLis = lis.slice(idx + 1);

                        var afterOl = null;
                        if (afterLis.length) {
                          afterOl = document.createElement('ol');
                          afterOl.className = ol.className;
                          afterOl.style.cssText = ol.style.cssText || '';
                          // counterReset 안 함 (자동 병합 observer가 처리)
                          afterLis.forEach(function(x){ afterOl.appendChild(x); });
                        }

                        var plain = document.createElement('div');
                        plain.innerHTML = '<br>';

                        li.remove();

                        if (!ol.querySelector('li')) {
                          if (afterOl) {
                            parent.insertBefore(plain, ol.nextSibling);
                            parent.insertBefore(afterOl, plain.nextSibling);
                          } else {
                            parent.insertBefore(plain, ol.nextSibling);
                          }
                          ol.remove();
                        } else {
                          parent.insertBefore(plain, ol.nextSibling);
                          if (afterOl) parent.insertBefore(afterOl, plain.nextSibling);
                        }

                        var exitRange = document.createRange();
                        exitRange.selectNodeContents(plain);
                        exitRange.collapse(true);
                        sel.removeAllRanges();
                        sel.addRange(exitRange);
                        return;
                      }

                      // 일반 항목에서 Enter → caret 위치 기준으로 현재 li를 둘로 분할
                      var beforeRange = document.createRange();
                      beforeRange.selectNodeContents(li);
                      beforeRange.setEnd(rangeNow.startContainer, rangeNow.startOffset);

                      var afterRange = document.createRange();
                      afterRange.selectNodeContents(li);
                      afterRange.setStart(rangeNow.startContainer, rangeNow.startOffset);

                      var beforeHtml = _cleanEndBr(_htmlOfRange(beforeRange));
                      var afterHtml = _cleanStartBr(_htmlOfRange(afterRange));

                      li.innerHTML = beforeHtml || '<br>';

                      var newLi = document.createElement('li');
                      newLi.innerHTML = afterHtml || '<br>';

                      if (li.nextSibling) {
                        ol.insertBefore(newLi, li.nextSibling);
                      } else {
                        ol.appendChild(newLi);
                      }

                      var newRange = document.createRange();
                      newRange.selectNodeContents(newLi);
                      newRange.collapse(true);
                      sel.removeAllRanges();
                      sel.addRange(newRange);
                      return;
                    }
                                        // Enter는 <br> 삽입으로 강제 (자동 <li> 생성 방지)
                    e.preventDefault();
                    document.execCommand('insertLineBreak');
                  }
                });
                // 붙여넣기 시 서식 제거 옵션 (선택적)
                editBox.addEventListener('paste', function(e){
                  // 브라우저 기본 붙여넣기 유지
                });
              }

              const bar = document.createElement('div');
              bar.className = 'ov-diff-bar';
              bar.style.cssText = 'display:flex;gap:8px;margin-top:8px;justify-content:flex-end;';
              bar.innerHTML = '<button type="button" class="ov-diff-cancel" style="background:#fff;color:#475467;border:1px solid #D0D5DD;border-radius:6px;padding:6px 14px;font-size:12px;font-weight:600;cursor:pointer;">취소</button>'
                + '<button type="button" class="ov-diff-apply" style="background:#7C3AED;color:#fff;border:0;border-radius:6px;padding:6px 14px;font-size:12px;font-weight:700;cursor:pointer;">적용</button>';

              blockEl.querySelectorAll('.ov-diff-bar').forEach(function(b){ b.remove(); });
              blockEl.appendChild(bar);

              bar.querySelector('.ov-diff-apply').addEventListener('click', function(){
                if (!editBox) { blk.body = polished; _markSectionsDirty(); window._renderManualSections(); return; }
                // diff 태그 벗기기: <del>(삭제) 제거, <ins>(추가) 는 내용만 남김
                const clone = editBox.cloneNode(true);
                clone.querySelectorAll('del').forEach(function(el){ el.remove(); });
                clone.querySelectorAll('ins').forEach(function(el){
                  const span = document.createElement('span');
                  span.innerHTML = el.innerHTML;
                  el.replaceWith(...span.childNodes);
                });
                // 오염 제거: <style> 태그, AI 안내 div 제거
                clone.querySelectorAll('style').forEach(function(el){ el.remove(); });
                clone.querySelectorAll('div').forEach(function(el){
                  var txt = (el.textContent || '').trim();
                  if (txt.indexOf('🤖 AI 정리 결과') === 0 || txt.indexOf('AI 정리 결과 (빨강=삭제') !== -1) {
                    el.remove();
                  }
                });
                const finalHtml = String(clone.innerHTML || '').trim();
                blk.body = finalHtml || polished;
                _markSectionsDirty();
                window._renderManualSections();
              });
              bar.querySelector('.ov-diff-cancel').addEventListener('click', function(){
                target.innerHTML = originalHtml;
                target.setAttribute('contenteditable', 'true');
                bar.remove();
              });
            } catch (e) {
              alert('AI 다듬기 실패: ' + String((e && e.message) || e));
            } finally {
              btn.textContent = oldText;
              btn.disabled = false;
            }
          });
        });
        sectionsRoot.querySelectorAll('.ov-block-del').forEach(function(btn){
          btn.addEventListener('click', function(){
            const si = parseInt(btn.getAttribute('data-sec-idx'), 10);
            const bi = parseInt(btn.getAttribute('data-blk-idx'), 10);
            if (!sectionsState[si] || !sectionsState[si].blocks) return;
            if (!confirm('이 블록을 삭제하시겠습니까?')) return;
            sectionsState[si].blocks.splice(bi, 1);
            _markSectionsDirty();
            window._renderManualSections();
          });
        });
        // 원본 편집박스에 Tab + undo/redo 이식 (idempotent, 재렌더 안전)
        window._attachOriginalTextBehavior = function(editBox){
          if (!editBox || editBox.__originalAttached) return;
          editBox.__originalAttached = true;
          if (window._aiSplitTextByNewlines) window._aiSplitTextByNewlines(editBox);

          // 기존 오염 body 청소 (한 번만)
          if (!editBox.__origBodyCleaned) {
            editBox.__origBodyCleaned = true;
            var changed = false;
            editBox.querySelectorAll('style').forEach(function(el){ el.remove(); changed = true; });
            editBox.querySelectorAll('div').forEach(function(el){
              var txt = (el.textContent || '').trim();
              if (txt.indexOf('🤖 AI 정리 결과') === 0 || txt.indexOf('AI 정리 결과 (빨강=삭제') !== -1) {
                el.remove();
                changed = true;
              }
            });
            if (changed) {
              try {
                var si = parseInt(editBox.getAttribute('data-sec-idx'), 10);
                var bi = parseInt(editBox.getAttribute('data-blk-idx'), 10);
                if (sectionsState[si] && sectionsState[si].blocks && sectionsState[si].blocks[bi]) {
                // data-indent 속성을 HTML에 확실히 반영
                editBox.querySelectorAll('[data-indent]').forEach(function(el) {
                  if (el.dataset.indent && el.dataset.indent !== '0') {
                    el.setAttribute('data-indent', el.dataset.indent);
                  }
                });
                                  sectionsState[si].blocks[bi].body = editBox.innerHTML;
                  _markSectionsDirty();
                }
              } catch(e) {}
            }
          }

          // === undo/redo 스택 (AI 편집박스와 동일 로직) ===
          editBox._origUndoStack = [];
          editBox._origRedoStack = [];
          editBox._origUndoDebounce = null;
          editBox._origUndoMax = 50;

          editBox._origGetCursorOffset = function(){
            var sel = window.getSelection();
            if (!sel.rangeCount) return -1;
            var range = sel.getRangeAt(0);
            if (!editBox.contains(range.startContainer)) return -1;
            var pre = range.cloneRange();
            pre.selectNodeContents(editBox);
            pre.setEnd(range.startContainer, range.startOffset);
            return pre.toString().length;
          };

          editBox._origSetCursorOffset = function(offset){
            if (offset < 0) return;
            var walker = document.createTreeWalker(editBox, NodeFilter.SHOW_TEXT, null, false);
            var node, count = 0;
            while ((node = walker.nextNode())) {
              var len = node.textContent.length;
              if (count + len >= offset) {
                var range = document.createRange();
                range.setStart(node, Math.max(0, offset - count));
                range.collapse(true);
                var sel = window.getSelection();
                sel.removeAllRanges();
                sel.addRange(range);
                return;
              }
              count += len;
            }
          };

          editBox._origPushUndo = function(){
            var snap = { html: editBox.innerHTML, offset: editBox._origGetCursorOffset() };
            var top = editBox._origUndoStack[editBox._origUndoStack.length - 1];
            if (top && top.html === snap.html) return;
            editBox._origUndoStack.push(snap);
            if (editBox._origUndoStack.length > editBox._origUndoMax) editBox._origUndoStack.shift();
            editBox._origRedoStack = [];
          };

          editBox._origPushUndoDebounced = function(){
            if (editBox._origUndoDebounce) clearTimeout(editBox._origUndoDebounce);
            editBox._origUndoDebounce = setTimeout(function(){
              editBox._origPushUndo();
              editBox._origUndoDebounce = null;
            }, 500);
          };

          editBox._origUndo = function(){
            if (editBox._origUndoStack.length === 0) return;
            var current = { html: editBox.innerHTML, offset: editBox._origGetCursorOffset() };
            var prev = editBox._origUndoStack.pop();
            editBox._origRedoStack.push(current);
            editBox.innerHTML = prev.html;
            editBox._origSetCursorOffset(prev.offset);
            // dirty 플래그 및 body 동기화
            try {
              var si = parseInt(editBox.getAttribute('data-sec-idx'), 10);
              var bi = parseInt(editBox.getAttribute('data-blk-idx'), 10);
              if (sectionsState[si] && sectionsState[si].blocks && sectionsState[si].blocks[bi]) {
                // data-indent 속성을 HTML에 확실히 반영
                editBox.querySelectorAll('[data-indent]').forEach(function(el) {
                  if (el.dataset.indent && el.dataset.indent !== '0') {
                    el.setAttribute('data-indent', el.dataset.indent);
                  }
                });
                                sectionsState[si].blocks[bi].body = editBox.innerHTML;
                _markSectionsDirty();
              }
            } catch(e) {}
          };

          editBox._origRedo = function(){
            if (editBox._origRedoStack.length === 0) return;
            var current = { html: editBox.innerHTML, offset: editBox._origGetCursorOffset() };
            var next = editBox._origRedoStack.pop();
            editBox._origUndoStack.push(current);
            editBox.innerHTML = next.html;
            editBox._origSetCursorOffset(next.offset);
            try {
              var si = parseInt(editBox.getAttribute('data-sec-idx'), 10);
              var bi = parseInt(editBox.getAttribute('data-blk-idx'), 10);
              if (sectionsState[si] && sectionsState[si].blocks && sectionsState[si].blocks[bi]) {
                // data-indent 속성을 HTML에 확실히 반영
                editBox.querySelectorAll('[data-indent]').forEach(function(el) {
                  if (el.dataset.indent && el.dataset.indent !== '0') {
                    el.setAttribute('data-indent', el.dataset.indent);
                  }
                });
                                sectionsState[si].blocks[bi].body = editBox.innerHTML;
                _markSectionsDirty();
              }
            } catch(e) {}
          };

          // 초기 스냅샷
          setTimeout(function(){ editBox._origPushUndo(); }, 10);

          // 타이핑 시 debounced push
          editBox.addEventListener('input', function(){
            editBox._origPushUndoDebounced();
          });

          // Ctrl+Z / Cmd+Z / Ctrl+Y 인터셉트
          editBox.addEventListener('keydown', function(ev){
            var isMod = ev.metaKey || ev.ctrlKey;
            if (!isMod) return;
            if (ev.key === 'z' || ev.key === 'Z') {
              ev.preventDefault();
              if (ev.stopImmediatePropagation) ev.stopImmediatePropagation();
              if (ev.shiftKey) editBox._origRedo();
              else editBox._origUndo();
            } else if (ev.key === 'y' || ev.key === 'Y') {
              ev.preventDefault();
              if (ev.stopImmediatePropagation) ev.stopImmediatePropagation();
              editBox._origRedo();
            }
          }, true);

          // Tab / Shift+Tab (AI 편집박스와 동일 로직 - 우리 헬퍼 재사용)
          editBox.addEventListener('keydown', function(ev){
            if (ev.key !== 'Tab') return;
            if (!window._aiCollectSelectedNodes) return;
            var collected = window._aiCollectSelectedNodes(editBox);
            if (collected.lis.length === 0 && collected.blocks.length === 0) return;

            ev.preventDefault();
            if (ev.stopImmediatePropagation) ev.stopImmediatePropagation();
            if (editBox._origPushUndo) editBox._origPushUndo();

            var savedOffset = editBox._origGetCursorOffset ? editBox._origGetCursorOffset() : -1;

            var lis = collected.lis.slice();
            var blocks = collected.blocks.slice();

            if (ev.shiftKey) {
              // Shift+Tab: outdent / unlist
              var sortedLisOut = lis.slice().sort(function(a, b){
                var pos = a.compareDocumentPosition(b);
                if (pos & Node.DOCUMENT_POSITION_FOLLOWING) return 1;
                if (pos & Node.DOCUMENT_POSITION_PRECEDING) return -1;
                return 0;
              });
              sortedLisOut.forEach(function(li){
                if (!li.parentNode) return;
                var parentOl = li.parentNode;
                var ovr = parentOl.dataset.depthOverride;
                if (ovr != null && ovr !== '') {
                  var currentOvr = parseInt(ovr, 10);
                  var currentDep = window._aiGetListDepth(parentOl, editBox);
                  if (currentOvr > currentDep) {
                    var newOvr = currentOvr - 1;
                    if (newOvr <= currentDep) delete parentOl.dataset.depthOverride;
                    else parentOl.dataset.depthOverride = String(newOvr);
                    return;
                  }
                }
                var hostLi = parentOl ? parentOl.parentNode : null;
                var isTopLevel = !hostLi || hostLi.tagName !== 'LI';
                if (isTopLevel) {
                  if (window._aiUnlistCurrentLi) window._aiUnlistCurrentLi(editBox, li);
                } else {
                  var grandOl = hostLi.parentNode;
                  if (grandOl && grandOl.tagName === 'OL') {
                    if (hostLi.nextSibling) grandOl.insertBefore(li, hostLi.nextSibling);
                    else grandOl.appendChild(li);
                    if (!parentOl.querySelector('li')) parentOl.remove();
                  }
                }
              });
              blocks.forEach(function(block){
                var cur = parseInt(block.dataset.indent || '0', 10) || 0;
                if (cur > 0) {
                  block.setAttribute('data-indent', String(cur - 1));
                  if (String(cur - 1) === '0') block.removeAttribute('data-indent');
                } else if (block.dataset.indent) {
                  delete block.dataset.indent;
                }
              });
            } else {
              // Tab: indent
              var sortedLis = lis.slice().sort(function(a, b){
                var pos = a.compareDocumentPosition(b);
                if (pos & Node.DOCUMENT_POSITION_FOLLOWING) return -1;
                if (pos & Node.DOCUMENT_POSITION_PRECEDING) return 1;
                return 0;
              });
              sortedLis.forEach(function(li, i){
                if (!li.parentNode) return;
                var parentOl = li.parentNode;
                var prevLi = li.previousElementSibling;
                var hasPrevLi = prevLi && prevLi.tagName === 'LI' && sortedLis.indexOf(prevLi) < 0;

                if (hasPrevLi) {
                  var childOl = window._aiEnsureChildOl(prevLi, window._aiStyleForDepth(window._aiGetListDepth(parentOl, editBox) + 1));
                  childOl.appendChild(li);
                  if (!parentOl.querySelector('li')) parentOl.remove();
                } else {
                  var siblings = Array.prototype.filter.call(parentOl.children, function(x){ return x.tagName === 'LI'; });
                  var targetOl = parentOl;
                  if (siblings.length > 1) {
                    var isolatedOl = document.createElement('ol');
                    isolatedOl.className = parentOl.className;
                    isolatedOl.dataset.start = '1';
                    isolatedOl.dataset.manualSplit = '1';
                    parentOl.dataset.manualSplit = '1';
                    isolatedOl.appendChild(li);
                    if (parentOl.nextSibling) parentOl.parentNode.insertBefore(isolatedOl, parentOl.nextSibling);
                    else parentOl.parentNode.appendChild(isolatedOl);
                    targetOl = isolatedOl;
                  }
                  var currentDepth = window._aiGetListDepth(targetOl, editBox);
                  var override = targetOl.dataset.depthOverride;
                  var currentEffective = (override != null && override !== '') ? parseInt(override, 10) : currentDepth;
                  var newDepth = currentEffective + 1;
                  if (newDepth <= 5) targetOl.dataset.depthOverride = String(newDepth);
                }
              });
              blocks.forEach(function(block){
                var cur = parseInt(block.dataset.indent || '0', 10) || 0;
                if (cur < 5) block.setAttribute('data-indent', String(cur + 1));
              });
            }

            if (window._aiRefreshAllLists) window._aiRefreshAllLists(editBox);
            if (savedOffset >= 0 && editBox._origSetCursorOffset) editBox._origSetCursorOffset(savedOffset);

            // body 동기화 + dirty
            try {
              var si2 = parseInt(editBox.getAttribute('data-sec-idx'), 10);
              var bi2 = parseInt(editBox.getAttribute('data-blk-idx'), 10);
              if (sectionsState[si2] && sectionsState[si2].blocks && sectionsState[si2].blocks[bi2]) {
                // data-indent 속성을 실제 HTML 속성으로 동기화 (innerHTML에 포함되도록)
                editBox.querySelectorAll('[data-indent]').forEach(function(el) {
                  if (el.dataset.indent && el.dataset.indent !== '0') {
                    el.setAttribute('data-indent', el.dataset.indent);
                  }
                });
                // HTML → items 변환 (구조화된 저장)
                var items = window._htmlToItems(editBox.innerHTML);
                sectionsState[si2].blocks[bi2].items = items;  // HTML 대신 items 저장
                sectionsState[si2].blocks[bi2].body = editBox.innerHTML;  // 기존 호환용
                _markSectionsDirty();
              }
            } catch(e) {}
          }, true);
        };

        // 텍스트 블록 인라인 편집: input 이벤트로 body 동기화 (rich text 이므로 innerHTML 저장)
        sectionsRoot.querySelectorAll('.ov-block-text-body').forEach(function(el){
          el.addEventListener('input', function(){
            const si = parseInt(el.getAttribute('data-sec-idx'), 10);
            const bi = parseInt(el.getAttribute('data-blk-idx'), 10);
            if (!sectionsState[si] || !sectionsState[si].blocks || !sectionsState[si].blocks[bi]) return;
            sectionsState[si].blocks[bi].body = el.innerHTML;
            _markSectionsDirty();
          });
        });
        // Rich text 툴바 바인딩
        function _getListInEditor(editor){
          const sel = window.getSelection();
          let node = sel && sel.anchorNode;
          while (node && node !== editor) {
            if (node.nodeType === 1 && (node.tagName === 'UL' || node.tagName === 'OL')) return node;
            node = node.parentNode;
          }
          return null;
        }
        function _ensureList(editor, ordered){
          editor.focus();
          try { document.execCommand(ordered ? 'insertOrderedList' : 'insertUnorderedList', false, null); } catch(e) {}
          return _getListInEditor(editor);
        }
        function _syncEditorState(editor){
          const si = parseInt(editor.getAttribute('data-sec-idx'), 10);
          const bi = parseInt(editor.getAttribute('data-blk-idx'), 10);
          if (sectionsState[si] && sectionsState[si].blocks && sectionsState[si].blocks[bi]) {
            sectionsState[si].blocks[bi].body = editor.innerHTML;
            _markSectionsDirty();
          }
        }
        function _applyListStyle(editor, listType, styleValue){
          let listEl = _getListInEditor(editor);
          if (!listEl || (listType === 'ul' && listEl.tagName !== 'UL') || (listType === 'ol' && listEl.tagName !== 'OL')) {
            listEl = _ensureList(editor, listType === 'ol');
          }
          if (!listEl) return;
          Array.from(listEl.classList).forEach(function(cls){
            if (cls.indexOf('ov-list-') === 0) listEl.classList.remove(cls);
          });
          listEl.removeAttribute('type');
          listEl.style.listStyleType = '';
          if (String(styleValue).indexOf('ov-') === 0) {
            listEl.classList.add('ov-list-' + styleValue.slice(3));
            listEl.style.listStyleType = 'none';
          } else {
            listEl.style.listStyleType = styleValue;
            if (listType === 'ol') {
              if (styleValue === 'decimal') listEl.setAttribute('type', '1');
              else if (styleValue === 'lower-alpha') listEl.setAttribute('type', 'a');
              else if (styleValue === 'upper-alpha') listEl.setAttribute('type', 'A');
              else if (styleValue === 'lower-roman') listEl.setAttribute('type', 'i');
              else if (styleValue === 'upper-roman') listEl.setAttribute('type', 'I');
            }
          }
          _syncEditorState(editor);
        }
        function _findPreviousOrderedList(root, currentOl){
          const all = Array.from(root.querySelectorAll('ol'));
          const idx = all.indexOf(currentOl);
          if (idx > 0) return all[idx - 1];
          return null;
        }
        function _orderedCount(ol){
          return ol.querySelectorAll('li').length || 0;
        }
        function _continueOrderedList(editor){
          let ol = _getListInEditor(editor);
          if (!ol || ol.tagName !== 'OL') ol = _ensureList(editor, true);
          if (!ol) return;
          const prev = _findPreviousOrderedList(editor.closest('.ov-section') || document, ol);
          let nextStart = 1;
          if (prev) {
            const prevStart = parseInt(prev.getAttribute('start') || '1', 10) || 1;
            const prevCount = _orderedCount(prev);
            nextStart = prevStart + Math.max(prevCount, 1);
          } else {
            const raw = prompt('이어쓰기 시작 번호', ol.getAttribute('start') || '1');
            if (!raw) return;
            nextStart = parseInt(raw, 10);
            if (!Number.isFinite(nextStart) || nextStart < 1) return;
          }
          ol.setAttribute('start', String(nextStart));
          _syncEditorState(editor);
        }
        function _restartOrderedList(editor){
          let ol = _getListInEditor(editor);
          if (!ol || ol.tagName !== 'OL') ol = _ensureList(editor, true);
          if (!ol) return;
          const raw = prompt('번호 시작 값', '1');
          if (!raw) return;
          const start = parseInt(raw, 10);
          if (!Number.isFinite(start) || start < 1) return;
          ol.setAttribute('start', String(start));
          _syncEditorState(editor);
        }

        sectionsRoot.querySelectorAll('.ov-block-text-body').forEach(function(el){
          el.addEventListener('input', function(){
            _syncEditorState(el);
          });
        });

        sectionsRoot.querySelectorAll('.ov-rt-toolbar').forEach(function(bar){

          const editorId = bar.getAttribute('data-editor');
          const editor = document.getElementById(editorId);
          function exec(cmd, val){
            if (!editor) return;
            editor.focus();
            try { document.execCommand(cmd, false, val || null); } catch(e){ console.warn(e); }
            _syncEditorState(editor);
          }
          bar.querySelectorAll('button[data-cmd]').forEach(function(btn){
            btn.addEventListener('mousedown', function(e){ e.preventDefault(); });
            btn.addEventListener('click', function(){ exec(btn.getAttribute('data-cmd')); });
          });
          const sel = bar.querySelector('select[data-cmd]');
          if (sel) {
            sel.addEventListener('change', function(){
              let v = sel.value;
              if (v === 'div') v = 'p';
              exec('formatBlock', v);
            });
          }

    /* === WORD_COLOR_POPOVER_BEGIN === */
    window.__ovSavedRange = window.__ovSavedRange || null;
    window.__ovColorPopover = window.__ovColorPopover || null;
    window.__ovColorAnchor = window.__ovColorAnchor || null;
    window.__ovColorDocBound = window.__ovColorDocBound || false;

    window.getActiveEditor = window.getActiveEditor || function(){
      return document.querySelector('[contenteditable="true"].is-editing, [contenteditable="true"][data-editor-active="1"], [contenteditable="true"]');
    };

    window.saveEditorSelection = function(){
      var sel = window.getSelection();
      if (!sel || sel.rangeCount === 0) return;
      var editor = window.getActiveEditor && window.getActiveEditor();
      var node = sel.anchorNode;
      if (!editor || !node || !editor.contains(node)) return;
      window.__ovSavedRange = sel.getRangeAt(0).cloneRange();
    };

    window.restoreEditorSelection = function(){
      if (!window.__ovSavedRange) return false;
      var sel = window.getSelection();
      if (!sel) return false;
      sel.removeAllRanges();
      sel.addRange(window.__ovSavedRange);
      return true;
    };

    window.closeWordColorPopover = function(){
      if (window.__ovColorPopover && window.__ovColorPopover.parentNode) {
        window.__ovColorPopover.parentNode.removeChild(window.__ovColorPopover);
      }
      window.__ovColorPopover = null;
      window.__ovColorAnchor = null;
    };

    function _hexToRgb(hex){
      var h = (hex || '').replace('#', '').trim();
      if (h.length === 3) h = h.split('').map(function(x){ return x + x; }).join('');
      var n = parseInt(h, 16);
      return { r:(n >> 16) & 255, g:(n >> 8) & 255, b:n & 255 };
    }

    function _rgbToHex(r, g, b){
      return '#' + [r,g,b].map(function(v){
        var s = Math.max(0, Math.min(255, Math.round(v))).toString(16);
        return s.length === 1 ? '0' + s : s;
      }).join('').toUpperCase();
    }

    function _mixHex(a, b, ratio){
      var c1 = _hexToRgb(a), c2 = _hexToRgb(b);
      return _rgbToHex(
        c1.r + (c2.r - c1.r) * ratio,
        c1.g + (c2.g - c1.g) * ratio,
        c1.b + (c2.b - c1.b) * ratio
      );
    }

    function _themeRows(){
      var base = ['#000000','#FFFFFF','#44546A','#E7E6E6','#5B9BD5','#ED7D31','#A5A5A5','#FFC000','#4472C4','#70AD47'];
      return [
        base,
        base.map(function(c){ return _mixHex(c, '#FFFFFF', 0.80); }),
        base.map(function(c){ return _mixHex(c, '#FFFFFF', 0.60); }),
        base.map(function(c){ return _mixHex(c, '#FFFFFF', 0.40); }),
        base.map(function(c){ return _mixHex(c, '#000000', 0.25); }),
        base.map(function(c){ return _mixHex(c, '#000000', 0.50); })
      ];
    }

    function _standardColors(){
      return ['#C00000','#FF0000','#FFC000','#FFFF00','#92D050','#00B050','#00B0F0','#0070C0','#002060','#7030A0'];
    }

    function _makeColorCell(color, onPick){
      var btn = document.createElement('button');
      btn.type = 'button';
      btn.title = color;
      btn.style.width = '18px';
      btn.style.height = '18px';
      btn.style.padding = '0';
      btn.style.border = '1px solid #D1D5DB';
      btn.style.borderRadius = '3px';
      btn.style.background = color;
      btn.style.cursor = 'pointer';
      btn.addEventListener('mousedown', function(e){
        e.preventDefault();
        e.stopPropagation();
        onPick(color);
      });
      return btn;
    }

    window.applyEditorColor = function(cmd, color, triggerBtn){
      var editor = window.getActiveEditor && window.getActiveEditor();
      if (!editor) return;
      editor.focus();
      if (window.restoreEditorSelection) window.restoreEditorSelection();
      try { document.execCommand('styleWithCSS', false, true); } catch(e) {}
      try {
        if (cmd === 'foreColor') document.execCommand('foreColor', false, color || '#000000');
        else document.execCommand('hiliteColor', false, color || '#FFF3B0');
      } catch(e) {}
      if (window.saveEditorSelection) window.saveEditorSelection();

      if (triggerBtn) {
        var preview = triggerBtn.querySelector('.word-color-preview');
        if (preview) preview.style.background = color;
      }
      if (window.closeWordColorPopover) window.closeWordColorPopover();
    };

    window.openWordColorPopover = function(anchorBtn, cmd, defaultColor){
      if (window.__ovColorPopover && window.__ovColorAnchor === anchorBtn) {
        window.closeWordColorPopover();
        return;
      }
      window.closeWordColorPopover();

      var pop = document.createElement('div');
      pop.style.position = 'absolute';
      pop.style.zIndex = '99999';
      pop.style.width = '232px';
      pop.style.padding = '10px';
      pop.style.background = '#FFFFFF';
      pop.style.border = '1px solid #D1D5DB';
      pop.style.borderRadius = '10px';
      pop.style.boxShadow = '0 12px 28px rgba(0,0,0,.14)';

      var rect = anchorBtn.getBoundingClientRect();
      pop.style.left = (window.scrollX + rect.left) + 'px';
      pop.style.top = (window.scrollY + rect.bottom + 8) + 'px';

      function sectionTitle(text){
        var el = document.createElement('div');
        el.textContent = text;
        el.style.fontSize = '11px';
        el.style.fontWeight = '700';
        el.style.color = '#6B7280';
        el.style.margin = '6px 0 6px';
        return el;
      }

      var autoBtn = document.createElement('button');
      autoBtn.type = 'button';
      autoBtn.textContent = (cmd === 'foreColor') ? '자동' : '채우기 없음';
      autoBtn.style.width = '100%';
      autoBtn.style.height = '28px';
      autoBtn.style.marginBottom = '8px';
      autoBtn.style.border = '1px solid #D1D5DB';
      autoBtn.style.borderRadius = '6px';
      autoBtn.style.background = '#fff';
      autoBtn.style.cursor = 'pointer';
      autoBtn.addEventListener('mousedown', function(e){
        e.preventDefault();
        e.stopPropagation();
        window.applyEditorColor(cmd, cmd === 'foreColor' ? '#000000' : '#FFFFFF', anchorBtn);
      });
      pop.appendChild(autoBtn);

      pop.appendChild(sectionTitle('테마 색'));
      var themeWrap = document.createElement('div');
      themeWrap.style.display = 'grid';
      themeWrap.style.gridTemplateColumns = 'repeat(10, 18px)';
      themeWrap.style.gap = '4px';
      _themeRows().forEach(function(row){
        row.forEach(function(color){
          themeWrap.appendChild(_makeColorCell(color, function(picked){
            window.applyEditorColor(cmd, picked, anchorBtn);
          }));
        });
      });
      pop.appendChild(themeWrap);

      pop.appendChild(sectionTitle('표준 색'));
      var stdWrap = document.createElement('div');
      stdWrap.style.display = 'grid';
      stdWrap.style.gridTemplateColumns = 'repeat(10, 18px)';
      stdWrap.style.gap = '4px';
      _standardColors().forEach(function(color){
        stdWrap.appendChild(_makeColorCell(color, function(picked){
          window.applyEditorColor(cmd, picked, anchorBtn);
        }));
      });
      pop.appendChild(stdWrap);

      var moreRow = document.createElement('div');
      moreRow.style.marginTop = '8px';

      var custom = document.createElement('input');
      custom.type = 'color';
      custom.value = defaultColor || '#000000';
      custom.style.width = '100%';
      custom.style.height = '30px';
      custom.style.border = '1px solid #D1D5DB';
      custom.style.borderRadius = '6px';
      custom.style.background = '#fff';
      custom.style.cursor = 'pointer';
      custom.addEventListener('input', function(){
        window.applyEditorColor(cmd, custom.value, anchorBtn);
      });

      moreRow.appendChild(custom);
      pop.appendChild(moreRow);

      if (window.EyeDropper) {
        var dropRow = document.createElement('div');
        dropRow.style.marginTop = '6px';

        var dropBtn = document.createElement('button');
        dropBtn.type = 'button';
        dropBtn.innerHTML = ''
          + '<svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="#111827" stroke-width="2" stroke-linecap="round" stroke-linejoin="round" style="flex:0 0 auto;">'
          +   '<path d="m2 22 1-1h3l9-9"/>'
          +   '<path d="M3 21v-3l9-9"/>'
          +   '<path d="m15 6 3.4-3.4a2.1 2.1 0 1 1 3 3L18 9l.4.4a2.1 2.1 0 1 1-3 3l-3.8-3.8a2.1 2.1 0 1 1 3-3l.4.4Z"/>'
          + '</svg>'
          + '<span style="margin-left:8px;font-weight:500;">스포이드</span>';
        dropBtn.style.width = '100%';
        dropBtn.style.height = '30px';
        dropBtn.style.border = '1px solid #D1D5DB';
        dropBtn.style.borderRadius = '6px';
        dropBtn.style.background = '#fff';
        dropBtn.style.cursor = 'pointer';
        dropBtn.style.fontSize = '12px';
        dropBtn.style.color = '#111827';
        dropBtn.style.display = 'inline-flex';
        dropBtn.style.alignItems = 'center';
        dropBtn.style.justifyContent = 'flex-start';
        dropBtn.style.padding = '0 10px';

        dropBtn.addEventListener('mousedown', function(e){
          e.preventDefault();
          e.stopPropagation();
          try {
            var eyeDropper = new window.EyeDropper();
            eyeDropper.open().then(function(result){
              if (result && result.sRGBHex) {
                window.applyEditorColor(cmd, result.sRGBHex, anchorBtn);
              }
            }).catch(function(){ /* 사용자 취소 */ });
          } catch(err) {
            console.warn('EyeDropper failed', err);
          }
        });

        dropRow.appendChild(dropBtn);
        pop.appendChild(dropRow);
      }

      document.body.appendChild(pop);
      window.__ovColorPopover = pop;
      window.__ovColorAnchor = anchorBtn;
    };

    if (!window.__ovColorDocBound) {
      document.addEventListener('mousedown', function(e){
        var pop = window.__ovColorPopover;
        var anchor = window.__ovColorAnchor;
        if (!pop) return;
        if (pop.contains(e.target)) return;
        if (anchor && anchor.contains(e.target)) return;
        window.closeWordColorPopover();
      }, true);

      document.addEventListener('selectionchange', function(){
        var editor = window.getActiveEditor && window.getActiveEditor();
        var sel = window.getSelection();
        if (!editor || !sel || sel.rangeCount === 0) return;
        var node = sel.anchorNode;
        if (node && editor.contains(node)) window.saveEditorSelection();
      });

      window.__ovColorDocBound = true;
    }
    /* === WORD_COLOR_POPOVER_END === */
          bar.querySelectorAll('.word-color-btn').forEach(function(btn){
            btn.addEventListener('mousedown', function(e){
              e.preventDefault();
              e.stopPropagation();
              if (window.saveEditorSelection) window.saveEditorSelection();
              if (window.openWordColorPopover) {
                window.openWordColorPopover(btn, btn.dataset.colorCmd, btn.dataset.defaultColor || '#000000');
              }
            });
          });
          bar.querySelectorAll('select[data-list-style="ul"]').forEach(function(sel2){
            sel2.addEventListener('change', function(){
              _applyListStyle(editor, 'ul', sel2.value);
            });
          });
          bar.querySelectorAll('select[data-list-style="ol"]').forEach(function(sel3){
            sel3.addEventListener('change', function(){
              _applyListStyle(editor, 'ol', sel3.value);
            });
          });
        });

        // 모든 원본 편집박스에 Tab + undo/redo 부착
        sectionsRoot.querySelectorAll('.ov-block-text-body').forEach(function(editor){
          if (window._attachOriginalTextBehavior) window._attachOriginalTextBehavior(editor);
        });

        // OL 안에서 우클릭 시 커스텀 메뉴 (이어/재시작)
        sectionsRoot.querySelectorAll('.ov-block-text-body').forEach(function(editor){
          editor.addEventListener('contextmenu', function(e){
            // AI 편집 박스 내부 우클릭은 전용 메뉴가 처리
            if (e.target && e.target.closest && e.target.closest('.ov-ai-diff-edit')) return;

            // 커서/클릭 위치가 OL 안에 있는지 확인
            let node = e.target;
            let inOl = null;
            while (node && node !== editor) {
              if (node.nodeType === 1 && node.tagName === 'OL') { inOl = node; break; }
              node = node.parentNode;
            }
            if (!inOl) return; // 기본 브라우저 메뉴 유지
            e.preventDefault();
            // 기존 메뉴 있으면 제거
            const prev = document.getElementById('ov-ol-ctxmenu');
            if (prev) prev.remove();
            const menu = document.createElement('div');
            menu.id = 'ov-ol-ctxmenu';
            menu.style.cssText = 'position:fixed;z-index:9999;background:#fff;border:1px solid #D9E3F1;border-radius:8px;box-shadow:0 4px 12px rgba(0,0,0,0.12);padding:4px;min-width:160px;font-size:13px;';
            menu.style.left = e.clientX + 'px';
            menu.style.top = e.clientY + 'px';
            menu.innerHTML = ''
              + '<div class="ov-ctx-item" data-action="continue" style="padding:8px 12px;cursor:pointer;border-radius:4px;">이전 번호에서 이어쓰기</div>'
              + '<div class="ov-ctx-item" data-action="restart" style="padding:8px 12px;cursor:pointer;border-radius:4px;">번호 다시 시작</div>';
            document.body.appendChild(menu);
            menu.querySelectorAll('.ov-ctx-item').forEach(function(item){
              item.addEventListener('mouseenter', function(){ item.style.background = '#EEF4FB'; });
              item.addEventListener('mouseleave', function(){ item.style.background = 'transparent'; });
              item.addEventListener('click', function(){
                const action = item.getAttribute('data-action');
                if (action === 'continue') _continueOrderedList(editor);
                else if (action === 'restart') _restartOrderedList(editor);
                menu.remove();
              });
            });
            // 바깥 클릭 시 제거
            setTimeout(function(){
              document.addEventListener('click', function _closeOnce(){
                menu.remove();
                document.removeEventListener('click', _closeOnce);
              }, { once: true });
            }, 0);
          });
        });

        sectionsRoot.querySelectorAll('.ov-sec-rename-btn').forEach(function(btn){
          btn.addEventListener('click', function(){
            const i = parseInt(btn.getAttribute('data-sec-idx'), 10);
            window._renameSectionFlow(i);
          });
        });
        sectionsRoot.querySelectorAll('.ov-sec-tab-rename-btn').forEach(function(btn){
          btn.addEventListener('click', function(e){
            e.preventDefault();
            e.stopPropagation();
            const i = parseInt(btn.getAttribute('data-sec-idx'), 10);
            window._renameSectionFlow(i);
          });
        });

        sectionsRoot.querySelectorAll('.ov-sec-del-btn').forEach(function(btn){
          btn.addEventListener('click', function(){
            const i = parseInt(btn.getAttribute('data-sec-idx'), 10);
            window._deleteSectionFlow(i);
          });
        });
      };

      window._addSectionFlow = function(){
        const raw = prompt('추가할 섹션명을 입력하세요\\n(예: 현황, 주차별 출하실적, 주차별 계획, 주요내용, 이슈/리스크, 요청사항)');
        if (raw === null) return;
        const name = String(raw || '').trim();
        if (!name) { alert('섹션명을 입력해주세요'); return; }

        // 완전 동일한 이름이 이미 있으면 차단
        for (let i=0;i<sectionsState.length;i++){
          if (_normSectionName(sectionsState[i].title) === _normSectionName(name)) {
            alert('이미 같은 이름의 섹션이 있습니다: ' + sectionsState[i].title);
            return;
          }
        }

        // 유사 이름 alias 검사
        const canonical = _findAliasCanonical(name);
        if (canonical && _normSectionName(canonical) !== _normSectionName(name)) {
          // 이미 canonical 섹션이 있는지 확인
          const existingIdx = sectionsState.findIndex(function(s){
            return _normSectionName(s.title) === _normSectionName(canonical);
          });
          if (existingIdx >= 0) {
            const ok = confirm('이미 "' + canonical + '" 섹션이 있습니다.\\n이 섹션으로 대체하시겠습니까?\\n\\n[확인] 기존 섹션 사용 (추가 안 함)\\n[취소] 그래도 "' + name + '"으로 추가');
            if (ok) return;
          } else {
            const useCanonical = confirm('"' + name + '" 은(는) 표준명 "' + canonical + '"와 유사합니다.\\n표준명으로 추가할까요?\\n\\n[확인] "' + canonical + '"로 추가\\n[취소] 입력한 그대로 "' + name + '"으로 추가');
            if (useCanonical) {
              sectionsState.push({id: _genSecId(), title: canonical, kind: 'mixed', blocks: []});
              _markSectionsDirty();
              window._renderManualSections();
              return;
            }
          }
        }

        sectionsState.push({id: _genSecId(), title: name, kind: 'mixed', blocks: []});
        _markSectionsDirty();
        window._renderManualSections();
      };

      window._renameSectionFlow = function(idx){
        if (idx < 0 || idx >= sectionsState.length) return;
        const cur = sectionsState[idx].title || '';
        const raw = prompt('새 섹션명을 입력하세요', cur);
        if (raw === null) return;
        const name = String(raw || '').trim();
        if (!name) { alert('섹션명을 입력해주세요'); return; }
        if (_normSectionName(name) === _normSectionName(cur)) return;
        for (let i=0;i<sectionsState.length;i++){
          if (i === idx) continue;
          if (_normSectionName(sectionsState[i].title) === _normSectionName(name)) {
            alert('이미 같은 이름의 섹션이 있습니다: ' + sectionsState[i].title);
            return;
          }
        }
        sectionsState[idx].title = name;
        _markSectionsDirty();
        window._renderManualSections();
      };

      window._deleteSectionFlow = function(idx){
        if (idx < 0 || idx >= sectionsState.length) return;
        const cur = sectionsState[idx].title || '';
        if (!confirm('"' + cur + '" 섹션을 삭제하시겠습니까?')) return;
        sectionsState.splice(idx, 1);
        _markSectionsDirty();
        window._renderManualSections();
      };

      // 엑셀 미리보기 로더
      // 자동 리스트 동작: "1)" "1." "-" "*" + Space → 리스트 시작, Enter → 다음 번호, 빈줄 Enter → 종료

      window._loadXlsxPreviews = function(){
        const previews = document.querySelectorAll('.ov-xlsx-preview');
        previews.forEach(async function(el){
          if (el.getAttribute('data-loaded') === '1') return;
          const doc = el.getAttribute('data-doc') || '';
          const file = el.getAttribute('data-file') || '';
          if (!doc || !file) {
            el.innerHTML = '<div style="color:#b42318;font-size:12px;">미리보기 정보 누락</div>';
            return;
          }
          try {
            const r = await fetch(
              '/admin/manual-files/' + encodeURIComponent(doc) + '/' + encodeURIComponent(file) + '/xlsx-preview',
              { credentials: 'same-origin' }
            );
            if (!r.ok) throw new Error('HTTP ' + r.status);
            const j = await r.json();
            const rows = Array.isArray(j.rows) ? j.rows : [];
            if (!rows.length) {
              el.innerHTML = '<div style="color:#667085;font-size:12px;">비어 있는 시트입니다</div>';
              el.setAttribute('data-loaded', '1');
              return;
            }
            const nRows = j.n_rows || rows.length;
            const nCols = j.n_cols || 0;
            let html = '';
            html += '<div style="font-size:12px;font-weight:700;color:#344054;margin-bottom:8px;">📊 ' + (j.sheet || '첫 시트') + ' (' + nRows + '행 × ' + nCols + '열)</div>';
            html += '<div style="overflow:auto;border:1px solid #D9E3F1;border-radius:8px;background:#fff;">';
            html += '<table style="border-collapse:separate;border-spacing:0;width:100%;font-size:12px;">';
            rows.forEach(function(row){
              html += '<tr>';
              (row || []).forEach(function(cell){
                let rawText = String(cell.text == null ? '' : cell.text);
                let safe = rawText
                  .replace(/&/g,'&amp;')
                  .replace(/</g,'&lt;')
                  .replace(/>/g,'&gt;');
                safe = safe.split(String.fromCharCode(13,10)).join('<br>');
                safe = safe.split(String.fromCharCode(13)).join('<br>');
                safe = safe.split(String.fromCharCode(10)).join('<br>');

                // PO증감/월별 문자열 강제 개행 (정규식 리터럴 금지: 파이썬 템플릿 깨짐 방지)
                const monthTokens = ['3월:', '4월:', '5월:', '6월:', '7월:', '8월:', '9월:', '10월:', '11월:', '12월:'];
                let tokens = [];
                for (let ti = 0; ti < monthTokens.length; ti++) {
                  tokens.push(monthTokens[ti]);
                }
                for (let wi = 1; wi <= 53; wi++) {
                  tokens.push('W' + wi + ':');
                }
                for (let i = 1; i < tokens.length; i++) {
                  safe = safe.split(' ' + tokens[i]).join('<br>' + tokens[i]);
                }
                safe = safe.split('<br> ').join('<br>');
                const rs = cell.rowspan && cell.rowspan > 1 ? ' rowspan="' + cell.rowspan + '"' : '';
                const cs = cell.colspan && cell.colspan > 1 ? ' colspan="' + cell.colspan + '"' : '';
                let style = 'padding:8px 12px;text-align:center;white-space:pre-line;min-width:60px;';
                // 셀별 border (색상 지정 있으면 그거, 없으면 기본 연회색)
                const _bd = cell.borders || {};
                const _bstyle = function(s){
                  const m = { thin: '1px solid', medium: '2px solid', thick: '3px solid', dashed: '1px dashed', dotted: '1px dotted', double: '3px double' };
                  return m[s] || '1px solid';
                };
                // 색상 border는 box-shadow inset으로 그림 (이웃 셀에 안 가려짐)
                const _shadows = [];
                ['top','right','bottom','left'].forEach(function(side){
                  if (_bd[side]) {
                    // border-width 픽셀 값
                    const s = _bd[side].style;
                    const w = (s === 'medium') ? 2 : (s === 'thick' || s === 'double') ? 3 : 1;
                    // inset shadow: inset X Y blur color, side 방향에 따라 offset 조정
                    const off = {top:'0 '+w+'px', right:'-'+w+'px 0', bottom:'0 -'+w+'px', left:w+'px 0'}[side];
                    _shadows.push('inset ' + off + ' 0 ' + _bd[side].color);
                    style += 'border-' + side + ':' + w + 'px solid ' + _bd[side].color + ';';
                  } else {
                    style += 'border-' + side + ':1px solid #D9E3F1;';
                  }
                });
                if (_shadows.length) {
                  style += 'box-shadow:' + _shadows.join(',') + ';position:relative;z-index:2;';
                }
                // 배경색이 있을 때만 폰트색도 반영 (흰색 폰트 + 흰 배경 방지)
                if (cell.bg) {
                  style += 'background:' + cell.bg + ';';
                  if (cell.fg) style += 'color:' + cell.fg + ';';
                  else style += 'color:#fff;';
                } else {
                  style += 'color:#344054;';
                }
                if (cell.bold) style += 'font-weight:700;';
                html += '<td' + rs + cs + ' style="' + style + '">' + safe + '</td>';
              });
              html += '</tr>';
            });
            html += '</table></div>';
            el.innerHTML = html;
            el.style.padding = '10px';
            el.style.border = '1px solid #E5E7EB';
            el.style.background = '#F9FAFB';
            el.setAttribute('data-loaded', '1');
          } catch (e) {
            el.innerHTML = '<div style="color:#b42318;font-size:12px;">엑셀 미리보기 실패: ' + String(e.message || e) + '</div>';
          }
        });
      };

      // 렌더 후 매번 preview 로드
      const _origRender = window._renderManualSections;
      window._renderManualSections = function(){
        _origRender();
        setTimeout(function(){ window._loadXlsxPreviews(); }, 0);
      };

      window._renderManualSections();
    } else {
      // PPT 프로젝트: 하드코딩 3탭 스크롤 바인딩
      const pptTabs = document.getElementById('ov-ppt-tabs');
      if (pptTabs) {
        pptTabs.querySelectorAll('[data-ppt-jump]').forEach(function(tab){
          tab.addEventListener('click', function(){
            const i = parseInt(tab.getAttribute('data-ppt-jump'), 10);
            const allSecs = c.querySelectorAll('.ov-section');
            if (allSecs[i]) {
              allSecs[i].scrollIntoView({behavior: 'smooth', block: 'start'});
              pptTabs.querySelectorAll('[data-ppt-jump]').forEach(function(t){ t.classList.remove('active'); });
              tab.classList.add('active');
            }
          });
        });
      }
    }

    // 4) 상단 버튼 바인딩
    const backBtn = document.getElementById('ov-back-btn');
    if (backBtn) backBtn.addEventListener('click', window.backToReportList);

    // HTML → items 변환 함수 (구조화된 저장)
    window._htmlToItems = function(html) {
      var parser = new DOMParser();
      var doc = parser.parseFromString(html, 'text/html');
      var items = [];
      
      function getDepth(el) {
        // data-indent 속성 또는 부모 요소로 depth 계산
        var indent = el.getAttribute('data-indent');
        if (indent) return parseInt(indent) || 0;
        return 0;
      }
      
      function processNode(node, depth) {
        if (node.nodeType === 3) {  // 텍스트 노드
          var text = node.textContent.trim();
          if (text) {
            items.push({
              type: depth > 0 ? 'sub' : 'bullet',
              indent: depth,
              text: text
            });
          }
        } else if (node.nodeType === 1) {  // 요소 노드
          var tag = node.tagName.toLowerCase();
          
          if (tag === 'div' || tag === 'p') {
            var itemDepth = getDepth(node) || depth;
            var hasBlockChild = node.querySelector('ul, ol, div, p, blockquote');
            
            if (!hasBlockChild) {
              var text = node.textContent.trim();
              if (text) {
                items.push({
                  type: itemDepth > 0 ? 'sub' : 'bullet',
                  indent: itemDepth,
                  text: text
                });
              }
            } else {
              // 자식 노드 재귀 처리
              for (var i = 0; i < node.childNodes.length; i++) {
                processNode(node.childNodes[i], itemDepth);
              }
            }
          } else if (tag === 'ul' || tag === 'ol') {
            var listDepth = getDepth(node) || depth;
            var lis = node.querySelectorAll(':scope > li');
            lis.forEach(function(li, idx) {
              var text = li.textContent.trim();
              if (text) {
                items.push({
                  type: listDepth > 0 ? 'sub' : 'bullet',
                  indent: listDepth,
                  text: text,
                  marker: tag === 'ul' ? '•' : (idx + 1) + '.'
                });
              }
            });
          } else if (tag === 'blockquote') {
            // blockquote 내부는 depth +1
            for (var i = 0; i < node.childNodes.length; i++) {
              processNode(node.childNodes[i], depth + 1);
            }
          } else if (tag === 'br') {
            // 줄바꿈은 무시 (텍스트로 처리됨)
          } else {
            // 기타 태그는 텍스트만 추출
            var text = node.textContent.trim();
            if (text) {
              items.push({
                type: depth > 0 ? 'sub' : 'bullet',
                indent: depth,
                text: text
              });
            }
          }
        }
      }
      
      for (var i = 0; i < doc.body.childNodes.length; i++) {
        processNode(doc.body.childNodes[i], 0);
      }
      
      return items;
    };
    
    // 4-b) 인라인 프로젝트명 편집 + 저장 버튼 활성화
    const labelEl = document.getElementById('ov-project-label');
    const saveBtn = document.getElementById('ov-save-btn');
    const initialLabel = labelEl ? labelEl.textContent.trim() : '';
    let dirty = false;
    function markDirty(){
      dirty = true;
      if (saveBtn) {
        saveBtn.disabled = false;
        saveBtn.style.opacity = '1';
        saveBtn.style.cursor = 'pointer';
      }
    }
    function markClean(){
      dirty = false;
      if (saveBtn) {
        saveBtn.disabled = true;
        saveBtn.style.opacity = '0.4';
        saveBtn.style.cursor = 'not-allowed';
      }
    }
    if (labelEl) {
      labelEl.addEventListener('focus', function(){ labelEl.style.borderBottom = '2px dashed #4A6FA5'; });
      labelEl.addEventListener('blur', function(){ labelEl.style.borderBottom = '2px dashed transparent'; });
      labelEl.addEventListener('input', function(){ if (labelEl.textContent.trim() !== initialLabel) markDirty(); });
      labelEl.addEventListener('keydown', function(e){
        if (e.key === 'Enter') { e.preventDefault(); labelEl.blur(); }
      });
    }
    if (saveBtn) {
      saveBtn.addEventListener('click', async function(){
        if (!dirty || saveBtn.disabled) return;
        const newLabel = labelEl ? labelEl.textContent.trim() : '';
        if (!newLabel) { alert('프로젝트명이 비어있습니다'); return; }
        // display_title에서 프로젝트명만 추출: " · W## 주간보고" 앞까지
        let newProjectName = newLabel;
        let newWeek = null;
        const m = newLabel.match(/^(.*?)(?:\s*·\s*W(\d+)\s*주간보고)?\s*$/);
        if (m) {
          newProjectName = (m[1] || '').trim();
          if (m[2]) newWeek = parseInt(m[2], 10);
        }
        // 원본 프로젝트명(파일명 파싱값)을 키로 사용
        const origProj = (target && target._split_project) || newProjectName;
        const overrides = {};
        overrides[origProj] = newProjectName;
        saveBtn.disabled = true;
        saveBtn.textContent = '저장 중...';
        try {
          const payload = { project_overrides: overrides };
          if (newWeek) payload.week_override = newWeek;
          // 수기 프로젝트: sections 도 함께 저장 (항상 최신 상태 전송)
          if (isManual) {
            payload.products = [{ sections: sectionsState }];
          }
          const res = await fetch('/admin/reports/' + encodeURIComponent(docId), {
            method: 'PUT',
            headers: { 'Content-Type': 'application/json' },
            body: JSON.stringify(payload)
          });
          if (!res.ok) { alert('저장 실패'); saveBtn.textContent = '저장'; markDirty(); return; }
          saveBtn.textContent = '저장됨 ✓';
          markClean();
          setTimeout(function(){ saveBtn.textContent = '저장'; }, 1200);
        } catch(e) {
          alert('저장 오류: ' + e.message);
          saveBtn.textContent = '저장';
          markDirty();
        }
      });
    }

    // 5) 현황 (headline) 채우기
    const cs = document.getElementById('ov-current-status');
    if (cs) {
      if (headline) {
        cs.className = '';
        cs.innerHTML = ''
          + '<div style="background:#F8FBFF;border:1px solid #E6EBF2;border-radius:14px;padding:16px 18px;">'
          + '  <div style="font-size:13px;color:#6E7785;font-weight:700;margin-bottom:6px;">헤드라인</div>'
          + '  <div style="font-size:15px;color:#111827;font-weight:700;">' + headline + '</div>'
          + '</div>';
      }
    }

    // 6) KPI 카드 렌더 (프로젝트가 KPI 지원하는 경우만)
    const slot = document.getElementById('ov-kpi-slot');
    const issuesSlot = document.getElementById('ov-issues-slot');
    let projectKey = '';
    if (projects.length) {
      const first = (projects[0] || '').toLowerCase();
      if (first.indexOf('메이저') >= 0 || first.indexOf('major') >= 0) projectKey = 'major_module';
    }

    if (slot) {
      if (projectKey) {
        try {
          const pres = await fetch('/projects/' + projectKey);
          const pdata = await pres.json();
          const kpi = pdata.kpi_card;
          const issues = pdata.issue_lines || [];
          slot.innerHTML = window.renderKpiCardInline(kpi, issues);
          window._currentKpiContext = { projectKey: projectKey, card: kpi, issues: issues };
          // 편집은 엑셀 업로드 방식으로 전환 예정 (Phase 2)

          if (issuesSlot && Array.isArray(issues) && issues.length > 0) {
            let ih = '<div class="ov-issue-list">';
            issues.forEach(function(ii, i){
              const text = ii.text || '';
              const dueDate = ii.due_date || '';
              let ddayText = '';
              if (dueDate && ii.show_dday !== false) {
                try {
                  const d = new Date(dueDate);
                  const today = new Date();
                  const diff = Math.round((d - today) / (1000*60*60*24));
                  const mmdd = (d.getMonth()+1) + '/' + d.getDate();
                  if (diff === 0) ddayText = 'D-day · ' + mmdd;
                  else if (diff > 0) ddayText = 'D-' + diff + ' · ' + mmdd;
                  else ddayText = 'D+' + Math.abs(diff) + ' 지남 · ' + mmdd;
                } catch(_){}
              }
              ih += ''
                + '<div class="ov-issue-row">'
                + '  <div class="ov-issue-text">' + (i+1) + ') ' + text + '</div>'
                + (ddayText ? '<div class="ov-issue-dday">' + ddayText + '</div>' : '')
                + '</div>';
            });
            ih += '</div>';
            issuesSlot.innerHTML = ih;
          }
        } catch(e){
          console.error(e);
          slot.innerHTML = '<div style="color:#B8302E;font-size:13px;text-align:center;padding:12px;">KPI 데이터 로드 실패</div>';
        }
      } else {
        slot.innerHTML = '<div style="color:#8593A6;font-size:13px;text-align:center;padding:12px;">이 프로젝트는 KPI 데이터가 아직 없습니다.</div>';
      }
    }
  };

    window.backToReportList = function(){
    const c = document.getElementById('v2-content-inner') || document.getElementById('v2-content');
    if (!c) return;
    c.innerHTML = window.renderReportPageHTML();
    window.loadAdminV2Reports();
    window.bindAdminV2Upload();
  };


  // ============================================================
  // Admin v2 · KPI 셀 편집 (모달 팝업)
  // ============================================================
  window._currentKpiContext = { projectKey: '', card: null, issues: [] };

  window.bindKpiCellEditors = function(){
    document.querySelectorAll('.ov-kpi-editable').forEach(function(cell){
      cell.addEventListener('click', function(){
        const kind = cell.getAttribute('data-kpi-kind');
        const key = cell.getAttribute('data-kpi-key');
        const label = cell.getAttribute('data-kpi-label');
        const type = cell.getAttribute('data-kpi-type');
        window.openKpiEditModal(kind, key, label, type);
      });
    });
  };

  window.openKpiEditModal = function(kind, key, label, type){
    const ctx = window._currentKpiContext || {};
    const card = ctx.card || {};

    // 현재 값 찾기
    let cur = null;
    const list = kind === 'month' ? (card.months || []) : (card.weeks || []);
    for (let i = 0; i < list.length; i++) {
      const it = list[i];
      const itKey = kind === 'month' ? (it.month_key || it.key || it.month) : (it.week_key || it.key || it.week);
      if (itKey === key || it.month === label || it.week === label) { cur = it; break; }
    }

    const efemCur = cur && cur.raw_efem != null ? cur.raw_efem : (cur && cur.efem_qty != null ? cur.efem_qty : '');
    const vtmCur = cur && cur.raw_vtm != null ? cur.raw_vtm : (cur && cur.vtm_qty != null ? cur.vtm_qty : '');

    // 모달 생성
    const backdrop = document.createElement('div');
    backdrop.id = 'ov-kpi-modal-backdrop';
    backdrop.style.cssText = 'position:fixed;top:0;left:0;right:0;bottom:0;background:rgba(15,44,89,0.35);z-index:9999;display:flex;align-items:center;justify-content:center;';

    const modal = document.createElement('div');
    modal.style.cssText = 'background:#fff;border-radius:22px;padding:26px 28px;min-width:360px;max-width:90vw;box-shadow:0 20px 60px rgba(0,0,0,0.2);';
    modal.innerHTML = ''
      + '<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:16px;">'
      + '  <div style="font-size:18px;font-weight:800;color:#0F2C59;">' + label + ' 편집</div>'
      + '  <button type="button" id="ov-kpi-modal-close" style="border:0;background:transparent;font-size:20px;cursor:pointer;color:#8593A6;">×</button>'
      + '</div>'
      + '<div style="font-size:13px;color:#6E7785;margin-bottom:16px;">' + (kind === 'month' ? '월별' : '주차별') + ' KPI 수량을 입력하세요 (판가 자동 계산)</div>'
      + '<div style="margin-bottom:14px;">'
      + '  <label style="display:block;font-size:12px;color:#6E7785;font-weight:700;margin-bottom:6px;">타입</label>'
      + '  <select id="ov-kpi-type" style="width:100%;padding:10px;border:1px solid #D9E0EA;border-radius:10px;font-size:14px;">'
      + '    <option value="plan"' + (type === 'plan' ? ' selected' : '') + '>계획 (plan)</option>'
      + '    <option value="actual"' + (type === 'actual' ? ' selected' : '') + '>실적 (actual)</option>'
      + '  </select>'
      + '</div>'
      + '<div style="display:grid;grid-template-columns:1fr 1fr;gap:12px;margin-bottom:16px;">'
      + '  <div>'
      + '    <label style="display:block;font-size:12px;color:#6E7785;font-weight:700;margin-bottom:6px;">EFEM 수량</label>'
      + '    <input type="number" id="ov-kpi-efem" value="' + efemCur + '" style="width:100%;padding:10px;border:1px solid #D9E0EA;border-radius:10px;font-size:14px;box-sizing:border-box;" />'
      + '  </div>'
      + '  <div>'
      + '    <label style="display:block;font-size:12px;color:#6E7785;font-weight:700;margin-bottom:6px;">VTM 수량</label>'
      + '    <input type="number" id="ov-kpi-vtm" value="' + vtmCur + '" style="width:100%;padding:10px;border:1px solid #D9E0EA;border-radius:10px;font-size:14px;box-sizing:border-box;" />'
      + '  </div>'
      + '</div>'
      + '<div id="ov-kpi-modal-status" style="font-size:13px;font-weight:600;margin-bottom:12px;min-height:18px;"></div>'
      + '<div style="display:flex;gap:8px;justify-content:flex-end;">'
      + '  <button type="button" id="ov-kpi-modal-cancel" style="border:1px solid #DCE4EF;background:#fff;color:#35527C;border-radius:12px;padding:10px 18px;font-size:14px;font-weight:700;cursor:pointer;">취소</button>'
      + '  <button type="button" id="ov-kpi-modal-save" style="border:0;background:#0F2C59;color:#fff;border-radius:12px;padding:10px 20px;font-size:14px;font-weight:800;cursor:pointer;">저장</button>'
      + '</div>';

    backdrop.appendChild(modal);
    document.body.appendChild(backdrop);

    function closeModal(){ document.body.removeChild(backdrop); }

    document.getElementById('ov-kpi-modal-close').addEventListener('click', closeModal);
    document.getElementById('ov-kpi-modal-cancel').addEventListener('click', closeModal);
    backdrop.addEventListener('click', function(e){ if (e.target === backdrop) closeModal(); });

    document.getElementById('ov-kpi-modal-save').addEventListener('click', async function(){
      const efem = parseFloat(document.getElementById('ov-kpi-efem').value || '0');
      const vtm = parseFloat(document.getElementById('ov-kpi-vtm').value || '0');
      const newType = document.getElementById('ov-kpi-type').value || 'plan';
      const statusEl = document.getElementById('ov-kpi-modal-status');
      statusEl.style.color = '#2E5B94';
      statusEl.textContent = '⏳ 저장 중...';

      const endpoint = '/admin/kpi/' + ctx.projectKey + (kind === 'month' ? '/months' : '/weeks');
      const payload = kind === 'month'
        ? { month: key, efem: efem, vtm: vtm, type: newType, source: 'admin_v2_edit' }
        : { week: key, efem: efem, vtm: vtm, type: newType, source: 'admin_v2_edit' };

      try {
        const res = await fetch(endpoint, {
          method: 'POST',
          headers: { 'Content-Type': 'application/json' },
          credentials: 'same-origin',
          body: JSON.stringify(payload)
        });
        const data = await res.json().catch(function(){ return {}; });
        if (res.ok) {
          statusEl.style.color = '#117A52';
          statusEl.textContent = '✅ 저장 완료';
          setTimeout(function(){
            closeModal();
            window.reloadKpiSlot();
          }, 500);
        } else if (res.status === 401) {
          statusEl.style.color = '#C1272D';
          statusEl.textContent = '❌ 인증 필요';
        } else {
          statusEl.style.color = '#C1272D';
          statusEl.textContent = '❌ 저장 실패: ' + (data.detail || res.status);
        }
      } catch(e){
        statusEl.style.color = '#C1272D';
        statusEl.textContent = '❌ 네트워크 오류';
      }
    });
  };

  window.reloadKpiSlot = async function(){
    const ctx = window._currentKpiContext || {};
    if (!ctx.projectKey) return;
    const slot = document.getElementById('ov-kpi-slot');
    if (!slot) return;
    try {
      const res = await fetch('/projects/' + ctx.projectKey);
      const data = await res.json();
      const kpi = data.kpi_card;
      const issues = data.issue_lines || [];
      slot.innerHTML = window.renderKpiCardInline(kpi, issues);
      window._currentKpiContext = { projectKey: ctx.projectKey, card: kpi, issues: issues };
      window.bindKpiCellEditors();
    } catch(e){ console.error(e); }
  };

  // ============================================================
  // Admin v2 · PPT 업로드 처리
  // ============================================================
  window.bindAdminV2Upload = function(){
    const drop = document.getElementById('ov-upload-drop');
    const input = document.getElementById('ov-upload-input');
    const btn = document.getElementById('ov-upload-btn');
    const statusEl = document.getElementById('ov-upload-status');
    if (!drop || !input || !btn || !statusEl) return;

    function setStatus(msg, color){
      statusEl.textContent = msg || '';
      statusEl.style.color = color || '#415064';
    }

    function highlight(on){
      drop.style.background = on ? '#EAF3FF' : '#F7FBFF';
      drop.style.borderColor = on ? '#2E5B94' : '#BFD2EA';
    }

    async function doUpload(file){
      if (!file) return;
      const name = (file.name || '').toLowerCase();
      if (!name.endsWith('.pptx') && !name.endsWith('.ppt')) {
        setStatus('❌ .pptx / .ppt 파일만 업로드 가능합니다.', '#C1272D');
        return;
      }
      const sizeMB = file.size / (1024 * 1024);
      if (sizeMB > 50) {
        setStatus('❌ 파일 크기가 50MB 를 초과합니다.', '#C1272D');
        return;
      }

      setStatus('⏳ 업로드 중... ' + file.name + ' (' + sizeMB.toFixed(1) + 'MB)', '#2E5B94');
      btn.disabled = true;
      btn.style.opacity = '0.6';

      const fd = new FormData();
      fd.append('file', file);
      fd.append('report_family', 'default');
      fd.append('allow_duplicate', 'false');

      try {
        const res = await fetch('/upload', { method: 'POST', body: fd, credentials: 'same-origin' });
        let data = {};
        try { data = await res.json(); } catch(_) {}

        if (res.status === 409) {
          const ok = confirm('이미 업로드된 파일과 동일합니다. 그래도 계속 업로드할까요?');
          if (ok) {
            const fd2 = new FormData();
            fd2.append('file', file);
            fd2.append('report_family', 'default');
            fd2.append('allow_duplicate', 'true');
            const res2 = await fetch('/upload', { method: 'POST', body: fd2, credentials: 'same-origin' });
            const data2 = await res2.json().catch(function(){ return {}; });
            if (res2.ok) {
              setStatus('✅ 업로드 완료 (중복 허용): ' + (data2.filename || file.name), '#117A52');
              window.loadAdminV2Reports && window.loadAdminV2Reports();
            } else {
              setStatus('❌ 업로드 실패: ' + (data2.detail || res2.status), '#C1272D');
            }
          } else {
            setStatus('취소됨 (중복 파일)', '#7A8595');
          }
        } else if (res.ok) {
          setStatus('✅ 업로드 완료: ' + (data.filename || file.name), '#117A52');
          window.loadAdminV2Reports && window.loadAdminV2Reports();
        } else if (res.status === 401) {
          setStatus('❌ 인증이 필요합니다. 다시 로그인해주세요.', '#C1272D');
        } else {
          setStatus('❌ 업로드 실패: ' + (data.detail || res.status), '#C1272D');
        }
      } catch (e) {
        setStatus('❌ 네트워크 오류: ' + e.message, '#C1272D');
      } finally {
        btn.disabled = false;
        btn.style.opacity = '1';
        input.value = '';
      }
    }

    btn.addEventListener('click', function(){ input.click(); });
    input.addEventListener('change', function(e){
      const f = e.target.files && e.target.files[0];
      if (f) doUpload(f);
    });

    ['dragenter','dragover'].forEach(function(evt){
      drop.addEventListener(evt, function(e){
        e.preventDefault(); e.stopPropagation();
        highlight(true);
      });
    });
    ['dragleave','drop'].forEach(function(evt){
      drop.addEventListener(evt, function(e){
        e.preventDefault(); e.stopPropagation();
        highlight(false);
      });
    });
    drop.addEventListener('drop', function(e){
      const dt = e.dataTransfer;
      const f = dt && dt.files && dt.files[0];
      if (f) doUpload(f);
    });
  };

  // 최초 진입 시 report 페이지 즉시 렌더
  document.addEventListener('DOMContentLoaded', function(){
    const root = document.getElementById('ov-report-root');
    if (root) {
      root.outerHTML = '<div id="v2-content-inner"></div>';
      const inner = document.getElementById('v2-content-inner');
      if (inner) {
        inner.innerHTML = window.renderReportPageHTML();
        window.loadAdminV2Reports();
        window.bindAdminV2Upload();
      }
    }
  });

})();
</script>
<script>
</script>

<script>
// 사업부 드롭다운 자동 채우기 (페이지 로드 시 1회)
(function() {
  function fillDivisions() {
    fetch('/divisions', { credentials: 'same-origin' })
      .then(function(r) { return r.ok ? r.json() : null; })
      .then(function(j) {
        if (!j || !j.divisions) return;
        var sel = document.getElementById('v2-division-select');
        if (!sel) return;
        sel.innerHTML = '';
        j.divisions.forEach(function(d) {
          var opt = document.createElement('option');
          opt.value = d.id;
          opt.textContent = d.label;
          sel.appendChild(opt);
        });
        var saved = localStorage.getItem('v2_selected_division') || (j.divisions[0] && j.divisions[0].id) || '';
        if (saved) {
          sel.value = saved;
          sel.dispatchEvent(new Event('change', { bubbles: true }));
        }
        sel.addEventListener('change', function(){
      try { localStorage.setItem('v2_division', this.value || ''); } catch(e){}
      if (typeof window.renderAdminV2ByDivision === 'function') {
        window.renderAdminV2ByDivision();
      }
    });
        console.log('✅ 사업부 ' + j.divisions.length + '개 로드됨');
      })
      .catch(function(e) { console.error('사업부 로드 실패:', e); });
  }
  if (document.readyState === 'loading') {
    document.addEventListener('DOMContentLoaded', fillDivisions);
  } else {
    fillDivisions();
  }
})();

// ═══════════════════════════════════════════════
// 모델 관리 페이지
// ═══════════════════════════════════════════════
(function() {
  'use strict';

  // 스타일 주입
  var style = document.createElement('style');
  style.textContent = [
    '.mdl-wrap { max-width: 800px; }',
    '.mdl-header { display: flex; align-items: center; gap: 12px; margin-bottom: 20px; }',
    '.mdl-header h2 { font-size: 20px; font-weight: 800; color: #111827; margin: 0; }',
    '.mdl-proj-sel { padding: 8px 12px; border: 1px solid #D1D5DB; border-radius: 8px; font-size: 14px; background: #fff; min-width: 200px; }',
    '.mdl-table { width: 100%; border-collapse: collapse; background: #fff; border-radius: 12px; overflow: hidden; border: 1px solid #E5E7EB; }',
    '.mdl-table th { background: #F9FAFB; padding: 10px 14px; font-size: 12px; font-weight: 700; color: #6B7280; text-align: left; border-bottom: 1px solid #E5E7EB; }',
    '.mdl-table td { padding: 10px 14px; font-size: 13px; border-bottom: 1px solid #F3F4F6; vertical-align: middle; }',
    '.mdl-table tr:last-child td { border-bottom: none; }',
    '.mdl-input { padding: 6px 10px; border: 1px solid #D1D5DB; border-radius: 6px; font-size: 13px; width: 100%; box-sizing: border-box; }',
    '.mdl-input:focus { outline: none; border-color: #3B82F6; }',
    '.mdl-sel { padding: 6px 8px; border: 1px solid #D1D5DB; border-radius: 6px; font-size: 13px; background: #fff; }',
    '.mdl-num { width: 70px; text-align: center; }',
    '.mdl-badge { display: inline-block; padding: 2px 8px; border-radius: 10px; font-size: 11px; font-weight: 700; }',
    '.mdl-badge-prod { background: #DBEAFE; color: #1D4ED8; }',
    '.mdl-badge-dev { background: #FEF3C7; color: #92400E; }',
    '.mdl-btn { padding: 6px 14px; border: none; border-radius: 6px; font-size: 13px; font-weight: 600; cursor: pointer; }',
    '.mdl-btn-add { background: #0F2C59; color: #fff; }',
    '.mdl-btn-add:hover { background: #12356F; }',
    '.mdl-btn-del { background: #FEE2E2; color: #DC2626; padding: 4px 10px; font-size: 12px; }',
    '.mdl-btn-del:hover { background: #FECACA; }',
    '.mdl-btn-save { background: #059669; color: #fff; padding: 10px 24px; font-size: 14px; }',
    '.mdl-btn-save:hover { background: #047857; }',
    '.mdl-actions { display: flex; gap: 8px; margin-top: 16px; align-items: center; }',
    '.mdl-status { font-size: 13px; margin-left: auto; }',
    '.mdl-status.ok { color: #059669; }',
    '.mdl-status.err { color: #DC2626; }',
    '.mdl-empty { text-align: center; padding: 40px; color: #9CA3AF; font-size: 14px; }',
    '.mdl-group-sep td { background: #F9FAFB; font-weight: 700; font-size: 12px; color: #374151; padding: 8px 14px; }',
    '.mdl-weekly-section { margin-top: 24px; padding-top: 20px; border-top: 2px solid #E5E7EB; }',
    '.mdl-weekly-section h3 { font-size: 16px; font-weight: 700; color: #111827; margin: 0 0 12px 0; }',
    '.mdl-weekly-status { font-size: 13px; color: #6B7280; margin-bottom: 12px; }',
    '.mdl-weekly-status.ok { color: #059669; }',
    '.mdl-weekly-status.err { color: #DC2626; }',
    '.mdl-weekly-actions { display: flex; gap: 8px; margin-bottom: 16px; }',
    '.mdl-weekly-preview { background: #F9FAFB; border-radius: 8px; padding: 12px; max-height: 300px; overflow: auto; }',
    '.mdl-weekly-preview table { width: 100%; border-collapse: collapse; font-size: 12px; }',
    '.mdl-weekly-preview th { background: #E5E7EB; padding: 6px 8px; text-align: left; font-weight: 600; }',
    '.mdl-weekly-preview td { padding: 6px 8px; border-bottom: 1px solid #E5E7EB; min-width: 40px; }',
    '.mdl-weekly-preview td:empty { background: #F3F4F6; }',
    '.mdl-wrap { padding: 4px 0; }',
    '.mdl-head { margin-bottom: 16px; }',
    '.mdl-title-row { display: flex; align-items: center; gap: 10px; margin-bottom: 14px; }',
    '.mdl-mark { color: #0F2C59; font-size: 14px; }',
    '.mdl-title { font-size: 18px; font-weight: 700; color: #111827; margin: 0; }',
    '.mdl-select { margin-left: auto; padding: 8px 12px; border: 1px solid #D1D5DB; border-radius: 8px; font-size: 13px; background: #fff; min-width: 160px; }',
    '.mdl-tabs { display: flex; gap: 4px; border-bottom: 2px solid #E5E7EB; }',
    '.mdl-tab { padding: 10px 18px; border: none; background: none; font-size: 14px; font-weight: 600; color: #6B7280; cursor: pointer; border-bottom: 2px solid transparent; margin-bottom: -2px; }',
    '.mdl-tab.is-active { color: #0F2C59; border-bottom-color: #0F2C59; }',
    '.mdl-badge { display: inline-block; min-width: 20px; padding: 1px 6px; background: #EEF2FF; color: #0F2C59; border-radius: 10px; font-size: 11px; font-weight: 700; margin-left: 4px; }',
    '.mdl-pane { padding-top: 16px; }',
    '.mdl-actions { display: flex; align-items: center; gap: 10px; margin-bottom: 12px; }',
    '.mdl-save-msg { font-size: 13px; font-weight: 600; }',
    '.mdl-hint { font-size: 12px; color: #9CA3AF; margin-top: 10px; }',
    '.mdl-input { width: 110px; padding: 6px 8px; border: 1px solid #D1D5DB; border-radius: 6px; font-size: 13px; text-align: right; }',
    '.mdl-td-name { font-weight: 600; color: #111827; }',
    '.mdl-td-ratio { font-weight: 700; color: #0F2C59; }',
    '.mdl-group-row td { background: #F9FAFB; padding: 8px 14px !important; }',
    '.mdl-group-badge { display: inline-block; padding: 2px 10px; border-radius: 10px; font-size: 12px; font-weight: 700; }',
    '.mdl-group-mass { background: #DBEAFE; color: #1D4ED8; }',
    '.mdl-group-dev { background: #FEF3C7; color: #B45309; }',
    '.mdl-group-cnt { margin-left: 8px; font-size: 12px; color: #6B7280; }',
    '.mdl-empty { padding: 40px; text-align: center; color: #9CA3AF; background: #F9FAFB; border-radius: 8px; }',
    '.mdl-loading { padding: 40px; text-align: center; color: #9CA3AF; }',
    '.mdl-plan-card { background: #fff; border: 1px solid #E5E7EB; border-radius: 12px; padding: 16px; }',
      /* ═══ 디자인 시안 기준 재정의 (덮어쓰기) ═══ */
    '.mdl-wrap { background: transparent; }',
    '.mdl-table { width: 100%; border-collapse: separate; border-spacing: 0 8px; background: transparent; border: none; }',
    '.mdl-table th { background: transparent; padding: 4px 14px; font-size: 12px; font-weight: 600; color: #64748B; text-align: left; border: none; }',
    '.mdl-table td { background: #FFFFFF; padding: 12px 14px; font-size: 13px; border-top: 1px solid #E4E7EC; border-bottom: 1px solid #E4E7EC; vertical-align: middle; box-shadow: 0 4px 12px rgba(0,0,0,0.03); }',
    '.mdl-table td:first-child { border-left: 1px solid #E4E7EC; border-radius: 8px 0 0 8px; }',
    '.mdl-table td:last-child { border-right: 1px solid #E4E7EC; border-radius: 0 8px 8px 0; }',
    '.mdl-table tr:last-child td { border-bottom: 1px solid #E4E7EC; }',
    '.mdl-group-row td { background: transparent !important; border: none !important; box-shadow: none !important; padding: 4px 6px !important; }',
    '.mdl-group-badge { display: inline-block; padding: 4px 12px; border-radius: 8px; font-size: 12px; font-weight: 700; }',
    '.mdl-group-mass { background: #EFF6FF; color: #1E40AF; }',
    '.mdl-group-dev { background: #FEF3C7; color: #92400E; }',
    '.mdl-group-cnt { margin-left: 8px; font-size: 12px; color: #64748B; font-weight: 600; }',
    '.mdl-input { width: 110px; padding: 8px 12px; border: 1px solid #E2E8F0; border-radius: 6px; font-size: 13px; text-align: right; background: #FFFFFF; }',
    '.mdl-input:focus { outline: none; border-color: #2563EB; }',
    '.mdl-table select[data-field="group"] { padding: 6px 10px; border: none; border-radius: 6px; font-size: 12px; font-weight: 600; background: #EFF6FF; color: #2563EB; cursor: pointer; }',
    '.mdl-devtype { padding: 6px 12px; border: none; border-radius: 6px; font-size: 12px; font-weight: 700; cursor: pointer; }',
    '.mdl-dt-hvm { background: #F5F3FF; color: #7C3AED; }',
    '.mdl-dt-rpm { background: #ECFDF5; color: #059669; }',
    '.mdl-btn-add { background: #0F172A; color: #FFFFFF; padding: 10px 18px; border-radius: 6px; font-size: 13px; font-weight: 700; }',
    '.mdl-btn-add:hover { background: #1E293B; }',
    '.mdl-btn-save { background: #10B981; color: #FFFFFF; padding: 10px 24px; border-radius: 6px; font-size: 14px; font-weight: 700; }',
    '.mdl-btn-save:hover { background: #059669; }',
    '.mdl-btn-del { background: #FEF2F2; color: #EF4444; padding: 6px 12px; font-size: 12px; border-radius: 6px; font-weight: 600; }',
    '.mdl-btn-del:hover { background: #FEE2E2; }',
    '.mdl-btn-process { background: #F5F3FF; color: #7C3AED; padding: 6px 12px; font-size: 12px; border-radius: 6px; font-weight: 600; margin-right: 6px; }',
    '.mdl-btn-process:hover { background: #EDE9FE; }',
    '.mdl-btn-note-save { background: #10B981; color: #FFFFFF; padding: 6px 16px; font-size: 12px; border-radius: 6px; font-weight: 700; }',
    '.mdl-btn-note-save:hover { background: #059669; }',
    '.mdl-note-card { background: #FFFFFF; border: 1px solid #E4E7EC; border-radius: 12px; padding: 16px 20px; margin-bottom: 16px; box-shadow: 0 4px 12px rgba(0,0,0,0.03); }',
    '.mdl-note-row { display: flex; align-items: center; justify-content: space-between; margin-bottom: 10px; }',
    '.mdl-note-label { font-size: 14px; font-weight: 700; color: #1E293B; }',
    '.mdl-note-card textarea { width: 100%; border: 1px solid #E2E8F0; border-radius: 8px; padding: 12px; font-size: 13px; resize: vertical; box-sizing: border-box; font-family: inherit; color: #1E293B; }',
    '.mdl-note-card textarea:focus { outline: none; border-color: #2563EB; }',
    '.mdl-plan-card { background: #FFFFFF; border: 1px solid #E4E7EC; border-radius: 12px; padding: 20px; box-shadow: 0 4px 12px rgba(0,0,0,0.03); }',
    '.mdl-empty { padding: 40px; text-align: center; color: #94A3B8; background: #FFFFFF; border-radius: 12px; border: 1px solid #E4E7EC; }',
    '.mdl-loading { padding: 40px; text-align: center; color: #94A3B8; }',
    /* ── 프로세스 에디터 시안 ── */
    '.mdl-proc { max-width: 1080px; }',
    '.mdl-proc-head { display: flex; align-items: center; gap: 12px; margin-bottom: 20px; }',
    '.mdl-proc-title { font-size: 20px; font-weight: 700; color: #1E293B; margin-right: auto; }',
    '.mdl-btn-back { background: #FFFFFF; color: #475569; border: 1px solid #E2E8F0; padding: 8px 14px; border-radius: 8px; font-size: 13px; font-weight: 600; }',
    '.mdl-btn-back:hover { background: #F8FAFC; }',
    '.mdl-proc-head .mdl-select { padding: 8px 12px; border: none; border-radius: 6px; font-size: 12px; font-weight: 700; background: #F5F3FF; color: #7C3AED; min-width: 90px; }',
    '.mdl-proc-summary { display: grid; grid-template-columns: 1fr 1.4fr 1fr; gap: 12px; margin-bottom: 20px; }',
    '.mdl-proc-card { background: #FFFFFF; border: 1px solid #E2E8F0; border-radius: 12px; padding: 18px 20px; box-shadow: 0 4px 12px rgba(0,0,0,0.03); }',
    '.mdl-proc-card .lbl { font-size: 12px; color: #64748B; margin-bottom: 6px; font-weight: 600; }',
    '.mdl-proc-card .val { font-size: 18px; font-weight: 800; color: #1E293B; }',
    '.mdl-proc-card .sub { font-size: 11px; color: #94A3B8; margin-top: 4px; }',
    '.mdl-proc-bar { height: 8px; background: #F1F5F9; border-radius: 4px; margin-top: 10px; overflow: hidden; }',
    '.mdl-proc-fill { height: 100%; background: #10B981; border-radius: 4px; transition: width 0.3s; }',
    '.mdl-proc-group { background: #FFFFFF; border: 1px solid #E4E7EC; border-radius: 12px; margin-bottom: 16px; overflow: hidden; box-shadow: 0 4px 12px rgba(0,0,0,0.03); }',
    '.mdl-proc-group-head { padding: 0 16px; height: 40px; font-size: 14px; font-weight: 700; display: flex; align-items: center; justify-content: space-between; }',
    '.mdl-proc-group:nth-of-type(2) .mdl-proc-group-head { background: #EEF2FF !important; color: #312E81 !important; }',
    '.mdl-proc-cnt { font-size: 11px; font-weight: 700; opacity: 0.85; padding: 2px 10px; border-radius: 10px; background: rgba(255,255,255,0.5); }',
    '.mdl-proc-row { display: flex; align-items: center; gap: 14px; padding: 0 16px; height: 48px; border-top: 1px solid #F1F5F9; background: #FFFFFF; }',
    '.mdl-proc-row.current { background: #F0F7FF; }',
    '.mdl-proc-icon { width: 20px; height: 20px; border-radius: 50%; display: inline-flex; align-items: center; justify-content: center; font-size: 11px; font-weight: 700; flex-shrink: 0; }',
    '.mdl-proc-icon.done { background: #22C55E; color: #FFFFFF; }',
    '.mdl-proc-icon.doing { background: #2563EB; color: #FFFFFF; }',
    '.mdl-proc-icon.todo { background: #E2E8F0; color: #94A3B8; }',
    '.mdl-proc-name { font-size: 13px; font-weight: 700; color: #1E293B; min-width: 180px; }',
    '.mdl-proc-dates { display: flex; gap: 18px; margin-left: auto; align-items: center; }',
    '.mdl-proc-dates label { font-size: 11px; color: #64748B; display: flex; align-items: center; gap: 6px; }',
    '.mdl-proc-dates input[type=date] { border: none; background: transparent; font-size: 11px; font-weight: 600; color: #2563EB; padding: 2px 0; cursor: pointer; }',
    '.mdl-proc-dates input[type=date]:focus { outline: none; }',
    '.mdl-proc-status { font-size: 11px; font-weight: 700; color: #94A3B8; width: 44px; text-align: right; }',
    '.mdl-proc-status.done { color: #64748B; }',
    '.mdl-proc-status.doing { color: #2563EB; }',
  ].join(' ');
  document.head.appendChild(style);

  var _modelsData = [];
  var _currentProjectKey = '';

  function _esc(s) {
    var d = document.createElement('div');
    d.textContent = s || '';
    return d.innerHTML;
  }

  function _genId(name) {
    return (name || '').toLowerCase().replace(/[^a-z0-9가-힣]/g, '-').replace(/-+/g, '-').replace(/^-|-$/g, '') || 'model-' + Date.now();
  }

  // 프로젝트 목록 로드 (드롭다운용)
  function loadProjectOptions(sel) {
    var divSel = document.getElementById('v2-division-select');
    var divId = divSel ? divSel.value : '';
    if (!divId) {
      sel.innerHTML = '<option value="">사업부를 먼저 선택하세요</option>';
      return;
    }
    fetch('/admin/config/projects?division_id=' + encodeURIComponent(divId), { credentials: 'same-origin' })
      .then(function(r) { return r.json(); })
      .then(function(d) {
        var projects = (d && d.projects) || [];
        sel.innerHTML = '<option value="">프로젝트 선택...</option>';
        projects.forEach(function(p) {
          var opt = document.createElement('option');
          opt.value = p.id || p.label;
          opt.textContent = p.label || p.id;
          sel.appendChild(opt);
        });
      })
      .catch(function() {
        sel.innerHTML = '<option value="">로드 실패</option>';
      });
  }

  // 모델 테이블 렌더링
  function syncDomToData(container) {
    container.querySelectorAll('tr[data-model-id]').forEach(function(r) {
      const m = _modelsData.find(function(x) { return x.id === r.dataset.modelId; });
      if (!m) return;
      const g = r.querySelector('[data-field="group"]');
      if (g) m.group = g.value;
      const dt = r.querySelector('[data-field="dev_type"]');
      if (dt) m.dev_type = dt.value;
      const st = r.querySelector('[data-field="status"]');
      if (st) m.status = st.value;
      m.price = parseInt((r.querySelector('[data-field="price"]').value || '0').replace(/,/g, ''), 10) || 0;
      m.material_cost = parseInt((r.querySelector('[data-field="material_cost"]').value || '0').replace(/,/g, ''), 10) || 0;
      const pg = r.querySelector('[data-field="progress"]');
      if (pg) m.progress = parseInt(pg.value || '0', 10) || 0;
    });
  }

  function renderTable(container) {
    const byGroup = { '양산': [], '개발': [] };
    _modelsData.forEach(function(m) {
      const g = (m.group === '개발') ? '개발' : '양산';
      byGroup[g].push(m);
    });
    const tabCount = document.getElementById('mdl-tab-count');
    if (tabCount) tabCount.textContent = _modelsData.length;

    let html = '<table class="mdl-table"><thead><tr>' +
      '<th>모델명</th><th>구분</th><th>개발 유형</th><th>판가($)</th><th>재료비($)</th><th>재료비율</th><th>관리</th>' +
      '</tr></thead><tbody>';

    ['양산', '개발'].forEach(function(g) {
      const list = byGroup[g];
      if (!list.length) return;
      html += '<tr class="mdl-group-row"><td colspan="7">' +
        '<span class="mdl-group-badge mdl-group-' + (g === '양산' ? 'mass' : 'dev') + '">' + g + '</span>' +
        '<span class="mdl-group-cnt">' + list.length + '개</span></td></tr>';
      list.forEach(function(m) {
        const price = m.price || 0;
        const mcost = m.material_cost || 0;
        const ratio = price > 0 ? ((mcost / price) * 100).toFixed(1) + '%' : '-';
        const isDev = g === '개발';
        const devType = (m.dev_type || 'HVM').toUpperCase();
        html += '<tr data-model-id="' + m.id + '">' +
          '<td class="mdl-td-name">' + m.name + '</td>' +
          '<td><select data-field="group">' +
            '<option value="양산"' + (m.group === '양산' ? ' selected' : '') + '>양산</option>' +
            '<option value="개발"' + (isDev ? ' selected' : '') + '>개발</option>' +
          '</select></td>' +
          '<td>' + (isDev
            ? '<select data-field="dev_type" class="mdl-devtype mdl-dt-' + devType.toLowerCase() + '">' +
              '<option value="HVM"' + (devType === 'HVM' ? ' selected' : '') + '>HVM</option>' +
              '<option value="RPM"' + (devType === 'RPM' ? ' selected' : '') + '>RPM</option></select>'
            : '<span class="mdl-dash">-</span>') + '</td>' +
          '<td><input data-field="price" type="number" min="0" value="' + price + '" class="mdl-input"></td>' +
          '<td><input data-field="material_cost" type="number" min="0" value="' + mcost + '" class="mdl-input"></td>' +
          '<td data-field="ratio" class="mdl-td-ratio">' + ratio + '</td>' +
          '<td class="mdl-td-actions">' +
            '<input data-field="status" type="hidden" value="' + (m.status || '정상') + '">' +
            '<input data-field="progress" type="hidden" value="' + (m.progress || 0) + '">' +
            (isDev ? '<button type="button" class="mdl-btn mdl-btn-process mdl-row-process">Process 입력</button>' : '') +
            '<button type="button" class="mdl-btn mdl-btn-del mdl-row-del">삭제</button>' +
          '</td></tr>';
      });
    });

    html += '</tbody></table>';
    if (!_modelsData.length) html = '<div class="mdl-empty">등록된 모델이 없습니다.</div>';
    container.innerHTML = html;

    container.querySelectorAll('input, select').forEach(function(el) {
      el.addEventListener('change', function() {
        const tr = el.closest('tr');
        if (!tr) return;
        if (el.dataset.field === 'group') {
          syncDomToData(container);
          renderTable(container);
          return;
        }
        if (el.dataset.field === 'dev_type') {
          el.className = 'mdl-devtype mdl-dt-' + el.value.toLowerCase();
        }
        if (el.dataset.field === 'price' || el.dataset.field === 'material_cost') {
          const pr = parseInt(tr.querySelector('[data-field="price"]').value || '0', 10) || 0;
          const mc = parseInt(tr.querySelector('[data-field="material_cost"]').value || '0', 10) || 0;
          const cell = tr.querySelector('[data-field="ratio"]');
          if (cell) cell.textContent = pr > 0 ? ((mc / pr) * 100).toFixed(1) + '%' : '-';
        }
      });
    });

    container.querySelectorAll('.mdl-row-del').forEach(function(btn) {
      btn.addEventListener('click', function() {
        const tr = btn.closest('tr');
        const mid = tr ? tr.dataset.modelId : null;
        if (!mid) return;
        if (confirm('[' + mid + '] 삭제하시겠습니까? (저장 버튼을 눌러야 반영됩니다)')) {
          _modelsData = _modelsData.filter(function(m) { return m.id !== mid; });
          renderTable(container);
        }
      });
    });

    container.querySelectorAll('.mdl-row-process').forEach(function(btn) {
      btn.addEventListener('click', function() {
        const tr = btn.closest('tr');
        if (tr && tr.dataset.modelId) openProcessEditor(tr.dataset.modelId);
      });
    });
  }

  function openProcessEditor(modelId) {
    if (!_currentProjectKey) { alert('프로젝트를 먼저 선택하세요'); return; }
    fetch('/projects/' + encodeURIComponent(_currentProjectKey) + '/models/' + encodeURIComponent(modelId) + '/process', { credentials: 'same-origin' })
      .then(function(r) { return r.json(); })
      .then(function(d) {
        if (d.detail) { alert(d.detail); return; }
        renderProcessEditor(d);
      })
      .catch(function(e) { alert('프로세스 로드 실패: ' + e.message); });
  }

  function renderProcessEditor(d) {
    const tableBox = document.getElementById('mdl-table-box');
    if (!tableBox) return;
    const groups = [['발주', '#DBEAFE', '#1D4ED8'], ['제작·검사', '#FEF3C7', '#B45309'], ['승인', '#EDE9FE', '#7C3AED']];

    let html = '<div class="mdl-proc">' +
      '<div class="mdl-proc-head">' +
        '<button type="button" class="mdl-btn mdl-btn-back" id="mdl-proc-back">← 목록으로</button>' +
        '<span class="mdl-proc-title">' + d.model_name + ' · 개발 Process 입력</span>' +
        '<select id="mdl-proc-devtype" class="mdl-select" style="min-width:90px">' +
          '<option value="HVM"' + (d.dev_type === 'HVM' ? ' selected' : '') + '>HVM</option>' +
          '<option value="RPM"' + (d.dev_type === 'RPM' ? ' selected' : '') + '>RPM</option>' +
        '</select>' +
        '<button type="button" class="mdl-btn mdl-btn-save" id="mdl-proc-save">저장</button>' +
        '<span id="mdl-proc-msg" class="mdl-save-msg"></span>' +
      '</div>' +
      '<div class="mdl-proc-summary">' +
        '<div class="mdl-proc-card"><div class="lbl">개발 유형</div><div class="val">' + d.dev_type + '</div></div>' +
        '<div class="mdl-proc-card"><div class="lbl">진행률</div><div class="val">' + d.done + ' / ' + d.total + ' 단계</div>' +
          '<div class="mdl-proc-bar"><div class="mdl-proc-fill" style="width:' + d.progress + '%"></div></div></div>' +
        '<div class="mdl-proc-card"><div class="lbl">다음 단계</div><div class="val">' + d.current_stage + '</div>' +
          '<div class="sub">' + (d.current_expected ? '예상 ' + d.current_expected : '예상일 미정') + '</div></div>' +
      '</div>';

    let num = 0;
    groups.forEach(function(g) {
      const steps = d.steps.filter(function(s) { return s.group === g[0]; });
      if (!steps.length) return;
      html += '<div class="mdl-proc-group">' +
        '<div class="mdl-proc-group-head" style="background:' + g[1] + ';color:' + g[2] + '">' +
          g[0] + '<span class="mdl-proc-cnt">' + steps.length + '단계</span></div>';
      steps.forEach(function(s) {
        num++;
        const isDone = (s.actual || '').trim() !== '';
        const isCurrent = !isDone && d.current_stage === s.name;
        html += '<div class="mdl-proc-row' + (isCurrent ? ' current' : '') + '" data-step-key="' + s.key + '">' +
          '<span class="mdl-proc-icon ' + (isDone ? 'done' : (isCurrent ? 'doing' : 'todo')) + '">' +
            (isDone ? '✓' : (isCurrent ? '●' : '○')) + '</span>' +
          '<span class="mdl-proc-name">' + ('0' + num).slice(-2) + ' ' + s.name + '</span>' +
          '<span class="mdl-proc-dates">' +
            '<label>예상 <input type="date" data-f="expected" value="' + (s.expected || '') + '"></label>' +
            '<label>실제 <input type="date" data-f="actual" value="' + (s.actual || '') + '"></label>' +
          '</span>' +
          '<span class="mdl-proc-status ' + (isDone ? 'done' : (isCurrent ? 'doing' : '')) + '">' +
            (isDone ? '완료' : (isCurrent ? '진행중' : '대기')) + '</span>' +
        '</div>';
      });
      html += '</div>';
    });
    html += '</div>';
    tableBox.innerHTML = html;

    document.getElementById('mdl-proc-back').addEventListener('click', function() {
      renderTable(tableBox);
    });

    document.getElementById('mdl-proc-save').addEventListener('click', function() {
      const steps = [];
      tableBox.querySelectorAll('.mdl-proc-row').forEach(function(row) {
        steps.push({
          key: row.dataset.stepKey,
          expected: row.querySelector('[data-f="expected"]').value || '',
          actual: row.querySelector('[data-f="actual"]').value || ''
        });
      });
      const msg = document.getElementById('mdl-proc-msg');
      fetch('/admin/projects/' + encodeURIComponent(_currentProjectKey) + '/models/' + encodeURIComponent(d.model_id) + '/process', {
        method: 'PUT', credentials: 'same-origin',
        headers: { 'Content-Type': 'application/json' },
        body: JSON.stringify({ dev_type: document.getElementById('mdl-proc-devtype').value, steps: steps })
      })
        .then(function(r) { return r.json(); })
        .then(function(res) {
          if (res.ok) {
            msg.textContent = '✅ 저장 완료 (진행률 ' + res.progress + '%)';
            msg.style.color = '#059669';
            const m = _modelsData.find(function(x) { return x.id === d.model_id; });
            if (m) { m.dev_type = res.dev_type; m.progress = res.progress; }
            // 저장 직후 서버 데이터로 화면 즉시 갱신 (체크 상태/진행률/다음 단계 반영)
            fetch('/projects/' + encodeURIComponent(_currentProjectKey) + '/models/' + encodeURIComponent(d.model_id) + '/process', { credentials: 'same-origin' })
              .then(function(r2) { return r2.json(); })
              .then(function(d2) {
                if (!d2.detail) renderProcessEditor(d2);
                // 상단 메시지는 리렌더 후에도 유지되도록 다시 표시
                setTimeout(function() {
                  const msg2 = document.getElementById('mdl-proc-msg');
                  if (msg2) { msg2.textContent = '✅ 저장 완료 (진행률 ' + res.progress + '%)'; msg2.style.color = '#059669'; }
                }, 50);
              });
          } else {
            msg.textContent = '❌ 실패: ' + (res.detail || '');
            msg.style.color = '#DC2626';
          }
        })
        .catch(function(e) {
          msg.textContent = '❌ ' + e.message;
          msg.style.color = '#DC2626';
        });
    });
  }

  function loadModels(projectKey, container) {
    if (!projectKey) { _modelsData = []; renderTable(container); return; }
    _currentProjectKey = projectKey;
    fetch('/admin/projects/' + encodeURIComponent(projectKey) + '/models', { credentials: 'same-origin' })
      .then(function(r) { return r.json(); })
      .then(function(d) {
        _modelsData = (d && d.models) || [];
        renderTable(container);
      })
      .catch(function() { _modelsData = []; renderTable(container); });
  }

  // 모델 저장
  function saveModels(statusEl) {
    if (!_currentProjectKey) {
      if (statusEl) { statusEl.textContent = '프로젝트를 먼저 선택하세요'; statusEl.className = 'mdl-status err'; }
      return;
    }
    if (statusEl) { statusEl.textContent = '저장 중...'; statusEl.className = 'mdl-status'; }
    fetch('/admin/projects/' + encodeURIComponent(_currentProjectKey) + '/models', {
      method: 'PUT',
      headers: { 'Content-Type': 'application/json' },
      credentials: 'same-origin',
      body: JSON.stringify({ models: _modelsData })
    })
    .then(function(r) {
      if (!r.ok) throw new Error('HTTP ' + r.status);
      return r.json();
    })
    .then(function(d) {
      if (statusEl) { statusEl.textContent = '✅ 저장 완료 (' + (d.count || 0) + '개)'; statusEl.className = 'mdl-status ok'; }
      // 서버 정렬 결과(양산→개발)를 화면에 즉시 반영
      var tableBox = document.getElementById('mdl-table-box');
      if (tableBox) loadModels(_currentProjectKey, tableBox);
    })
    .catch(function(e) {
      if (statusEl) { statusEl.textContent = '❌ 저장 실패: ' + e.message; statusEl.className = 'mdl-status err'; }
    });
  }

  // 사업부 select가 비어있으면 자동 fallback: 첫 번째 사업부 선택
  function ensureDivisionSelected() {
    var divSel = document.getElementById('v2-division-select');
    if (!divSel) return false;
    if (divSel.value) return true;
    if (!divSel.options || divSel.options.length <= 1) return false;
    var first = divSel.options[1];  // 0번은 placeholder
    if (first && first.value) {
      divSel.value = first.value;
      try { localStorage.setItem('v2_division', divSel.value); } catch(e){}
      divSel.dispatchEvent(new Event('change', { bubbles: true }));
      return true;
    }
    return false;
  }

  // 페이지 렌더링 (data-page="models" 클릭 시)
  function renderModelsPage() {
    const content = document.getElementById('v2-content');
    if (!content) return;

    content.innerHTML =
      '<div class="mdl-wrap">' +
        '<div class="mdl-head">' +
          '<div class="mdl-title-row">' +
            '<span class="mdl-mark">◆</span>' +
            '<h2 class="mdl-title">모델 관리</h2>' +
            '<select id="mdl-proj-select" class="mdl-select"></select>' +
          '</div>' +
          '<div class="mdl-tabs">' +
            '<button type="button" class="mdl-tab is-active" data-mtab="list">📋 모델 목록 <span id="mdl-tab-count" class="mdl-badge">0</span></button>' +
            '<button type="button" class="mdl-tab" data-mtab="plan">📅 주차별 계획</button>' +
          '</div>' +
        '</div>' +

        '<div class="mdl-pane" data-mpane="list">' +
          '<div class="mdl-note-card">' +
            '<div class="mdl-note-row">' +
              '<div class="mdl-note-label">프로젝트 현황</div>' +
              '<button type="button" class="mdl-btn mdl-btn-note-save" id="mdl-note-save">저장</button>' +
            '</div>' +
            '<textarea id="mdl-status-note" rows="3" placeholder="프로젝트 현황을 입력하세요 (앱에 표시됩니다)"></textarea>' +
            '<div id="mdl-note-msg" class="mdl-save-msg"></div>' +
          '</div>' +
          '<div class="mdl-actions">' +
            '<button type="button" class="mdl-btn mdl-btn-add" id="mdl-add-btn">＋ 모델 추가</button>' +
            '<button type="button" class="mdl-btn mdl-btn-save" id="mdl-save-btn">저장</button>' +
            '<span id="mdl-save-msg" class="mdl-save-msg"></span>' +
          '</div>' +
          '<div id="mdl-table-box"><div class="mdl-loading">불러오는 중...</div></div>' +
          '<p class="mdl-hint">판가·재료비는 달러($) 기준이며, 재료비율은 자동 계산됩니다 (재료비 ÷ 판가 × 100).</p>' +
        '</div>' +

        '<div class="mdl-pane" data-mpane="plan" hidden>' +
          '<div class="mdl-plan-card">' +
            '<div id="mdl-weekly-status" class="mdl-weekly-status">로딩 중...</div>' +
            '<div class="mdl-weekly-actions">' +
              '<input type="file" id="mdl-weekly-file" accept=".xlsx,.xls,.xlsm" style="display:none">' +
              '<button type="button" class="mdl-btn mdl-btn-add" id="mdl-weekly-upload">📁 엑셀 업로드</button>' +
              '<button type="button" class="mdl-btn mdl-btn-del" id="mdl-weekly-delete" style="display:none">🗑️ 삭제</button>' +
            '</div>' +
            '<div id="mdl-weekly-preview" class="mdl-weekly-preview"></div>' +
          '</div>' +
        '</div>' +
      '</div>';

    // 탭 전환
    const tabs = content.querySelectorAll('.mdl-tab');
    const panes = content.querySelectorAll('.mdl-pane');
    tabs.forEach(function(t) {
      t.addEventListener('click', function() {
        tabs.forEach(function(x) { x.classList.remove('is-active'); });
        t.classList.add('is-active');
        panes.forEach(function(pn) { pn.hidden = pn.dataset.mpane !== t.dataset.mtab; });
      });
    });

    const sel = document.getElementById('mdl-proj-select');
    const tableBox = document.getElementById('mdl-table-box');
    const weeklyStatus = document.getElementById('mdl-weekly-status');
    const weeklyPreview = document.getElementById('mdl-weekly-preview');
    const weeklyFileInput = document.getElementById('mdl-weekly-file');
    const weeklyUploadBtn = document.getElementById('mdl-weekly-upload');
    const weeklyDeleteBtn = document.getElementById('mdl-weekly-delete');
    const saveBtn = document.getElementById('mdl-save-btn');
    const saveMsg = document.getElementById('mdl-save-msg');
    const addBtn = document.getElementById('mdl-add-btn');

    function loadStatusNote(projectKey) {
      if (!projectKey) return;
      fetch('/projects/' + encodeURIComponent(projectKey) + '/models/detail', { credentials: 'same-origin' })
        .then(function(r) { return r.json(); })
        .then(function(d) {
          const ta = document.getElementById('mdl-status-note');
          if (ta) ta.value = d.status_note || '';
        })
        .catch(function() {});
    }

    function loadWeeklyPlan(projectKey) {
      if (!projectKey) {
        weeklyStatus.textContent = '프로젝트를 선택하세요';
        weeklyStatus.className = 'mdl-weekly-status';
        weeklyPreview.innerHTML = '';
        weeklyDeleteBtn.style.display = 'none';
        return;
      }
      weeklyStatus.textContent = '불러오는 중...';
      weeklyStatus.className = 'mdl-weekly-status';
      fetch('/projects/' + encodeURIComponent(projectKey) + '/weekly-plan', { credentials: 'same-origin' })
        .then(function(r) { return r.json(); })
        .then(function(d) {
          if (d.has_plan && d.url) {
            weeklyStatus.textContent = '✅ ' + (d.file_name || '엑셀 파일') + ' (' + (d.uploaded_at || '').substring(0, 10) + ')';
            weeklyStatus.className = 'mdl-weekly-status ok';
            weeklyDeleteBtn.style.display = 'inline-block';
            weeklyPreview.innerHTML = '<img src="' + d.url + '" style="max-width:100%;border:1px solid #E5E7EB;border-radius:8px;background:#fff;cursor:zoom-in;" alt="weekly plan" onerror="this.parentNode.textContent=&#39;이미지 로드 실패&#39;;this.remove();">';
          } else {
            weeklyStatus.textContent = '등록된 주차별 계획이 없습니다';
            weeklyStatus.className = 'mdl-weekly-status';
            weeklyPreview.innerHTML = '';
            weeklyDeleteBtn.style.display = 'none';
          }
        })
        .catch(function(e) {
          weeklyStatus.textContent = '❌ 로드 실패: ' + e.message;
          weeklyStatus.className = 'mdl-weekly-status err';
        });
    }

    weeklyUploadBtn.addEventListener('click', function() {
      if (!sel.value) { alert('프로젝트를 먼저 선택하세요'); return; }
      weeklyFileInput.click();
    });

    weeklyFileInput.addEventListener('change', function() {
      const file = weeklyFileInput.files[0];
      if (!file) return;
      const projectKey = sel.value;
      if (!projectKey) { alert('프로젝트를 먼저 선택하세요'); weeklyFileInput.value = ''; return; }
      weeklyStatus.textContent = '⏳ 업로드 중...';
      weeklyStatus.className = 'mdl-weekly-status';
      const fd = new FormData();
      fd.append('file', file);
      fetch('/admin/projects/' + encodeURIComponent(projectKey) + '/weekly-plan', {
        method: 'POST', credentials: 'same-origin', body: fd
      })
        .then(function(r) { return r.json().then(function(d) { return { ok: r.ok, d: d }; }); })
        .then(function(res) {
          if (res.ok && res.d.ok) {
            weeklyStatus.textContent = '✅ ' + (res.d.file_name || file.name);
            weeklyStatus.className = 'mdl-weekly-status ok';
            loadWeeklyPlan(projectKey);
          } else {
            weeklyStatus.textContent = '❌ ' + (res.d.detail || '업로드 실패');
            weeklyStatus.className = 'mdl-weekly-status err';
          }
        })
        .catch(function(e) {
          weeklyStatus.textContent = '❌ ' + e.message;
          weeklyStatus.className = 'mdl-weekly-status err';
        })
        .finally(function() { weeklyFileInput.value = ''; });
    });

    weeklyDeleteBtn.addEventListener('click', function() {
      const projectKey = sel.value;
      if (!projectKey) return;
      if (!confirm('주차별 계획을 삭제하시겠습니까?')) return;
      fetch('/admin/projects/' + encodeURIComponent(projectKey) + '/weekly-plan', {
        method: 'DELETE', credentials: 'same-origin'
      })
        .then(function(r) { return r.json(); })
        .then(function(d) {
          if (d.ok) loadWeeklyPlan(projectKey);
          else alert('삭제 실패: ' + (d.detail || ''));
        })
        .catch(function(e) { alert('삭제 실패: ' + e.message); });
    });

    saveBtn.addEventListener('click', function() {
      syncDomToData(tableBox);
      saveModels(saveMsg);
    });

    const noteSaveBtn = document.getElementById('mdl-note-save');
    if (noteSaveBtn) {
      noteSaveBtn.addEventListener('click', function() {
        if (!_currentProjectKey) { alert('프로젝트를 먼저 선택하세요'); return; }
        const noteEl = document.getElementById('mdl-status-note');
        const msg = document.getElementById('mdl-note-msg');
        fetch('/admin/projects/' + encodeURIComponent(_currentProjectKey) + '/status-note', {
          method: 'PUT', credentials: 'same-origin',
          headers: { 'Content-Type': 'application/json' },
          body: JSON.stringify({ note: noteEl ? noteEl.value : '' })
        })
          .then(function(r) { return r.json(); })
          .then(function(d) {
            if (d.ok) { msg.textContent = '✅ 현황 저장 완료'; msg.style.color = '#059669'; }
            else { msg.textContent = '❌ 저장 실패'; msg.style.color = '#DC2626'; }
          })
          .catch(function(e) { msg.textContent = '❌ ' + e.message; msg.style.color = '#DC2626'; });
      });
    }

    addBtn.addEventListener('click', function() {
      const name = prompt('모델명 (예: CUP-100)');
      if (!name || !name.trim()) return;
      const trimmed = name.trim();
      if (_modelsData.some(function(m) { return m.id === trimmed; })) {
        alert('이미 존재하는 모델명입니다');
        return;
      }
      _modelsData.push({ id: trimmed, name: trimmed, group: '양산', status: '정상', progress: 0, price: 0, material_cost: 0 });
      renderTable(tableBox);
    });

    fetch('/admin/config/projects?division_id=semiconductor', { credentials: 'same-origin' })
      .then(function(r) { return r.json(); })
      .then(function(data) {
        const projects = data.projects || [];
        sel.innerHTML = '';
        if (!projects.length) {
          sel.innerHTML = '<option value="">프로젝트 없음</option>';
          tableBox.innerHTML = '<div class="mdl-empty">프로젝트가 없습니다</div>';
          return;
        }
        projects.forEach(function(pj) {
          const opt = document.createElement('option');
          opt.value = pj.id;
          opt.textContent = pj.label || pj.id;
          sel.appendChild(opt);
        });
        loadModels(sel.value, tableBox);
        loadWeeklyPlan(sel.value); loadStatusNote(sel.value);
      })
      .catch(function(e) {
        sel.innerHTML = '<option value="">오류: ' + e.message + '</option>';
        tableBox.innerHTML = '<div class="mdl-empty">프로젝트 목록 로드 실패</div>';
      });

    sel.addEventListener('change', function() {
      loadModels(sel.value, tableBox);
      loadWeeklyPlan(sel.value); loadStatusNote(sel.value);
    });
  }

  document.addEventListener('click', function(e) {
    var navItem = e.target.closest('.nav-item[data-page="models"]');
    if (navItem) {
      // 기존 active 해제
      document.querySelectorAll('.nav-item').forEach(function(n) { n.classList.remove('active'); });
      navItem.classList.add('active');
      var crumb = document.getElementById('v2-crumb-page');
      if (crumb) crumb.textContent = '모델 관리';
      renderModelsPage();
    }
  });
})();
</script>
</body>
</html>
"""

@app.get("/admin/upload", response_class=HTMLResponse)
def admin_upload_page(admin_auth: Optional[str] = Cookie(default=None)):
    if not _verify_session(admin_auth):
        return RedirectResponse(url="/admin/login?next=/admin/upload", status_code=302)
    return HTMLResponse(content=_ADMIN_UPLOAD_HTML)


@app.get("/admin/uploads/history")
def admin_uploads_history(_admin: int = Depends(get_admin_session)):
    """업로드된 PPT 목록 (삭제 UI용)"""
    items = _read_json(LATEST_FILE, [])
    history = []
    for item in items:
        doc_id = item.get("doc_id", "")
        history.append({
            "doc_id": doc_id,
            "file_name": item.get("file_name") or item.get("filename") or "",
            "upload_timestamp": item.get("upload_timestamp") or item.get("uploaded_at") or "",
            "product_count": len(item.get("products", [])),
            "slide_count": len(item.get("slide_images", {})),
        })
    # 최신 업로드 먼저
    history.sort(key=lambda x: x.get("upload_timestamp", ""), reverse=True)
    return {"items": history}


@app.get("/admin/reports/all")
def admin_list_reports_all(_admin: int = Depends(get_admin_session)):
    # Admin 전용: hidden 항목까지 모두 반환
    raw = _read_json(LATEST_FILE, [])
    enriched = []
    for r in (raw or []):
        try:
            meta = (r or {}).get("report_meta") or {}
            fname = (r or {}).get("file_name") or ""
            parsed = _parse_report_filename(fname, meta.get("date", ""))
            manual_projs = (r or {}).get("manual_projects")
            if manual_projs:
                parsed = dict(parsed)
                parsed["projects"] = list(manual_projs)
            mw = (r or {}).get("week_override")
            if mw:
                parsed["week"] = mw
            if parsed.get("projects"):
                _wk = parsed.get("week")
                parsed["display_title"] = ", ".join(parsed["projects"]) + (" · W" + str(_wk) + " 주간보고" if _wk else "")
            classified = _classify_report_status(
                (r or {}).get("upload_timestamp", ""),
                parsed.get("date") or meta.get("date", "")
            )
            r_enriched = dict(r or {})
            r_enriched["parsed"] = parsed
            r_enriched["display_title"] = parsed.get("display_title", "")
            r_enriched["report_status"] = classified.get("status", "published")
            r_enriched["d_day"] = classified.get("d_day")
            r_enriched["hidden"] = bool((r or {}).get("hidden", False))
            enriched.append(r_enriched)
        except Exception:
            enriched.append(r)
    split_cards = []
    for e in enriched:
        try:
            split_cards.extend(_split_report_by_project(e))
        except Exception:
            split_cards.append(e)
    split_cards = _dedupe_cards_by_project(split_cards)
    return {"reports": split_cards}


@app.post("/admin/reports")
def admin_create_manual_report(payload: dict, _admin: int = Depends(get_admin_session)):
    # 수기 프로젝트 생성. PPT 없이 리포트 카드만 추가.
    project_name = (payload.get("project_name") or "").strip()
    division_id = (payload.get("division_id") or "semiconductor").strip() or "semiconductor" 

    # 중복 프로젝트명 검사 (정규화 후 비교)
    _norm_new = _normalize_project_name(project_name)
    if _norm_new:
        try:
            _existing = _load_latest_reports()
        except Exception:
            _existing = []
        for _r in (_existing or []):
            # 기존 리포트가 가진 모든 프로젝트명 (override 포함)
            _parsed = (_r or {}).get("parsed") or {}
            _projs = list(_parsed.get("projects") or [])
            _manual = (_r or {}).get("manual_projects") or []
            _projs.extend(_manual)
            _overrides = (_r or {}).get("project_overrides") or {}
            _all_names = []
            for _pn in _projs:
                _all_names.append(_overrides.get(_pn, _pn))
            for _pn in _all_names:
                if _normalize_project_name(_pn) == _norm_new:
                    raise HTTPException(status_code=400, detail=f"이미 존재하는 프로젝트명입니다: {project_name}")
    week = payload.get("week")
    headline = (payload.get("headline") or "").strip()
    if not project_name:
        raise HTTPException(status_code=400, detail="project_name required")
    try:
        week_int = int(week)
        if week_int < 1 or week_int > 53:
            raise ValueError
    except Exception:
        raise HTTPException(status_code=400, detail="valid week required (1-53)")

    import uuid, datetime as _dt
    doc_id = "manual_" + uuid.uuid4().hex[:12]
    now_iso = _dt.datetime.utcnow().isoformat()
    # 해당 주차의 월요일 날짜 계산 (ISO week)
    today = _dt.date.today()
    year = today.year
    try:
        report_date = _dt.date.fromisocalendar(year, week_int, 1).isoformat()
    except Exception:
        report_date = today.isoformat()

    display_title = project_name + " · W" + str(week_int) + " 주간보고"
    default_sections = [
        {"id": "sec_" + uuid.uuid4().hex[:8], "title": "현황", "kind": "text", "blocks": []},
        {"id": "sec_" + uuid.uuid4().hex[:8], "title": "주차별 계획", "kind": "mixed", "blocks": []},
        {"id": "sec_" + uuid.uuid4().hex[:8], "title": "주요내용", "kind": "mixed", "blocks": []},
    ]
    new_report = {
        "doc_id": doc_id,
        "file_name": "",
        "upload_timestamp": now_iso,
        "is_manual": True,
        "division_id": division_id,
        "report_meta": {"date": report_date},
        "manual_projects": [project_name],
        "manual_projects": [project_name],
        "products": [{
            "name": project_name,
            "category": "",
            "status": "",
            "headline": headline,
            "summary_bullets": [],
            "sections": default_sections,
            "division_id": division_id
        }],
        "project_overrides": {},
        "week_override": week_int,
        "manual_projects": [project_name],
    }
    items = _read_json(LATEST_FILE, [])
    items.insert(0, new_report)
    _write_json(LATEST_FILE, items)
    try:
        _sync_report_to_notes(new_report)
    except Exception as _e:
        print("[WARN] create single sync failed:", _e)
    return {"ok": True, "doc_id": doc_id, "display_title": display_title}


@app.delete("/admin/reports/{doc_id}")
def admin_delete_manual_report(doc_id: str, _admin: int = Depends(get_admin_session)):
    # 수기 리포트 삭제. 안전을 위해 is_manual=True 인 것만 허용.
    items = _read_json(LATEST_FILE, [])
    target = None
    for it in items:
        if it.get("doc_id") == doc_id:
            target = it
            break
    if not target:
        raise HTTPException(status_code=404, detail="not found")
    if not target.get("is_manual"):
        raise HTTPException(status_code=403, detail="only manual reports can be deleted here")
    items = [it for it in items if it.get("doc_id") != doc_id]
    _write_json(LATEST_FILE, items)
    try:
        _resync_notes_from_latest()
    except Exception as _e:
        print("[WARN] delete resync failed:", _e)
    return {"ok": True, "doc_id": doc_id}



# ─── reports → notes.json 자동 동기화 (수기 doc용) ───


def _extract_due_date(text: str) -> tuple:
    """(clean_text, due_iso or None) 반환.
    
    ★ 원문 그대로 유지 정책 (2026-07-27~):
      - 사용자가 적은 문장을 그대로 앱에 표시하기 위해 텍스트는 건드리지 않음.
      - due_date 만 추출 (화살표 인식형 파서 사용).
      - 명시적 대괄호 [YYYY-MM-DD] / [YY-MM-DD] 는 우선순위로 처리.
    """
    import re as _re
    from datetime import date as _date
    if not text or not isinstance(text, str):
        return (text or "", None)

    due_iso = None

    # 우선순위 1: 대괄호 안의 명시적 ISO 날짜
    m = _re.search(r'\[\s*(20\d{2})[-./](\d{1,2})[-./](\d{1,2})\s*\]', text)
    if m:
        try:
            due_iso = _date(int(m.group(1)), int(m.group(2)), int(m.group(3))).isoformat()
        except Exception:
            pass
    else:
        m = _re.search(r'\[\s*(\d{2})[-./](\d{1,2})[-./](\d{1,2})\s*\]', text)
        if m:
            try:
                due_iso = _date(2000 + int(m.group(1)), int(m.group(2)), int(m.group(3))).isoformat()
            except Exception:
                pass

    # 우선순위 2: 새 파서 (화살표 인식 + 마지막 세그먼트 우선)
    if not due_iso:
        due_iso = _extract_due_date_from_text(text) or None

    # ★ text 는 정리하지 않고 원문 그대로 반환
    return (text, due_iso)

def _collect_inline_runs(node, inherited_style: dict) -> list:
    """BeautifulSoup 노드를 재귀 순회하며 [(text, style_dict), ...] 반환.
    style_dict keys: color, bold, italic, underline, size_scale.
    부모의 스타일은 자식에게 상속됨.
    """
    from bs4 import NavigableString
    if isinstance(node, NavigableString):
        return [(str(node), dict(inherited_style))]

    style = dict(inherited_style)
    name = getattr(node, 'name', '') or ''
    name = name.lower()

    # ★ 블록 태그(div/p/br)는 줄바꿈 run으로 처리 — 저장 시 줄바꿈 보존
    if name in ('div', 'p'):
        children_runs = []
        for child in node.children:
            children_runs.extend(_collect_inline_runs(child, style))
        return children_runs + [('\n', {})]  # 블록 끝에 줄바꿈
    if name == 'br':
        return [('\n', {})]  # <br>은 줄바꿈

    # 태그별 스타일 부여
    if name == 'font':
        c = node.get('color')
        if c:
            style['color'] = c.strip()
        sz = node.get('size')
        if sz:
            try:
                # <font size="1..7"> → 스케일 (3이 기본)
                n = int(str(sz).strip())
                _map = {1: 0.7, 2: 0.85, 3: 1.0, 4: 1.15, 5: 1.3, 6: 1.5, 7: 1.75}
                style['size_scale'] = _map.get(n, 1.0)
            except ValueError:
                pass
    elif name in ('b', 'strong'):
        style['bold'] = True
    elif name in ('i', 'em'):
        style['italic'] = True
    elif name == 'u':
        style['underline'] = True

    # style 속성은 태그 종류에 관계없이 파싱 (b/strong/i/em/u/span/div/p 등 모두)
    # 예: <b style="color:red">, <span style="font-weight:bold">, <div style="color:blue">
    s_attr = node.get('style') if hasattr(node, 'get') else None
    if s_attr:
        import re as _re
        for m in _re.finditer(r'([\w-]+)\s*:\s*([^;]+)', s_attr):
            k = m.group(1).strip().lower()
            v = m.group(2).strip()
            if k == 'color':
                style['color'] = v
            elif k == 'font-weight':
                if v.lower() in ('bold', 'bolder') or (v.isdigit() and int(v) >= 600):
                    style['bold'] = True
            elif k == 'font-style' and v.lower() == 'italic':
                style['italic'] = True
            elif k == 'text-decoration' and 'underline' in v.lower():
                style['underline'] = True
            elif k == 'font-size':
                # "14px", "1.2em", "120%"
                _m2 = _re.match(r'(\d+(?:\.\d+)?)\s*(px|em|%)?', v)
                if _m2:
                    num = float(_m2.group(1))
                    unit = (_m2.group(2) or 'px').lower()
                    if unit == 'px':
                        style['size_scale'] = num / 14.0  # 14px 기준
                    elif unit == 'em':
                        style['size_scale'] = num
                    elif unit == '%':
                        style['size_scale'] = num / 100.0

    runs = []
    for child in node.children:
        runs.extend(_collect_inline_runs(child, style))
    return runs


def _runs_to_output(runs: list) -> tuple:
    """runs [(text, style), ...] → (plain_text, text_runs 또는 None).
    연속된 동일 style은 병합. 스타일이 하나도 없으면 text_runs=None 반환.
    """
    if not runs:
        return ('', None)
    # 앞뒤 공백을 유지하되 완전히 빈 run은 제외
    filtered = [(t, st) for (t, st) in runs if t]
    if not filtered:
        return ('', None)
    # 병합 — 단, 줄바꿈 run은 항상 별도 유지
    merged = []
    for t, st in filtered:
        if t == '\n':
            merged.append([t, {}])  # 줄바꿈은 스타일 없이 별도
        elif merged and merged[-1][1] == st and merged[-1][0] != '\n':
            merged[-1] = (merged[-1][0] + t, st)
        else:
            merged.append([t, st])
    plain = ''.join(t for t, _ in merged)
    # 스타일이 하나라도 있는지
    has_style = any(bool(st) for _, st in merged)
    if not has_style:
        return (plain, None)
    # text_runs 포맷 (스키마 slim)
    text_runs = []
    for t, st in merged:
        run = {'text': t}
        if st.get('color'):
            run['color'] = st['color']
        if st.get('bold'):
            run['bold'] = True
        if st.get('italic'):
            run['italic'] = True
        if st.get('underline'):
            run['underline'] = True
        if 'size_scale' in st and st['size_scale'] != 1.0:
            run['size_scale'] = round(st['size_scale'], 2)
        text_runs.append(run)
    return (plain, text_runs)


def _html_body_to_items(html: str) -> list:
    """HTML(관리자 편집기 결과)을 모바일 앱이 이해하는 items 배열로 변환.
    
    우선 items가 이미 있으면 그대로 반환 (구조화된 저장 포맷).
    없으면 HTML 파싱 (기존 방식, 호환용).
    """
    # items가 이미 있으면 그대로 반환
    if isinstance(html, list):
        return html
    if isinstance(html, dict) and 'items' in html:
        return html['items']
    print(f"[DEBUG _html_body_to_items] input length={len(html)}, first 500 chars: {html[:500]}")
    """
    print(f"[DEBUG _html_body_to_items] input length={len(html)}, first 500 chars: {html[:500]}")
    
    

    # 중첩된 <div> 평탄화 (브라우저 contenteditable 버그 대응)
    # 예: <div>외부<div>내부1</div><div>내부2</div></div> → <div>외부</div><div>내부1</div><div>내부2</div>
    from bs4 import BeautifulSoup as _BS
    _soup = _BS(html, 'html.parser')
    for _div in list(_soup.find_all('div')):
        _nested = _div.find_all('div', recursive=False)
        if _nested:
            for _child in _nested:
                _div.insert_after(_child)
    html = str(_soup)

    변환 규칙:
    - <ol depth 0> → {"type": "bullet", "text": "1. ..."}
    - <ol depth 1> → {"type": "sub", "text": "1) ..."}
    - <ol depth 2+> → {"type": "sub", "text": "① ..."} (원형숫자)
    - <div>*text*</div> 또는 <div>★text</div> → highlight
    - <div data-indent="N">text</div> (N>=1) → sub
    - <div>일반</div> → bullet
    - <del>...</del> 완전 제거
    - <ins>text</ins> → 텍스트만 유지
    """
    from bs4 import BeautifulSoup, NavigableString
    
    if not html or not isinstance(html, str):
        return []
    
    # <del> 제거, <ins>는 텍스트만 남김
    soup = BeautifulSoup(html, 'html.parser')
    for del_tag in soup.find_all('del'):
        del_tag.decompose()
    for ins_tag in soup.find_all('ins'):
        ins_tag.unwrap()
    
    # depth별 마커 함수
    CIRCLED = ['①','②','③','④','⑤','⑥','⑦','⑧','⑨','⑩',
               '⑪','⑫','⑬','⑭','⑮','⑯','⑰','⑱','⑲','⑳']
    
    def marker_for_depth(depth: int, idx: int, is_ul: bool = False) -> str:
        """depth+index+is_ul로 실제 표시 문자 생성."""
        if is_ul:
            # ul(불릿): depth 0=•, 1=-, 2=*, 3+=·
            if depth == 0:
                return "•"
            elif depth == 1:
                return "-"
            elif depth == 2:
                return "*"
            else:
                return "·"
        else:
            # ol(숫자): depth 0=1., 1=1), 2=①...
            style = depth % 3
            if style == 0:
                return f"{idx + 1}."
            elif style == 1:
                return f"{idx + 1})"
            else:  # circled
                return CIRCLED[idx] if idx < len(CIRCLED) else f"{idx + 1}."
    
    def get_ol_depth(ol) -> int:
        """ol의 effective depth 계산 (data-depth-override 우선)."""
        override = ol.get('data-depth-override')
        if override:
            try:
                return int(override)
            except ValueError:
                pass
        # DOM nesting으로 계산
        depth = 0
        cur = ol.parent
        while cur is not None:
            if getattr(cur, 'name', None) == 'ol':
                ovr = cur.get('data-depth-override') if hasattr(cur, 'get') else None
                if ovr:
                    try:
                        return int(ovr) + depth + 1
                    except ValueError:
                        pass
                depth += 1
            cur = getattr(cur, 'parent', None)
        return depth
    
    items = []
    
    def process_ol(ol, extra_prefix=""):
        """ol의 li들을 순회하며 items에 추가 (nested ol도 재귀). data-depth-override 우선."""
        print(f"[DEBUG process_ol] called, tag={ol.name}, children={len(list(ol.children))}, is_ul={ol.name == 'ul'}")
        # data-depth-override 속성이 있으면 우선 사용 (에디터 Tab 들여쓰기)
        _override = ol.get('data-depth-override')
        if _override is not None:
            try:
                depth = int(_override)
            except (ValueError, TypeError):
                depth = get_ol_depth(ol)
        else:
            depth = get_ol_depth(ol)
        # start 값 반영
        try:
            start = int(ol.get('data-start') or ol.get('start') or 1)
        except (ValueError, TypeError):
            start = 1
        
        lis = [c for c in ol.children if getattr(c, 'name', None) == 'li']
        for i, li in enumerate(lis):
            # li 안의 인라인 runs 수집 (색상/굵기 등 스타일 보존)
            li_runs = []
            nested_ols = []
            for child in li.children:
                if getattr(child, 'name', None) in ('ol', 'ul'):
                    nested_ols.append(child)
                else:
                    li_runs.extend(_collect_inline_runs(child, {}))
            
            li_plain, li_text_runs = _runs_to_output(li_runs)
            li_plain = li_plain.rstrip()
            print(f"[DEBUG process_ol] li[{i}] tag_name={repr(ol.name)}, plain='{li_plain[:50]}...'")
            if li_plain:
                # ul은 bullet(•), ol은 숫자 마커
                if ol.name == 'ul':
                    marker = '•'
                else:
                    marker = marker_for_depth(depth, i + (start - 1), is_ul=(ol.name == 'ul'))
                item_type = 'bullet' if depth == 0 else 'sub'
                clean_text, due_iso = _extract_due_date(li_plain)
                full_text = f"{marker} {clean_text}" if clean_text else f"{marker}"
                _item = {
                    'type': item_type,
                    'text': full_text,
                }
                if due_iso:
                    _item['due_date'] = due_iso
                # text_runs가 있고 due_date 추출로 텍스트가 안 바뀐 경우만 유지
                if li_text_runs and clean_text == li_plain:
                    prefixed_runs = [{'text': f"{marker} "}] + li_text_runs
                    _item['text_runs'] = prefixed_runs
                items.append(_item)
            
            # nested ol 재귀 처리
            for nested in nested_ols:
                process_ol(nested)
    
    def emit_runs(runs: list):
        """[(text, style), ...] 리스트 → bullet 아이템으로 emit.
        text_runs가 있으면 함께 저장. '*'/'★'/'※' 특별 취급 없음."""
        plain, text_runs = _runs_to_output(runs)
        plain = plain.rstrip()
        if not plain:
            return
        # text_runs가 있으면 각 run의 앞뒤 공백/양쪽 트림 동기화
        if text_runs:
            # plain의 정리에 맞춰 text_runs도 정리
            # 앞 공백 제거 (앱이 sub 들여쓰기를 대신하므로 이중 적용 방지)
            while text_runs and text_runs[0]['text'] and text_runs[0]['text'] != text_runs[0]['text'].lstrip():
                text_runs[0]['text'] = text_runs[0]['text'].lstrip()
                if not text_runs[0]['text']:
                    text_runs.pop(0)
                else:
                    break
            # 뒤 공백 제거
            while text_runs and text_runs[-1]['text'] and text_runs[-1]['text'].rstrip() != text_runs[-1]['text']:
                text_runs[-1]['text'] = text_runs[-1]['text'].rstrip()
                if not text_runs[-1]['text']:
                    text_runs.pop()
                else:
                    break
        clean, due_iso = _extract_due_date(plain)
        _bul = {'type': 'bullet', 'text': clean}
        if due_iso:
            _bul['due_date'] = due_iso
        if text_runs and clean == plain:
            # due_date 추출로 text가 바뀌지 않은 경우만 text_runs 유지 (안전)
            _bul['text_runs'] = text_runs
        items.append(_bul)

    def process_element(el):
        """editor 직속 자식 처리 (div, ol, ul 등). 인라인 태그는 상위 flush에서 처리됨.
        '*' / '★' / '※' 특별 취급 없이 원문 그대로 유지."""
        name = getattr(el, 'name', None)
        if name in ('ol', 'ul'):
            process_ol(el)
        elif name in ('div', 'p'):
            # 인라인 runs 수집 (색상/굵기 등 보존)
            # 단, 직접 자식 중 block 요소(ul, ol, div, p)는 제외 — 중첩 구조에서 텍스트 합침 방지
            BLOCK_TAGS = {'ul', 'ol', 'div', 'p', 'li'}
            _el_runs = []
            _has_block_child = False
            for _c in el.children:
                _c_name = getattr(_c, 'name', None)
                if _c_name in BLOCK_TAGS:
                    _has_block_child = True
                    continue  # block 자식은 건너뛰고, 나중에 재귀 처리
                _el_runs.extend(_collect_inline_runs(_c, {}))
            _el_plain, _el_text_runs = _runs_to_output(_el_runs)
            text = (_el_plain or '').rstrip()
            
            # block 자식(ul/ol)은 1단계(find_all)에서 이미 처리됨 → 여기선 중첩 div/p만 재귀
            if _has_block_child:
                # 자기 텍스트 먼저 flush (block 자식보다 위에 위치해야 함)
                if text:
                    clean, due_iso = _extract_due_date(text)
                    _bul = {'type': 'bullet', 'text': clean}
                    if _el_text_runs and clean == text:
                        _bul['text_runs'] = _el_text_runs
                    if due_iso:
                        _bul['due_date'] = due_iso
                    items.append(_bul)
                    text = ''  # 아래 공통 emit 방지
                for _c in el.children:
                    _c_name = getattr(_c, 'name', None)
                    if _c_name in ('div', 'p'):
                        process_element(_c)
                    elif _c_name in ('ul', 'ol'):
                        process_ol(_c)
                # div 자체의 텍스트도 있으면 추가로 item 생성
                if not text:
                    return
            # 들여쓰기 체크: data-indent 속성 우선, 없으면 텍스트 앞 공백 개수로 계산
            _indent_attr = el.get('data-indent')
            try:
                indent = int(_indent_attr) if _indent_attr else 0
            except (ValueError, TypeError):
                indent = 0
            
            # data-indent가 없으면 앞 공백 개수로 depth 계산 (2칸=1depth)
            if indent == 0:
                _leading = _el_plain or ''
                # 앞쪽 공백/탭/nbsp 개수 세기
                _stripped = _leading.lstrip(' \t\u00a0')
                _space_count = len(_leading) - len(_stripped)
                indent = _space_count // 2  # 2칸당 1depth
                if indent > 0:
                    print(f"[DEBUG indent] text='{text[:30]}...', space_count={_space_count}, depth={indent}")
            else:
                print(f"[DEBUG indent] text='{text[:30]}...', data-indent={_indent_attr}, depth={indent}")
            
            print(f"[DEBUG FINAL] text='{text[:40]}', indent={indent}, _has_block_child={_has_block_child}")
            if indent >= 1:
                clean, due_iso = _extract_due_date(text)
                _sub = {'type': 'sub', 'text': clean}
                if _el_text_runs and clean == text:
                    _sub['text_runs'] = _el_text_runs
                if due_iso:
                    _sub['due_date'] = due_iso
                items.append(_sub)
            else:
                clean, due_iso = _extract_due_date(text)
                _bul = {'type': 'bullet', 'text': clean}
                if _el_text_runs and clean == text:
                    _bul['text_runs'] = _el_text_runs
                if due_iso:
                    _bul['due_date'] = due_iso
                items.append(_bul)
    
    # 최상위 자식들 순회.
    # 연속된 NavigableString + 인라인 태그(font/span/b/i/u/em/strong/sub/sup 등)는 하나의 텍스트로 병합.
    # <br> 만나면 flush(하나의 아이템으로 emit).
    INLINE_TAGS = {'font', 'span', 'b', 'i', 'u', 'em', 'strong', 'sub', 'sup', 'small', 'mark', 'a', 'code'}
    runs_buf = []  # [(text, style_dict), ...]

    def flush_buf():
        if runs_buf:
            emit_runs(list(runs_buf))
            runs_buf.clear()

    # 처리된 리스트 추적 (중복 방지)
    _processed_lists = set()
    
    for child in soup.children:
        name = getattr(child, 'name', None)
        if isinstance(child, NavigableString):
            runs_buf.append((str(child), {}))
        elif name == 'br':
            flush_buf()
        elif name in INLINE_TAGS:
            runs_buf.extend(_collect_inline_runs(child, {}))
        elif name in ('ul', 'ol'):
            # 제자리에서 바로 처리 (문서 순서 유지)
            flush_buf()
            process_ol(child)
        else:
            # 블록 요소(div, p 등): 지금까지 쌓인 인라인 flush 후 별도 처리
            flush_buf()
            process_element(child)
    flush_buf()
    
    return items


def _products_to_note_cards(products: list, week_override=None) -> list:
    """products[] → notes.json 형식의 cards[] 로 변환."""
    if not isinstance(products, list):
        return []
    
    cards = []
    for prod in products:
        if not isinstance(prod, dict):
            continue
        title = (prod.get('name') or '').strip()
        if not title:
            continue
        
        card_sections = []
        for sec in (prod.get('sections') or []):
            if not isinstance(sec, dict):
                continue
            sec_title = (sec.get('title') or '').strip()
            blocks = sec.get('blocks') or []
            
            items = []
            for blk in blocks:
                if not isinstance(blk, dict):
                    continue
                kind = blk.get('kind', 'text')
                if kind == 'text':
                    body = blk.get('body', '')
                    items.extend(_html_body_to_items(body))
                elif kind == 'file':
                    # 파일 첨부는 photo item으로. xlsx/xlsm 이면 PNG 자동 변환.
                    fname = blk.get('file_name', '') or ''
                    url = blk.get('url', '') or ''
                    if fname or url:
                        _lower = (fname + ' ' + url).lower()
                        _photo_ref = url  # 기본값: URL 그대로
                        if '.xlsx' in _lower or '.xlsm' in _lower:
                            try:
                                local_path = _xlsx_url_to_local_path(url)
                                if local_path and local_path.exists():
                                    # division_id 유추 (파서는 report 전체를 안 봄 → product 이름으로)
                                    _div_id = ''
                                    try:
                                        _div_id = _cl.derive_division_from_project(title) or ''
                                    except Exception:
                                        _div_id = ''
                                    if not _div_id:
                                        _div_id = 'semiconductor'
                                    png_bytes = _xlsx_file_to_png_bytes(local_path)
                                    _photo_ref = _save_note_photo(_div_id, png_bytes, ext='png')
                                    print(f'[parser] xlsx→png 자동 변환: {fname} → {_photo_ref}')
                                else:
                                    print(f'[parser] xlsx 파일 찾지 못함, URL 그대로: {url}')
                            except Exception as _e:
                                print(f'[parser] xlsx→png 변환 실패 (URL 그대로 사용): {_e}')
                        items.append({
                            'type': 'photo',
                            'text': fname or url,
                            'photo_ref': _photo_ref,
                        })
            
            if items:
                new_sec = {
                    'title': sec_title,
                    'items': items,
                }
                # sales_visible 이 명시적으로 지정돼 있으면 유지 (기본 True 는 저장 안 함)
                if isinstance(sec.get('sales_visible'), bool):
                    new_sec['sales_visible'] = sec['sales_visible']
                card_sections.append(new_sec)
        
        if card_sections:
            card = {
                'title': title,
                'sections': card_sections,
            }
            # summary/headline 등도 복사
            if prod.get('headline'):
                card['headline'] = prod.get('headline')
            if prod.get('status'):
                card['status'] = prod.get('status')
            cards.append(card)
    
    return cards


def _sync_report_to_notes(it: dict) -> None:
    """수기 doc 하나를 notes.json에 upsert.
    
    - division_id 유추: it['division_id'] > product.name derive > semiconductor 폴백
    - cards는 title 기준 merge
    - _normalize_note_cards로 후처리
    """
    try:
        if not isinstance(it, dict):
            return
        if not it.get('is_manual'):
            return
        
        products = it.get('products') or []
        cards = _products_to_note_cards(products, week_override=it.get('week_override'))
        if not cards:
            return
        
        # division_id 결정
        division_id = (it.get('division_id') or '').strip()
        if not division_id:
            # product.name으로 유추 시도
            for prod in products:
                if isinstance(prod, dict) and prod.get('name'):
                    try:
                        division_id = _cl.derive_division_from_project(prod['name']) or ''
                    except Exception:
                        pass
                    if division_id:
                        break
        if not division_id:
            division_id = 'semiconductor'  # 폴백
        
        # 정규화
        parsed = {'cards': cards}
        _normalize_note_cards(parsed)
        cards = parsed['cards']
        
        # sales_input 자동 계산 (주차별 계획 섹션)
        try:
            products = it.get('products') or []
            for card in cards:
                for sec in card.get('sections', []) or []:
                    # products에서 같은 title의 section 찾아 sales_input 가져옴
                    sec_title = (sec.get('title') or '').strip()
                    for prod in products:
                        for psec in prod.get('sections', []) or []:
                            if (psec.get('title') or '').strip() == sec_title:
                                # [SALES v2] 우선순위:
                                #   1) 섹션에 첨부된 xlsx + sales_prices → 자동 파싱/계산
                                #   2) sales_data (구조화 JSON)
                                #   3) sales_input (구식 텍스트)
                                sd = psec.get('sales_data')
                                si = psec.get('sales_input') or ''
                                prices = psec.get('sales_prices') or {}

                                # 1) xlsx 파싱 (섹션 내 첨부 파일 자동 인식)
                                #    reports_latest.json 은 blocks[] 를 씀 → blocks 우선 탐색
                                xlsx_ref = None
                                for _blk in (psec.get('blocks') or []):
                                    if not isinstance(_blk, dict):
                                        continue
                                    _kind = (_blk.get('kind') or '').lower()
                                    _url = _blk.get('url') or _blk.get('file_url') or ''
                                    _fname = _blk.get('file_name') or ''
                                    _all = (_url + ' ' + _fname).lower()
                                    if '.xlsx' in _all or '.xlsm' in _all:
                                        xlsx_ref = _url or _fname
                                        break
                                # fallback: items[] (notes 스키마)
                                if not xlsx_ref:
                                    xlsx_ref = _find_xlsx_in_section_items(psec.get('items') or [])
                                if not xlsx_ref:
                                    src_meta = psec.get('sales_source') or {}
                                    xlsx_ref = src_meta.get('file_url')

                                parsed_ok = False
                                if xlsx_ref:
                                    local_path = _xlsx_url_to_local_path(xlsx_ref)
                                    if local_path and local_path.exists():
                                        try:
                                            parsed = parse_sales_excel(local_path)
                                            sec['sales_source'] = {
                                                'file_name': local_path.name,
                                                'file_url': str(xlsx_ref),
                                                'parsed_at': __import__('datetime').datetime.now().isoformat(timespec='seconds'),
                                                'parsed': parsed,
                                            }
                                            sec['sales_prices'] = prices
                                            if prices:
                                                r = compute_sales_from_parsed(parsed, prices)
                                                sec['sales_summary'] = r['sales_summary']
                                                sec['sales_summary_data'] = r['sales_summary_data']
                                                sec['sales_computed_at'] = r['sales_computed_at']
                                            parsed_ok = True
                                        except Exception as _pe:
                                            print(f'[sales_xlsx parse ERROR] {xlsx_ref} :: {_pe}')

                                # 2) fallback: sales_data (JSON)
                                if not parsed_ok and isinstance(sd, dict) and (sd.get('prices') or sd.get('weeks')):
                                    sec['sales_data'] = sd
                                    r = compute_sales_from_data(sd)
                                    sec['sales_summary'] = r['sales_summary']
                                    sec['sales_computed_at'] = r['sales_computed_at']
                                    parsed_ok = True

                                # 3) fallback: sales_input (구식 텍스트)
                                if not parsed_ok and si:
                                    sec['sales_input'] = si
                                    r = compute_sales_from_input(si)
                                    sec['sales_summary'] = r['sales_summary']
                                    sec['sales_computed_at'] = r['sales_computed_at']
                                break
        except Exception as _e:
            print(f'[sales_compute in sync ERROR] {_e}')
        
        # notes.json upsert
        from datetime import datetime
        data = _load_notes()
        notes_map = data.setdefault('notes', {})
        existing = notes_map.get(division_id) or {}
        existing_cards = existing.get('cards') or []
        
        def _norm_title(t):
            return re.sub(r'\s+', '', (t or '').strip()).lower()
        
        by_title = {}
        order = []
        for c in existing_cards:
            if not isinstance(c, dict):
                continue
            tkey = _norm_title(c.get('title'))
            if not tkey or tkey in by_title:
                continue
            by_title[tkey] = c
            order.append(tkey)
        
        def _merge_card_preserve_manual_state(old_card, new_card):
            try:
                def _norm_text(x):
                    return re.sub(r'\s+', ' ', (x or '').strip()).lower()

                old_secs = old_card.get('sections') or []
                new_secs = new_card.get('sections') or []
                old_item_map = {}

                for os_ in old_secs:
                    os_title = (os_.get('title') or '').strip()
                    for oi in (os_.get('items') or []):
                        if not isinstance(oi, dict):
                            continue
                        keys = []
                        oi_id = (oi.get('item_id') or '').strip()
                        oi_text = (oi.get('text') or '').strip()
                        if oi_id:
                            keys.append(('id', os_title, oi_id))
                        if oi_text:
                            keys.append(('text', os_title, _norm_text(oi_text)))
                            keys.append(('text-any', '', _norm_text(oi_text)))
                        for k in keys:
                            old_item_map[k] = oi

                for ns in new_secs:
                    ns_title = (ns.get('title') or '').strip()
                    for ni in (ns.get('items') or []):
                        if not isinstance(ni, dict):
                            continue
                        match = None
                        ni_id = (ni.get('item_id') or '').strip()
                        ni_text = (ni.get('text') or '').strip()
                        if ni_id:
                            match = old_item_map.get(('id', ns_title, ni_id))
                        if not match and ni_text:
                            match = old_item_map.get(('text', ns_title, _norm_text(ni_text)))
                        if not match and ni_text:
                            match = old_item_map.get(('text-any', '', _norm_text(ni_text)))
                        if match:
                            if match.get('auto_due_hidden'):
                                ni['auto_due_hidden'] = True
                            if (match.get('due_date_override') or '').strip():
                                ni['due_date_override'] = (match.get('due_date_override') or '').strip()
                            if (match.get('due_date_auto') or '').strip() and not (ni.get('due_date_auto') or '').strip():
                                ni['due_date_auto'] = (match.get('due_date_auto') or '').strip()
                            try:
                                _apply_effective_due_date(ni)
                            except Exception:
                                pass
                return new_card
            except Exception:
                return new_card

        # 기존 notes 전체에서 수동 상태 맵 생성
        manual_state_map = {}
        for _oc in existing_cards:
            if not isinstance(_oc, dict):
                continue
            _ocard_title = (_oc.get('title') or '').strip()
            for _os in (_oc.get('sections') or []):
                if not isinstance(_os, dict):
                    continue
                _osec_title = (_os.get('title') or '').strip()
                for _oi in (_os.get('items') or []):
                    if not isinstance(_oi, dict):
                        continue
                    _otext = re.sub(r'\s+', ' ', (_oi.get('text') or '').strip()).lower()
                    if not _otext:
                        continue
                    _key = (_ocard_title, _osec_title, _otext)
                    manual_state_map[_key] = {
                        'auto_due_hidden': bool(_oi.get('auto_due_hidden')),
                        'due_date_override': (_oi.get('due_date_override') or '').strip(),
                        'due_date_auto': (_oi.get('due_date_auto') or '').strip(),
                    }

        def _apply_manual_state(card_obj):
            try:
                _ct = (card_obj.get('title') or '').strip()
                for _ns in (card_obj.get('sections') or []):
                    if not isinstance(_ns, dict):
                        continue
                    _st = (_ns.get('title') or '').strip()
                    for _ni in (_ns.get('items') or []):
                        if not isinstance(_ni, dict):
                            continue
                        _tx = re.sub(r'\s+', ' ', (_ni.get('text') or '').strip()).lower()
                        if not _tx:
                            continue
                        _m = manual_state_map.get((_ct, _st, _tx))
                        if not _m:
                            continue
                        if _m.get('auto_due_hidden'):
                            _ni['auto_due_hidden'] = True
                        if _m.get('due_date_override'):
                            _ni['due_date_override'] = _m['due_date_override']
                        if _m.get('due_date_auto') and not (_ni.get('due_date_auto') or '').strip():
                            _ni['due_date_auto'] = _m['due_date_auto']
                        try:
                            _apply_effective_due_date(_ni)
                        except Exception:
                            pass
                return card_obj
            except Exception:
                return card_obj

        for c in cards:
            tkey = _norm_title(c.get('title'))
            if not tkey:
                continue
            c = _apply_manual_state(c)
            if tkey in by_title:
                by_title[tkey] = _merge_card_preserve_manual_state(by_title[tkey], c)
            else:
                by_title[tkey] = c
                order.append(tkey)
        
        merged_cards = [by_title[k] for k in order]
        
        report_date = ''
        meta = it.get('report_meta') or {}
        if isinstance(meta, dict):
            report_date = meta.get('date', '') or ''
        
        existing_out = notes_map.get(division_id) or {}
        existing_cards_out = list(existing_out.get('cards') or [])

        # 현재 동기화 대상 카드 title 집합
        _target_titles = set()
        for _c in cards:
            _t = _norm_title((_c or {}).get('title'))
            if _t:
                _target_titles.add(_t)

        _final_cards = []
        for _c in existing_cards_out:
            if not isinstance(_c, dict):
                _final_cards.append(_c)
                continue
            _t = _norm_title((_c or {}).get('title'))
            if _t and _t in _target_titles:
                continue
            _final_cards.append(_c)

        # 현재 doc 카드만 추가 (기존 순서 유지 + 신규 카드 append)
        for _c in merged_cards:
            if isinstance(_c, dict):
                _final_cards.append(_c)

        notes_map[division_id] = {
            'report_date': report_date or existing.get('report_date', ''),
            'updated_at': datetime.now().isoformat(),
            'raw_text': existing.get('raw_text', ''),
            'cards': _final_cards,
        }
        _save_notes(data)
    except Exception as e:
        import traceback
        print(f'[sync_report_to_notes ERROR] {e}')
        traceback.print_exc()


@app.get("/admin/projects/{project_key}/models")
def admin_get_project_models(project_key: str, _admin: int = Depends(get_admin_session)):
    """admin용 모델 목록 조회"""
    models = _get_project_models(project_key)
    return {
        "project_key": project_key,
        "models": models,
    }


@app.post("/admin/projects/{project_key}/models")
def admin_add_model(project_key: str, payload: dict, _admin: int = Depends(get_admin_session)):
    """모델 단건 추가. 같은 id가 이미 있으면 409 반환.
    payload: {"id": "cup-100", "name": "CUP-100", "group": "양산", "progress": 35, "status": "지연"}
    """
    _alias = {
        "havaplate": "hrva_plate",
        "hrvaplate": "hrva_plate",
        "hrva-plate": "hrva_plate",
    }
    _key = _alias.get(project_key.strip().lower(), project_key.strip())
    data = _load_models()
    proj = data.setdefault("projects", {}).setdefault(_key, {"models": []})
    existing_ids = {m.get("id") for m in proj.get("models", [])}
    norm = _normalize_model(payload or {}, existing_ids)
    if norm is None:
        raise HTTPException(status_code=400, detail="모델명이 필요하거나 ID가 중복됩니다.")
    proj.setdefault("models", []).append(norm)
    proj["models"].sort(key=lambda m: 0 if m.get("group") == "양산" else 1)
    _save_models(data)
    print(f"[models] POST {_key}: +{norm['id']} ({len(proj['models'])} models)")
    return {"ok": True, "project_key": _key, "model": norm, "total": len(proj["models"])}


@app.patch("/admin/projects/{project_key}/models/{model_id}")
def admin_update_model(project_key: str, model_id: str, payload: dict, _admin: int = Depends(get_admin_session)):
    """모델 단건 수정 (부분 업데이트)
    payload: 변경할 필드만 {"name": "...", "progress": 50, ...}"""
    from urllib.parse import unquote
    model_id = unquote(model_id)
    _alias = {
        "havaplate": "hrva_plate",
        "hrvaplate": "hrva_plate",
        "hrva-plate": "hrva_plate",
    }
    _key = _alias.get(project_key.strip().lower(), project_key.strip())
    data = _load_models()
    proj = data.setdefault("projects", {}).setdefault(_key, {"models": []})
    target = None
    for m in proj.get("models", []):
        if m.get("id") == model_id:
            target = m
            break
    if target is None:
        raise HTTPException(status_code=404, detail="모델을 찾을 수 없습니다.")
    if "name" in payload and str(payload["name"]).strip():
        target["name"] = str(payload["name"]).strip()
    if "group" in payload:
        g = str(payload["group"]).strip()
        if g in ("양산", "개발"):
            target["group"] = g
    if "progress" in payload:
        try:
            p = int(payload["progress"])
        except (ValueError, TypeError):
            p = target.get("progress", 0)
        target["progress"] = max(0, min(100, p))
    if "status" in payload:
        s = str(payload["status"]).strip()
        if s in ("정상", "주의", "지연"):
            target["status"] = s
    if "price" in payload:
        try:
            target["price"] = max(0, int(payload["price"]))
        except (ValueError, TypeError):
            pass
    if "material_cost" in payload:
        try:
            target["material_cost"] = max(0, int(payload["material_cost"]))
        except (ValueError, TypeError):
            pass
    _save_models(data)
    print(f"[models] PATCH {_key}/{model_id}")
    return {"ok": True, "project_key": _key, "model": target}


@app.delete("/admin/projects/{project_key}/models/{model_id}")
def admin_delete_model(project_key: str, model_id: str, _admin: int = Depends(get_admin_session)):
    """모델 단건 삭제"""
    from urllib.parse import unquote
    model_id = unquote(model_id)
    _alias = {
        "havaplate": "hrva_plate",
        "hrvaplate": "hrva_plate",
        "hrva-plate": "hrva_plate",
    }
    _key = _alias.get(project_key.strip().lower(), project_key.strip())
    data = _load_models()
    proj = data.setdefault("projects", {}).setdefault(_key, {"models": []})
    before = len(proj.get("models", []))
    proj["models"] = [m for m in proj.get("models", []) if m.get("id") != model_id]
    after = len(proj["models"])
    if before == after:
        raise HTTPException(status_code=404, detail="모델을 찾을 수 없습니다.")
    _save_models(data)
    print(f"[models] DELETE {_key}/{model_id}")
    return {"ok": True, "project_key": _key, "removed": before - after, "total": after}


@app.put("/admin/projects/{project_key}/models")
def admin_put_project_models(project_key: str, payload: dict, _admin: int = Depends(get_admin_session)):
    """admin용 모델 목록 저장 (전체 교체) - dev_type/process/weekly_plan/status_note 보존"""
    _alias = {
        "havaplate": "hrva_plate",
        "hrvaplate": "hrva_plate",
        "hrva-plate": "hrva_plate",
    }
    _key = _alias.get(project_key.strip().lower(), project_key.strip())

    raw_models = payload.get("models") or []
    if not isinstance(raw_models, list):
        raise HTTPException(status_code=400, detail="models는 배열이어야 합니다.")

    data = _load_models()
    projects = data.setdefault("projects", {})
    proj = projects.setdefault(_key, {})
    old_map = {m.get("id"): m for m in proj.get("models", []) if isinstance(m, dict)}

    normalized = []
    seen_ids = set()
    for m in raw_models:
        if not isinstance(m, dict):
            continue
        mid = str(m.get("id") or "").strip()
        name = str(m.get("name") or "").strip()
        if not mid or not name:
            continue
        if mid in seen_ids:
            continue
        seen_ids.add(mid)
        group = str(m.get("group") or "양산").strip()
        if group not in ("양산", "개발"):
            group = "양산"
        try:
            progress = int(m.get("progress") or 0)
        except (ValueError, TypeError):
            progress = 0
        progress = max(0, min(100, progress))
        status = str(m.get("status") or "정상").strip()
        try:
            price = int(m.get("price") or 0)
        except (ValueError, TypeError):
            price = 0
        try:
            material_cost = int(m.get("material_cost") or 0)
        except (ValueError, TypeError):
            material_cost = 0
        entry = {
            "id": mid,
            "name": name,
            "group": group,
            "progress": progress,
            "status": status,
            "price": max(0, price),
            "material_cost": max(0, material_cost),
        }
        old = old_map.get(mid) or {}
        if group == "개발":
            dev_type = str(m.get("dev_type") or old.get("dev_type") or "HVM").strip().upper()
            if dev_type not in ("HVM", "RPM"):
                dev_type = "HVM"
            entry["dev_type"] = dev_type
            proc = m.get("process") if isinstance(m.get("process"), list) else old.get("process")
            entry["process"] = proc if isinstance(proc, list) and len(proc) == 13 else _default_process()
        normalized.append(entry)

    normalized.sort(key=lambda m: 0 if m.get("group") == "양산" else 1)
    proj["models"] = normalized
    _save_models(data)
    print(f"[models] saved {_key}: {len(normalized)} models")
    return {"ok": True, "project_key": _key, "count": len(normalized)}


# ─── 주차별 계획 (프로젝트당 1개) ───

@app.get("/projects/{project_key}/weekly-plan")
def get_weekly_plan(project_key: str):
    """프로젝트 주차별 계획 조회 (앱용)"""
    _alias = {
        "havaplate": "hrva_plate",
        "hrvaplate": "hrva_plate",
        "hrva-plate": "hrva_plate",
    }
    _key = _alias.get(project_key.strip().lower(), project_key.strip())
    data = _load_models()
    proj = (data.get("projects") or {}).get(_key) or {}
    plan = proj.get("weekly_plan") or {}
    _ref = plan.get("photo_ref")
    return {
        "project_key": _key,
        "has_plan": bool(_ref),
        "file_name": plan.get("file_name"),
        "uploaded_at": plan.get("uploaded_at"),
        "photo_ref": _ref,
        "url": f"/note_photos/{_ref}" if _ref else None,
    }


@app.post("/admin/projects/{project_key}/weekly-plan")
async def admin_upload_weekly_plan(
    project_key: str,
    file: UploadFile = File(...),
    _admin: int = Depends(get_admin_session),
):
    """주차별 계획 엑셀 업로드 → 첫 시트를 PNG로 변환 → 저장 (보고서 엑셀과 동일 방식)"""
    _alias = {
        "havaplate": "hrva_plate",
        "hrvaplate": "hrva_plate",
        "hrva-plate": "hrva_plate",
    }
    _key = _alias.get(project_key.strip().lower(), project_key.strip())

    orig_name = file.filename or "weekly_plan.xlsx"
    lower = orig_name.lower()
    if not (lower.endswith(".xlsx") or lower.endswith(".xlsm")):
        raise HTTPException(status_code=400, detail="xlsx/xlsm 파일만 허용")
    raw = await file.read()
    if not raw:
        raise HTTPException(status_code=400, detail="빈 파일")

    try:
        import openpyxl
        from io import BytesIO
        import base64
        wb = openpyxl.load_workbook(BytesIO(raw), data_only=True)
        ws = wb.worksheets[0]
        data_url = _excel_sheet_to_preview_data_url(ws)
        b64 = data_url.split(",", 1)[1]
        png_bytes = base64.b64decode(b64)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"엑셀 변환 실패: {e}")

    asset_id = "plan_" + _new_asset_id(_key)
    out_path = _photo_path(_key, asset_id, "png")
    out_path.write_bytes(png_bytes)
    photo_ref = f"{_key}/{asset_id}.png"

    data = _load_models()
    proj = data.setdefault("projects", {}).setdefault(_key, {"models": []})
    # 기존 계획 이미지가 있으면 삭제
    old_ref = (proj.get("weekly_plan") or {}).get("photo_ref")
    if old_ref:
        try:
            _delete_note_photo(old_ref)
        except Exception:
            pass
    proj["weekly_plan"] = {
        "file_name": orig_name,
        "uploaded_at": __import__("datetime").datetime.now().isoformat(),
        "photo_ref": photo_ref,
    }
    _save_models(data)
    print(f"[weekly_plan] saved {_key}: {orig_name} -> {photo_ref}")
    return {
        "ok": True,
        "project_key": _key,
        "file_name": orig_name,
        "photo_ref": photo_ref,
        "url": f"/note_photos/{photo_ref}",
    }


@app.delete("/admin/projects/{project_key}/weekly-plan")
def admin_delete_weekly_plan(project_key: str, _admin: int = Depends(get_admin_session)):
    """주차별 계획 삭제"""
    _alias = {
        "havaplate": "hrva_plate",
        "hrvaplate": "hrva_plate",
        "hrva-plate": "hrva_plate",
    }
    _key = _alias.get(project_key.strip().lower(), project_key.strip())
    data = _load_models()
    proj = data.setdefault("projects", {}).setdefault(_key, {"models": []})
    if "weekly_plan" in proj:
        _ref = (proj.get("weekly_plan") or {}).get("photo_ref")
        if _ref:
            try:
                _delete_note_photo(_ref)
            except Exception:
                pass
        del proj["weekly_plan"]
        _save_models(data)
    return {"ok": True, "project_key": _key}


@app.put("/admin/reports/{doc_id}")
def admin_update_report(doc_id: str, payload: dict, _admin: int = Depends(get_admin_session)):
    # 수기 편집 저장
    # payload 예: {"project_overrides": {"메이져모듈": "메이저모듈"}, "products": [...]}
    items = _read_json(LATEST_FILE, [])
    # FCM alarm hook: 저장 전 상태 스냅샷
    _fcm_before_snapshot = _snapshot_notes_status()
    edit_mode = payload.get("edit_mode", False)
    if edit_mode:
        # 편집 모드: 알람 억제, 세션에 snapshot 저장
        _save_edit_session(doc_id, _fcm_before_snapshot)
        _cleanup_expired_sessions()
        print(f"[FCM] edit mode ON for {doc_id}: alarm suppressed")
    found = False
    for it in items:
        if it.get("doc_id") != doc_id:
            continue
        found = True
        try:
            _now = datetime.datetime.utcnow()
            _iso = _now.isocalendar()
            _week_now = int(_iso[1])
            _date_now = _now.date().isoformat()

            _meta = it.get("report_meta") or {}
            _meta["week"] = _week_now
            _meta["date"] = _date_now
            it["report_meta"] = _meta
            it["week_override"] = _week_now

            _parsed = it.get("parsed") or {}
            _projects = list(_parsed.get("projects") or [])
            if not _projects and (it.get("products") or []):
                _p0 = (it.get("products") or [])[0] or {}
                _nm = (_p0.get("name") or "").strip()
                if _nm:
                    _projects = [_nm]
            _projects = [str(x).strip() for x in _projects if str(x).strip()]
            if _projects:
                _parsed["projects"] = _projects
            _parsed["week"] = _week_now
            _parsed["date"] = _date_now
            _parsed["display_title"] = ", ".join(_projects) + (" · W" + str(_week_now) + " 주간보고" if _week_now else "")
            it["parsed"] = _parsed
            it["display_title"] = _parsed["display_title"]
        except Exception as _e:
            print("[WARN] week refresh failed:", _e)
        overrides = payload.get("project_overrides")
        # 이름 변경 시 다른 리포트와 중복되는지 검사
        if isinstance(overrides, dict) and overrides:
            _new_names = set()
            for _v in overrides.values():
                _n = _normalize_project_name(_v)
                if _n:
                    _new_names.add(_n)
            for _r in items:
                if (_r or {}).get("doc_id") == doc_id:
                    continue
                _parsed = (_r or {}).get("parsed") or {}
                _projs = list(_parsed.get("projects") or [])
                _manual = (_r or {}).get("manual_projects") or []
                _projs.extend(_manual)
                _overrides_o = (_r or {}).get("project_overrides") or {}
                for _pn in _projs:
                    _final = _overrides_o.get(_pn, _pn)
                    if _normalize_project_name(_final) in _new_names:
                        raise HTTPException(status_code=400, detail="이미 존재하는 프로젝트명입니다: " + str(_final))
        if isinstance(overrides, dict):
            existing_ov = it.get("project_overrides") or {}
            existing_ov.update({str(k): str(v) for k, v in overrides.items() if v})
            it["project_overrides"] = existing_ov

        # N-1: 저장 시점의 현재 주차로 강제 업데이트 (수정 저장 기준)
        from datetime import date
        current_week = date.today().isocalendar()[1]
        it["week_override"] = current_week
        
        # parsed.display_title도 업데이트
        if it.get("parsed") and it["parsed"].get("projects"):
            projects_str = ", ".join(it["parsed"]["projects"])
            it["parsed"]["week"] = current_week
            it["parsed"]["display_title"] = projects_str + " · W" + str(current_week) + " 주간보고"
            it["display_title"] = it["parsed"]["display_title"]

        new_products = payload.get("products")
        if isinstance(new_products, list):
            existing = it.get("products", [])
            merged = []
            for i, np in enumerate(new_products):
                base = dict(existing[i]) if i < len(existing) else {}
                for key in ("name", "headline", "category", "status", "summary_bullets", "sections", "sales_input", "sales_data", "sales_source", "sales_prices", "sales_summary", "sales_summary_data", "sales_computed_at", "sales_visible"):
                    if key in np:
                        base[key] = np[key]
                merged.append(base)
            it["products"] = merged
        break
    if not found:
        raise HTTPException(status_code=404, detail="doc_id not found")
    _write_json(LATEST_FILE, items)
    # 수기 doc이면 notes.json에는 현재 doc만 자동 동기화 (모바일 앱 반영용)
    for _it in items:
        if _it.get("doc_id") == doc_id and _it.get("is_manual"):
            try:
                _sync_report_to_notes(_it)
            except Exception as _e:
                print("[WARN] update single sync failed:", _e)
            break

    # FCM alarm hook: 저장 후 상태 diff → RED/ORANGE 신규 전이만 알람
    # edit_mode=True면 스킵 (세션에 저장됨, edit_done에서 처리)
    if not edit_mode:
        try:
            _fcm_after_snapshot = _snapshot_notes_status()
            _fcm_events = _detect_status_transitions(_fcm_before_snapshot, _fcm_after_snapshot)
            if _fcm_events:
                print(f"[FCM] 상태 전이 감지: {_fcm_events}")
                _fire_status_alarms(_fcm_events)
        except Exception as _fcm_e:
            print(f"[FCM] hook 오류: {_fcm_e}")
    else:
        print(f"[FCM] edit mode: alarm deferred to edit_done")

    return {"ok": True, "doc_id": doc_id, "edit_mode": edit_mode}




@app.post("/admin/reports/{doc_id}/edit_done")
def admin_report_edit_done(doc_id: str, _admin: int = Depends(get_admin_session)):
    """편집 완료: 저장된 before snapshot과 현재 상태 비교 → 알람 1회 발송"""
    _before_snapshot = _clear_edit_session(doc_id)
    if not _before_snapshot:
        return {"ok": True, "doc_id": doc_id, "message": "no edit session found"}
    
    try:
        _after_snapshot = _snapshot_notes_status()
        _events = _detect_status_transitions(_before_snapshot, _after_snapshot)
        if _events:
            print(f"[FCM] edit_done: 상태 전이 감지 → {_events}")
            _fire_status_alarms(_events)
            return {"ok": True, "doc_id": doc_id, "alarms_fired": len(_events)}
        else:
            print(f"[FCM] edit_done: 변경 없음 (no alarms)")
            return {"ok": True, "doc_id": doc_id, "alarms_fired": 0}
    except Exception as e:
        print(f"[FCM] edit_done 오류: {e}")
        return {"ok": False, "doc_id": doc_id, "error": str(e)}

@app.post("/admin/reports/{doc_id}/section-file")
async def admin_upload_section_file(
    doc_id: str,
    file: UploadFile = File(...),
    _admin: int = Depends(get_admin_session)
):
    # 수기 프로젝트의 섹션 안 파일 블록용 업로드.
    # uploads/manual/{doc_id}/ 아래에 저장하고 다운로드 URL 반환.
    import re as _re, uuid as _uuid, shutil as _shutil
    items = _read_json(LATEST_FILE, [])
    target = None
    for it in items:
        if it.get("doc_id") == doc_id:
            target = it
            break
    if not target:
        raise HTTPException(status_code=404, detail="doc_id not found")
    if not target.get("is_manual"):
        raise HTTPException(status_code=403, detail="only manual reports support section files")

    if not file.filename:
        raise HTTPException(status_code=400, detail="file required")

    # 안전한 파일명: 원본 이름 유지하되 경로 문자 제거
    safe_name = _re.sub(r"[\\/]", "_", file.filename).strip()
    if not safe_name:
        safe_name = "file"
    # 중복 시 앞에 랜덤 프리픽스
    manual_dir = UPLOAD_DIR / "manual" / doc_id
    manual_dir.mkdir(parents=True, exist_ok=True)
    dest = manual_dir / safe_name
    if dest.exists():
        stem = dest.stem
        suffix = dest.suffix
        safe_name = stem + "_" + _uuid.uuid4().hex[:6] + suffix
        dest = manual_dir / safe_name

    data = await file.read()
    dest.write_bytes(data)

    # xlsx/xlsm 이면 PNG로도 자동 변환해서 photo_ref 확보 (앱 렌더용)
    photo_ref = None
    lower = safe_name.lower()
    if lower.endswith(('.xlsx', '.xlsm')):
        try:
            division_id = _derive_division_id_from_report(target)
            png_bytes = _xlsx_file_to_png_bytes(dest)
            photo_ref = _save_note_photo(division_id, png_bytes, ext='png')
            print(f"[section-file] xlsx→png 자동 변환 OK: {safe_name} (div={division_id}) → {photo_ref}")
        except Exception as _e:
            print(f"[section-file] xlsx→png 변환 실패 (원본은 저장됨): {_e}")

    return {
        "ok": True,
        "doc_id": doc_id,
        "file_name": safe_name,
        "size": len(data),
        "url": "/admin/manual-files/" + doc_id + "/" + safe_name,
        "photo_ref": photo_ref,
    }



@app.post("/admin/reports/{doc_id}/section-parse-xlsx")
async def admin_parse_section_xlsx(
    doc_id: str,
    payload: dict,
    _admin: int = Depends(get_admin_session)
):
    """섹션에 첨부된 xlsx 파일을 파싱해서 models/weeks 리턴 (판가 입력용)."""
    import re as _re
    filename = (payload or {}).get('filename') or ''
    filename = filename.strip()
    if not filename:
        raise HTTPException(status_code=400, detail='filename required')
    if _re.search(r'[\\/]|\.\.', filename):
        raise HTTPException(status_code=400, detail='invalid filename')
    fp = UPLOAD_DIR / 'manual' / doc_id / filename
    if not fp.exists() or not fp.is_file():
        raise HTTPException(status_code=404, detail='xlsx not found')
    if not filename.lower().endswith(('.xlsx', '.xlsm')):
        raise HTTPException(status_code=400, detail='not an excel file')
    try:
        parsed = parse_sales_excel(fp)
        return {'ok': True, 'filename': filename, 'parsed': parsed}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f'parse error: {e}')

@app.get("/admin/manual-files/{doc_id}/{filename}")
def admin_download_section_file(
    doc_id: str,
    filename: str,
    _admin: int = Depends(get_admin_session)
):
    # 수기 프로젝트 섹션 파일 다운로드
    import re as _re
    from fastapi.responses import FileResponse
    if _re.search(r"[\\/]|\.\.", filename):
        raise HTTPException(status_code=400, detail="invalid filename")
    fp = UPLOAD_DIR / "manual" / doc_id / filename
    if not fp.exists() or not fp.is_file():
        raise HTTPException(status_code=404, detail="not found")
    return FileResponse(str(fp), filename=filename)


@app.get("/admin/manual-files/{doc_id}/{filename}/xlsx-preview")
def admin_xlsx_preview(
    doc_id: str,
    filename: str,
    _admin: int = Depends(get_admin_session)
):
    # 수기 프로젝트 섹션 엑셀 미리보기: 병합 셀 + 스타일 유지, 빈 행/열 제거
    import re as _re
    if _re.search(r"[\\/]|\.\.", filename):
        raise HTTPException(status_code=400, detail="invalid filename")
    fp = UPLOAD_DIR / "manual" / doc_id / filename
    if not fp.exists() or not fp.is_file():
        raise HTTPException(status_code=404, detail="not found")
    if not filename.lower().endswith((".xlsx", ".xls")):
        raise HTTPException(status_code=400, detail="not an excel file")
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(fp), data_only=True)
        ws = wb.active
        sheet_name = ws.title if ws else ""

        MAX_ROWS = 50
        MAX_COLS = 20

        # 1) 병합 범위 파악: (min_row, min_col) → (rowspan, colspan)
        merged_map = {}
        skip_cells = set()
        for mr in ws.merged_cells.ranges:
            r1, c1, r2, c2 = mr.min_row, mr.min_col, mr.max_row, mr.max_col
            merged_map[(r1, c1)] = (r2 - r1 + 1, c2 - c1 + 1)
            for rr in range(r1, r2 + 1):
                for cc in range(c1, c2 + 1):
                    if (rr, cc) != (r1, c1):
                        skip_cells.add((rr, cc))

        # 2) 실제 데이터 범위 (빈 행/열 제거 + 시작 여백 trim)
        max_r = min(ws.max_row or 0, MAX_ROWS)
        max_c = min(ws.max_column or 0, MAX_COLS)

        non_empty = []
        for r in range(1, max_r + 1):
            for c in range(1, max_c + 1):
                v = ws.cell(row=r, column=c).value
                if v is not None and str(v).strip() != "":
                    non_empty.append((r, c))

        if not non_empty:
            wb.close()
            return {"ok": True, "sheet": sheet_name, "rows": [], "n_rows": 0, "n_cols": 0}

        real_min_r = min(r for r, _ in non_empty)
        real_max_r = max(r for r, _ in non_empty)
        real_min_c = min(c for _, c in non_empty)
        real_max_c = max(c for _, c in non_empty)

        # 3) 행 데이터 구성 (실제 데이터 bounding box만)
        rows = []
        for r in range(real_min_r, real_max_r + 1):
            row_cells = []
            for c in range(real_min_c, real_max_c + 1):
                if (r, c) in skip_cells:
                    continue
                cell = ws.cell(row=r, column=c)
                v = cell.value
                text = "" if v is None else str(v)

                # 병합 정보
                rspan, cspan = merged_map.get((r, c), (1, 1))

                # 현재 crop 범위를 벗어나는 병합은 잘라냄
                if r + rspan - 1 > real_max_r:
                    rspan = real_max_r - r + 1
                if c + cspan - 1 > real_max_c:
                    cspan = real_max_c - c + 1

                # 스타일 정보 (rgb / theme / indexed 모두 처리)
                bg = ""
                fg = ""
                bold = False

                # openpyxl theme index → 근사 RGB (Office 기본 테마)
                _THEME_RGB = {
                    0: "FFFFFF", 1: "000000", 2: "E7E6E6", 3: "44546A",
                    4: "4472C4", 5: "ED7D31", 6: "A5A5A5", 7: "FFC000",
                    8: "5B9BD5", 9: "70AD47",
                }

                def _apply_tint(hex6: str, tint) -> str:
                    """openpyxl tint(-1.0~1.0)를 근사 반영."""
                    try:
                        t = float(tint or 0)
                    except Exception:
                        t = 0.0
                    if not t:
                        return hex6
                    r = int(hex6[0:2], 16)
                    g = int(hex6[2:4], 16)
                    b = int(hex6[4:6], 16)
                    if t < 0:
                        f = 1 + t
                        r = int(r * f); g = int(g * f); b = int(b * f)
                    else:
                        r = int(r + (255 - r) * t)
                        g = int(g + (255 - g) * t)
                        b = int(b + (255 - b) * t)
                    r = max(0, min(255, r)); g = max(0, min(255, g)); b = max(0, min(255, b))
                    return f"{r:02X}{g:02X}{b:02X}"

                def _color_to_hex(color_obj):
                    if color_obj is None:
                        return ""
                    try:
                        ctype = getattr(color_obj, "type", None)
                        if ctype == "rgb":
                            rgb = getattr(color_obj, "rgb", None)
                            if rgb and isinstance(rgb, str) and rgb != "00000000":
                                return "#" + rgb[-6:]
                        elif ctype == "theme":
                            theme_idx = getattr(color_obj, "theme", None)
                            if isinstance(theme_idx, int) and theme_idx in _THEME_RGB:
                                base = _THEME_RGB[theme_idx]
                                tint = getattr(color_obj, "tint", 0)
                                return "#" + _apply_tint(base, tint)
                        elif ctype == "indexed":
                            # 흔한 인덱스만 처리 (필요시 확장)
                            _INDEXED = {
                                0: "000000", 1: "FFFFFF", 2: "FF0000", 3: "00FF00",
                                4: "0000FF", 5: "FFFF00", 6: "FF00FF", 7: "00FFFF",
                                64: "000000",  # system foreground
                            }
                            idx = getattr(color_obj, "indexed", None)
                            if isinstance(idx, int) and idx in _INDEXED:
                                return "#" + _INDEXED[idx]
                    except Exception:
                        return ""
                    return ""

                try:
                    if cell.fill and cell.fill.patternType:
                        bg = _color_to_hex(getattr(cell.fill, "fgColor", None))
                    if cell.font and cell.font.color:
                        fg = _color_to_hex(cell.font.color)
                    if cell.font and cell.font.bold:
                        bold = True
                except Exception:
                    pass

                # border 정보 (색상이 지정된 경우만 - default 검정선은 스킵)
                borders = {}
                try:
                    b = cell.border
                    for side_name in ("left", "right", "top", "bottom"):
                        side = getattr(b, side_name, None)
                        if not side or not side.style:
                            continue
                        col = side.color
                        if not col:
                            continue
                        # rgb 값 안전하게 추출 (openpyxl 버전 차이 대응)
                        rgb_val = None
                        try:
                            ctype = getattr(col, "type", None)
                            if ctype == "rgb":
                                raw = col.__dict__.get("rgb") if hasattr(col, "__dict__") else None
                                if isinstance(raw, str) and len(raw) >= 6:
                                    rgb_val = raw
                        except Exception:
                            pass
                        # default 검정선은 제외 (파일 크기 절약)
                        if not rgb_val:
                            continue
                        if rgb_val.upper() in ("00000000", "FF000000", "000000"):
                            continue
                        borders[side_name] = {
                            "color": "#" + rgb_val[-6:],
                            "style": side.style,
                        }
                except Exception:
                    pass

                row_cells.append({
                    "text": text,
                    "rowspan": rspan,
                    "colspan": cspan,
                    "bg": bg,
                    "fg": fg,
                    "bold": bold,
                    "borders": borders,
                })
            rows.append(row_cells)

        wb.close()
        return {
            "ok": True,
            "sheet": sheet_name,
            "rows": rows,
            "n_rows": real_max_r - real_min_r + 1,
            "n_cols": real_max_c - real_min_c + 1,
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"excel parse error: {e}")



@app.post("/admin/reports/{doc_id}/polish-text")
def admin_polish_text(
    doc_id: str,
    payload: dict,
    _admin: int = Depends(get_admin_session)
):
    # 수기 프로젝트 텍스트 블록 AI 다듬기
    import re as _re
    from html import unescape as _html_unescape

    if client is None:
        raise HTTPException(status_code=500, detail="OPENAI_API_KEY not configured")

    raw_html = str((payload or {}).get("text") or "").strip()
    section_title = str((payload or {}).get("section_title") or "").strip()

    if not raw_html:
        raise HTTPException(status_code=400, detail="text is required")

    def _html_to_text(s: str) -> str:
        s = _re.sub(r'<br\s*/?>', '\n', s, flags=_re.I)
        s = _re.sub(r'</(p|div|li|h[1-6])>', '\n', s, flags=_re.I)
        s = _re.sub(r'<[^>]+>', '', s)
        s = _html_unescape(s)
        s = _re.sub(r'\n{3,}', '\n\n', s)
        s = _re.sub(r'[ \t]+', ' ', s)
        return s.strip()

    # few-shot 예시 수집
    items = _read_json(LATEST_FILE, [])
    examples = []
    seen = set()

    def _push_example(txt: str):
        t = _html_to_text(txt or "")
        t = _re.sub(r'\s+', ' ', t).strip()
        if len(t) < 25:
            return
        if t in seen:
            return
        seen.add(t)
        examples.append(t[:700])

    for it in items:
        if (it or {}).get("doc_id") == doc_id:
            continue
        for prod in ((it or {}).get("products") or []):
            if prod.get("headline"):
                _push_example(str(prod.get("headline")))
            for b in (prod.get("summary_bullets") or []):
                _push_example(str(b))
            for sec in (prod.get("sections") or []):
                for blk in (sec.get("blocks") or []):
                    if (blk or {}).get("type") == "text" and (blk or {}).get("body"):
                        _push_example(str(blk.get("body")))
            if len(examples) >= 3:
                break
        if len(examples) >= 3:
            break

    fewshot = "\n\n".join(
        [f"[예시 {i+1}]\n{ex}" for i, ex in enumerate(examples[:3])]
    ) or "(예시 없음)"

    system_prompt = """
너는 한국어 제조/개발 주간보고 문체 편집기다.
원본을 최대한 보존하면서 표기와 문법만 다듬는 게 유일한 목적이다.

[절대 규칙 - 위반 시 실패]
1. 원본에 있는 문장/항목/번호를 절대 삭제하지 마라.
2. 원본이 N줄이면 결과도 정확히 N줄이어야 한다.
3. 원본에 없는 내용을 추가하지 마라.
4. 숫자/날짜/주차/수량/고유명사(제품명, 파트명, 부서명, 사람이름)를 절대 바꾸지 마라.
5. 항목 번호(1. 2. 3. / 1) 2) 3) / 가) 나) / ① ②)와 순서를 절대 바꾸지 마라.
6. 원본의 bullet, sub-bullet, 들여쓰기 계층을 그대로 유지하라.
7. 요약/축약/의역 절대 금지. 표현이 어색해도 원본 단어를 다른 단어로 대체하지 마라.

[표기 정돈 - 반드시 적용]
1. 화살표 통일: `-->`, `->`, `=>`, `⇒`, `~>` → 모두 `→`
2. 라인 시작의 하이픈 bullet: `- ` → `• ` (라인 첫 문자가 하이픈일 때만)
3. 줄 앞 불필요한 공백 정리 (단, 들여쓰기 계층은 유지):
   - 라인 시작의 tab은 스페이스 4개로
   - 라인 시작 스페이스는 4의 배수로 정규화
4. 콜론 뒤 공백 하나 확보: `단어:다음` → `단어: 다음`
5. 괄호 앞뒤 공백 정리: `단어(내용)` → `단어 (내용)`은 하지 마라. 원본 그대로.
6. 마침표/쉼표 뒤 공백 하나 확보

[문법 확인 - 반드시 적용]
1. 맞춤법/띄어쓰기 오류 교정
2. 오탈자 교정
3. 어색한 조사 수정 (예: "을" ↔ "를", "이" ↔ "가", "은" ↔ "는")
4. 잘못된 어미 자연스럽게 수정
5. 반드시 최소 1개 이상의 실질적 텍스트 수정을 수행할 것

[중요]
- 공백/줄바꿈만 바꾸는 건 수정으로 치지 않음
- 실제 글자를 바꿔야 수정임
- 원문이 완벽하면 그대로 두되, 표기 정돈은 무조건 적용

[하지 말 것]
- 문장 삭제 금지
- 항목 삭제 금지
- 요약 금지
- 표현 대체 금지 (원본 단어를 다른 단어로 바꾸지 마)
- HTML 태그 임의 제거 금지
- <ol>, <ul>, <li> 새로 만들지 마라 (원본에 없으면 만들지 마)

[출력]
- HTML fragment만 (설명 없이)
- markdown code fence(```) 사용 금지
- 원본 HTML 태그(<p>, <span>, <b>, <br>, <ul>, <li> 등) 그대로 유지
- 원본에 <br>이 N개면 결과도 <br>이 N개

[리스트 규칙 - 매우 중요]
- 원본에 "1)" "2)" "가)" "1." "-" 같은 수동 번호매기기가 있으면 그대로 텍스트로 유지 (bullet 규칙에 따라 `-`는 `•`로만 변환)
- 절대로 <ol>, <ul>, <li> 태그로 감싸지 마라. 원본이 <ol>/<ul>이 아니면 만들지 마라.
- 원본이 평문(줄바꿈만 있는 텍스트)이면 결과도 평문으로 유지
- 원본의 들여쓰기(공백, 탭)를 최대한 그대로 유지
""".strip()

    user_prompt = f"""
[섹션]
{section_title or "일반"}

[참고 문체 예시]
{fewshot}

[원본 HTML]
{raw_html}
""".strip()

    try:
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            temperature=0.2,
            max_tokens=4000,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
        )
        polished = (resp.choices[0].message.content or "").strip()
        polished = _re.sub(r"^```[a-zA-Z0-9_-]*\s*", "", polished)
        polished = _re.sub(r"\s*```$", "", polished)
        if not polished:
            raise HTTPException(status_code=500, detail="empty model output")

        # === 후처리: 표기 통일 강제 (프롬프트가 놓친 경우 대비) ===
        # 화살표 통일: --> / -> / => / ⇒ / ~> → →
        # 주의: <br>, </li>, <br/> 등 HTML 안의 - 는 건드리지 않음
        polished = _re.sub(r"-->", "→", polished)
        polished = _re.sub(r"=>", "→", polished)
        polished = _re.sub(r"⇒", "→", polished)
        polished = _re.sub(r"~>", "→", polished)
        # -> 는 HTML 태그(<->, </->)와 겹치지 않게 좌우 공백 있는 경우만
        polished = _re.sub(r"(?<=\s)->(?=\s)", "→", polished)
        polished = _re.sub(r"^->(?=\s)", "→", polished, flags=_re.M)

        # 라인 시작 하이픈 bullet: `- xxx` → `• xxx`
        # HTML 태그 안이 아니라 텍스트 라인 시작만 처리
        # <br> 직후, <p> 직후, <div> 직후, 문자열 시작
        polished = _re.sub(r"(^|<br\s*/?>|<p[^>]*>|<div[^>]*>|<li[^>]*>)(\s*)-\s", r"\1\2• ", polished, flags=_re.I)

        # 인라인 diff HTML 생성 (플레인 텍스트 비교 후 마크업)
        import difflib as _difflib
        orig_text = _html_to_text(raw_html)
        new_text = _html_to_text(polished)
        sm = _difflib.SequenceMatcher(None, orig_text, new_text, autojunk=False)
        diff_parts = []
        for tag, i1, i2, j1, j2 in sm.get_opcodes():
            a_seg = orig_text[i1:i2]
            b_seg = new_text[j1:j2]
            def _esc(t):
                return (t.replace("&", "&amp;")
                         .replace("<", "&lt;")
                         .replace(">", "&gt;")
                         .replace("\n", "<br>"))
            if tag == "equal":
                diff_parts.append(_esc(a_seg))
            elif tag == "delete":
                diff_parts.append('<del style="background:#FEE2E2;color:#B42318;text-decoration:line-through;padding:0 2px;border-radius:3px;">' + _esc(a_seg) + '</del>')
            elif tag == "insert":
                diff_parts.append('<ins style="background:#D1FAE5;color:#065F46;text-decoration:none;padding:0 2px;border-radius:3px;">' + _esc(b_seg) + '</ins>')
            elif tag == "replace":
                diff_parts.append('<del style="background:#FEE2E2;color:#B42318;text-decoration:line-through;padding:0 2px;border-radius:3px;">' + _esc(a_seg) + '</del>')
                diff_parts.append('<ins style="background:#D1FAE5;color:#065F46;text-decoration:none;padding:0 2px;border-radius:3px;">' + _esc(b_seg) + '</ins>')
        diff_html = "".join(diff_parts)

        return {
            "ok": True,
            "doc_id": doc_id,
            "section_title": section_title,
            "original": raw_html,
            "polished": polished,
            "diff_html": diff_html,
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"ai polish error: {e}")


@app.post("/admin/reports/{doc_id}/hide")
def admin_toggle_hide(doc_id: str, _admin: int = Depends(get_admin_session)):
    # 숨김/숨김해제 토글. hidden=True 이면 /reports 에서 제외됨.
    items = _read_json(LATEST_FILE, [])
    found = False
    new_state = False
    for it in items:
        if it.get("doc_id") == doc_id:
            new_state = not bool(it.get("hidden", False))
            it["hidden"] = new_state
            found = True
            break
    if not found:
        raise HTTPException(status_code=404, detail="doc_id not found")
    _write_json(LATEST_FILE, items)
    _resync_notes_from_latest()
    return {"ok": True, "doc_id": doc_id, "hidden": new_state}



def _resync_notes_from_latest() -> None:
    """reports_latest.json 기준으로 notes.json 자동동기화 카드를 재생성"""
    try:
        latest_items = _read_json(LATEST_FILE, [])
        notes_data = _load_notes()
        notes_map = notes_data.get("notes", {})
        for div_id, payload in list(notes_map.items()):
            if not isinstance(payload, dict):
                continue
            cards = payload.get("cards", []) or []
            preserved = [c for c in cards if isinstance(c, dict) and c.get("note_only") is True]
            payload["cards"] = preserved
        _save_notes(notes_data)
        for it in latest_items:
            try:
                _sync_report_to_notes(it)
            except Exception:
                import traceback
                traceback.print_exc()
    except Exception:
        import traceback
        traceback.print_exc()


@app.delete("/admin/doc/{doc_id}")
def admin_delete_doc(doc_id: str, _admin: int = Depends(get_admin_session)):
    """PPT 파일 단위 삭제: reports_latest.json + 관련 폴더/파일 전부 정리"""
    import shutil
    # 1) reports_latest.json 에서 제거
    items = _read_json(LATEST_FILE, [])
    before = len(items)
    items = [it for it in items if it.get("doc_id") != doc_id]
    after = len(items)
    if before == after:
        raise HTTPException(status_code=404, detail="해당 doc_id를 찾을 수 없습니다.")
    _write_json(LATEST_FILE, items)

    # 2) 파일 시스템 정리
    removed = []
    pptx_path = UPLOAD_DIR / f"{doc_id}.pptx"
    if pptx_path.exists():
        pptx_path.unlink()
        removed.append(str(pptx_path.name))

    for folder_dir in [SLIDES_DIR, SLIDE_IMAGES_DIR, CROPPED_DIR]:
        target = folder_dir / doc_id
        if target.exists() and target.is_dir():
            shutil.rmtree(target)
            removed.append(f"{folder_dir.name}/{doc_id}/")

    # 3) upload_hashes.json 에서 제거
    try:
        idx = _load_upload_hash_index()
        new_idx = {k: v for k, v in idx.items() if v.get("doc_id") != doc_id}
        if len(new_idx) != len(idx):
            _save_upload_hash_index(new_idx)
            removed.append("upload_hash entry")
    except Exception:
        pass

    return {"ok": True, "doc_id": doc_id, "removed": removed, "remaining": after}


@app.post("/admin/docs/delete-batch")
def admin_delete_docs_batch(
    payload: dict,
    _admin: int = Depends(get_admin_session),
):
    """PPT doc_id 여러 개 일괄 삭제"""
    doc_ids = payload.get("doc_ids", []) if isinstance(payload, dict) else []
    if not isinstance(doc_ids, list):
        raise HTTPException(status_code=400, detail="doc_ids must be a list")

    removed = []
    failed = []
    for did in doc_ids:
        if not isinstance(did, str):
            continue
        try:
            admin_delete_doc(did, _admin=_admin)
            removed.append(did)
        except HTTPException as e:
            failed.append({"doc_id": did, "status": e.status_code, "detail": e.detail})
        except Exception as e:
            failed.append({"doc_id": did, "detail": str(e)})

    return {
        "ok": True,
        "removed_count": len(removed),
        "removed": removed,
        "failed": failed,
    }


@app.delete("/admin/card")
def admin_delete_card(
    doc_id: str,
    product_idx: int,
    _admin: int = Depends(get_admin_session),
):
    """카드(product) 1장 단위 삭제"""
    items = _read_json(LATEST_FILE, [])
    target = None
    for it in items:
        if it.get("doc_id") == doc_id:
            target = it
            break
    if not target:
        raise HTTPException(status_code=404, detail="해당 doc_id를 찾을 수 없습니다.")

    products = target.get("products", [])
    if product_idx < 0 or product_idx >= len(products):
        raise HTTPException(status_code=404, detail=f"product_idx 범위 초과 (0~{len(products)-1})")

    removed_product = products.pop(product_idx)
    target["products"] = products
    _write_json(LATEST_FILE, items)
    return {
        "ok": True,
        "doc_id": doc_id,
        "removed_product_name": removed_product.get("name") or removed_product.get("product"),
        "remaining_products": len(products),
    }


@app.post("/admin/reset")
def admin_reset(password: str = Form(...)):
    _has_session = bool(_verify_session(admin_auth))
    if not _has_session and password != UPLOAD_PASSWORD:
        raise HTTPException(status_code=401, detail="비밀번호가 틀립니다.")
    _write_json(LATEST_FILE, [])
    return {"ok": True, "message": "최신 데이터가 초기화되었습니다."}


# =========================================================
# 어드민 업로드 페이지 HTML
# =========================================================
_ADMIN_UPLOAD_HTML = """
<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>사업부 보고 관리자</title>
<style>
  * { box-sizing: border-box; margin: 0; padding: 0; }
  body {
    font-family: -apple-system, BlinkMacSystemFont, "Apple SD Gothic Neo", sans-serif;
    background: #f5f7fa;
    color: #1f2937;
    line-height: 1.6;
  }
  header {
    background: #1E3A5F;
    color: white;
    padding: 20px 32px;
  }
  header h1 { font-size: 20px; font-weight: 700; }
  header .meta { font-size: 12px; color: #cbd5e1; margin-top: 4px; }

  nav.tabs {
    display: flex;
    background: white;
    border-bottom: 1px solid #e5e7eb;
    padding: 0 32px;
    gap: 8px;
    overflow-x: auto;
  }
  nav.tabs button {
    padding: 14px 20px;
    border: none;
    background: transparent;
    cursor: pointer;
    font-size: 15px;
    color: #6b7280;
    border-bottom: 3px solid transparent;
    transition: all 0.15s;
    font-weight: 500;
    white-space: nowrap;
  }
  nav.tabs button:hover { color: #1E3A5F; background: #f9fafb; }
  nav.tabs button.active {
    color: #1E3A5F;
    border-bottom-color: #1E3A5F;
    font-weight: 700;
  }

  main {
    max-width: 960px;
    margin: 0 auto;
    padding: 32px 20px;
  }
  .tab-content { display: none; }
  .tab-content.active { display: block; }

  .card {
    background: white;
    padding: 24px;
    border-radius: 12px;
    margin-bottom: 20px;
    box-shadow: 0 2px 8px rgba(0,0,0,0.05);
  }
  .card h2 { margin: 0 0 16px; font-size: 18px; color: #1f2937; }

  label { display: block; margin: 14px 0 6px; font-weight: 600; font-size: 13px; color: #333; }
  input, select, textarea {
    width: 100%;
    padding: 10px;
    border: 1px solid #ddd;
    border-radius: 8px;
    font-size: 14px;
    background: white;
    font-family: inherit;
  }
  button {
    background: #1E3A5F;
    color: white;
    border: none;
    padding: 12px 24px;
    border-radius: 8px;
    font-size: 15px;
    font-weight: 600;
    cursor: pointer;
    margin-top: 16px;
  }
  button:hover { background: #2c4d75; }
  button.delete {
    background: #ff3b30;
    padding: 6px 12px;
    font-size: 13px;
    margin-top: 0;
  }
  button.delete:hover { background: #d92e26; }

  .image-item {
    display: flex;
    align-items: center;
    gap: 16px;
    padding: 12px;
    border-bottom: 1px solid #eee;
  }
  .image-item:last-child { border-bottom: none; }
  .image-item img {
    width: 80px;
    height: 60px;
    object-fit: cover;
    border-radius: 6px;
    border: 1px solid #ddd;
  }
  .image-item .info { flex: 1; font-size: 13px; color: #555; }
  .image-item .info strong { color: #1f2937; }
  .image-item .badges { margin-top: 4px; }
  .badge {
    display: inline-block;
    padding: 2px 8px;
    border-radius: 4px;
    font-size: 11px;
    background: #e5e7eb;
    color: #374151;
    margin-right: 4px;
  }
  .badge.section { background: #fff3cd; color: #856404; }

  #uploadMsg {
    margin-top: 12px;
    padding: 10px;
    border-radius: 6px;
    font-size: 14px;
  }
  #uploadMsg.success { background: #d4edda; color: #155724; }
  #uploadMsg.error { background: #f8d7da; color: #721c24; }

  /* 노트 첨부 (표/사진) 버튼 */
  .note-item-row { display: flex; align-items: flex-start; gap: 6px; }
  .note-item-text { flex: 1; }
  .note-attach-btn {
    background: transparent; border: 1px solid #d1d5db; color: #6b7280;
    padding: 2px 8px; font-size: 11px; border-radius: 4px; cursor: pointer;
    margin: 0; line-height: 1.4; white-space: nowrap;
  }
  .note-attach-btn:hover { background: #f3f4f6; color: #1E3A5F; border-color: #1E3A5F; }
  .note-attach-btn.has-attachment { background: #fef3c7; color: #92400e; border-color: #f59e0b; }
  .note-mini-table {
    margin: 6px 0 6px 28px; border-collapse: collapse; font-size: 12px;
    background: white; border: 1px solid #e5e7eb; border-radius: 4px;
  }
  .note-mini-table th, .note-mini-table td {
    border: 1px solid #e5e7eb; padding: 4px 8px; text-align: left;
  }
  .note-mini-table th { background: #f9fafb; font-weight: 600; color: #374151; }
  .note-mini-photo {
    margin: 6px 0 6px 28px; max-width: 200px; max-height: 120px;
    border-radius: 6px; border: 1px solid #e5e7eb; cursor: pointer; display: block;
  }
  .note-attach-controls { display: flex; gap: 4px; margin: 4px 0 0 28px; }
  .note-attach-controls button {
    font-size: 11px; padding: 2px 6px; margin: 0;
    background: #f3f4f6; color: #4b5563; border: 1px solid #d1d5db; border-radius: 4px;
  }
  .note-attach-controls button.danger { background: #fee2e2; color: #991b1b; border-color: #fca5a5; }

  /* 표 편집 모달 */
  .modal-overlay {
    display: none; position: fixed; inset: 0; background: rgba(0,0,0,0.5);
    z-index: 1000; justify-content: center; align-items: center;
  }
  .modal-overlay.show { display: flex; }
  .modal-box {
    background: white; border-radius: 12px; padding: 24px;
    max-width: 90vw; max-height: 90vh; overflow: auto; min-width: 520px;
  }
  .modal-box h3 { margin: 0 0 12px; color: #1E3A5F; font-size: 18px; }
  .modal-box .modal-input {
    width: 100%; padding: 8px 10px; border: 1px solid #d1d5db; border-radius: 6px;
    font-size: 13px; margin-bottom: 10px; box-sizing: border-box;
  }
  .modal-box textarea.modal-input { min-height: 80px; font-family: monospace; }
  .modal-box .modal-table {
    border-collapse: collapse; width: 100%; margin-bottom: 12px; font-size: 13px;
  }
  .modal-box .modal-table th, .modal-box .modal-table td {
    border: 1px solid #d1d5db; padding: 4px;
  }
  .modal-box .modal-table input {
    border: none; width: 100%; padding: 4px; font-size: 13px; box-sizing: border-box;
  }
  .modal-box .modal-table input:focus { outline: 2px solid #1E3A5F; }
  .modal-actions { display: flex; gap: 8px; justify-content: flex-end; margin-top: 12px; }
  .modal-actions button { margin: 0; padding: 8px 16px; font-size: 13px; }
  .modal-actions button.cancel { background: #6b7280; }
  .modal-actions button.danger { background: #dc2626; }

  /* 사진 확대 모달 */
  .photo-overlay {
    display: none; position: fixed; inset: 0; background: rgba(0,0,0,0.9);
    z-index: 1100; justify-content: center; align-items: center; cursor: pointer;
  }
  .photo-overlay.show { display: flex; }
  .photo-overlay img { max-width: 95vw; max-height: 95vh; }


  /* 매핑 관리 대시보드 */
  .map-chip {
    padding: 6px 12px; border: 1px solid #d1d5db; border-radius: 16px;
    background: white; color: #4b5563; font-size: 13px; cursor: pointer;
    transition: all 0.15s;
  }
  .map-chip:hover { background: #f3f4f6; border-color: #1E3A5F; color: #1E3A5F; }
  .map-chip.active { background: #1E3A5F; color: white; border-color: #1E3A5F; font-weight: 600; }
  .map-chip .count {
    display: inline-block; margin-left: 6px; padding: 1px 6px; border-radius: 8px;
    font-size: 11px; background: #e5e7eb; color: #374151;
  }
  .map-chip.active .count { background: rgba(255,255,255,0.25); color: white; }

  .map-pcard {
    background: white; border: 1px solid #e5e7eb; border-radius: 8px;
    padding: 12px 14px; transition: all 0.15s; cursor: default;
  }
  .map-pcard:hover { border-color: #1E3A5F; box-shadow: 0 2px 8px rgba(30,58,95,0.08); }
  .map-pcard .ptitle { font-size: 14px; font-weight: 600; color: #1f2937; margin-bottom: 8px; }
  .map-pcard .pmeta { display: flex; gap: 10px; font-size: 12px; color: #6b7280; }
  .map-pcard .pmeta .badge-ppt { color: #1E3A5F; }
  .map-pcard .pmeta .badge-note { color: #10b981; }
  .map-pcard .pmeta .badge-empty { color: #d1d5db; }

  .placeholder {
    background: white;
    border: 2px dashed #d1d5db;
    border-radius: 12px;
    padding: 60px 40px;
    text-align: center;
    color: #6b7280;
  }
  .placeholder .icon { font-size: 48px; margin-bottom: 16px; }
  .placeholder h3 { font-size: 18px; margin-bottom: 8px; color: #374151; }
  .placeholder p { font-size: 14px; }

  .footer-info {
    text-align: center;
    color: #9ca3af;
    font-size: 12px;
    padding: 24px;
  }

  /* Dropzone */
  .dropzone {
    border: 2px dashed #cbd5e1;
    border-radius: 10px;
    padding: 28px 16px;
    text-align: center;
    cursor: pointer;
    transition: all 0.15s ease;
    background: #f8fafc;
  }
  .dropzone:hover { border-color: #1E3A5F; background: #f1f5f9; }
  .dropzone.dragover { border-color: #1E3A5F; background: #e0e7ff; transform: scale(1.01); }
  .dz-icon { font-size: 36px; margin-bottom: 8px; }
  .dz-title { font-size: 14px; font-weight: 600; color: #1E3A5F; margin-bottom: 4px; }
  .dz-sub { font-size: 12px; color: #64748b; }
  .file-row {
    display: flex; justify-content: space-between; align-items: center;
    padding: 8px 12px; margin-top: 6px;
    background: #f8fafc; border-radius: 6px; border: 1px solid #e2e8f0;
  }
  .file-row-name { font-size: 13px; color: #1f2937; flex: 1; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; margin-right: 8px; }
  .file-row-size { font-size: 11px; color: #6b7280; margin-right: 8px; }
  .file-row-remove { background: transparent; color: #b91c1c; border: none; cursor: pointer; font-size: 16px; padding: 0 4px; }
  .file-row-remove:hover { color: #7f1d1d; }

  
  /* Progress bar */
  .pf-row { padding: 10px 12px; border-bottom: 1px solid #f1f5f9; font-size: 13px; }
  .pf-head { display: flex; align-items: center; gap: 10px; margin-bottom: 6px; }
  .pf-status { width: 22px; text-align: center; font-size: 15px; }
  .pf-name { font-weight: 600; color: #1f2937; flex: 1; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; }
  .pf-pct { color: #1E3A5F; font-weight: 700; min-width: 44px; text-align: right; }
  .pf-barwrap { background: #e5e7eb; border-radius: 999px; height: 8px; overflow: hidden; }
  .pf-bar { background: linear-gradient(90deg, #1E3A5F, #3b82f6); height: 100%; width: 0%; transition: width 0.18s ease; }
  .pf-bar.done { background: linear-gradient(90deg, #10b981, #059669); }
  .pf-bar.fail { background: linear-gradient(90deg, #ef4444, #b91c1c); }
  .pf-bar.warn { background: linear-gradient(90deg, #f59e0b, #d97706); }
  .pf-detail { margin-top: 6px; color: #6b7280; font-size: 12px; }
  .pf-total {
    display: flex; align-items: center; gap: 10px;
    background: #f8fafc; border: 1px solid #e2e8f0; border-radius: 8px;
    padding: 10px 14px; margin-bottom: 12px; font-size: 13px;
  }
  .pf-total .pf-barwrap { flex: 1; height: 10px; }
  .pf-total-label { font-weight: 700; color: #1E3A5F; }
  .pf-total-pct { font-weight: 700; color: #1E3A5F; min-width: 44px; text-align: right; }

  </style>
</head>
<body>
<header style="display:flex;justify-content:space-between;align-items:flex-start;gap:16px;flex-wrap:wrap;">
  <div>
    <h1 style="margin:0;">📊 사업부 보고 관리자</h1>
    <div class="meta" id="header-meta">로딩 중...</div>
  </div>
  <div id="adminSessionBar" style="display:flex;align-items:center;gap:8px;background:#1E3A5F;color:#fff;padding:8px 12px;border-radius:8px;font-size:13px;white-space:nowrap;">
    <span>🔐 관리자 세션 유지 중</span>
    <span style="background:#0f2342;padding:3px 10px;border-radius:999px;">남은 시간 <strong id="adminSessionRemain">확인 중...</strong></span>
    <button type="button" id="adminExtendBtn" style="background:#3b82f6;color:#fff;border:none;border-radius:6px;padding:6px 10px;cursor:pointer;">8시간 연장</button>
    <button type="button" id="adminLogoutBtn" style="background:#ef4444;color:#fff;border:none;border-radius:6px;padding:6px 10px;cursor:pointer;">로그아웃</button>
  </div>
</header>

<nav class="tabs">
  <button class="tab-btn active" data-tab="ppt">📤 PPT 업로드</button>
  <button class="tab-btn" data-tab="notes">📝 주간 보고 업로드</button>
  <button class="tab-btn" data-tab="mapping">⚙️ 매핑 관리</button>
</nav>

<main>

  <!-- ============================== 탭 1: PPT 업로드 ============================== -->
  <section class="tab-content active" id="tab-ppt">
    <div class="card">
      <h2>📤 새 PPT 업로드</h2>
      <form id="pptUploadForm">
        <label for="pptPassword">업로드 비밀번호</label>
        <input type="password" id="pptPassword" placeholder="관리자 비밀번호" required />

        <label>PPT 파일 (여러 개 가능 · 드래그 & 드롭 지원)</label>
        <div id="pptDropzone" class="dropzone">
          <div class="dz-icon">📂</div>
          <div class="dz-title">여기로 PPT 파일을 드래그하거나 클릭해서 선택</div>
          <div class="dz-sub">.pptx / .ppt · 여러 번 추가 가능 · 중복 자동 제거</div>
          <input type="file" id="pptFiles" accept=".pptx,.ppt" multiple hidden />
        </div>
        <div id="pptFileListHeader" style="display:none; margin-top:12px; display:flex; justify-content:space-between; align-items:center;">
          <div id="pptFileCount" style="font-size:13px; color:#374151; font-weight:600;"></div>
          <button type="button" id="pptClearAll" style="background:#fee2e2; color:#b91c1c; border:none; padding:6px 12px; border-radius:6px; cursor:pointer; font-size:12px;">전체 비우기</button>
        </div>
        <div id="pptFileList" style="margin-top: 8px; font-size: 13px;"></div>

        <label for="reportFamily">리포트 패밀리 (선택)</label>
        <input type="text" id="reportFamily" placeholder="default" value="default" />

        <!-- 중복 허용은 업로드 시작 버튼 옆으로 이동 -->

        <details style="margin-top: 20px;">
          <summary style="cursor: pointer; font-weight: 600; color: #1E3A5F; padding: 10px 0;">
            ▼ 사전 매핑 (선택) — 모든 파일에 동일 적용
          </summary>
          <div style="padding: 12px 0; border-top: 1px solid #eee; margin-top: 8px;">
            <p style="font-size: 12px; color: #6b7280; margin-bottom: 12px;">
              자동 분류가 실패한 카드만 여기서 정한 값으로 폴백됩니다. 자동 분류 잘 되는 카드는 영향 없음.
            </p>

            <label for="pptDivision">사업부</label>
            <select id="pptDivision">
              <option value="">— 자동 분류 —</option>
            </select>

            <label for="pptProject">프로젝트</label>
            <select id="pptProject" disabled>
              <option value="">— 사업부 먼저 선택 —</option>
            </select>
          </div>
        </details>

        <div style="margin-top:16px;">
          <div style="display:flex; align-items:center; gap:12px;">
            <button type="submit" id="pptUploadBtn">업로드 시작</button>
            <label style="display:inline-flex; align-items:center; gap:6px; font-size:13px; padding:6px 10px; background:#f3f4f6; border-radius:6px; cursor:pointer; white-space:nowrap;">
              <input type="checkbox" id="allowDuplicateUpload" />
              <span>중복 허용</span>
            </label>
          </div>
          <div style="font-size:11px; color:#9ca3af; margin-top:6px;">
            기본 차단 · 같은 파일 다시 올릴 때만 체크
          </div>
        </div>
      </form>
    </div>

    <div class="card" id="uploadHistoryCard">
      <h2 style="display:flex; align-items:center; justify-content:space-between;">
        <span>🗂️ 업로드 내역</span>
        <span style="font-size:12px; font-weight:normal;">
          <label style="margin-right:8px;"><input type="checkbox" id="historySelectAll" /> 전체 선택</label>
          <button type="button" id="historyDeleteBtn" style="background:#fee2e2; color:#b91c1c; border:none; padding:6px 12px; border-radius:6px; cursor:pointer; font-size:12px;">선택 삭제</button>
        </span>
      </h2>
      <p style="color:#6b7280;font-size:13px;margin-top:-4px;">PPT 파일 단위로 삭제하면 관련 카드가 모두 제거됩니다.</p>
      <div id="uploadHistoryArea" style="overflow-x:auto;">
        <div style="color:#999;padding:16px;">불러오는 중...</div>
      </div>
    </div>
    <div class="card" id="pptProgressCard" style="display: none;">
      <h2>📊 업로드 진행</h2>
      <div id="pptProgressSummary" style="font-size: 14px; color: #374151; margin-bottom: 12px;"></div>
      <div id="pptProgressList"></div>
    </div>

  </section>

  <!-- ============================== 탭 2: 이미지 업로드 (기존 폼 그대로) ============================== -->

  <!-- ============================== 탭 3: 매핑 관리 ============================== -->
  <section class="tab-content" id="tab-notes">
    <div class="card">
      <h2>📝 주간 보고 업로드</h2>
      <p style="color:#6b7280;font-size:13px;margin-top:-4px;">담당자가 받은 주간 보고 텍스트를 그대로 붙여넣고, AI 정리 후 저장하면 사장님 앱에 표시됩니다.</p>

      <div style="display:flex;gap:12px;align-items:flex-end;margin:16px 0;flex-wrap:wrap;">
        <div>
          <label style="display:block;font-size:13px;color:#374151;margin-bottom:4px;">사업부</label>
          <select id="noteDivision" style="padding:8px 12px;border:1px solid #d1d5db;border-radius:6px;font-size:14px;min-width:180px;">
            <option value="">로딩 중...</option>
          </select>
        </div>
        <div>
          <label style="display:block;font-size:13px;color:#374151;margin-bottom:4px;">보고일자</label>
          <input type="date" id="noteReportDate" style="padding:8px 12px;border:1px solid #d1d5db;border-radius:6px;font-size:14px;" />
        </div>
        <button type="button" id="noteLoadBtn" style="padding:8px 14px;background:#e5e7eb;color:#374151;border:none;border-radius:6px;cursor:pointer;font-size:13px;">📥 기존 노트 불러오기</button>
      </div>

      <label style="display:block;font-size:13px;color:#374151;margin-bottom:4px;">원본 텍스트</label>
      <textarea id="noteRawText" placeholder="<페러데이 4T>&#10;1. LPM&#10; - 5호기, 6호기-테스트완료(4/3)&#10;..." style="width:100%;height:280px;padding:12px;border:1px solid #d1d5db;border-radius:6px;font-size:13px;font-family:'Menlo','Monaco',monospace;line-height:1.5;"></textarea>

      <div style="display:flex;gap:8px;margin-top:12px;flex-wrap:wrap;">
        <button type="button" id="noteAiParseBtn" style="padding:10px 18px;background:#3b82f6;color:white;border:none;border-radius:6px;cursor:pointer;font-size:14px;font-weight:500;">🤖 AI로 정리하기</button>

        <span id="noteStatus" style="align-self:center;font-size:13px;color:#6b7280;"></span>
      </div>
    </div>

    <div class="card" id="notePreviewCard" style="display:none;">
      <h2>✨ 미리보기 (앱에서 이렇게 보입니다)</h2>
      <div id="notePreviewArea" style="background:#f5f7fb;padding:16px;border-radius:8px;"></div>
    </div>
  </section>

  <section class="tab-content" id="tab-mapping">
    <div class="card">
      <h2>📊 사업부 ↔ 프로젝트 통합 대시보드</h2>
      <p style="color:#6b7280;font-size:13px;margin-bottom:14px;">사업부 칩을 클릭하면 해당 사업부의 프로젝트 리스트가 표시됩니다.</p>
      <div id="mappingDivisionChips" style="display:flex;flex-wrap:wrap;gap:8px;margin-bottom:14px;"></div>
      <div id="mappingDivisionMeta" style="font-size:13px;color:#6b7280;margin-bottom:10px;"></div>
      <div id="mappingProjectGrid" style="display:grid;grid-template-columns:repeat(auto-fill,minmax(220px,1fr));gap:12px;"></div>
      <div id="mappingEmpty" style="color:#9ca3af;padding:24px;text-align:center;display:none;">프로젝트가 등록되지 않은 사업부입니다.</div>
    </div>
  </section>


</main>

<div class="footer-info">
  v2 — 5단계 진행 중 ·
  <a href="/docs" target="_blank" style="color: #6b7280;">API Docs</a>
</div>

<script>
  // 헤더 시간 표시
  function updateHeaderTime() {
    const now = new Date();
    const wd = ['일','월','화','수','목','금','토'][now.getDay()];
    const ampm = now.getHours() < 12 ? '오전' : '오후';
    const h12 = now.getHours() % 12 || 12;
    const m = String(now.getMinutes()).padStart(2, '0');
    document.getElementById('header-meta').textContent =
      `마지막 새로고침: ${now.getMonth()+1}/${now.getDate()} (${wd}) ${ampm} ${h12}:${m}`;
  }
  updateHeaderTime();

  // 탭 전환
  // ─── PPT → 주간 보고 초안 생성 ───
  let _draftDocId = null;

  function openDraftModal(docId, fileName) {
    _draftDocId = docId;
    document.getElementById('draftModalFile').textContent = fileName || docId;
    document.getElementById('draftModalStatus').textContent = '';
    const wrap = document.getElementById('draftDivChips');
    wrap.innerHTML = '<span style="color:#999;">사업부 로딩 중...</span>';
    document.getElementById('draftDivModal').classList.add('show');

    fetch('/divisions', { credentials: 'same-origin' })
      .then(r => r.json())
      .then(data => {
        const divs = data.divisions || [];
        wrap.innerHTML = '';
        divs.forEach(d => {
          const chip = document.createElement('button');
          chip.type = 'button';
          chip.textContent = (d.badge_short_label || d.label) + ' (' + (d.projects ? d.projects.length : 0) + ')';
          chip.className = 'map-chip';
          chip.style.cursor = 'pointer';
          chip.onclick = () => runPptxDraft(d.id, d.label);
          wrap.appendChild(chip);
        });
      })
      .catch(e => {
        wrap.innerHTML = '<span style="color:#c00;">사업부 로드 실패: ' + e.message + '</span>';
      });
  }

  function closeDraftModal() {
    document.getElementById('draftDivModal').classList.remove('show');
    _draftDocId = null;
  }

  async function runPptxDraft(divisionId, divisionLabel) {
    const status = document.getElementById('draftModalStatus');
    status.textContent = '⏳ PPT 분석 중... (10-30초)';
    status.style.color = '#6b7280';
    try {
      const r = await fetch('/admin/notes/from_pptx', {
        method: 'POST',
        headers: { 'Content-Type': 'application/json' },
        credentials: 'same-origin',
        body: JSON.stringify({ doc_id: _draftDocId, division_id: divisionId }),
      });
      if (!r.ok) {
        let detail = r.status;
        try { const err = await r.json(); detail = err.detail || detail; } catch(_){}
        status.textContent = '❌ 실패: ' + detail;
        status.style.color = '#c00';
        return;
      }
      const data = await r.json();
      const cards = data.cards || [];
      if (cards.length === 0) {
        status.textContent = '⚠️ 추출된 카드가 없습니다';
        status.style.color = '#c00';
        return;
      }
      // 1) 주간 보고 탭으로 전환
      document.querySelectorAll('.tab-btn').forEach(b => b.classList.remove('active'));
      document.querySelectorAll('.tab-content').forEach(c => c.classList.remove('active'));
      const noteBtn = document.querySelector('.tab-btn[data-tab="notes"]');
      if (noteBtn) noteBtn.classList.add('active');
      document.getElementById('tab-notes').classList.add('active');

      // 2) 사업부 셀렉트 + 날짜 자동 설정
      const sel = document.getElementById('noteDivision');
      if (sel) sel.value = divisionId;
      const dateEl = document.getElementById('noteReportDate');
      if (dateEl && !dateEl.value) {
        const today = new Date();
        dateEl.value = today.toISOString().slice(0, 10);
      }

      // 3) 미리보기 자동 주입
      _noteParsedCards = cards;
      renderNotePreview(_noteParsedCards);
      const noteStatus = document.getElementById('noteStatus');
      if (noteStatus) {
        noteStatus.textContent = '✅ PPT에서 ' + cards.length + '개 카드 자동 생성 — 확인 후 저장하세요';
        noteStatus.style.color = '#16a34a';
      }
      closeDraftModal();
    } catch (e) {
      status.textContent = '❌ 오류: ' + e.message;
      status.style.color = '#c00';
    }
  }

  // ESC로 초안 모달 닫기
  document.addEventListener('keydown', (e) => {
    if (e.key === 'Escape') {
      const m = document.getElementById('draftDivModal');
      if (m && m.classList.contains('show')) closeDraftModal();
    }
  });

  document.querySelectorAll('.tab-btn').forEach(btn => {
    btn.addEventListener('click', () => {
      const target = btn.dataset.tab;
      document.querySelectorAll('.tab-btn').forEach(b => b.classList.remove('active'));
      document.querySelectorAll('.tab-content').forEach(c => c.classList.remove('active'));
      btn.classList.add('active');
      document.getElementById('tab-' + target).classList.add('active');
    });
  });

  // ============================== 기존 이미지 업로드 로직 (그대로 유지) ==============================

  // 사업부 목록 로드
  async function loadProjects() {
    try {
      const res = await fetch('/projects');
      const data = await res.json();
      const sel = document.getElementById('projectKey');
      if (!sel) return;
      data.projects.forEach(p => {
        const opt = document.createElement('option');
        opt.value = p.key;
        opt.textContent = p.label || p.key;
        sel.appendChild(opt);
      });
    } catch (e) {
      console.error('사업부 목록 로드 실패:', e);
    }
  }

  // 등록된 이미지 목록 로드
  

  
  // ============================================================
  // 업로드 내역 (PPT doc_id 단위 삭제)
  // ============================================================
  async function loadUploadHistory() {
    const area = document.getElementById('uploadHistoryArea');
    if (!area) return;
    try {
      const r = await fetch('/admin/uploads/history', { credentials: 'same-origin' });
      if (!r.ok) {
        area.innerHTML = '<div style="color:#c00;padding:16px;">권한이 없습니다. 다시 로그인하세요.</div>';
        return;
      }
      const j = await r.json();
      const items = j.items || [];
      if (items.length === 0) {
        area.innerHTML = '<div style="color:#999;padding:16px;">업로드된 PPT가 없습니다.</div>';
        return;
      }
      let html = '<table style="width:100%;border-collapse:collapse;font-size:13px;min-width:600px;">';
      html += '<thead><tr style="background:#f3f4f6;text-align:left;">';
      html += '<th style="padding:10px 8px;width:32px;text-align:center;"></th>';
      html += '<th style="padding:10px 8px;">업로드 일시</th>';
      html += '<th style="padding:10px 8px;">파일명 / doc_id</th>';
      html += '<th style="padding:10px 8px;text-align:center;">카드</th>';
      html += '<th style="padding:10px 8px;text-align:center;">슬라이드</th>';
      html += '<th style="padding:10px 8px;text-align:center;">삭제</th>';
      html += '</tr></thead><tbody>';
      items.forEach(it => {
        let tsLabel = '-';
        if (it.upload_timestamp) {
          try {
            tsLabel = new Date(it.upload_timestamp).toLocaleString('ko-KR', {
              year: '2-digit', month: '2-digit', day: '2-digit',
              hour: '2-digit', minute: '2-digit'
            });
          } catch (_) { tsLabel = it.upload_timestamp; }
        }
        const safeFname = (it.file_name || '(이름 없음)').replace(/</g, '&lt;');
        html += '<tr style="border-bottom:1px solid #e5e7eb;">';
        html += `<td style="padding:10px 8px;text-align:center;"><input type="checkbox" class="history-row-check" data-doc-id="${it.doc_id}"></td>`;
        html += `<td style="padding:10px 8px;white-space:nowrap;">${tsLabel}</td>`;
        html += `<td style="padding:10px 8px;">${safeFname}<br><small style="color:#999;">${it.doc_id}</small></td>`;
        html += `<td style="padding:10px 8px;text-align:center;">${it.product_count}</td>`;
        html += `<td style="padding:10px 8px;text-align:center;">${it.slide_count}</td>`;
        html += `<td style="padding:10px 8px;text-align:center;white-space:nowrap;">`
              + `<button data-doc-id="${it.doc_id}" data-file-name="${(it.file_name || '').replace(/"/g, '&quot;').replace(/</g, '&lt;')}" class="draft-btn" style="background:#3b82f6;color:#fff;border:none;padding:6px 10px;border-radius:6px;cursor:pointer;font-size:12px;margin-right:6px;">📝 초안</button>`
              + `<button onclick="deleteDoc('${it.doc_id}', ${it.product_count})" style="background:#ef4444;color:#fff;border:none;padding:6px 10px;border-radius:6px;cursor:pointer;font-size:12px;">삭제</button>`
              + `</td>`;
        html += '</tr>';
      });
      html += '</tbody></table>';
      area.innerHTML = html;
      // 초안 버튼 이벤트 바인딩
      area.querySelectorAll('.draft-btn').forEach(btn => {
        btn.addEventListener('click', () => {
          openDraftModal(btn.dataset.docId, btn.dataset.fileName);
        });
      });
    } catch (e) {
      area.innerHTML = `<div style="color:#c00;padding:16px;">불러오기 실패: ${e.message}</div>`;
    }
  }

  async function deleteDoc(docId, productCount) {
    if (!confirm(`정말 삭제하시겠습니까?\n\n관련 카드 ${productCount}개가 모두 제거되고\nPPT 파일과 슬라이드 이미지도 삭제됩니다.`)) return;
    try {
      const r = await fetch(`/admin/doc/${docId}`, {
        method: 'DELETE',
        credentials: 'same-origin'
      });
      if (!r.ok) {
        let detail = r.status;
        try { const err = await r.json(); detail = err.detail || detail; } catch (_) {}
        alert('삭제 실패: ' + detail);
        return;
      }
      const result = await r.json();
      alert(`✅ 삭제 완료\n남은 PPT: ${result.remaining}개`);
      loadUploadHistory();
    } catch (e) {
      alert('삭제 오류: ' + e.message);
    }
  }



  // ============================================================
  // 관리자 세션 UI
  // ============================================================
  let adminSessionExpiresAt = null;

  function fmtRemain(sec) {
    if (sec < 0) sec = 0;
    const h = Math.floor(sec / 3600);
    const m = Math.floor((sec % 3600) / 60);
    const s = sec % 60;
    if (h > 0) return `${h}시간 ${m}분`;
    if (m > 0) return `${m}분 ${s}초`;
    return `${s}초`;
  }

  async function refreshAdminSessionInfo() {
    try {
      const r = await fetch('/admin/session', { credentials: 'same-origin' });
      if (!r.ok) {
        location.href = '/admin/login?next=/admin/upload';
        return;
      }
      const j = await r.json();
      adminSessionExpiresAt = j.expires_at;

      // password input 숨김 + required 해제
      document.querySelectorAll('input[type="password"]').forEach(el => {
        el.required = false;
        const label = el.id ? document.querySelector(`label[for="${el.id}"]`) : null;
        if (label) label.style.display = 'none';
        el.style.display = 'none';
      });
    } catch (e) {
      console.error('session load fail', e);
    }
  }

  function tickAdminSessionRemain() {
    const el = document.getElementById('adminSessionRemain');
    if (!el || !adminSessionExpiresAt) return;
    const remain = adminSessionExpiresAt - Math.floor(Date.now() / 1000);
    el.textContent = fmtRemain(remain);
    if (remain <= 0) {
      location.href = '/admin/login?next=/admin/upload';
    }
  }

  document.getElementById('adminExtendBtn')?.addEventListener('click', async () => {
    const r = await fetch('/admin/extend', { method: 'POST', credentials: 'same-origin' });
    if (r.ok) {
      const j = await r.json();
      adminSessionExpiresAt = j.expires_at;
      tickAdminSessionRemain();
      alert('세션이 8시간 연장되었습니다.');
    } else {
      location.href = '/admin/login?next=/admin/upload';
    }
  });

  document.getElementById('adminLogoutBtn')?.addEventListener('click', async () => {
    await fetch('/admin/logout', { method: 'POST', credentials: 'same-origin' });
    location.href = '/admin/login?next=/admin/upload';
  });

  refreshAdminSessionInfo().then(() => {
    tickAdminSessionRemain();
    setInterval(tickAdminSessionRemain, 1000);
    // 5초마다 백엔드 핑(서버측 만료 즉시 감지)
    setInterval(refreshAdminSessionInfo, 5000);
  });

  // 탭/창 복귀 시 즉시 재검증
  document.addEventListener('visibilitychange', () => {
    if (document.visibilityState === 'visible') {
      refreshAdminSessionInfo();
    }
  });
  window.addEventListener('focus', refreshAdminSessionInfo);


  // 초기 로드
  loadProjects();
  loadUploadHistory();

  // ============================================================
  // 5-2b-2: PPT 업로드 폼 (다중 파일 + 사전 매핑 + 진행 표시)
  // ============================================================

  // 사업부 dropdown 채우기
  async function loadDivisionsForPpt() {
    try {
      const res = await fetch('/admin/config/divisions');
      const data = await res.json();
      const sel = document.getElementById('pptDivision');
      data.divisions.forEach(d => {
        const opt = document.createElement('option');
        opt.value = d.id;
        opt.textContent = d.label;
        sel.appendChild(opt);
      });
    } catch (e) {
      console.error('사업부 로드 실패:', e);
    }
  }

  // 사업부 선택 시 프로젝트 dropdown 채우기
  document.getElementById('pptDivision').addEventListener('change', async (e) => {
    const divId = e.target.value;
    const projSel = document.getElementById('pptProject');
    projSel.innerHTML = '<option value="">— 자동 분류 —</option>';

    if (!divId) {
      projSel.disabled = true;
      projSel.innerHTML = '<option value="">— 사업부 먼저 선택 —</option>';
      return;
    }

    try {
      const res = await fetch('/admin/config/projects?division_id=' + encodeURIComponent(divId));
      const data = await res.json();
      data.projects.forEach(p => {
        const opt = document.createElement('option');
        opt.value = p.id;
        opt.textContent = p.label + (p.group ? ' (' + p.group + ')' : '');
        projSel.appendChild(opt);
      });
      projSel.disabled = false;
    } catch (e) {
      console.error('프로젝트 로드 실패:', e);
    }
  });

  // 파일 선택 시 목록 표시

  // 한 파일 업로드 (XHR 기반 진행률 % 반환)
  function uploadSinglePpt(file, password, reportFamily, divisionId, projectId, allowDuplicate, itemEl, onProgress) {
    return new Promise((resolve) => {
      const fd = new FormData();
      fd.append('file', file);
      fd.append('password', password);
      fd.append('report_family', reportFamily || 'default');
      if (divisionId) fd.append('division_id', divisionId);
      if (projectId) fd.append('project_id', projectId);
      if (allowDuplicate) fd.append('allow_duplicate', 'true');

      // 상태: 시작
      const statusEl = itemEl.querySelector('.pf-status');
      const pctEl = itemEl.querySelector('.pf-pct');
      const barEl = itemEl.querySelector('.pf-bar');
      const detailEl = itemEl.querySelector('.pf-detail');
      statusEl.textContent = '⬆️';
      detailEl.textContent = '업로드 중...';

      const xhr = new XMLHttpRequest();
      xhr.open('POST', '/upload', true);

      // 업로드 바이트 진행률 (0~90% 구간 사용, 서버 처리는 90~100% 마무리에서 채움)
      xhr.upload.onprogress = (ev) => {
        if (!ev.lengthComputable) return;
        const pct = Math.round((ev.loaded / ev.total) * 90);
        barEl.style.width = pct + '%';
        pctEl.textContent = pct + '%';
        if (typeof onProgress === 'function') onProgress(pct);
      };

      // 업로드 완료 → 서버 처리 단계로 전환
      xhr.upload.onload = () => {
        statusEl.textContent = '⚙️';
        detailEl.textContent = '서버 처리 중...';
        barEl.style.width = '95%';
        pctEl.textContent = '95%';
        if (typeof onProgress === 'function') onProgress(95);
      };

      xhr.onload = () => {
        let data = {};
        try { data = JSON.parse(xhr.responseText || '{}'); } catch (e) {}
        if (xhr.status === 409 && data && data.code === 'duplicate_upload') {
          barEl.style.width = '100%';
          barEl.classList.add('warn');
          pctEl.textContent = '중복';
          statusEl.textContent = '⚠️';
          detailEl.innerHTML = '<span style="color:#b45309;">' + (data.detail || '이미 업로드된 파일') + '</span>';
          if (typeof onProgress === 'function') onProgress(100);
          resolve({ ok: false, duplicate: true, file: file.name, file_key: fileKey(file), error: data.detail || 'duplicate upload' });
          return;
        }
        if (xhr.status >= 200 && xhr.status < 300 && data && data.ok !== false) {
          barEl.style.width = '100%';
          barEl.classList.add('done');
          pctEl.textContent = '100%';
          statusEl.textContent = '✅';
          detailEl.textContent = '카드 ' + (data.product_count || 0) + '개 추출 (슬라이드 ' + (data.slide_count || 0) + '장)';
          if (typeof onProgress === 'function') onProgress(100);
          resolve({ ok: true, file: file.name, file_key: fileKey(file), data });
        } else {
          const msg = (data && (data.detail || data.message)) || ('HTTP ' + xhr.status);
          barEl.classList.add('fail');
          statusEl.textContent = '❌';
          detailEl.innerHTML = '<span style="color:#dc2626;">' + msg + '</span>';
          if (typeof onProgress === 'function') onProgress(100);
          resolve({ ok: false, file: file.name, file_key: fileKey(file), error: msg });
        }
      };

      xhr.onerror = () => {
        barEl.classList.add('fail');
        statusEl.textContent = '❌';
        detailEl.innerHTML = '<span style="color:#dc2626;">네트워크 오류</span>';
        if (typeof onProgress === 'function') onProgress(100);
        resolve({ ok: false, file: file.name, file_key: fileKey(file), error: 'network error' });
      };

      xhr.send(fd);
    });
  }

  // 업로드 폼 submit

  // ============================================================
  // 5-2b-2: 누적 + 드래그&드롭 + 개별 삭제 + 전체 비우기
  // ============================================================
  let pptSelectedFiles = []; // 누적된 File 객체 배열

  function formatBytes(b) {
    if (b < 1024) return b + ' B';
    if (b < 1024*1024) return (b/1024).toFixed(1) + ' KB';
    return (b/1024/1024).toFixed(1) + ' MB';
  }

  function fileKey(f) { return f.name + '_' + f.size; }

  function renderPptFileList() {
    const listEl = document.getElementById('pptFileList');
    const headerEl = document.getElementById('pptFileListHeader');
    const countEl = document.getElementById('pptFileCount');
    const btn = document.getElementById('pptUploadBtn');

    listEl.innerHTML = '';
    if (pptSelectedFiles.length === 0) {
      headerEl.style.display = 'none';
      btn.textContent = '업로드 시작';
      btn.disabled = true;
      btn.style.opacity = '0.5';
      return;
    }

    headerEl.style.display = 'flex';
    countEl.textContent = pptSelectedFiles.length + '개 파일 선택됨';
    btn.textContent = pptSelectedFiles.length + '개 업로드';
    btn.disabled = false;
    btn.style.opacity = '1';

    pptSelectedFiles.forEach((f, idx) => {
      const row = document.createElement('div');
      row.className = 'file-row';
      row.innerHTML =
        '<div class="file-row-name">📄 ' + f.name + '</div>' +
        '<div class="file-row-size">' + formatBytes(f.size) + '</div>' +
        '<button type="button" class="file-row-remove" data-idx="' + idx + '" title="제거">✕</button>';
      listEl.appendChild(row);
    });

    listEl.querySelectorAll('.file-row-remove').forEach(b => {
      b.addEventListener('click', (e) => {
        const i = parseInt(e.currentTarget.getAttribute('data-idx'));
        pptSelectedFiles.splice(i, 1);
        renderPptFileList();
      });
    });
  }

  function addPptFiles(fileList) {
    const allowed = ['.pptx', '.ppt'];
    const existing = new Set(pptSelectedFiles.map(fileKey));
    let added = 0, skipped = 0;
    Array.from(fileList).forEach(f => {
      const lower = f.name.toLowerCase();
      if (!allowed.some(ext => lower.endsWith(ext))) { skipped++; return; }
      if (existing.has(fileKey(f))) { skipped++; return; }
      pptSelectedFiles.push(f);
      existing.add(fileKey(f));
      added++;
    });
    renderPptFileList();
    if (skipped > 0) console.log('skip(중복/비PPT):', skipped);
  }

  // 파일 선택 input change
  document.getElementById('pptFiles').addEventListener('change', (e) => {
    addPptFiles(e.target.files);
    e.target.value = ''; // 같은 파일 재선택 가능하도록
  });

  // 드롭존 클릭 → input 열기
  const dz = document.getElementById('pptDropzone');
  dz.addEventListener('click', () => document.getElementById('pptFiles').click());

  // 드래그 & 드롭
  ['dragenter','dragover'].forEach(ev => dz.addEventListener(ev, (e) => {
    e.preventDefault(); e.stopPropagation();
    dz.classList.add('dragover');
  }));
  ['dragleave','drop'].forEach(ev => dz.addEventListener(ev, (e) => {
    e.preventDefault(); e.stopPropagation();
    dz.classList.remove('dragover');
  }));
  dz.addEventListener('drop', (e) => {
    if (e.dataTransfer && e.dataTransfer.files) {
      addPptFiles(e.dataTransfer.files);
    }
  });

  // 페이지 전체 드래그 기본동작 방지
  ['dragover','drop'].forEach(ev => {
    window.addEventListener(ev, (e) => { e.preventDefault(); }, false);
  });

  // 전체 비우기
  document.getElementById('pptClearAll').addEventListener('click', () => {
    pptSelectedFiles = [];
    renderPptFileList();
  });

  // 초기 렌더 (버튼 비활성화)
  renderPptFileList();

  document.getElementById('pptUploadForm').addEventListener('submit', async (e) => {
    e.preventDefault();

    const files = pptSelectedFiles.slice();
    if (files.length === 0) {
      alert('PPT 파일을 1개 이상 선택해주세요.');
      return;
    }

    const password = document.getElementById('pptPassword').value;
    const reportFamily = document.getElementById('reportFamily').value || 'default';
    const divisionId = document.getElementById('pptDivision').value;
    const projectId = document.getElementById('pptProject').value;
    const allowDuplicate = document.getElementById('allowDuplicateUpload').checked;

    // 진행 영역 표시 + 초기화
    const progressCard = document.getElementById('pptProgressCard');
    const progressList = document.getElementById('pptProgressList');
    const progressSummary = document.getElementById('pptProgressSummary');
    // 진행카드를 업로드내역 위로 이동 (업로드 중에만 보임)
    const historyCard = document.getElementById('uploadHistoryCard');
    if (historyCard && progressCard && historyCard.parentNode) {
      historyCard.parentNode.insertBefore(progressCard, historyCard);
    }
    progressCard.style.display = 'block';
    progressList.innerHTML = '';
    progressSummary.textContent = `총 ${files.length}개 파일 처리 시작...`;

    // 전체 진행률 영역
    progressSummary.innerHTML = `
      <div class="pf-total">
        <span class="pf-total-label">전체 진행률</span>
        <span id="pfTotalText">0 / ${files.length}</span>
        <div class="pf-barwrap"><div class="pf-bar" id="pfTotalBar"></div></div>
        <span class="pf-total-pct" id="pfTotalPct">0%</span>
      </div>
    `;
    const totalBar = document.getElementById('pfTotalBar');
    const totalPct = document.getElementById('pfTotalPct');
    const totalText = document.getElementById('pfTotalText');

    // 각 파일별 진행 항목 생성
    const fileProgress = new Array(files.length).fill(0);
    const itemEls = files.map((f, idx) => {
      const div = document.createElement('div');
      div.className = 'pf-row';
      div.innerHTML = `
        <div class="pf-head">
          <span class="pf-status">⏳</span>
          <span class="pf-name">📄 ${f.name}</span>
          <span class="pf-pct">0%</span>
        </div>
        <div class="pf-barwrap"><div class="pf-bar"></div></div>
        <div class="pf-detail">대기 중</div>
      `;
      progressList.appendChild(div);
      return div;
    });

    function updateTotal() {
      const sum = fileProgress.reduce((a,b)=>a+b, 0);
      const pct = Math.round(sum / files.length);
      totalBar.style.width = pct + '%';
      totalPct.textContent = pct + '%';
      const doneCnt = fileProgress.filter(p => p >= 100).length;
      totalText.textContent = doneCnt + ' / ' + files.length;
    }

    // 업로드 버튼 비활성화
    const btn = document.getElementById('pptUploadBtn');
    btn.disabled = true;
    btn.textContent = '업로드 중...';

    // 순차 업로드 (파일별 % → 전체 진행률 합산)
    const results = [];
    for (let i = 0; i < files.length; i++) {
      const result = await uploadSinglePpt(
        files[i], password, reportFamily, divisionId, projectId, allowDuplicate, itemEls[i],
        (pct) => { fileProgress[i] = pct; updateTotal(); }
      );
      fileProgress[i] = 100;
      updateTotal();
      results.push(result);
    }

    // 완료 후 통계 (전체 진행률 바 아래에 결과 줄 추가)
    const ok = results.filter(r => r.ok).length;
    const fail = results.length - ok;
    const totalCards = results.filter(r => r.ok).reduce((s, r) => s + (r.data.product_count || 0), 0);
    const doneLine = document.createElement('div');
    doneLine.style.cssText = 'margin-top:4px; font-size:13px;';
    doneLine.innerHTML = '<strong>완료:</strong> ' + ok + '/' + files.length + '개 성공' +
      (fail > 0 ? ', <span style="color:#dc2626;">' + fail + '개 실패</span>' : '') +
      ' · 총 카드 <strong>' + totalCards + '개</strong> 추출';
    progressSummary.appendChild(doneLine);

    // 성공한 파일만 선택 목록에서 제거
    const successKeys = new Set(results.filter(r => r.ok).map(r => r.file_key));
    pptSelectedFiles = pptSelectedFiles.filter(f => !successKeys.has(fileKey(f)));
    renderPptFileList();

    // 버튼 복구
    btn.disabled = false;
    btn.textContent = '업로드 시작';

    // 이미지 탭의 사업부 목록 갱신 (새로 업로드된 PPT 의 프로젝트가 잡힐 수 있게)
    try { loadProjects(); } catch (e) {}

    // 업로드 내역 즉시 갱신
    try { loadUploadHistory(); } catch (e) {}

    // 5초 뒤 진행카드 숨김
    setTimeout(() => {
      const pc = document.getElementById('pptProgressCard');
      if (pc) pc.style.display = 'none';
    }, 5000);
  });

  // 업로드 내역 다중 선택 핸들러 (전체 선택 / 선택 삭제)
  document.addEventListener('change', (e) => {
    if (e.target && e.target.id === 'historySelectAll') {
      const checked = e.target.checked;
      document.querySelectorAll('.history-row-check').forEach(cb => { cb.checked = checked; });
    }
  });
  document.addEventListener('click', async (e) => {
    if (e.target && e.target.id === 'historyDeleteBtn') {
      const ids = Array.from(document.querySelectorAll('.history-row-check:checked'))
        .map(cb => cb.dataset.docId).filter(Boolean);
      if (ids.length === 0) { alert('삭제할 항목을 선택하세요.'); return; }
      if (!confirm(ids.length + '개의 업로드를 삭제할까요? PPT 파일과 관련 카드가 모두 제거됩니다.')) return;
      e.target.disabled = true;
      try {
        const r = await fetch('/admin/docs/delete-batch', {
          method: 'POST',
          credentials: 'include',
          headers: {'Content-Type': 'application/json'},
          body: JSON.stringify({doc_ids: ids})
        });
        if (r.status === 401) { location.href = '/admin/login'; return; }
        const d = await r.json();
        alert((d.removed_count || 0) + '개 삭제 완료');
        try { loadUploadHistory(); } catch(_){}
      } catch (err) {
        alert('삭제 실패: ' + err);
      } finally {
        e.target.disabled = false;
      }
    }
  });




  // PPT 탭 진입 시 dropdown 초기 채우기
  loadDivisionsForPpt();

  // ============================== 노트 관리 (사업부 보고) ==============================
  let _noteParsedCards = null;

  async function loadNoteDivisions() {
    try {
      const r = await fetch('/admin/config/divisions', { credentials: 'same-origin' });
      if (!r.ok) throw new Error('fetch failed');
      const data = await r.json();
      const sel = document.getElementById('noteDivision');
      if (!sel) return;
      sel.innerHTML = '';
      (data.divisions || []).forEach(d => {
        const opt = document.createElement('option');
        opt.value = d.id;
        opt.textContent = d.label;
        sel.appendChild(opt);
      });
    } catch (e) {
      console.error('사업부 로드 실패', e);
    }
  }

  // 보고일자 기본값 = 오늘
  function initNoteDate() {
    const el = document.getElementById('noteReportDate');
    if (!el) return;
    const today = new Date();
    const y = today.getFullYear();
    const m = String(today.getMonth() + 1).padStart(2, '0');
    const d = String(today.getDate()).padStart(2, '0');
    el.value = `${y}-${m}-${d}`;
  }

  

function showUploadDonePopup(cardCount, onClose) {
  // 기존 팝업 있으면 제거
  const oldBg = document.getElementById('uploadDoneBg');
  if (oldBg) oldBg.remove();

  const bg = document.createElement('div');
  bg.id = 'uploadDoneBg';
  bg.style.cssText = 'position:fixed;top:0;left:0;width:100%;height:100%;background:rgba(0,0,0,0.45);z-index:9999;display:flex;align-items:center;justify-content:center;';

  const box = document.createElement('div');
  box.style.cssText = 'background:#fff;border-radius:14px;padding:28px 32px;min-width:300px;max-width:90%;box-shadow:0 10px 30px rgba(0,0,0,0.2);text-align:center;font-family:inherit;';

  box.innerHTML = '<div style="font-size:42px;margin-bottom:10px;">✅</div>'
    + '<div style="font-size:18px;font-weight:600;color:#111;margin-bottom:6px;">업로드 완료되었습니다</div>'
    + '<div style="font-size:14px;color:#6b7280;margin-bottom:20px;">' + cardCount + '개 카드가 저장되었습니다</div>'
    + '<button id="uploadDoneOk" style="background:#10b981;color:#fff;border:none;border-radius:8px;padding:10px 28px;font-size:15px;font-weight:600;cursor:pointer;">확인</button>';

  bg.appendChild(box);
  document.body.appendChild(bg);

  const ok = document.getElementById('uploadDoneOk');
  function close() {
    bg.remove();
    if (typeof onClose === 'function') onClose();
  }
  ok.addEventListener('click', close);
  bg.addEventListener('click', function(e){ if (e.target === bg) close(); });
}

function renderNotePreview(cards) {
  const area = document.getElementById('notePreviewArea');
  const card = document.getElementById('notePreviewCard');
  if (!area || !card) return;

  if (!cards || cards.length === 0) {
    card.style.display = 'none';
    area.innerHTML = '';
    return;
  }

  const esc = (s) => String(s == null ? '' : s)
    .replace(/&/g, '&amp;')
    .replace(/</g, '&lt;')
    .replace(/>/g, '&gt;')
    .replace(/"/g, '&quot;');

  const renderAttachBtns = (ci, si, ii) => {
    // disabled by request — 미리보기에서 표/사진 버튼 숨김
    return '';
  };

  const renderMiniTable = (table) => {
    if (!table) return '';
    const headers = Array.isArray(table.headers) ? table.headers : [];
    const rows = Array.isArray(table.rows) ? table.rows : [];
    const previewRows = rows.slice(0, 3);
    let html = `
      <div style="margin-top:8px;padding:10px;border:1px solid #d1d5db;border-radius:8px;background:#fafafa;">
        <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:6px;">
          <strong style="font-size:13px;">📊 ${esc(table.title || '표')}</strong>
          <span style="font-size:11px;color:#6b7280;">탭/클릭으로 확대·수정</span>
        </div>
        <div style="overflow-x:auto;">
          <table style="border-collapse:collapse;width:100%;font-size:12px;background:#fff;">
    `;
    if (headers.length) {
      html += '<thead><tr>';
      headers.forEach(h => {
        html += `<th style="border:1px solid #e5e7eb;padding:6px 8px;background:#f3f4f6;text-align:left;">${esc(h)}</th>`;
      });
      html += '</tr></thead>';
    }
    html += '<tbody>';
    previewRows.forEach(r => {
      html += '<tr>';
      (Array.isArray(r) ? r : []).forEach(c => {
        html += `<td style="border:1px solid #e5e7eb;padding:6px 8px;">${esc(c)}</td>`;
      });
      html += '</tr>';
    });
    html += '</tbody></table></div>';
    if (rows.length > 3) {
      html += `<div style="margin-top:4px;font-size:11px;color:#6b7280;">… 외 ${rows.length - 3}행</div>`;
    }
    html += '</div>';
    return html;
  };

  const renderDueChip = (dueRaw) => {
    if (!dueRaw) return '';
    const m = /^(\d{4})-(\d{2})-(\d{2})/.exec(dueRaw);
    if (!m) return '';
    const due = new Date(parseInt(m[1]), parseInt(m[2])-1, parseInt(m[3]));
    const today = new Date();
    today.setHours(0,0,0,0);
    const diff = Math.round((due - today) / (1000*60*60*24));
    let bg = '#dbeafe', fg = '#1e40af', label;
    if (diff < 0) { bg = '#fee2e2'; fg = '#b91c1c'; label = `D+${-diff} 지남`; }
    else if (diff === 0) { bg = '#fef3c7'; fg = '#92400e'; label = 'D-0 오늘'; }
    else if (diff <= 3) { bg = '#fed7aa'; fg = '#9a3412'; label = `D-${diff}`; }
    else { label = `D-${diff}`; }
    const dateLabel = `${m[2]}/${m[3]}`;
    return `<span style="display:inline-block;margin-left:6px;padding:1px 7px;border-radius:10px;background:${bg};color:${fg};font-size:11px;font-weight:600;">${label} (${dateLabel})</span>`;
  };

  const renderPhoto = (photoRef) => {
    if (!photoRef) return '';
    const url = '/note_photos/' + photoRef;
    const isExcel = /\/xls_/.test(photoRef);
    if (isExcel) {
      return `
        <div style="margin:8px 0 14px 0;">
          <div style="background:#fff;border:1px solid #e5e7eb;border-radius:10px;padding:8px;overflow:auto;">
            <img src="${esc(url)}"
                 onclick="showPhotoOverlay('${esc(url)}')"
                 style="display:block;width:100%;height:auto;border-radius:6px;cursor:pointer;" />
          </div>
        </div>
      `;
    }
    return `
      <div style="margin-top:8px;">
        <img src="${esc(url)}"
             onclick="showPhotoOverlay('${esc(url)}')"
             style="max-width:240px;max-height:150px;border-radius:8px;border:1px solid #d1d5db;cursor:pointer;object-fit:cover;" />
      </div>
    `;
  };

  // 일반 항목 한 줄 렌더 (group 묶음 안/밖 공용)
  const renderOneItem = (it, ci, si, ii) => {
    if (typeof it === 'string') {
      it = { type: 'bullet', text: it };
    }
    if (!it || typeof it !== 'object') return '';
    const type = it.type || 'bullet';
    const text = esc(it.text || '');
    const hasTable = !!it.table_ref;
    const hasPhoto = !!it.photo_ref;
    const tableData = it.table_data || null;
    const photoRef = it.photo_ref || '';

    if (type === 'table') {
      if (tableData && tableData.preview_image_data) {
        const _safeTitle = String(it.text || tableData.title || '엑셀').replace(/[&<>"']/g, function(ch){
          return {'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'}[ch];
        });
        return `
          <div style="margin:8px 0 14px 0;">
            <div style="font-size:13px;font-weight:700;color:#334155;margin-bottom:6px;">📊 ${_safeTitle}</div>
            <div style="background:#fff;border:1px solid #e5e7eb;border-radius:10px;padding:8px;overflow:auto;">
              <img src="${tableData.preview_image_data}" style="display:block;max-width:100%;height:auto;border-radius:6px;" />
            </div>
          </div>
        `;
      }
      return ``;
    }
    if (type === 'photo') {
      return `<div style="margin:8px 0;">${renderPhoto(photoRef)}</div>`;
    }

    // 색상 매핑 (item.color > 타입 기본색)
    const _colorMap = { red:'#dc2626', blue:'#1e88e5', orange:'#ef6c00' };
    const itemColor = _colorMap[String(it.color || '').toLowerCase()] || '';

    let prefix = '•';
    let textStyle = 'font-size:14px;color:#111827;';
    let leftPad = '';
    const rawText = String(it.text || '').trim();
    const startsWithStar = rawText.startsWith('※');

    if (type === 'highlight') {
      textStyle = 'font-size:14px;color:#dc2626;font-weight:700;';
    } else if (type === 'sub') {
      // 화살표 중복 방지: 본문이 이미 → / ↳ / => 로 시작하면 prefix 비움
      prefix = /^(?:→|↳|=>)\s*/.test(rawText) ? '' : '→';
      textStyle = 'font-size:13px;color:#374151;';
      leftPad = 'padding-left:18px;';
    }

    // ※ 줄은 굵게 + 기본 검정 (color 지정 있으면 그게 우선)
    if (startsWithStar) {
      textStyle = 'font-size:14px;color:#111827;font-weight:700;';
      prefix = '';
    }

    // item.color 가 있으면 최종적으로 그 색으로 덮어쓰기 (font-weight 유지)
    if (itemColor) {
      textStyle = textStyle.replace(/color:[^;]+;?/, '') + 'color:' + itemColor + ';';
    }

    if (hasTable && !tableData && typeof loadTableDataForItem === 'function') {
      setTimeout(() => {
        try { loadTableDataForItem(it, () => renderNotePreview(_noteParsedCards)); } catch (_) {}
      }, 0);
    }

    const dueChip = renderDueChip(it.due_date || '');
    return `
      <div style="display:flex;align-items:flex-start;justify-content:space-between;gap:8px;margin:6px 0;${leftPad}">
        <div style="flex:1;min-width:0;">
          <div style="${textStyle}">${prefix} ${text} ${dueChip}</div>
          ${hasTable ? renderMiniTable(tableData) : ''}
          ${hasPhoto ? renderPhoto(photoRef) : ''}
        </div>
        ${renderAttachBtns(ci, si, ii)}
      </div>
    `;
  };

  // group_id 같은 항목들을 연속 구간으로 묶어 노란 박스로 렌더
  const renderItemsWithGroups = (items, ci, si) => {
    let html = '';
    let i = 0;
    while (i < items.length) {
      const raw = items[i];
      const it = (typeof raw === 'string') ? { type: 'bullet', text: raw } : (raw || {});
      const gid = it.group_id || '';

      if (gid) {
        // 같은 group_id 가진 연속 구간 모으기
        const groupItems = [];
        let j = i;
        while (j < items.length) {
          const r = items[j];
          const x = (typeof r === 'string') ? { type: 'bullet', text: r } : (r || {});
          if (x.group_id !== gid) break;
          groupItems.push({ raw: r, idx: j });
          j++;
        }

        // 노란 박스 시작
        html += `
          <div style="margin:8px 0;padding:10px 12px;border-left:4px solid #d97706;background:#fffbeb;border-radius:6px;">
        `;

        // group_note 와 일반 항목 분리
        const normals = groupItems.filter(g => {
          const o = (typeof g.raw === 'string') ? { type: 'bullet' } : (g.raw || {});
          return o.type !== 'group_note';
        });
        const notes = groupItems.filter(g => {
          const o = (typeof g.raw === 'string') ? { type: 'bullet' } : (g.raw || {});
          return o.type === 'group_note';
        });

        // 일반 항목 먼저
        normals.forEach(g => {
          html += renderOneItem(g.raw, ci, si, g.idx);
        });

        // group_note (있으면 박스 안 아래에)
        notes.forEach(g => {
          const o = (typeof g.raw === 'string') ? { type: 'group_note', text: g.raw } : g.raw;
          const text = esc(o.text || '');
          html += `
            <div style="margin:6px 0 0 4px;padding-top:6px;border-top:1px dashed #fcd34d;">
              <div style="font-size:13px;color:#92400e;font-style:italic;">↪ ${text} ${renderDueChip(o.due_date || '')}</div>
            </div>
          `;
        });

        html += `</div>`;
        i = j;
      } else {
        // group 아닌 일반 항목
        if (it.type === 'group_note') {
          // group_id 없는 group_note 도 노란 박스로 단독 표시
          const text = esc(it.text || '');
          html += `
            <div style="margin:8px 0 8px 18px;padding:8px 10px;border-left:4px solid #d97706;background:#fffbeb;border-radius:6px;">
              <div style="font-size:13px;color:#92400e;font-style:italic;">↪ ${text}</div>
            </div>
          `;
        } else {
          html += renderOneItem(raw, ci, si, i);
        }
        i++;
      }
    }
    return html;
  };

  let html = '';
  cards.forEach((c, ci) => {
    const title = esc(c.title || '');
    const sections = Array.isArray(c.sections) ? c.sections : [];
    html += `
      <div style="background:#fff;border:1px solid #e5e7eb;border-radius:12px;padding:16px;margin-bottom:16px;">
        <div style="font-weight:800;font-size:18px;margin-bottom:12px;">${title}</div>
    `;

    sections.forEach((s, si) => {
      const stitle = esc(s.title || '');
      const items = Array.isArray(s.items) ? s.items : [];
      html += `
        <div style="margin-bottom:12px;">
          <div style="background:#f5f7fb;border-radius:8px;padding:8px 10px;font-weight:700;margin-bottom:8px;">${stitle}</div>
      `;

      html += renderItemsWithGroups(items, ci, si);

      // 섹션 레벨 표 (구버전 호환)
      if (s.table_ref && s.table_data) {
        html += renderMiniTable(s.table_data);
      } else if (s.table_ref && typeof loadTableDataForItem === 'function') {
        setTimeout(() => {
          try { loadTableDataForItem(s, () => renderNotePreview(_noteParsedCards)); } catch (_) {}
        }, 0);
      }

      html += `</div>`;
    });

    html += `</div>`;
  });

  // 미리보기 맨 아래 저장 버튼
  html += `
    <div style="display:flex;justify-content:flex-end;margin-top:16px;">
      <button type="button" id="noteSaveBtn"
        style="padding:12px 24px;background:#10b981;color:white;border:none;border-radius:8px;cursor:pointer;font-size:15px;font-weight:600;">
        💾 저장
      </button>
    </div>
  `;
  area.innerHTML = html;
  // 새로 생긴 저장 버튼에 이벤트 바인딩
  const newSaveBtn = area.querySelector('#noteSaveBtn');
  if (newSaveBtn) newSaveBtn.addEventListener('click', noteSave);
  card.style.display = 'block';
}
  window.renderNotePreview = renderNotePreview;




  async function loadTableDataForItem(it, cb) {
    if (!it.table_ref || it._table_loading) return;
    it._table_loading = true;
    try {
      const r = await fetch('/notes/table/' + it.table_ref);
      if (r.ok) {
        const d = await r.json();
        it._table_data = d.table;
      }
    } catch (e) {}
    it._table_loading = false;
    if (cb) cb();
  }

  function showPhotoOverlay(url) {
    document.getElementById('photoOverlayImg').src = url;
    document.getElementById('photoOverlay').classList.add('show');
  }

  async function noteAiParse() {
    const text = document.getElementById('noteRawText').value.trim();
    const status = document.getElementById('noteStatus');
    const saveBtn = document.getElementById('noteSaveBtn');
    if (!text) {
      status.textContent = '⚠️ 텍스트를 입력하세요';
      status.style.color = '#dc2626';
      return;
    }
    status.textContent = '🤖 AI 정리 중...';
    status.style.color = '#6b7280';
    if (saveBtn) { saveBtn.disabled = true; saveBtn.style.opacity = '0.5'; }
    const _divId = document.getElementById('noteDivision').value;
    try {
      const r = await fetch('/admin/notes/ai_parse', {
        method: 'POST',
        credentials: 'same-origin',
        headers: { 'Content-Type': 'application/json' },
        body: JSON.stringify({ text, division_id: _divId }),
      });
      if (!r.ok) throw new Error(`HTTP ${r.status}`);
      const data = await r.json();
      _noteParsedCards = data.cards || [];
      window._noteParsedCards = _noteParsedCards;
      renderNotePreview(_noteParsedCards);
      status.textContent = `✅ ${_noteParsedCards.length}개 카드 정리 완료`;
      status.style.color = '#10b981';
      if (saveBtn) { saveBtn.disabled = false; saveBtn.style.opacity = '1'; }
    } catch (e) {
      status.textContent = '❌ AI 정리 실패: ' + e.message;
      status.style.color = '#dc2626';
      _noteParsedCards = null;
    }
  }
  window.noteAiParse = noteAiParse;


  async function noteSave() {
    if (!_noteParsedCards) {
      alert('먼저 AI로 정리해주세요');
      return;
    }
    const divisionId = document.getElementById('noteDivision').value;
    const reportDate = document.getElementById('noteReportDate').value;
    if (!divisionId) { alert('사업부를 선택하세요'); return; }
    if (!reportDate) { alert('보고일자를 선택하세요'); return; }

    const status = document.getElementById('noteStatus');
    status.textContent = '💾 저장 중...';
    status.style.color = '#6b7280';
    try {
      const r = await fetch('/admin/notes', {
        method: 'POST',
        credentials: 'same-origin',
        headers: { 'Content-Type': 'application/json' },
        body: JSON.stringify({
          division_id: divisionId,
          report_date: reportDate,
          raw_text: document.getElementById('noteRawText').value || '',
          cards: _noteParsedCards,
        }),
      });
      if (!r.ok) {
        const err = await r.text();
        throw new Error(err || `HTTP ${r.status}`);
      }
      const data = await r.json();
      status.textContent = `✅ 업로드 완료 (${data.card_count}개 카드)`;
      status.style.color = '#10b981';

      // 미리보기 숨기고 원본 텍스트만 남기기
      _noteParsedCards = null;
      const previewCard = document.getElementById('notePreviewCard');
      const previewArea = document.getElementById('notePreviewArea');
      if (previewArea) previewArea.innerHTML = '';
      if (previewCard) previewCard.style.display = 'none';

      const saveBtn = document.getElementById('noteSaveBtn');
      if (saveBtn) { saveBtn.disabled = true; saveBtn.style.opacity = '0.5'; }
    } catch (e) {
      status.textContent = '❌ 저장 실패: ' + e.message;
      status.style.color = '#dc2626';
    }
  }

  // cards JSON → 텍스트 변환 (편집 가능 포맷)
  function _cardsToText(cards) {
    if (!cards || cards.length === 0) return '';
    const out = [];
    cards.forEach((c, ci) => {
      out.push(`<${c.title || '프로젝트'}>`);
      const sections = Array.isArray(c.sections) ? c.sections : [];
      sections.forEach((s, si) => {
        if (s.title) out.push(`[${s.title}]`);
        const items = Array.isArray(s.items) ? s.items : [];
        // group_id 모아서 } 표기 만들기
        let i = 0;
        while (i < items.length) {
          const it = items[i] || {};
          const gid = it.group_id || '';
          if (gid) {
            // 같은 group_id 연속 묶기
            const group = [];
            let j = i;
            while (j < items.length && (items[j] || {}).group_id === gid) {
              group.push(items[j]);
              j++;
            }
            const normals = group.filter(g => (g.type || '') !== 'group_note');
            const notes = group.filter(g => (g.type || '') === 'group_note');
            normals.forEach((g, idx) => {
              const isLast = idx === normals.length - 1;
              const t = (g.type || 'bullet');
              let prefix = '- ';
              if (t === 'highlight') prefix = '*';
              else if (t === 'sub') prefix = '  → ';
              const tail = (isLast && notes.length > 0) ? ' } ' + notes.map(n => n.text || '').join(', ') : '';
              out.push(prefix + (g.text || '') + tail);
            });
            // 단일 group_note 만 있고 normal 없으면 그냥 출력
            if (normals.length === 0 && notes.length > 0) {
              notes.forEach(n => out.push('} ' + (n.text || '')));
            }
            i = j;
          } else {
            const t = (it.type || 'bullet');
            if (t === 'highlight') out.push('*' + (it.text || ''));
            else if (t === 'sub') out.push('  → ' + (it.text || ''));
            else if (t === 'group_note') out.push('} ' + (it.text || ''));
            else if (t === 'table') out.push('[표]');
            else if (t === 'photo') out.push('[사진]');
            else out.push('- ' + (it.text || ''));
            i++;
          }
        }
      });
      if (ci < cards.length - 1) out.push('');
    });
    return out.join(String.fromCharCode(10));
  }

  async function noteLoadExisting() {
    const divisionId = document.getElementById('noteDivision').value;
    if (!divisionId) { alert('사업부를 먼저 선택하세요'); return; }
    const status = document.getElementById('noteStatus');
    status.textContent = '📥 불러오는 중...';
    status.style.color = '#6b7280';
    try {
      const r = await fetch(`/notes?division_id=${encodeURIComponent(divisionId)}`);
      if (!r.ok) throw new Error(`HTTP ${r.status}`);
      const data = await r.json();
      const cards = data.cards || [];
      if (data.report_date) {
        document.getElementById('noteReportDate').value = data.report_date;
      }
      if (cards.length > 0) {
        const txt = _cardsToText(cards);
        const ta = document.getElementById('noteRawText');
        if (ta) ta.value = txt;
        // 미리보기 닫기 (사용자가 텍스트 먼저 보고 수정하도록)
        _noteParsedCards = null;
        const card = document.getElementById('notePreviewCard');
        if (card) card.style.display = 'none';
        status.textContent = `✅ ${cards.length}개 카드를 텍스트로 불러옴 — 수정 후 "🤖 AI 정리" 를 눌러주세요`;
        status.style.color = '#10b981';
      } else {
        status.textContent = '저장된 노트가 없습니다';
        status.style.color = '#6b7280';
      }
    } catch (e) {
      status.textContent = '❌ 불러오기 실패: ' + e.message;
      status.style.color = '#dc2626';
    }
  }

  // ─── 표/사진 첨부 ───
  let _editingTable = null; // {cardIdx, secIdx, itemIdx, headers, rows, title}

  function _getItem(ci, si, ii) {
    if (!_noteParsedCards) return null;
    const c = _noteParsedCards[ci]; if (!c) return null;
    const s = (c.sections || [])[si]; if (!s) return null;
    return (s.items || [])[ii] || null;
  }

  async function openTableModal(ci, si, ii) {
    const it = _getItem(ci, si, ii); if (!it) return;
    let tableData = { title: '', headers: ['구분', '값'], rows: [['', '']] };
    let isExisting = false;
    if (it.table_ref) {
      try {
        const r = await fetch('/notes/table/' + it.table_ref);
        if (r.ok) {
          const d = await r.json();
          tableData = d.table || tableData;
          isExisting = true;
        }
      } catch (e) {}
    }
    _editingTable = {
      ci, si, ii,
      title: tableData.title || '',
      headers: (tableData.headers && tableData.headers.length) ? [...tableData.headers] : ['구분', '값'],
      rows: (tableData.rows && tableData.rows.length) ? tableData.rows.map(r => [...r]) : [['', '']],
    };
    document.getElementById('tableModalTitle').value = _editingTable.title;
    document.getElementById('tablePasteArea').value = '';
    document.getElementById('tableModalDeleteBtn').style.display = isExisting ? 'inline-block' : 'none';
    renderTableEditor();
    document.getElementById('tableModal').classList.add('show');
  }

  function closeTableModal() {
    document.getElementById('tableModal').classList.remove('show');
    _editingTable = null;
  }

  function renderTableEditor() {
    if (!_editingTable) return;
    const area = document.getElementById('tableEditorArea');
    area.innerHTML = '';
    const tbl = document.createElement('table');
    tbl.className = 'modal-table';

    // 헤더 행
    const thead = document.createElement('thead');
    const trh = document.createElement('tr');
    _editingTable.headers.forEach((h, ci) => {
      const th = document.createElement('th');
      const inp = document.createElement('input');
      inp.value = h;
      inp.style.fontWeight = '600';
      inp.style.background = '#f9fafb';
      inp.oninput = (e) => { _editingTable.headers[ci] = e.target.value; };
      th.appendChild(inp);
      trh.appendChild(th);
    });
    thead.appendChild(trh);
    tbl.appendChild(thead);

    // 데이터 행
    const tbody = document.createElement('tbody');
    _editingTable.rows.forEach((row, ri) => {
      const tr = document.createElement('tr');
      _editingTable.headers.forEach((_, ci) => {
        const td = document.createElement('td');
        const inp = document.createElement('input');
        inp.value = row[ci] !== undefined ? row[ci] : '';
        inp.oninput = (e) => {
          while (row.length < _editingTable.headers.length) row.push('');
          row[ci] = e.target.value;
        };
        td.appendChild(inp);
        tr.appendChild(td);
      });
      tbody.appendChild(tr);
    });
    tbl.appendChild(tbody);
    area.appendChild(tbl);
  }

  function addTableRow() {
    if (!_editingTable) return;
    _editingTable.rows.push(_editingTable.headers.map(() => ''));
    renderTableEditor();
  }

  function removeTableRow() {
    if (!_editingTable) return;
    if (_editingTable.rows.length <= 1) { alert('최소 1행은 있어야 합니다'); return; }
    _editingTable.rows.pop();
    renderTableEditor();
  }

  function addTableCol() {
    if (!_editingTable) return;
    _editingTable.headers.push('항목' + (_editingTable.headers.length + 1));
    _editingTable.rows.forEach(r => r.push(''));
    renderTableEditor();
  }

  function removeTableCol() {
    if (!_editingTable) return;
    if (_editingTable.headers.length <= 1) { alert('최소 1열은 있어야 합니다'); return; }
    _editingTable.headers.pop();
    _editingTable.rows.forEach(r => r.pop());
    renderTableEditor();
  }

  // 클립보드 붙여넣기 자동 파싱 (Tab/쉼표)
  document.addEventListener('DOMContentLoaded', () => {
    const pa = document.getElementById('tablePasteArea');
    if (pa) {
      pa.addEventListener('paste', (e) => {
        setTimeout(() => {
          const text = pa.value;
          if (!text.trim() || !_editingTable) return;
          const _LF = String.fromCharCode(10);
          const _CR = String.fromCharCode(13);
          const _TAB = String.fromCharCode(9);
          const _norm = String(text).split(_CR + _LF).join(_LF).split(_CR).join(_LF);
          const lines = _norm.split(_LF).filter(function(l){ return l.trim(); });
          if (lines.length === 0) return;
          const sep = lines[0].indexOf(_TAB) >= 0 ? _TAB : ',';
          const parsed = lines.map(l => l.split(sep).map(c => c.trim()));
          if (parsed.length >= 1) {
            _editingTable.headers = parsed[0];
            _editingTable.rows = parsed.slice(1).length ? parsed.slice(1) : [parsed[0].map(() => '')];
            renderTableEditor();
            pa.value = '';
          }
        }, 50);
      });
    }
  });

  async function saveTableFromModal() {
    if (!_editingTable) return;
    const divisionId = document.getElementById('noteDivision').value;
    if (!divisionId) { alert('사업부를 선택하세요'); return; }
    const it = _getItem(_editingTable.ci, _editingTable.si, _editingTable.ii); if (!it) return;

    const tableData = {
      title: document.getElementById('tableModalTitle').value.trim(),
      headers: _editingTable.headers,
      rows: _editingTable.rows,
    };

    try {
      const body = { division_id: divisionId, table: tableData };
      if (it.table_ref) body.table_ref = it.table_ref;
      const r = await fetch('/admin/notes/table', {
        method: 'POST',
        credentials: 'same-origin',
        headers: { 'Content-Type': 'application/json' },
        body: JSON.stringify(body),
      });
      if (!r.ok) throw new Error('HTTP ' + r.status);
      const d = await r.json();
      it.table_ref = d.table_ref;
      it._table_data = tableData;
      closeTableModal();
      renderNotePreview(_noteParsedCards);
      // 카드 변경 자동 반영을 위해 노트 저장 권장 메시지
      const st = document.getElementById('noteStatus');
      st.textContent = '✅ 표 저장됨 — 노트 변경사항이 있다면 "💾 저장" 버튼을 눌러주세요';
      st.style.color = '#10b981';
    } catch (e) {
      alert('표 저장 실패: ' + e.message);
    }
  }

  async function deleteTableFromModal() {
    if (!_editingTable) return;
    const it = _getItem(_editingTable.ci, _editingTable.si, _editingTable.ii); if (!it) return;
    if (!it.table_ref) { closeTableModal(); return; }
    if (!confirm('이 표를 삭제하시겠습니까?')) return;
    try {
      const r = await fetch('/admin/notes/table/' + it.table_ref, {
        method: 'DELETE', credentials: 'same-origin',
      });
      // 404여도 클라이언트는 정리
      delete it._table_data;
      delete it.table_ref;
      closeTableModal();
      renderNotePreview(_noteParsedCards);
      const st = document.getElementById('noteStatus');
      st.textContent = '🗑️ 표 삭제됨 — "💾 저장" 버튼을 눌러주세요';
      st.style.color = '#f59e0b';
    } catch (e) {
      alert('표 삭제 실패: ' + e.message);
    }
  }

  // 사진 업로드
  let _pendingPhotoTarget = null; // {ci, si, ii}
  function openPhotoUpload(ci, si, ii) {
    const it = _getItem(ci, si, ii); if (!it) return;
    const divisionId = document.getElementById('noteDivision').value;
    if (!divisionId) { alert('사업부를 선택하세요'); return; }
    _pendingPhotoTarget = { ci, si, ii };
    document.getElementById('notePhotoFileInput').click();
  }

  document.addEventListener('DOMContentLoaded', () => {
    const inp = document.getElementById('notePhotoFileInput');
    if (inp) {
      inp.addEventListener('change', async (e) => {
        const file = e.target.files[0];
        if (!file || !_pendingPhotoTarget) return;
        const { ci, si, ii } = _pendingPhotoTarget;
        const it = _getItem(ci, si, ii);
        const divisionId = document.getElementById('noteDivision').value;
        if (!it || !divisionId) { inp.value = ''; return; }
        const fd = new FormData();
        fd.append('division_id', divisionId);
        fd.append('file', file);
        try {
          const r = await fetch('/admin/notes/photo', {
            method: 'POST', credentials: 'same-origin', body: fd,
          });
          if (!r.ok) throw new Error('HTTP ' + r.status);
          const d = await r.json();
          it.photo_ref = d.photo_ref;
          renderNotePreview(_noteParsedCards);
          const st = document.getElementById('noteStatus');
          st.textContent = '✅ 사진 업로드됨 — "💾 저장" 버튼을 눌러주세요';
          st.style.color = '#10b981';
        } catch (err) {
          alert('사진 업로드 실패: ' + err.message);
        } finally {
          inp.value = '';
          _pendingPhotoTarget = null;
        }
      });
    }
  });

  async function removeNotePhoto(ci, si, ii) {
    const it = _getItem(ci, si, ii); if (!it || !it.photo_ref) return;
    if (!confirm('이 사진을 삭제하시겠습니까?')) return;
    try {
      await fetch('/admin/notes/photo/' + it.photo_ref, {
        method: 'DELETE', credentials: 'same-origin',
      });
    } catch (e) {}
    delete it.photo_ref;
    renderNotePreview(_noteParsedCards);
    const st = document.getElementById('noteStatus');
    st.textContent = '🗑️ 사진 삭제됨 — "💾 저장" 버튼을 눌러주세요';
    st.style.color = '#f59e0b';
  }

  // ─── 매핑 관리 대시보드 ───
  let _mapDivisions = [];
  let _mapActiveDivId = null;
  let _mapReports = [];
  let _mapNotes = {};

  async function initMappingTab() {
    try {
      const divResp = await fetch('/divisions');
      if (!divResp.ok) throw new Error('divisions HTTP ' + divResp.status);
      const divData = await divResp.json();
      _mapDivisions = divData.divisions || [];

      try {
        const r = await fetch('/admin/uploads/history', { credentials: 'same-origin' });
        if (r.ok) {
          const d = await r.json();
          _mapReports = d.items || d.reports || [];
        }
      } catch (e) { _mapReports = []; }

      _mapNotes = {};
      for (const d of _mapDivisions) {
        try {
          const r = await fetch('/notes?division_id=' + encodeURIComponent(d.id));
          if (r.ok) {
            const data = await r.json();
            _mapNotes[d.id] = data;
          }
        } catch (e) {}
      }

      renderMappingChips();
      if (_mapDivisions.length > 0) {
        selectMappingDivision(_mapDivisions[0].id);
      }
    } catch (e) {
      const grid = document.getElementById('mappingProjectGrid');
      if (grid) grid.innerHTML = '<div style="color:#dc2626;padding:16px;">로드 실패: ' + e.message + '</div>';
    }
  }

  function renderMappingChips() {
    const wrap = document.getElementById('mappingDivisionChips');
    if (!wrap) return;
    wrap.innerHTML = '';
    _mapDivisions.forEach(d => {
      const chip = document.createElement('button');
      chip.className = 'map-chip' + (d.id === _mapActiveDivId ? ' active' : '');
      chip.dataset.divId = d.id;
      chip.innerHTML = d.label + '<span class="count">' + (d.projects || []).length + '</span>';
      chip.onclick = () => selectMappingDivision(d.id);
      wrap.appendChild(chip);
    });
  }

  function _countProjectMatches(divisionId, project) {
    let pptCardCount = 0;
    const aliases = (project.aliases || []).map(s => s.toLowerCase());
    const labelLc = (project.label || '').toLowerCase();
    const projectId = project.id;

    for (const rep of _mapReports) {
      const products = rep.products || [];
      for (const p of products) {
        const pkey = p.project_key || p.project || '';
        if (pkey === projectId) { pptCardCount += 1; continue; }
        const text = ((p.name || '') + ' ' + (p.headline || '')).toLowerCase();
        if (labelLc && text.includes(labelLc)) { pptCardCount += 1; continue; }
        for (const al of aliases) {
          if (al && text.includes(al)) { pptCardCount += 1; break; }
        }
      }
    }
    return pptCardCount;
  }

  function selectMappingDivision(divId) {
    _mapActiveDivId = divId;
    renderMappingChips();

    const div = _mapDivisions.find(d => d.id === divId);
    if (!div) return;

    const meta = document.getElementById('mappingDivisionMeta');
    const grid = document.getElementById('mappingProjectGrid');
    const empty = document.getElementById('mappingEmpty');

    const projects = div.projects || [];
    const noteData = _mapNotes[divId] || {};
    const noteCards = noteData.cards || [];

    if (meta) {
      meta.textContent = '▼ ' + div.label + ' — 프로젝트 ' + projects.length + '개 / 주간 보고 카드 ' + noteCards.length + '개';
    }

    if (projects.length === 0) {
      if (grid) grid.innerHTML = '';
      if (empty) empty.style.display = 'block';
      return;
    }
    if (empty) empty.style.display = 'none';

    if (grid) {
      grid.innerHTML = '';
      projects.forEach(prj => {
        const card = document.createElement('div');
        card.className = 'map-pcard';
        card.style.cursor = 'pointer';
        card.onclick = () => openMappingDetail(divId, prj.id);

        const labelLc = (prj.label || '').toLowerCase();
        const noteMatch = noteCards.some(nc => {
          const t = (nc.title || '').toLowerCase();
          return t && (t.includes(labelLc) || labelLc.includes(t));
        });

        const pptCount = _countProjectMatches(divId, prj);

        const pptCls = pptCount > 0 ? 'badge-ppt' : 'badge-empty';
        const noteCls = noteMatch ? 'badge-note' : 'badge-empty';
        const noteText = noteMatch ? '✅' : '—';

        card.innerHTML =
          '<div class="ptitle">' + (prj.label || '?') + '</div>' +
          '<div class="pmeta">' +
            '<span class="' + pptCls + '">📊 PPT ' + pptCount + '</span>' +
            '<span class="' + noteCls + '">📝 노트 ' + noteText + '</span>' +
          '</div>';

        grid.appendChild(card);
      });
    }
  }

  function openMappingDetail(divisionId, projectId) {
    const div = _mapDivisions.find(d => d.id === divisionId);
    if (!div) return;
    const prj = (div.projects || []).find(x => x.id === projectId);
    if (!prj) return;

    document.getElementById('mappingDetailTitle').textContent = prj.label + ' (' + div.label + ')';

    // PPT 카드 매칭
    const aliases = (prj.aliases || []).map(s => s.toLowerCase());
    const labelLc = (prj.label || '').toLowerCase();
    const pptMatches = [];

    for (const rep of _mapReports) {
      const fname = rep.file_name || '';
      const ts = rep.upload_timestamp || '';
      const products = rep.products || [];
      for (const pd of products) {
        const pkey = pd.project_key || pd.project || '';
        const text = ((pd.name || '') + ' ' + (pd.headline || '')).toLowerCase();
        let matched = false;
        if (pkey === projectId) matched = true;
        else if (labelLc && text.includes(labelLc)) matched = true;
        else {
          for (const al of aliases) {
            if (al && text.includes(al)) { matched = true; break; }
          }
        }
        if (matched) {
          pptMatches.push({ rep_file: fname, rep_ts: ts, product: pd });
        }
      }
    }

    // 노트 카드 매칭 (해당 사업부의 노트 중 카드 title이 프로젝트명 포함)
    const noteData = _mapNotes[divisionId] || {};
    const noteCards = noteData.cards || [];
    const noteMatches = noteCards.filter(nc => {
      const t = (nc.title || '').toLowerCase();
      return t && (t.includes(labelLc) || labelLc.includes(t));
    });

    // Subtitle
    const sub = document.getElementById('mappingDetailSubtitle');
    sub.textContent = 'PPT 카드 ' + pptMatches.length + '개 / 노트 카드 ' + noteMatches.length + '개';

    // PPT 영역
    const pptArea = document.getElementById('mappingDetailPpt');
    if (pptMatches.length === 0) {
      pptArea.innerHTML = '<div style="color:#9ca3af;padding:10px 0;">매핑된 PPT 카드가 없습니다.</div>';
    } else {
      pptArea.innerHTML = '';
      pptMatches.forEach(m => {
        const wrap = document.createElement('div');
        wrap.style.cssText = 'background:#f9fafb;border:1px solid #e5e7eb;border-radius:8px;padding:10px 12px;margin-bottom:8px;';
        const fname = m.rep_file || '(파일명 없음)';
        const ts = m.rep_ts ? new Date(m.rep_ts).toLocaleString('ko-KR') : '';
        const name = m.product.name || '';
        const headline = m.product.headline || '';
        const bullets = m.product.summary_bullets || [];

        let html = '<div style="font-size:11px;color:#9ca3af;margin-bottom:4px;">📄 ' + fname + (ts ? ' · ' + ts : '') + '</div>';
        html += '<div style="font-weight:600;color:#1f2937;margin-bottom:2px;">' + (name || '(이름 없음)') + '</div>';
        if (headline) html += '<div style="color:#4b5563;margin-bottom:6px;">' + headline + '</div>';
        if (bullets.length) {
          html += '<ul style="margin:4px 0 0 18px;padding:0;color:#374151;">';
          bullets.forEach(b => { html += '<li style="margin-bottom:2px;">' + b + '</li>'; });
          html += '</ul>';
        }
        wrap.innerHTML = html;
        pptArea.appendChild(wrap);
      });
    }

    // 노트 영역
    const noteArea = document.getElementById('mappingDetailNotes');
    if (noteMatches.length === 0) {
      noteArea.innerHTML = '<div style="color:#9ca3af;padding:10px 0;">매핑된 주간 보고 노트가 없습니다.</div>';
    } else {
      noteArea.innerHTML = '';
      const reportDate = noteData.report_date || '';
      noteMatches.forEach(nc => {
        const wrap = document.createElement('div');
        wrap.style.cssText = 'background:#f0fdf4;border:1px solid #bbf7d0;border-radius:8px;padding:10px 12px;margin-bottom:8px;';
        let html = '<div style="font-weight:600;color:#1E3A5F;margin-bottom:4px;">' + (nc.title || '(제목 없음)') + (reportDate ? ' · ' + reportDate : '') + '</div>';
        (nc.sections || []).forEach(sec => {
          html += '<div style="font-size:12px;font-weight:600;color:#374151;margin:6px 0 2px;padding:3px 8px;background:#dcfce7;border-radius:4px;">' + (sec.title || '') + '</div>';
          (sec.items || []).forEach(it => {
            const isHl = it.type === 'highlight';
            const isSub = it.type === 'sub';
            const dot = isSub ? '↳' : '•';
            const color = isHl ? '#dc2626' : '#1f2937';
            const fw = isHl ? '600' : '400';
            const ml = isSub ? '20px' : '6px';
            html += '<div style="font-size:12px;color:' + color + ';font-weight:' + fw + ';margin-left:' + ml + ';line-height:1.6;">' + dot + ' ' + (it.text || '') + '</div>';
          });
        });
        wrap.innerHTML = html;
        noteArea.appendChild(wrap);
      });
    }

    document.getElementById('mappingDetailModal').classList.add('show');
  }

  function closeMappingDetail() {
    document.getElementById('mappingDetailModal').classList.remove('show');
  }

  let _mapInitialized = false;
  document.addEventListener('DOMContentLoaded', () => {
    document.querySelectorAll('.tab-btn').forEach(btn => {
      if (btn.dataset.tab === 'mapping') {
        btn.addEventListener('click', () => {
          if (!_mapInitialized) {
            _mapInitialized = true;
            initMappingTab();
          }
        });
      }
    });
    // ESC로 모달 닫기
    document.addEventListener('keydown', (e) => {
      if (e.key === 'Escape') {
        const m = document.getElementById('mappingDetailModal');
        if (m && m.classList.contains('show')) m.classList.remove('show');
      }
    });
  });

  // 노트 탭 초기화
  document.addEventListener('DOMContentLoaded', () => {
    initNoteDate();
    loadNoteDivisions();
    const aiBtn = document.getElementById('noteAiParseBtn');
    if (aiBtn) aiBtn.addEventListener('click', noteAiParse);
    const saveBtn = document.getElementById('noteSaveBtn');
    if (saveBtn) saveBtn.addEventListener('click', noteSave);
    const loadBtn = document.getElementById('noteLoadBtn');
    if (loadBtn) loadBtn.addEventListener('click', function(){ if (window.openNoteLoadModal) { window.openNoteLoadModal(); } else { alert('note_loader.js 로딩 실패'); } });
  });

</script>

<!-- 표 편집 모달 -->
<div class="modal-overlay" id="tableModal">
  <div class="modal-box">
    <h3>📊 표 편집</h3>
    <input type="text" id="tableModalTitle" class="modal-input" placeholder="표 제목 (선택)" />
    <div style="font-size:12px;color:#6b7280;margin-bottom:6px;">
      💡 Excel에서 복사한 표를 아래에 붙여넣으면 자동 인식됩니다 (Tab/쉼표 구분).
    </div>
    <textarea id="tablePasteArea" class="modal-input" placeholder="여기에 붙여넣기 → 자동 파싱"></textarea>
    <div id="tableEditorArea"></div>
    <div style="display:flex; gap:6px; margin-bottom:12px;">
      <button onclick="addTableRow()" style="background:#e5e7eb;color:#374151;margin:0;padding:6px 12px;font-size:12px;">＋ 행 추가</button>
      <button onclick="addTableCol()" style="background:#e5e7eb;color:#374151;margin:0;padding:6px 12px;font-size:12px;">＋ 열 추가</button>
      <button onclick="removeTableRow()" style="background:#fee2e2;color:#991b1b;margin:0;padding:6px 12px;font-size:12px;">－ 마지막 행 삭제</button>
      <button onclick="removeTableCol()" style="background:#fee2e2;color:#991b1b;margin:0;padding:6px 12px;font-size:12px;">－ 마지막 열 삭제</button>
    </div>
    <div class="modal-actions">
      <button class="danger" id="tableModalDeleteBtn" onclick="deleteTableFromModal()">🗑️ 삭제</button>
      <button class="cancel" onclick="closeTableModal()">취소</button>
      <button onclick="saveTableFromModal()">💾 저장</button>
    </div>
  </div>
</div>

<!-- 매핑 관리: 프로젝트 상세 모달 -->
<div class="modal-overlay" id="mappingDetailModal">
  <div class="modal-box" style="min-width:600px;max-width:800px;">
    <div style="display:flex;justify-content:space-between;align-items:start;margin-bottom:12px;">
      <h3 id="mappingDetailTitle" style="margin:0;color:#1E3A5F;">프로젝트 상세</h3>
      <button onclick="closeMappingDetail()" style="background:transparent;color:#6b7280;font-size:24px;padding:0;margin:0;border:none;cursor:pointer;line-height:1;">✕</button>
    </div>

    <div style="font-size:13px;color:#6b7280;margin-bottom:14px;" id="mappingDetailSubtitle"></div>

    <h4 style="margin:14px 0 8px;color:#1E3A5F;font-size:15px;">📊 PPT 카드</h4>
    <div id="mappingDetailPpt" style="font-size:13px;"></div>

    <h4 style="margin:18px 0 8px;color:#1E3A5F;font-size:15px;">📝 주간 보고 노트</h4>
    <div id="mappingDetailNotes" style="font-size:13px;"></div>
  </div>
</div>

<!-- PPT 초안 생성 사업부 선택 모달 -->
<div class="modal-overlay" id="draftDivModal">
  <div class="modal-box" style="max-width:520px;">
    <div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:12px;">
      <h3 style="margin:0;">📝 주간 보고 초안 생성</h3>
      <button onclick="closeDraftModal()" style="background:transparent;color:#6b7280;font-size:24px;padding:0;margin:0;border:none;cursor:pointer;line-height:1;">✕</button>
    </div>
    <p style="color:#6b7280;font-size:13px;margin:0 0 10px;"><span id="draftModalFile" style="color:#374151;font-weight:600;"></span></p>
    <p style="color:#6b7280;font-size:13px;margin:0 0 14px;">이 PPT를 어느 사업부 주간 보고로 변환할까요?</p>
    <div id="draftDivChips" style="display:flex;flex-wrap:wrap;gap:8px;margin-bottom:14px;"></div>
    <div id="draftModalStatus" style="font-size:13px;color:#6b7280;margin-top:8px;min-height:18px;"></div>
  </div>
</div>

<!-- 사진 확대 모달 -->
<div class="photo-overlay" id="photoOverlay" onclick="this.classList.remove('show')">
  <img id="photoOverlayImg" src="" alt="확대 사진" />
</div>

<!-- 숨김 사진 업로드 input -->
<input type="file" id="notePhotoFileInput" accept="image/*" style="display:none;" />



<script src="/static/note_loader.js"></script>
<script src="/static/excel_drop.js"></script>
<script src="/static/photo_drop.js">
    
    </script>
</body>

</html>
"""

@app.post("/chat")
async def chat(payload: dict):
    """
    RAG 챗봇 엔드포인트
    입력: {"message": "질문 텍스트", "top_k": 5}
    출력: {"answer": "...", "sources": [{project_label, division_id, ...}]}
    """
    message = (payload or {}).get("message", "").strip()
    top_k = int((payload or {}).get("top_k", 5))
    if not message:
        return {"answer": "", "sources": [], "error": "empty message"}

    if not _vs.is_ready():
        return {
            "answer": "챗봇 인덱스가 준비되지 않았어요. 관리자에게 문의해주세요.",
            "sources": [],
            "error": "index_not_ready",
        }

    try:
        hits = _vs.search(message, top_k=top_k)
    except Exception as e:
        return {"answer": "", "sources": [], "error": f"search_failed: {e}"}

    if not hits:
        return {
            "answer": "관련된 프로젝트 정보를 찾지 못했어요.",
            "sources": [],
        }

    context_parts = []
    for i, h in enumerate(hits):
        context_parts.append(
            f"[문서 {i+1}] 사업부: {h.get('division_id')} / "
            f"프로젝트: {h.get('project_label')} / "
            f"보고일: {h.get('report_date')}\n"
            f"{h.get('text','')}"
        )
    context = "\n\n".join(context_parts)

    system_prompt = (
        "너는 반도체 사업부의 프로젝트 보고 어시스턴트다. "
        "주어진 [문서] 내용만을 근거로 사용자의 질문에 한국어로 간결하게 답한다. "
        "규칙: 1) 문서에 없는 내용은 추측하지 말 것 2) 숫자·일정·모델명은 원문 그대로 유지 "
        "3) 3~5문장 이내 4) 마지막에 '(근거: 프로젝트명)' 형태로 인용 표시 "
        "5) 여러 문서를 종합해야 하면 각 문서별 사실만 언급."
    )
    user_prompt = f"[문서]\n{context}\n\n[질문]\n{message}"

    try:
        if client is None:
            return {"answer": "OpenAI 키가 설정되지 않았어요.", "sources": hits}
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            temperature=0.2,
            max_tokens=500,
        )
        answer = (resp.choices[0].message.content or "").strip()
    except Exception as e:
        return {"answer": "", "sources": hits, "error": f"llm_failed: {e}"}

    return {
        "answer": answer,
        "sources": [
            {
                "division_id": h.get("division_id"),
                "project_label": h.get("project_label"),
                "project_key": h.get("project_key"),
                "report_date": h.get("report_date"),
                "score": round(h.get("score", 0.0), 3),
            }
            for h in hits
        ],
    }



# =========================================================
# KPI 카드: 히스토리 파일 기반 (프로젝트별 월/주차 누적)
# =========================================================
from datetime import datetime, date, timedelta
import calendar as _calendar

KPI_HISTORY_FILE = BASE_DIR / "kpi_history.json" if "BASE_DIR" in globals() else Path("kpi_history.json")
PROJECT_PROFILES_FILE = BASE_DIR / "project_profiles.json" if "BASE_DIR" in globals() else Path("project_profiles.json")


def _extract_tables_from_pptx(file_path) -> list:
    """PPT 파일에서 모든 표를 추출.

    반환 형식: [
      {
        "slide_index": int,   # 1부터 시작
        "slide_title": str,   # 첫 번째 텍스트 도형(있으면)
        "nearby_text": str,   # 슬라이드의 모든 텍스트를 이어붙인 것
        "rows": [[cell, cell, ...], ...]
      },
      ...
    ]
    """
    from pptx import Presentation as _Presentation
    tables = []
    try:
        prs = _Presentation(str(file_path))
    except Exception as e:
        return []

    for s_idx, slide in enumerate(prs.slides, start=1):
        # 슬라이드 전체 텍스트 수집
        slide_texts = []
        slide_title = ""
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    tf_text = "\n".join(
                        p.text for p in shape.text_frame.paragraphs if p.text
                    ).strip()
                    if tf_text:
                        slide_texts.append(tf_text)
                        if not slide_title:
                            slide_title = tf_text.split("\n")[0]
            except Exception:
                continue

        nearby_text = "\n".join(slide_texts)

        # 슬라이드 안의 표 추출
        for shape in slide.shapes:
            try:
                if not shape.has_table:
                    continue
                tbl = shape.table
                rows = []
                for row in tbl.rows:
                    cells = []
                    for cell in row.cells:
                        try:
                            txt = cell.text_frame.text if cell.text_frame else ""
                        except Exception:
                            txt = ""
                        cells.append((txt or "").strip())
                    rows.append(cells)
                tables.append({
                    "slide_index": s_idx,
                    "slide_title": slide_title,
                    "nearby_text": nearby_text,
                    "rows": rows,
                })
            except Exception:
                continue

    return tables


def _is_shipping_plan_table(table_dict: dict) -> bool:
    """표가 '출하 계획 및 실적' 관련인지 판정."""
    if not isinstance(table_dict, dict):
        return False
    keywords = [
        "모델 별 출하 계획 및 실적",
        "모델별 출하 계획 및 실적",
        "출하 계획 및 실적",
        "주차별 출하실적",
        "주차별 출하 실적",
    ]
    # 슬라이드 주변 텍스트에서 키워드 탐지
    nearby = (table_dict.get("nearby_text") or "").replace(" ", "")
    for kw in keywords:
        if kw.replace(" ", "") in nearby:
            return True
    # 표 헤더에서도 탐지 (첫 3행)
    rows = table_dict.get("rows") or []
    header_text = " ".join(
        " ".join(r) for r in rows[:3] if isinstance(r, list)
    ).replace(" ", "")
    for kw in ["W27", "W28", "PO", "누적합계", "누적 합계"]:
        if kw.replace(" ", "") in header_text:
            # 주차 헤더가 있는 표면 후보로 인정
            return True
    return False


def _find_shipping_plan_tables(file_path) -> list:
    """PPT에서 출하 계획 관련 표만 추출."""
    all_tables = _extract_tables_from_pptx(file_path)
    return [t for t in all_tables if _is_shipping_plan_table(t)]


def _pick_primary_shipping_table(tables: list, project_key: str) -> dict | None:
    """여러 표 중 진짜 출하 계획 표를 하나만 선택.

    판별 기준:
    1. R1 첫 셀에 '구분' 포함
    2. R1에 '누적 합계' 또는 '누적합계' 포함
    3. 데이터 행 첫 컬럼에 프로파일 등록 모델명 존재
    """
    if not tables:
        return None
    profile = _get_project_profile(project_key) or {}
    profile_models = {
        (m.get("display_name") or "").strip()
        for m in (profile.get("models") or [])
    }

    candidates = []
    for t in tables:
        rows = t.get("rows") or []
        if len(rows) < 3:
            continue
        r1 = rows[0]
        r1_txt = " ".join(r1)
        r1_first = (r1[0] or "").strip() if r1 else ""
        # 조건 1: 첫 셀 '구분'
        if "구분" not in r1_first:
            continue
        # 조건 2: 누적 합계 OR (PO + 실적/출하) 조합 존재
        r1_norm = r1_txt.replace(" ", "").replace("\n", "")
        has_cumulative = "누적" in r1_norm
        has_po_group = ("PO수량" in r1_norm or "PO" in r1_norm) and ("실적" in r1_norm or "출하" in r1_norm)
        if not (has_cumulative or has_po_group):
            continue
        # 조건 3: 데이터 행 첫 컬럼에 프로파일 모델 존재
        data_models = set()
        for row in rows[2:]:
            if not row:
                continue
            first = (row[0] or "").strip()
            if first:
                data_models.add(first)
        if profile_models and not (profile_models & data_models):
            # 프로파일 모델과 하나도 매칭 안 됨
            continue
        candidates.append(t)

    if not candidates:
        return None
    # 여러 개면 데이터 행이 가장 많은 것 선택
    candidates.sort(key=lambda x: len(x.get("rows") or []), reverse=True)
    return candidates[0]


def _parse_shipping_table_header(rows: list) -> dict:
    """헤더 3행(R1, R2, R3)을 분석해 각 컬럼의 의미를 매핑.

    처리 규칙:
    - R1: 대분류 (월/누적합계) - forward-fill (병합 셀)
    - R2: 서브그룹 (W27/W28.../합계) - forward-fill (계획/실적 쌍의 왼쪽 셀에만 값 있음)
    - R3: 세부 (계획/실적/PO/출하)
    """
    if len(rows) < 3:
        return {"columns": [], "detected_months": [], "current_month": None}

    r1 = rows[0]
    r2 = rows[1]
    r3 = rows[2]
    n_cols = max(len(r1), len(r2), len(r3))

    def _fill(row, n):
        """forward-fill: 빈 셀은 왼쪽 값으로 채움"""
        out = []
        last = ""
        for i in range(n):
            v = (row[i] if i < len(row) else "").strip()
            if v:
                last = v
            out.append(last)
        return out

    r1_filled = _fill(r1, n_cols)
    # R2는 forward-fill 하되, 다음 대분류가 바뀌면 리셋
    r2_filled = []
    last_r1 = ""
    last_r2 = ""
    for i in range(n_cols):
        cur_r1 = r1_filled[i]
        cur_r2 = (r2[i] if i < len(r2) else "").strip()
        if cur_r1 != last_r1:
            # 대분류가 바뀌면 R2 리셋
            last_r2 = ""
            last_r1 = cur_r1
        if cur_r2:
            last_r2 = cur_r2
        r2_filled.append(last_r2)

    columns = []
    detected_months = []
    months_with_weeks = []

    for i in range(1, n_cols):  # 0번은 '구분' 열
        top = r1_filled[i]
        mid = r2_filled[i]
        sub = (r3[i] if i < len(r3) else "").strip()

        if not top:
            continue

        top_norm = top.replace(" ", "").replace("\n", "")
        is_month = top.endswith("월") and "누적" not in top and top_norm not in ("PO수량", "출하실적", "잔량")
        is_cumulative = "누적" in top_norm
        is_direct_po = top_norm in ("PO수량", "PO",)
        is_direct_shipped = top_norm in ("출하실적", "출하",)
        is_direct_remaining = top_norm in ("잔량",)
        is_yearly = ("년실적" in top_norm) or ("년도실적" in top_norm)

        if is_month and top not in detected_months:
            detected_months.append(top)

        if is_direct_po:
            columns.append({"index": i, "kind": "cumulative_po"})
            continue
        if is_direct_shipped:
            columns.append({"index": i, "kind": "cumulative_shipped"})
            continue
        if is_direct_remaining:
            columns.append({"index": i, "kind": "cumulative_remaining"})
            continue
        if is_yearly:
            columns.append({"index": i, "kind": "yearly_actual", "label": top})
            continue

        if is_cumulative:
            if "PO" in sub.upper():
                columns.append({"index": i, "kind": "cumulative_po"})
            elif "출하" in sub:
                columns.append({"index": i, "kind": "cumulative_shipped"})
            else:
                columns.append({"index": i, "kind": "cumulative_other"})
            continue

        if is_month:
            # W##이 있으면 주차별
            if mid.startswith("W") and mid[1:].isdigit():
                if top not in months_with_weeks:
                    months_with_weeks.append(top)
                if sub == "계획":
                    columns.append({"index": i, "kind": "week_plan", "month": top, "week": mid})
                elif sub == "실적":
                    columns.append({"index": i, "kind": "week_actual", "month": top, "week": mid})
            elif mid == "합계":
                if sub == "계획":
                    columns.append({"index": i, "kind": "month_total_plan", "month": top})
                elif sub == "실적":
                    columns.append({"index": i, "kind": "month_total_actual", "month": top})
            else:
                # 완료 월 또는 다음 월 계획만
                if sub == "계획":
                    columns.append({"index": i, "kind": "month_plan", "month": top})
                elif sub == "실적":
                    columns.append({"index": i, "kind": "month_actual", "month": top})

    current_month = months_with_weeks[0] if months_with_weeks else None

    return {
        "columns": columns,
        "detected_months": detected_months,
        "current_month": current_month,
    }


def _shipping_table_to_kpi_data(table: dict, project_key: str, file_name: str = "") -> dict:
    """표 → KPI 스키마로 변환.

    반환:
    {
      "project_key": "major_module",
      "report_month": "2026-07",
      "report_year": 2026,
      "models": [
        {
          "model_id": "major_module::EFEM",
          "display_name": "EFEM",
          "months": {"2026-04": {"plan": 16, "actual": 9}, ...},
          "weeks": {"2026-W27": {"plan": 2, "actual": 2}, ...},
          "cumulative": {"po": 74, "shipped": 37}
        }
      ],
      "warnings": [...]
    }
    """
    warnings = []
    profile = _get_project_profile(project_key) or {}
    profile_models = {
        (m.get("display_name") or "").strip(): m
        for m in (profile.get("models") or [])
    }

    rows = table.get("rows") or []
    header_info = _parse_shipping_table_header(rows)
    columns = header_info["columns"]
    current_month_str = header_info["current_month"]  # 예: "7월"

    # 연도 판별: 파일명에서 추출
    year = None
    parsed = _parse_report_filename(file_name or "") if file_name else {}
    date_str = parsed.get("date") if parsed else None
    if date_str:
        try:
            year = int(date_str.split("-")[0])
        except Exception:
            pass
    if not year:
        from datetime import datetime as _dt
        year = _dt.now().year

    # 파일명 월 vs 표 현재 월 일치 확인
    file_month = None
    if date_str:
        try:
            file_month = int(date_str.split("-")[1])
        except Exception:
            pass
    table_month = None
    if current_month_str and current_month_str.endswith("월"):
        try:
            table_month = int(current_month_str[:-1])
        except Exception:
            pass
    if file_month and table_month and file_month != table_month:
        warnings.append(
            f"파일명 월({file_month}월)과 표의 현재 월({table_month}월)이 다릅니다"
        )

    def _month_key(month_str):
        if not month_str or not month_str.endswith("월"):
            return None
        try:
            m = int(month_str[:-1])
            return f"{year:04d}-{m:02d}"
        except Exception:
            return None

    def _week_key(week_str):
        if not week_str or not week_str.startswith("W"):
            return None
        num = week_str[1:]
        if not num.isdigit():
            return None
        return f"{year:04d}-W{int(num):02d}"

    def _to_int(v):
        if v is None:
            return None
        s = str(v).strip()
        if s in ("", "-", "·", "N/A"):
            return None
        try:
            return int(float(s.replace(",", "")))
        except Exception:
            return None

    # 데이터 행 순회 (R4부터)
    models_data = []
    for row in rows[2:]:
        if not row:
            continue
        model_name = (row[0] or "").strip()
        if not model_name:
            continue
        if model_name in ("합계", "총계", "TOTAL"):
            continue  # 합계 행은 파싱 대상 아님 (검증용, 이번엔 스킵)

        # 프로파일 매칭
        profile_model = profile_models.get(model_name)
        if not profile_model:
            warnings.append(f"프로파일에 없는 모델: '{model_name}' (project={project_key})")
            continue

        model_id = profile_model.get("model_id") or f"{project_key}::{model_name}"

        months_data = {}
        weeks_data = {}
        cumulative = {"po": None, "shipped": None}

        for col in columns:
            idx = col["index"]
            v = _to_int(row[idx]) if idx < len(row) else None
            kind = col["kind"]

            if kind == "month_plan":
                mk = _month_key(col["month"])
                if mk:
                    months_data.setdefault(mk, {})["plan"] = v
            elif kind == "month_actual":
                mk = _month_key(col["month"])
                if mk:
                    months_data.setdefault(mk, {})["actual"] = v
            elif kind == "week_plan":
                wk = _week_key(col["week"])
                if wk:
                    weeks_data.setdefault(wk, {})["plan"] = v
            elif kind == "week_actual":
                wk = _week_key(col["week"])
                if wk:
                    weeks_data.setdefault(wk, {})["actual"] = v
            elif kind == "month_total_plan":
                mk = _month_key(col["month"])
                if mk:
                    months_data.setdefault(mk, {})["plan"] = v
            elif kind == "month_total_actual":
                mk = _month_key(col["month"])
                if mk:
                    months_data.setdefault(mk, {})["actual"] = v
            elif kind == "cumulative_po":
                cumulative["po"] = v
            elif kind == "cumulative_shipped":
                cumulative["shipped"] = v
            elif kind == "cumulative_remaining":
                cumulative["remaining"] = v
            elif kind == "yearly_actual":
                cumulative["yearly_actual"] = v

        # === 주차 실적/계획을 월별로 합산해서 월 데이터 보정 ===
        # 규칙: 표 헤더에서 이미 "이 W##이 어느 월 컬럼 아래 있는지" 알고 있으므로,
        #       ISO 주차 계산 없이 헤더의 month 정보를 그대로 사용
        # 필요한 매핑: week_key(2026-W27) -> month_str("7월") -> month_key("2026-07")
        week_to_month = {}  # "2026-W27" -> "2026-07"
        for col in columns:
            if col.get("kind") in ("week_plan", "week_actual"):
                w = col.get("week")  # "W27"
                m = col.get("month")  # "7월"
                if w and m:
                    wk_full = _week_key(w)  # "2026-W27"
                    mk_full = _month_key(m)  # "2026-07"
                    if wk_full and mk_full:
                        week_to_month[wk_full] = mk_full

        # 주차 데이터를 월별로 합산
        week_by_month = {}
        for wk, wvals in weeks_data.items():
            mk = week_to_month.get(wk)
            if not mk:
                continue
            bucket = week_by_month.setdefault(mk, {"plan": [], "actual": []})
            if wvals.get("plan") is not None:
                bucket["plan"].append(wvals.get("plan"))
            if wvals.get("actual") is not None:
                bucket["actual"].append(wvals.get("actual"))

        # 월 데이터에 반영 (기존 값 우선, 없으면 합산값 채움)
        for mk, bucket in week_by_month.items():
            entry = months_data.setdefault(mk, {})
            if entry.get("actual") is None and bucket["actual"]:
                entry["actual"] = sum(bucket["actual"])
            if entry.get("plan") is None and bucket["plan"]:
                entry["plan"] = sum(bucket["plan"])

        models_data.append({
            "model_id": model_id,
            "display_name": model_name,
            "months": months_data,
            "weeks": weeks_data,
            "cumulative": cumulative,
        })

    return {
        "project_key": project_key,
        "report_month": f"{year:04d}-{table_month:02d}" if table_month else None,
        "report_year": year,
        "models": models_data,
        "warnings": warnings,
    }


def _load_project_profiles() -> dict:
    """프로젝트 프로파일 전체 로드."""
    try:
        import json as _json
        with open(PROJECT_PROFILES_FILE, "r", encoding="utf-8") as f:
            data = _json.load(f)
        if not isinstance(data, dict):
            return {"profiles": {}}
        if "profiles" not in data or not isinstance(data.get("profiles"), dict):
            return {"profiles": {}}
        return data
    except FileNotFoundError:
        return {"profiles": {}}
    except Exception:
        return {"profiles": {}}


def _get_project_profile(project_key: str) -> dict | None:
    """특정 프로젝트 프로파일 반환. 없으면 None."""
    if not project_key:
        return None
    data = _load_project_profiles()
    return (data.get("profiles") or {}).get(project_key)


def _save_project_profiles(data: dict) -> None:
    """프로젝트 프로파일 저장. updated_at 자동 갱신."""
    import json as _json
    from datetime import datetime as _dt
    if not isinstance(data, dict):
        raise ValueError("profiles data must be dict")
    now = _dt.now().isoformat()
    for pk, prof in (data.get("profiles") or {}).items():
        if isinstance(prof, dict):
            prof["updated_at"] = now
    with open(PROJECT_PROFILES_FILE, "w", encoding="utf-8") as f:
        _json.dump(data, f, ensure_ascii=False, indent=2)



def _auto_extract_kpi_from_pptx(
    saved_file_path,
    file_name: str,
    doc_id: str,
    project_id_hint=None,
    products=None,
) -> dict:
    """PPT 업로드 후 자동 KPI 추출 훅.

    실패해도 예외를 던지지 않음 (업로드는 성공 유지).
    프로파일이 등록된 프로젝트만 파싱 시도.
    """
    result = {
        "attempted": False,
        "project_key": None,
        "ok": False,
        "reason": None,
        "months_count": 0,
        "weeks_count": 0,
        "cumulative_updated": False,
    }
    try:
        # 1. project_key 결정
        project_key = None

        # 1-a. 명시적 힌트 우선
        if project_id_hint:
            project_key = str(project_id_hint).strip() or None

        # 1-b. products에서 자동 매칭
        if not project_key and products:
            try:
                from project_templates import _match_project_key
                for p in products:
                    if not isinstance(p, dict):
                        continue
                    name = (p.get("name") or p.get("product") or "").strip()
                    if name:
                        pk = _match_project_key(name)
                        if pk:
                            project_key = pk
                            break
                    category = (p.get("category") or "").strip()
                    if category:
                        pk = _match_project_key(category)
                        if pk:
                            project_key = pk
                            break
            except Exception as e:
                print(f"[KPI 자동 추출] products 매칭 실패: {e}")

        # 1-c. 파일명에서 추론
        if not project_key and file_name:
            try:
                from project_templates import _match_project_key
                parsed = _parse_report_filename(file_name)
                for proj_name in (parsed or {}).get("projects", []):
                    pk = _match_project_key(proj_name)
                    if pk:
                        project_key = pk
                        break
            except Exception as e:
                print(f"[KPI 자동 추출] 파일명 매칭 실패: {e}")

        if not project_key:
            result["reason"] = "project_key 결정 실패"
            print(f"[KPI 자동 추출] 스킵: {result['reason']} (file={file_name})")
            return result

        result["project_key"] = project_key

        # 2. 프로파일 조회
        profile = _get_project_profile(project_key)
        if not profile:
            result["reason"] = f"프로파일 없음: {project_key}"
            print(f"[KPI 자동 추출] 스킵: {result['reason']}")
            return result

        result["attempted"] = True

        # 3. 표 감지
        plan_tables = _find_shipping_plan_tables(saved_file_path)
        if not plan_tables:
            result["reason"] = "출하 계획 표를 찾지 못함"
            print(f"[KPI 자동 추출] 스킵: {result['reason']}")
            return result

        # 4. 진짜 출하 표 선택
        primary = _pick_primary_shipping_table(plan_tables, project_key)
        if not primary:
            result["reason"] = "출하 표 후보 없음 (프로파일 모델 미매칭)"
            print(f"[KPI 자동 추출] 스킵: {result['reason']}")
            return result

        # 5. KPI 스키마 변환
        kpi_data = _shipping_table_to_kpi_data(primary, project_key, file_name or "")
        if not kpi_data.get("models"):
            result["reason"] = "모델 데이터가 비어있음"
            print(f"[KPI 자동 추출] 스킵: {result['reason']}")
            return result

        # 6. 저장
        source_label = f"pptx:{doc_id}" if doc_id else f"pptx:{file_name}"
        apply_result = _apply_kpi_from_pptx(
            project_key,
            kpi_data,
            source_info={"label": source_label},
            dry_run=False,
        )
        result["ok"] = True
        result["months_count"] = apply_result.get("months_count", 0)
        result["weeks_count"] = apply_result.get("weeks_count", 0)
        result["cumulative_updated"] = apply_result.get("cumulative_updated", False)

        warnings = kpi_data.get("warnings") or []
        warn_str = f", warnings={len(warnings)}" if warnings else ""
        print(
            f"[KPI 자동 추출 성공] project={project_key}, "
            f"months={result['months_count']}, weeks={result['weeks_count']}, "
            f"cumulative={result['cumulative_updated']}{warn_str}"
        )
        return result
    except Exception as e:
        result["reason"] = f"예외 발생: {e}"
        print(f"[KPI 자동 추출 실패] {e}")
        import traceback
        traceback.print_exc()
        return result


def _apply_kpi_from_pptx(
    project_key: str,
    kpi_data: dict,
    source_info: dict = None,
    dry_run: bool = True,
) -> dict:
    """PPT에서 추출한 KPI 데이터를 kpi_history.json에 병합.

    - 기존 필드(efem, vtm, type, source, updated_at) 유지
    - 신규 필드(models, cumulative) 추가
    - dry_run=True: 실제 저장 안 하고 변경 예정 결과만 반환
    - dry_run=False: 실제 저장
    """
    from datetime import datetime as _dt
    source_info = source_info or {}
    now = _dt.now().isoformat()

    def _pick_value(model_dict, key_type):
        """수량 값 결정: actual 우선, 없으면 plan"""
        v = (model_dict or {}).get(key_type)
        return v

    def _pick_efem_vtm_and_type(period_dict):
        """월/주차 데이터에서 efem, vtm, type 결정.

        period_dict: {"major_module::EFEM": {"plan": x, "actual": y}, ...}
        """
        efem_val = None
        vtm_val = None
        chosen_type = None
        for mid, vals in period_dict.items():
            if not isinstance(vals, dict):
                continue
            actual = vals.get("actual")
            plan = vals.get("plan")
            # actual 있으면 actual, 없으면 plan
            if actual is not None:
                v = actual
                t = "actual"
            elif plan is not None:
                v = plan
                t = "plan"
            else:
                continue
            display = mid.split("::")[-1] if "::" in mid else mid
            if display == "EFEM":
                efem_val = v
                chosen_type = t if chosen_type is None else chosen_type
            elif display == "VTM":
                vtm_val = v
                chosen_type = t if chosen_type is None else chosen_type
        # 둘 다 actual이면 actual, 하나라도 plan이면 plan 우선? -> 우선 actual 우선
        # 위 로직은 먼저 나온 것 채택. 명시적으로:
        # 두 모델 중 하나라도 actual 있으면 -> actual
        actual_found = any(
            (period_dict.get(k) or {}).get("actual") is not None
            for k in period_dict
        )
        chosen_type = "actual" if actual_found else "plan"
        return efem_val, vtm_val, chosen_type

    # 로드
    data = _load_kpi_history()
    if not isinstance(data, dict):
        data = {"version": 1, "projects": {}}
    data.setdefault("projects", {})
    proj = data["projects"].setdefault(project_key, {})
    proj.setdefault("months", {})
    proj.setdefault("weeks", {})

    changes = {"months": [], "weeks": [], "cumulative": None}

    # 월별 병합
    # kpi_data["models"] 구조: [{model_id, months:{...}, weeks:{...}, cumulative:{}}]
    all_months = {}   # month_key -> {model_id: {plan, actual}}
    all_weeks = {}
    cumulative = {}
    for m in kpi_data.get("models") or []:
        mid = m.get("model_id")
        for mk, vals in (m.get("months") or {}).items():
            all_months.setdefault(mk, {})[mid] = vals
        for wk, vals in (m.get("weeks") or {}).items():
            all_weeks.setdefault(wk, {})[mid] = vals
        cum = m.get("cumulative") or {}
        if cum:
            cumulative[mid] = cum

    src_label = source_info.get("label") or "pptx"

    # 월 병합
    for mk in sorted(all_months.keys()):
        period_models = all_months[mk]
        efem_v, vtm_v, t = _pick_efem_vtm_and_type(period_models)
        old = proj["months"].get(mk) or {}
        new_entry = dict(old)
        if efem_v is not None:
            new_entry["efem"] = efem_v
        if vtm_v is not None:
            new_entry["vtm"] = vtm_v
        new_entry["type"] = t
        new_entry["source"] = src_label
        new_entry["updated_at"] = now
        new_entry["month"] = mk
        # 신규 필드
        new_entry["models"] = {
            mid: {"plan": v.get("plan"), "actual": v.get("actual")}
            for mid, v in period_models.items()
        }
        changes["months"].append({
            "key": mk, "before": old, "after": new_entry
        })
        proj["months"][mk] = new_entry

    # 주 병합
    for wk in sorted(all_weeks.keys()):
        period_models = all_weeks[wk]
        efem_v, vtm_v, t = _pick_efem_vtm_and_type(period_models)
        old = proj["weeks"].get(wk) or {}
        new_entry = dict(old)
        if efem_v is not None:
            new_entry["efem"] = efem_v
        if vtm_v is not None:
            new_entry["vtm"] = vtm_v
        new_entry["type"] = t
        new_entry["source"] = src_label
        new_entry["updated_at"] = now
        new_entry["week"] = wk
        new_entry["models"] = {
            mid: {"plan": v.get("plan"), "actual": v.get("actual")}
            for mid, v in period_models.items()
        }
        changes["weeks"].append({
            "key": wk, "before": old, "after": new_entry
        })
        proj["weeks"][wk] = new_entry

    # 누적 병합
    if cumulative:
        proj["cumulative"] = {
            "models": {
                mid: {"po": c.get("po"), "shipped": c.get("shipped")}
                for mid, c in cumulative.items()
            },
            "source": src_label,
            "updated_at": now,
        }
        changes["cumulative"] = proj["cumulative"]

    data["updated_at"] = now + "Z"

    if not dry_run:
        _save_kpi_history(data)

    return {
        "dry_run": dry_run,
        "project_key": project_key,
        "months_count": len(changes["months"]),
        "weeks_count": len(changes["weeks"]),
        "cumulative_updated": changes["cumulative"] is not None,
        "changes": changes,
    }


def _load_kpi_history() -> dict:
    """kpi_history.json 로드. 없으면 빈 구조."""
    try:
        with open(KPI_HISTORY_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {"version": 1, "projects": {}}


def _save_kpi_history(data: dict) -> None:
    data["updated_at"] = datetime.utcnow().isoformat() + "Z"
    with open(KPI_HISTORY_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def _iso_week_of(dt: date) -> str:
    """date -> 'YYYY-Wnn'."""
    y, w, _ = dt.isocalendar()
    return f"{y}-W{w:02d}"


def _month_key(dt: date) -> str:
    return f"{dt.year:04d}-{dt.month:02d}"


def _month_label(month_key: str) -> str:
    """'2026-07' -> '7월'."""
    try:
        _, m = month_key.split("-")
        return f"{int(m)}월"
    except Exception:
        return month_key


def _week_label(week_key: str) -> str:
    """'2026-W28' -> 'W28'."""
    try:
        return week_key.split("-")[1]
    except Exception:
        return week_key


def _slice_window(items: dict, key_fn_today, past: int, future: int) -> list:
    """
    items: dict[key -> row], key는 정렬 가능한 문자열
    key_fn_today: 오늘의 기준 key
    past 개 이전 ~ future 개 이후까지 반환. 실제 있는 것만.
    """
    today_key = key_fn_today()
    all_keys = sorted(items.keys())
    if not all_keys:
        return []
    # 오늘 key 앞뒤로 잘라내기 (없으면 가장 가까운 이전 key 기준)
    if today_key in all_keys:
        idx = all_keys.index(today_key)
    else:
        # today_key보다 작거나 같은 가장 마지막 key
        idx = 0
        for i, k in enumerate(all_keys):
            if k <= today_key:
                idx = i
    lo = max(0, idx - past)
    hi = min(len(all_keys), idx + future + 1)
    return [(k, items[k]) for k in all_keys[lo:hi]]


def _build_major_module_kpi_card() -> dict:
    """
    kpi_history.json 기반으로 메이저모듈 KPI 카드 생성.
    - 판가 기준 매출 계산
    - 월: 오늘 기준 직전 1개월(실적) + 향후 2개월 = 최대 3개
    - 주: 오늘 기준 직전 0주 + 향후 4주 = 최대 5개
    """
    hist = _load_kpi_history()
    proj = (hist.get("projects") or {}).get("major_module") or {}
    ue = proj.get("unit_economics") or {
        "EFEM": {"asp": 130000, "material": 107900},
        "VTM":  {"asp": 240000, "material": 196800},
    }
    asp_e = ue.get("EFEM", {}).get("asp", 130000)
    asp_v = ue.get("VTM",  {}).get("asp", 240000)

    def money_10k(efem_qty: float, vtm_qty: float) -> dict:
        efem = efem_qty * asp_e / 10000
        vtm  = vtm_qty  * asp_v / 10000
        return {"efem": round(efem, 2), "vtm": round(vtm, 2), "total": round(efem + vtm, 2)}

    today = date.today()

    months_raw = proj.get("months") or {}
    weeks_raw  = proj.get("weeks")  or {}

    m_slice = _slice_window(months_raw, lambda: _month_key(today), past=1, future=2)
    w_slice = _slice_window(weeks_raw,  lambda: _iso_week_of(today), past=1, future=4)

    months = []
    for mk, row in m_slice:
        e = row.get("efem", 0); v = row.get("vtm", 0)
        months.append({
            "month": _month_label(mk),
            "type": row.get("type", "plan"),
            **money_10k(e, v),
        })

    weeks = []
    for wk, row in w_slice:
        e = row.get("efem", 0); v = row.get("vtm", 0)
        weeks.append({
            "week": _week_label(wk),
            "type": row.get("type", "plan"),
            **money_10k(e, v),
        })

    # 주차별 타이틀 (가장 이른 주가 속한 달 표시)
    week_group_label = ""
    if w_slice:
        first_wk = w_slice[0][0]
        try:
            y, wn = first_wk.split("-W")
            monday = datetime.strptime(f"{y} {int(wn)} 1", "%G %V %u").date()
            week_group_label = f"{monday.month}월 주차별"
        except Exception:
            week_group_label = "주차별"

    return {
        "title": "월별/주차별 매출",
        "metric_mode": "revenue",
        "metric_note": "판가 기준 매출",
        "unit_label": "만불",
        "unit_economics": {
            "EFEM": {"asp": asp_e, "material": ue.get("EFEM", {}).get("material", 0)},
            "VTM":  {"asp": asp_v, "material": ue.get("VTM",  {}).get("material", 0)},
        },
        "week_group_label": week_group_label,
        "months": months,
        "weeks": weeks,
        "footnotes": [
            "※ 단위: 만불 (USD 10K)",
            "※ 매출 = 수량 × 판가 (ASP 기준)",
        ],
    }


def _build_major_module_issue_lines() -> list:
    """이슈 라인 로드.
    1) kpi_history.json 에 수동 등록된 값이 있으면 그걸 우선.
    2) 없으면 notes.json 의 메이저모듈 sections 에서 자동 생성:
       - type='sub' 이거나 due_date 있는 아이템을 추출
       - due_date 있으면 severity/show_dday 자동 계산
    """
    from datetime import date as _date

    # 1) 수동 등록 우선
    hist = _load_kpi_history()
    proj = (hist.get("projects") or {}).get("major_module") or {}
    manual = proj.get("issue_lines") or []
    if manual:
        return manual

    # 2) notes.json에서 자동 생성
    try:
        notes_data = _load_notes()
        semi = (notes_data.get("notes") or {}).get("semiconductor") or {}
        card = None
        for c in (semi.get("cards") or []):
            if (c.get("title") or "").strip() == "메이저모듈":
                card = c
                break
        if not card:
            return []

        today = _date.today()
        auto_lines = []
        for sec in (card.get("sections") or []):
            for it in (sec.get("items") or []):
                if not isinstance(it, dict):
                    continue
                itype = (it.get("type") or "").lower()
                text = (it.get("text") or "").strip()
                if not text:
                    continue
                # photo 는 스킵
                if itype == "photo":
                    continue
                # sub 또는 due_date 있는 것만
                due = (it.get("due_date") or it.get("due_date_auto") or "").strip()
                if itype != "sub" and not due:
                    continue

                line = {"text": text, "show_dday": False}
                if due:
                    try:
                        y, m, d = due[:10].split("-")
                        due_d = _date(int(y), int(m), int(d))
                        delta = (due_d - today).days
                        line["due_date"] = due[:10]
                        line["show_dday"] = True
                        if delta < 0:
                            line["severity"] = "high"
                        elif delta <= 3:
                            line["severity"] = "high"
                        elif delta <= 14:
                            line["severity"] = "mid"
                        else:
                            line["severity"] = "info"
                    except Exception:
                        line["severity"] = "info"
                else:
                    line["severity"] = "info"
                auto_lines.append(line)

        return auto_lines
    except Exception as _e:
        print(f"[issue_lines] 자동 생성 실패: {_e}")
        return []


# =========================================================
# KPI 관리 API (Phase 2): 월/주차/이슈 라인 upsert
# =========================================================
from pydantic import BaseModel


class KpiMonthPayload(BaseModel):
    month: str            # "2026-07"
    efem: float
    vtm: float
    type: str = "plan"    # "actual" | "plan"
    source: str = ""


class KpiWeekPayload(BaseModel):
    week: str             # "2026-W28"
    efem: float
    vtm: float
    type: str = "plan"
    source: str = ""


class KpiIssueLinesPayload(BaseModel):
    issue_lines: list


def _ensure_project(hist: dict, project_key: str) -> dict:
    hist.setdefault("projects", {})
    proj = hist["projects"].setdefault(project_key, {
        "unit_economics": {},
        "months": {},
        "weeks": {},
        "issue_lines": [],
    })
    proj.setdefault("months", {})
    proj.setdefault("weeks", {})
    proj.setdefault("issue_lines", [])
    return proj


def _should_overwrite(existing: dict, incoming_type: str) -> bool:
    """
    같은 key가 이미 있을 때 덮어쓸지 판단.
    - existing이 없으면 True
    - actual > plan (actual이 들어오면 무조건 덮음)
    - 같은 type이면 최신 upload가 덮음 (True)
    """
    if not existing:
        return True
    et = existing.get("type", "plan")
    if incoming_type == "actual" and et == "plan":
        return True
    if incoming_type == "plan" and et == "actual":
        return False
    return True  # 같은 type이면 최신이 덮음


@app.get("/admin/kpi/{project_key}")
def admin_kpi_get(project_key: str, _exp: int = Depends(get_admin_session)):
    hist = _load_kpi_history()
    proj = (hist.get("projects") or {}).get(project_key)
    if not proj:
        return {"project_key": project_key, "months": {}, "weeks": {}, "issue_lines": []}
    return {"project_key": project_key, **proj}


@app.post("/admin/kpi/{project_key}/months")
def admin_kpi_upsert_month(
    project_key: str,
    payload: KpiMonthPayload,
    _exp: int = Depends(get_admin_session),
):
    hist = _load_kpi_history()
    proj = _ensure_project(hist, project_key)
    existing = proj["months"].get(payload.month)
    if not _should_overwrite(existing, payload.type):
        return {"status": "skipped", "reason": "existing actual not overwritten by plan", "month": payload.month}
    proj["months"][payload.month] = {
        "efem": payload.efem,
        "vtm": payload.vtm,
        "type": payload.type,
        "source": payload.source,
        "uploaded_at": datetime.utcnow().isoformat() + "Z",
    }
    _save_kpi_history(hist)
    return {"status": "ok", "month": payload.month, "row": proj["months"][payload.month]}


@app.post("/admin/kpi/{project_key}/weeks")
def admin_kpi_upsert_week(
    project_key: str,
    payload: KpiWeekPayload,
    _exp: int = Depends(get_admin_session),
):
    hist = _load_kpi_history()
    proj = _ensure_project(hist, project_key)
    existing = proj["weeks"].get(payload.week)
    if not _should_overwrite(existing, payload.type):
        return {"status": "skipped", "reason": "existing actual not overwritten by plan", "week": payload.week}
    proj["weeks"][payload.week] = {
        "efem": payload.efem,
        "vtm": payload.vtm,
        "type": payload.type,
        "source": payload.source,
        "uploaded_at": datetime.utcnow().isoformat() + "Z",
    }
    _save_kpi_history(hist)
    return {"status": "ok", "week": payload.week, "row": proj["weeks"][payload.week]}


@app.post("/admin/kpi/{project_key}/issue_lines")

@app.post("/admin/kpi/{project_key}/upload_excel")
async def admin_kpi_upload_excel(
    project_key: str,
    file: UploadFile = File(...),
    _admin: int = Depends(get_admin_session),
):
    """
    출하실적 엑셀 업로드 → 자동 파싱 → KPI 히스토리 반영.

    엑셀 스키마 (major_module 기준):
    - R1: 빈 행
    - R2: 카테고리 헤더 (PO/실적/잔량/6월/7월/W27~W31 등)
    - R3: 주차 헤더 (W27, W28, ...)
    - R4: 계획/실적 서브헤더
    - R5: EFEM 데이터
    - R6: VTM 데이터
    - R7: 합계 (파생, 저장 안 함)
    """
    import openpyxl
    from io import BytesIO
    from datetime import datetime as _dt

    raw = await file.read()
    if not raw:
        raise HTTPException(status_code=400, detail="빈 파일")

    try:
        wb = openpyxl.load_workbook(BytesIO(raw), data_only=True)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"엑셀 파싱 실패: {e}")

    ws = wb.worksheets[0]

    # 헤더 인식: 주차 헤더 행에서 W## 위치 파악
    week_header_row = None
    plan_actual_row = None
    for r in range(1, min(10, ws.max_row) + 1):
        row_vals = [ws.cell(row=r, column=c).value for c in range(1, ws.max_column + 1)]
        row_str = [str(v) if v is not None else "" for v in row_vals]
        joined = "|".join(row_str)
        if "W27" in joined or "W28" in joined:
            week_header_row = r
        if ("계획" in joined and "실적" in joined) and week_header_row and r > week_header_row:
            plan_actual_row = r
            break

    if not week_header_row or not plan_actual_row:
        raise HTTPException(status_code=400, detail="주차/계획/실적 헤더를 찾을 수 없습니다")

    # 주차 컬럼 매핑 (W## → (plan_col, actual_col))
    week_cols = {}
    for c in range(1, ws.max_column + 1):
        v = ws.cell(row=week_header_row, column=c).value
        if v and isinstance(v, str) and v.strip().startswith("W"):
            wk = v.strip()
            # 오른쪽으로 스캔해서 계획/실적 위치 찾기
            plan_c = actual_c = None
            for cc in [c, c+1]:
                sub = ws.cell(row=plan_actual_row, column=cc).value
                if sub == "계획":
                    plan_c = cc
                elif sub == "실적":
                    actual_c = cc
            if plan_c is None or actual_c is None:
                continue
            week_cols[wk] = (plan_c, actual_c)

    # 월 컬럼 (6월/7월/8월) 파싱 - 카테고리 헤더 기준
    month_cols = {}
    if week_header_row > 2:
        cat_row = week_header_row - 1
        for c in range(1, ws.max_column + 1):
            v = ws.cell(row=cat_row, column=c).value
            if v and isinstance(v, str):
                sv = v.strip()
                if sv in ("6월", "7월", "8월"):
                    # 서브헤더에서 계획/실적 찾기
                    plan_c = actual_c = None
                    for cc in range(c, min(c + 3, ws.max_column + 1)):
                        sub = ws.cell(row=plan_actual_row, column=cc).value
                        if sub == "계획" and plan_c is None:
                            plan_c = cc
                        elif sub == "실적" and actual_c is None:
                            actual_c = cc
                    if plan_c is not None:
                        month_cols[sv] = (plan_c, actual_c)

    # 모델 데이터 행 찾기 (EFEM, VTM)
    data_rows = {}
    for r in range(plan_actual_row + 1, ws.max_row + 1):
        model = ws.cell(row=r, column=2).value  # B열
        if model and isinstance(model, str):
            key = model.strip().upper()
            if key in ("EFEM", "VTM"):
                data_rows[key] = r

    if not data_rows:
        raise HTTPException(status_code=400, detail="EFEM/VTM 데이터 행을 찾을 수 없습니다")

    def _to_num(v):
        if v is None or v == "":
            return None
        try:
            return float(v)
        except (TypeError, ValueError):
            return None

    # 파싱 결과
    result = {"project_key": project_key, "months": [], "weeks": [], "raw_summary": {}}

    # 현재 연도 (report_date 파일명 기준으로 나중에 개선 가능)
    year = _dt.now().year

    # 월 데이터
    month_num_map = {"6월": 6, "7월": 7, "8월": 8}
    for month_label, (plan_c, actual_c) in month_cols.items():
        mnum = month_num_map.get(month_label)
        if not mnum:
            continue
        month_key = f"{year}-{mnum:02d}"

        # 실적이 있으면 actual, 없으면 plan
        efem_actual = _to_num(ws.cell(row=data_rows["EFEM"], column=actual_c).value) if actual_c else None
        vtm_actual = _to_num(ws.cell(row=data_rows["VTM"], column=actual_c).value) if actual_c else None
        efem_plan = _to_num(ws.cell(row=data_rows["EFEM"], column=plan_c).value)
        vtm_plan = _to_num(ws.cell(row=data_rows["VTM"], column=plan_c).value)

        # 실적/계획 둘 다 저장 (실적이 있으면 actual 타입, 아니면 plan 타입)
        if efem_actual is not None or vtm_actual is not None:
            result["months"].append({
                "month": month_key,
                "type": "actual",
                "efem": efem_actual if efem_actual is not None else 0,
                "vtm": vtm_actual if vtm_actual is not None else 0,
            })
        if efem_plan is not None or vtm_plan is not None:
            result["months"].append({
                "month": month_key,
                "type": "plan",
                "efem": efem_plan if efem_plan is not None else 0,
                "vtm": vtm_plan if vtm_plan is not None else 0,
            })

    # 주차 데이터
    for wk, (plan_c, actual_c) in week_cols.items():
        try:
            wnum = int(wk.lstrip("W"))
        except Exception:
            continue
        week_key = f"{year}-W{wnum:02d}"

        efem_actual = _to_num(ws.cell(row=data_rows["EFEM"], column=actual_c).value)
        vtm_actual = _to_num(ws.cell(row=data_rows["VTM"], column=actual_c).value)
        efem_plan = _to_num(ws.cell(row=data_rows["EFEM"], column=plan_c).value)
        vtm_plan = _to_num(ws.cell(row=data_rows["VTM"], column=plan_c).value)

        if efem_actual is not None or vtm_actual is not None:
            result["weeks"].append({
                "week": week_key,
                "type": "actual",
                "efem": efem_actual if efem_actual is not None else 0,
                "vtm": vtm_actual if vtm_actual is not None else 0,
            })
        if efem_plan is not None or vtm_plan is not None:
            result["weeks"].append({
                "week": week_key,
                "type": "plan",
                "efem": efem_plan if efem_plan is not None else 0,
                "vtm": vtm_plan if vtm_plan is not None else 0,
            })

    # kpi_history.json 반영
    hist = _load_kpi_history()
    proj = hist.setdefault("projects", {}).setdefault(project_key, {
        "months": {}, "weeks": {}, "issue_lines": []
    })

    # 병합 규칙: actual > plan (같은 키에서 actual 우선), source 기록
    now_iso = _dt.now().isoformat()

    def _upsert(dst_map, key, entry):
        cur = dst_map.get(key)
        # actual 은 plan 을 덮어씀. 같은 type 이면 항상 최신 upload 로 교체.
        if cur is None:
            dst_map[key] = {**entry, "source": "excel_upload", "updated_at": now_iso}
        else:
            cur_type = cur.get("type", "plan")
            new_type = entry.get("type", "plan")
            if new_type == "actual" or cur_type == new_type:
                dst_map[key] = {**entry, "source": "excel_upload", "updated_at": now_iso}

    for m in result["months"]:
        _upsert(proj.setdefault("months", {}), m["month"], m)
    for w in result["weeks"]:
        _upsert(proj.setdefault("weeks", {}), w["week"], w)

    _save_kpi_history(hist)

    return {
        "status": "ok",
        "project_key": project_key,
        "months_parsed": len(result["months"]),
        "weeks_parsed": len(result["weeks"]),
        "detail": result,
    }

def admin_kpi_replace_issue_lines(
    project_key: str,
    payload: KpiIssueLinesPayload,
    _exp: int = Depends(get_admin_session),
):
    hist = _load_kpi_history()
    proj = _ensure_project(hist, project_key)
    proj["issue_lines"] = payload.issue_lines
    _save_kpi_history(hist)
    return {"status": "ok", "count": len(payload.issue_lines)}

# ============================================================
# 매출 자동 계산 헬퍼
# ============================================================
import re
from datetime import date, datetime
from typing import Optional

def _parse_sales_input(text: str) -> dict:
    """
    [모델/판가]
    EFEM=65
    VTM=24
    
    [주차별]
    W27: EFEM 2/2, VTM 2/2
    W30: EFEM 5/0, VTM 1/0
    
    → {"prices": {"EFEM": 65.0}, "weeks": {"W27": {"EFEM": {"plan":2,"actual":2}}}}
    """
    prices = {}
    weeks = {}
    section = None
    for raw in (text or "").splitlines():
        line = raw.strip()
        if not line:
            continue
        # 섹션 헤더
        m = re.match(r'^\[(.+?)\]\s*$', line)
        if m:
            tag = m.group(1).strip().lower()
            if '판가' in tag or 'price' in tag or '모델' in tag:
                section = 'prices'
            elif '주차' in tag or 'week' in tag:
                section = 'weeks'
            else:
                section = None
            continue
        if section == 'prices':
            # EFEM=65  또는  EFEM = 65.0
            m = re.match(r'^([^\s=]+)\s*=\s*([\d.]+)\s*$', line)
            if m:
                prices[m.group(1).strip()] = float(m.group(2))
            continue
        if section == 'weeks':
            # W30: EFEM 5/0, VTM 1/0
            m = re.match(r'^(W\d+)\s*[:：]?\s*(.+)$', line, re.IGNORECASE)
            if not m:
                continue
            wk = m.group(1).upper()
            payload = m.group(2)
            models = {}
            # "EFEM 5/0, VTM 1/0" → 콤마/공백 mix 처리
            for part in re.split(r'\s*,\s*', payload):
                part = part.strip()
                if not part:
                    continue
                # "EFEM 5/0"
                mm = re.match(r'^(\S+)\s+([\d]+)\s*/\s*([\d]+)\s*$', part)
                if mm:
                    models[mm.group(1)] = {
                        "plan": int(mm.group(2)),
                        "actual": int(mm.group(3)),
                    }
            if models:
                weeks[wk] = models
    return {"prices": prices, "weeks": weeks}


def _iso_week_number(d: date) -> int:
    return d.isocalendar()[1]


def _week_to_month(week_num: int, year: int) -> Optional[int]:
    """W-번호 → 그 주의 목요일이 속한 월 (ISO 규칙)"""
    try:
        # ISO week's Thursday belongs to that year's month
        thursday = datetime.strptime(f"{year}-W{week_num:02d}-4", "%G-W%V-%u").date()
        return thursday.month
    except Exception:
        return None


def _compute_sales_summary(parsed: dict, today: Optional[date] = None) -> str:
    """
    각 주차별로:
      - 과거 → 실적
      - 현재 → 실적 우선, 없으면 계획
      - 미래 → 계획
    매출 = 수량 × 판가
    반환: "7월 74.0만불 · 8월 407.0만불 · W30 89.0만불"
    """
    if today is None:
        today = date.today()
    cur_week = _iso_week_number(today)
    year = today.year
    
    prices = parsed.get("prices", {})
    weeks = parsed.get("weeks", {})
    
    if not prices or not weeks:
        return ""
    
    # 주차별 매출 계산
    week_sales = {}  # {W27: 189.0}
    for wk, models in weeks.items():
        m = re.match(r'^W(\d+)$', wk, re.IGNORECASE)
        if not m:
            continue
        wk_num = int(m.group(1))
        total = 0.0
        for model_name, qty in models.items():
            price = prices.get(model_name, 0.0)
            if wk_num < cur_week:
                # 과거: 실적
                use = qty.get("actual", 0)
            elif wk_num > cur_week:
                # 미래: 계획
                use = qty.get("plan", 0)
            else:
                # 현재: 실적 우선, 없으면 계획
                use = qty.get("actual", 0) or qty.get("plan", 0)
            total += use * price
        week_sales[wk_num] = total
    
    # 월별 집계
    month_sales = {}  # {7: 74.0, 8: 407.0}
    for wk_num, sales in week_sales.items():
        mo = _week_to_month(wk_num, year)
        if mo is None:
            continue
        month_sales[mo] = month_sales.get(mo, 0.0) + sales
    
    # 문자열 조립
    parts = []
    for mo in sorted(month_sales.keys()):
        parts.append(f"{mo}월 {month_sales[mo]:.1f}만불")
    # 현재 주차만 별도 강조 (있으면)
    if cur_week in week_sales:
        parts.append(f"W{cur_week} {week_sales[cur_week]:.1f}만불")
    
    return " · ".join(parts)



# ============================================================
# [SALES-XLSX] 출하계획 엑셀 자동 파싱 + 판가 기반 매출 계산
# ============================================================
def parse_sales_excel(file_path):
    """출하계획 xlsx를 읽어서 models / weeks(plan/actual) 구조로 반환."""
    from openpyxl import load_workbook
    import re as _re

    wb = load_workbook(str(file_path), data_only=True)
    ws = wb.active

    _WEEK_RE = _re.compile(r'W\s*(\d{1,2})', _re.I)

    def _cell_str(v):
        return '' if v is None else str(v).strip()

    def _safe_num(v):
        if v in (None, '', '-'):
            return 0.0
        try:
            return float(v)
        except Exception:
            s = str(v).replace(',', '').strip()
            try:
                return float(s) if s else 0.0
            except Exception:
                return 0.0

    def _month_from_label(raw):
        s = _cell_str(raw).replace(' ', '')
        m = _re.search(r'(\d{1,2})월', s)
        return int(m.group(1)) if m else None

    def _month_for_col(col):
        # 자기 자신 셀 or 왼쪽 셀들에서 "N월" 라벨 스캔 (병합셀 대응)
        for r in (1, 2):
            m = _month_from_label(ws.cell(r, col).value)
            if m:
                return m
        for c in range(col, 0, -1):
            for r in (1, 2):
                m = _month_from_label(ws.cell(r, c).value)
                if m:
                    return m
        return None

    def _build_model_name(row):
        b = _cell_str(ws.cell(row, 2).value)
        c = _cell_str(ws.cell(row, 3).value)
        if not b and not c:
            return None
        low_b = b.lower()
        low_c = c.lower()
        if low_b in ('합계', 'total', '계') or low_c in ('합계', 'total', '계'):
            return None
        # B+C 2단 (프레임처럼)
        if c and not _re.fullmatch(r'[\d,.\-]+', c):
            return f'{b} / {c}' if b else c
        return b or c

    # 주차 컬럼 스캔
    week_cols = []
    for c in range(1, ws.max_column + 1):
        label = _cell_str(ws.cell(3, c).value)
        m = _WEEK_RE.search(label)
        if not m:
            continue
        week_num = int(m.group(1))
        month_num = _month_for_col(c)
        week_cols.append((c, week_num, month_num))

    models = []
    weeks_map = {}

    for row in range(5, ws.max_row + 1):
        model = _build_model_name(row)
        if not model:
            continue

        row_has_value = False
        for col, week_num, month_num in week_cols:
            plan = _safe_num(ws.cell(row, col).value)
            actual = _safe_num(ws.cell(row, col + 1).value)
            if plan != 0 or actual != 0:
                row_has_value = True

            weeks_map.setdefault(week_num, {
                'week': week_num,
                'month': month_num,
                'models': {}
            })
            weeks_map[week_num]['models'][model] = {
                'plan': plan,
                'actual': actual,
            }

        if row_has_value and model not in models:
            models.append(model)

    return {
        'models': models,
        'weeks': [weeks_map[k] for k in sorted(weeks_map.keys())],
    }


def compute_sales_from_parsed(parsed, prices, today=None):
    """parsed(xlsx) + prices(만불) → 3박스 요약 문자열 + 구조화 데이터."""
    from datetime import date as _date, datetime as _dt

    today = today or _date.today()
    current_week = today.isocalendar()[1]
    current_month = today.month
    next_month = 1 if current_month == 12 else current_month + 1

    month_totals = {}
    week_totals = {}

    for w in (parsed.get('weeks') or []):
        wk = int(w.get('week') or 0)
        month_num = w.get('month')
        total = 0.0

        for model, qa in (w.get('models') or {}).items():
            price = float((prices or {}).get(model, 0) or 0)
            plan = float((qa or {}).get('plan', 0) or 0)
            actual = float((qa or {}).get('actual', 0) or 0)

            if wk < current_week:
                qty = actual
            elif wk == current_week:
                qty = actual if actual > 0 else plan
            else:
                qty = plan

            total += qty * price

        week_totals[wk] = total
        if month_num:
            month_totals[month_num] = month_totals.get(month_num, 0.0) + total

    cur_month_amt = round(month_totals.get(current_month, 0.0), 1)
    cur_week_amt = round(week_totals.get(current_week, 0.0), 1)
    next_month_amt = round(month_totals.get(next_month, 0.0), 1)

    delta_pct = 0.0
    if cur_month_amt > 0:
        delta_pct = round(((next_month_amt - cur_month_amt) / cur_month_amt) * 100.0, 1)

    summary_text = (
        f'{current_month}월 {cur_month_amt:.1f}만불 · '
        f'이번주 {cur_week_amt:.1f}만불 · '
        f'{next_month}월 {next_month_amt:.1f}만불'
    )
    if cur_month_amt > 0:
        arrow = '▲' if delta_pct >= 0 else '▼'
        summary_text += f' {arrow}{abs(delta_pct):.1f}%'

    return {
        'sales_summary': summary_text,
        'sales_summary_data': {
            'boxes': [
                {'key': 'current_month', 'label': f'{current_month}월',
                 'amount': cur_month_amt, 'unit': '만불'},
                {'key': 'current_week', 'label': '이번주',
                 'amount': cur_week_amt, 'unit': '만불', 'week': current_week},
                {'key': 'next_month', 'label': f'{next_month}월',
                 'amount': next_month_amt, 'unit': '만불', 'delta_pct': delta_pct},
            ]
        },
        'sales_computed_at': _dt.now().isoformat(timespec='seconds'),
    }


def _find_xlsx_in_section_items(items):
    """섹션 items에서 첫 번째 .xlsx 파일의 서버 URL 추출."""
    if not items:
        return None
    for it in items:
        if not isinstance(it, dict):
            continue
        for key in ('url', 'file_url', 'href', 'photo_ref', 'text'):
            v = it.get(key)
            if not v:
                continue
            s = str(v)
            if '.xlsx' in s.lower() or '.xlsm' in s.lower():
                return s
    return None


def _xlsx_url_to_local_path(url_str):
    """/admin/manual-files/{doc_id}/{filename} → 실제 파일 경로."""
    if not url_str:
        return None
    s = str(url_str)
    # /admin/manual-files/{doc_id}/{filename} 패턴
    import re as _re
    m = _re.search(r'/admin/manual-files/([^/]+)/([^/?#]+)', s)
    if m:
        doc_id = m.group(1)
        fname = m.group(2)
        return UPLOAD_DIR / 'manual' / doc_id / fname
    return None


def _is_sales_section_title(title):
    t = re.sub(r'\s+', '', (title or ''))
    if not t:
        return False
    return ('주차별' in t and ('계획' in t or '출하' in t)) or t == '주차별계획'

def compute_sales_from_input(sales_input: str, today: Optional[date] = None) -> dict:
    """외부에서 호출하는 진입점"""
    if not sales_input or not sales_input.strip():
        return {"sales_summary": "", "sales_computed_at": ""}
    parsed = _parse_sales_input(sales_input)
    summary = _compute_sales_summary(parsed, today=today)
    return {
        "sales_summary": summary,
        "sales_computed_at": datetime.now().isoformat(timespec='seconds'),
        "_debug": {
            "prices": parsed.get("prices"),
            "weeks_count": len(parsed.get("weeks", {})),
        },
    }



# ============================================================
# 매출 자동 계산 API
# ============================================================
@app.post("/admin/notes/section/sales_input")
def admin_set_sales_input(payload: dict, _admin: int = Depends(get_admin_session)):
    """
    body: { division_id, card_title, section_title, sales_input }
    → notes.json 의 해당 section 에 sales_input/sales_summary/sales_computed_at 저장
    → 계산 결과 반환
    """
    division_id = (payload.get("division_id") or "").strip()
    card_title = (payload.get("card_title") or "").strip()
    section_title = (payload.get("section_title") or "").strip()
    sales_input = payload.get("sales_input") or ""
    
    if not division_id or not card_title or not section_title:
        raise HTTPException(status_code=400, detail="division_id/card_title/section_title 필요")
    
    result = compute_sales_from_input(sales_input)
    
    # notes.json 갱신
    data = _load_notes()
    notes_map = data.get("notes", {}) or {}
    div = notes_map.get(division_id) or {}
    cards = div.get("cards", []) or []
    target_card = None
    for c in cards:
        if (c.get("title") or "").strip() == card_title:
            target_card = c
            break
    if not target_card:
        raise HTTPException(status_code=404, detail=f"카드 없음: {card_title}")
    
    sections = target_card.get("sections", []) or []
    target_sec = None
    for s in sections:
        if (s.get("title") or "").strip() == section_title:
            target_sec = s
            break
    if not target_sec:
        raise HTTPException(status_code=404, detail=f"섹션 없음: {section_title}")
    
    target_sec["sales_input"] = sales_input
    target_sec["sales_summary"] = result["sales_summary"]
    target_sec["sales_computed_at"] = result["sales_computed_at"]
    
    # 저장
    NOTES_FILE.parent.mkdir(parents=True, exist_ok=True)
    NOTES_FILE.write_text(
        json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    
    return {
        "ok": True,
        "sales_summary": result["sales_summary"],
        "sales_computed_at": result["sales_computed_at"],
        "debug": result.get("_debug", {}),
    }


def compute_sales_from_data(sales_data: dict, today: Optional[date] = None) -> dict:
    """구조화된 sales_data → summary 계산.
    sales_data = {
      "prices": {"EFEM": 65, ...},
      "weeks": [{"week": 27, "models": {"EFEM": {"plan": 2, "actual": 2}}}, ...]
    }
    """
    if not isinstance(sales_data, dict):
        return {"sales_summary": "", "sales_computed_at": ""}
    prices = sales_data.get("prices") or {}
    weeks_list = sales_data.get("weeks") or []
    if not prices or not weeks_list:
        return {"sales_summary": "", "sales_computed_at": ""}
    
    # sales_helpers.py 의 _compute_sales_summary 와 동일 규칙
    if today is None:
        today = date.today()
    cur_week = today.isocalendar()[1]
    year = today.year
    
    # weeks 리스트를 dict로 변환하여 _compute_sales_summary 재활용
    parsed = {"prices": {k: float(v) for k, v in prices.items()}, "weeks": {}}
    for w in weeks_list:
        wk_num = w.get("week")
        if not isinstance(wk_num, int) or wk_num <= 0:
            continue
        models = w.get("models") or {}
        parsed["weeks"][f"W{wk_num}"] = {
            m: {"plan": int(v.get("plan", 0) or 0), "actual": int(v.get("actual", 0) or 0)}
            for m, v in models.items() if isinstance(v, dict)
        }
    
    summary = _compute_sales_summary(parsed, today=today)
    return {
        "sales_summary": summary,
        "sales_computed_at": datetime.now().isoformat(timespec="seconds"),
    }


@app.post("/admin/notes/section/sales_recompute")
def admin_recompute_all_sales(_admin: int = Depends(get_admin_session)):
    """모든 section의 sales_input 재계산 (판가 정책 바뀔 때 등)"""
    data = _load_notes()
    notes_map = data.get("notes", {}) or {}
    count = 0
    for div_id, div in notes_map.items():
        for c in div.get("cards", []) or []:
            for s in c.get("sections", []) or []:
                si = s.get("sales_input")
                if not si:
                    continue
                r = compute_sales_from_input(si)
                s["sales_summary"] = r["sales_summary"]
                s["sales_computed_at"] = r["sales_computed_at"]
                count += 1
    NOTES_FILE.write_text(
        json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return {"ok": True, "recomputed": count}

@app.post("/device-tokens")
def register_device_token(payload: dict):
    """Flutter 앱이 앱 시작 시 FCM 토큰을 서버에 등록."""
    token = (payload or {}).get("token")
    if not token or not isinstance(token, str) or len(token) < 20:
        raise HTTPException(status_code=400, detail="invalid token")
    platform = (payload or {}).get("platform") or "android"
    debug = bool((payload or {}).get("debug", True))
    _add_device_token(token, platform=platform, debug=debug)
    return {"ok": True}




@app.post("/admin/push/update-notice")
def admin_push_update_notice(_admin: int = Depends(get_admin_session)):
    """업데이트 공지 푸시 발송 (관리자 전용)"""
    result = _send_fcm_to_all(
        "새 버전 업데이트 안내",
        "새 버전 2.1.2가 배포되었습니다. 최신 APK를 설치해주세요.",
        data={
            "type": "app_update",
            "url": "https://github.com/msk05317/project_summary/releases/tag/v2.1.2",
        },
    )
    return result


@app.post("/admin/fcm-test")
def admin_fcm_test(_admin: int = Depends(get_admin_session)):
    """알람 시스템 동작 테스트용 엔드포인트."""
    result = _send_fcm_to_all(
        title="🔔 테스트 알람",
        body="FCM 파이프라인 정상 동작 중",
        data={"type": "test"},
    )
    return result



# ─── 개발 승인 프로세스 (13단계 고정 템플릿) ───
DEV_PROCESS_STEPS = [
    ("fa_po", "FA PO", "발주"),
    ("material_order", "자재 발주", "발주"),
    ("incoming", "입고", "발주"),
    ("machining", "가공 (조립)", "제작·검사"),
    ("la_incoming", "LA 입고", "제작·검사"),
    ("lair_write", "LAIR 제출", "제작·검사"),
    ("lair_approval", "LAIR 승인", "제작·검사"),
    ("source_inspection", "Source Inspection", "승인"),
    ("fair_write", "FAIR 제출", "승인"),
    ("fair_approval", "FAIR 승인", "승인"),
    ("lap_test", "LAP TEST", "승인"),
    ("cdr", "CDR", "승인"),
    ("final_approval", "최종 승인 완료", "승인"),
]

def _default_process() -> list:
    return [{"key": k, "name": n, "group": g, "expected": "", "actual": ""} for k, n, g in DEV_PROCESS_STEPS]

def _model_key_alias(project_key: str) -> str:
    _alias = {"havaplate": "hrva_plate", "hrvaplate": "hrva_plate", "hrva-plate": "hrva_plate"}
    return _alias.get(project_key.strip().lower(), project_key.strip())

def _ensure_process(m: dict) -> list:
    proc = m.get("process")
    if not isinstance(proc, list) or len(proc) != 13:
        proc = _default_process()
        m["process"] = proc
    return proc

def _process_progress(proc: list) -> int:
    if not proc:
        return 0
    done = sum(1 for s in proc if str(s.get("actual") or "").strip())
    return round(done / len(proc) * 100)

def _process_current(proc: list):
    for s in proc:
        if not str(s.get("actual") or "").strip():
            return s.get("name") or "", str(s.get("expected") or "")
    return "최종 승인 완료", ""

def _enrich_model(m: dict) -> dict:
    out = {k: m.get(k) for k in ("id", "name", "group", "status", "progress", "price", "material_cost", "dev_type")}
    if m.get("group") == "개발":
        proc = _ensure_process(m)
        out["dev_type"] = m.get("dev_type") or "HVM"
        out["progress"] = _process_progress(proc)
        stage, expected = _process_current(proc)
        out["current_stage"] = stage
        out["current_expected"] = expected
        out["done_steps"] = sum(1 for s in proc if str(s.get("actual") or "").strip())
        out["total_steps"] = len(proc)
    return out

@app.get("/projects/{project_key}/models/detail")
def get_project_models_detail(project_key: str):
    """앱용 확장 모델 목록: dev_type, 자동 진행률, 현황(현재 단계), 완료예정일, status_note 포함"""
    _key = _model_key_alias(project_key)
    data = _load_models()
    proj = data.get("projects", {}).get(_key, {})
    changed = False
    enriched = []
    for m in proj.get("models", []):
        if not isinstance(m, dict):
            continue
        if m.get("group") == "개발" and (not isinstance(m.get("process"), list) or len(m.get("process")) != 13):
            changed = True
        enriched.append(_enrich_model(m))
    if changed:
        _save_models(data)
    enriched.sort(key=lambda x: 0 if x.get("group") == "양산" else 1)
    return {
        "project_key": _key,
        "has_models": len(enriched) > 0,
        "total": len(enriched),
        "status_note": proj.get("status_note") or "",
        "models": enriched,
    }

@app.get("/projects/{project_key}/models/{model_id}/process")
def get_model_process(project_key: str, model_id: str):
    """앱용 개발 승인 프로세스 조회"""
    from urllib.parse import unquote
    model_id = unquote(model_id)
    _key = _model_key_alias(project_key)
    data = _load_models()
    proj = data.get("projects", {}).get(_key, {})
    for m in proj.get("models", []):
        if isinstance(m, dict) and (m.get("id") or "").lower() == model_id.lower():
            if m.get("group") != "개발":
                raise HTTPException(status_code=400, detail="개발 모델만 프로세스가 있습니다.")
            proc = _ensure_process(m)
            _save_models(data)
            done = sum(1 for s in proc if str(s.get("actual") or "").strip())
            stage, expected = _process_current(proc)
            return {
                "model_id": model_id,
                "model_name": m.get("name") or model_id,
                "dev_type": m.get("dev_type") or "HVM",
                "steps": proc,
                "done": done,
                "total": len(proc),
                "progress": _process_progress(proc),
                "current_stage": stage,
                "current_expected": expected,
            }
    raise HTTPException(status_code=404, detail="모델을 찾을 수 없습니다.")

@app.put("/admin/projects/{project_key}/models/{model_id}/process")
def admin_put_model_process(project_key: str, model_id: str, payload: dict, _admin: int = Depends(get_admin_session)):
    """admin용 프로세스 저장: {"dev_type": "HVM", "steps": [{"key","expected","actual"}, ...]}"""
    from urllib.parse import unquote
    model_id = unquote(model_id)
    _key = _model_key_alias(project_key)
    data = _load_models()
    proj = data.setdefault("projects", {}).setdefault(_key, {"models": []})
    target = None
    for m in proj.get("models", []):
        if isinstance(m, dict) and (m.get("id") or "").lower() == model_id.lower():
            target = m
            break
    if target is None:
        raise HTTPException(status_code=404, detail="모델을 찾을 수 없습니다.")
    dt = str(payload.get("dev_type") or "").strip().upper()
    if dt in ("HVM", "RPM"):
        target["dev_type"] = dt
    proc = _ensure_process(target)
    incoming = payload.get("steps")
    if isinstance(incoming, list):
        by_key = {s.get("key"): s for s in incoming if isinstance(s, dict)}
        for step in proc:
            inc = by_key.get(step.get("key"))
            if inc is None:
                continue
            step["expected"] = str(inc.get("expected") or "").strip()[:10]
            step["actual"] = str(inc.get("actual") or "").strip()[:10]
    target["progress"] = _process_progress(proc)
    _save_models(data)
    return {"ok": True, "model_id": model_id, "progress": target["progress"], "dev_type": target.get("dev_type")}

@app.put("/admin/projects/{project_key}/status-note")
def admin_put_status_note(project_key: str, payload: dict, _admin: int = Depends(get_admin_session)):
    """프로젝트 현황 텍스트 저장: {"note": "..."}"""
    _key = _model_key_alias(project_key)
    data = _load_models()
    proj = data.setdefault("projects", {}).setdefault(_key, {"models": []})
    proj["status_note"] = str(payload.get("note") or "").strip()
    _save_models(data)
    return {"ok": True, "project_key": _key}



if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8080)
