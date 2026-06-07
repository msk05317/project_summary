from pptx import Presentation
from dataclasses import dataclass, field
from typing import List


@dataclass
class SlideTable:
    headers: List[str]
    rows: List[List[str]]

    def to_markdown(self) -> str:
        if not self.headers:
            return ""
        lines = ["| " + " | ".join(self.headers) + " |"]
        lines.append("|" + "|".join(["---"] * len(self.headers)) + "|")
        for r in self.rows:
            lines.append("| " + " | ".join(r) + " |")
        return "\n".join(lines)


@dataclass
class SlideContent:
    slide_number: int
    title: str
    body_text: str
    tables: List[SlideTable] = field(default_factory=list)
    issues: List[str] = field(default_factory=list)
    references: List[str] = field(default_factory=list)


@dataclass
class ParsedPPT:
    file_name: str
    author: str
    date: str
    title: str
    slides: List[SlideContent]

    def to_llm_input(self) -> str:
        out = [f"# {self.title}", f"작성자: {self.author} / 작성일: {self.date}\n"]
        for s in self.slides:
            out.append(f"\n## [슬라이드 {s.slide_number}] {s.title}")
            if s.body_text:
                out.append(s.body_text)
            for t in s.tables:
                out.append("\n" + t.to_markdown())
            if s.issues:
                out.append("\n[이슈사항]\n- " + "\n- ".join(s.issues))
            if s.references:
                out.append("\n[참고사항]\n- " + "\n- ".join(s.references))
        return "\n".join(out)


class PPTXParser:
    ISSUE_MARKERS = ["※ 이슈사항", "이슈사항", "리스크"]
    REF_MARKERS = ["※ 참고사항", "참고사항", "특이사항"]

    def parse(self, path: str) -> ParsedPPT:
        prs = Presentation(path)
        slides = [self._parse_slide(s, i + 1) for i, s in enumerate(prs.slides)]
        meta = self._parse_cover(prs.slides[0]) if prs.slides else {}
        return ParsedPPT(
            file_name=path.split("/")[-1],
            author=meta.get("author", ""),
            date=meta.get("date", ""),
            title=meta.get("title", ""),
            slides=slides,
        )

    def _parse_cover(self, slide) -> dict:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        full = "\n".join(texts)
        meta = {"title": "", "author": "", "date": ""}
        for line in full.split("\n"):
            line = line.strip()
            if not line:
                continue
            if "작성일" in line:
                meta["date"] = line.replace("작성일", "").strip()
            elif "작성자" in line:
                meta["author"] = line.replace("작성자", "").strip()
            elif not meta["title"] and "confidential" not in line.lower():
                meta["title"] = line
        return meta

    def _parse_slide(self, slide, num: int) -> SlideContent:
        title = ""
        body, issues, refs, tables = [], [], [], []

        for shape in slide.shapes:
            if shape.has_table:
                tables.append(self._parse_table(shape.table))
            elif shape.has_text_frame:
                text = shape.text_frame.text
                if not title and 0 < len(text.strip()) < 60 and "\n" not in text:
                    title = text.strip()
                    continue
                section = None
                for line in text.split("\n"):
                    s = line.strip()
                    if not s:
                        continue
                    if any(m in s for m in self.ISSUE_MARKERS):
                        section = "issue"
                        continue
                    if any(m in s for m in self.REF_MARKERS):
                        section = "ref"
                        continue
                    if section == "issue":
                        issues.append(s)
                    elif section == "ref":
                        refs.append(s)
                    else:
                        body.append(s)

        return SlideContent(
            slide_number=num,
            title=title,
            body_text="\n".join(body),
            tables=tables,
            issues=issues,
            references=refs,
        )

    def _parse_table(self, tbl) -> SlideTable:
        rows = [[c.text.strip() for c in r.cells] for r in tbl.rows]
        return SlideTable(headers=rows[0] if rows else [], rows=rows[1:])
