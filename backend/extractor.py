"""
PPT에서 키워드 기반으로 차트/표/이미지를 자동 추출하는 모듈
- 슬라이드 번호 의존 X
- 키워드로 섹션 자동 매칭
- 매주 PPT가 바뀌어도 자동 작동
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from PIL import Image
import re

# 실제 PPT 4개에서 뽑은 키워드 사전
# (파워박스, 메이저모듈, 반도체장비, 내재화 PPT 기준)
SECTION_KEYWORDS = {
    # ===== 공통 섹션 =====
    "양산": [
        "양산", "양산 ", "연간", "Aether GDX", "Aether",
        "양산종", "양산 14종", "양산 13종",
    ],
    "개발": [
        "개발", "개발종", "개발 18종", "개발 16종",
        "세이버", "SABRE", "승인 타겟", "내재화",
        "자재준비중", "CB완료", "양산타겟",
    ],
    "EMA": [
        "EMA", "ema", "조건부 승인", "조건부승인",
        "648-", "170-", "조립완료", "출고 대기", "출고대기",
        "순차 승인", "순차승인",
    ],
    "주차별 출하실적 및 계획": [
        "주차별", "출하실적", "출하계획", "출하 계획", "출하 실적",
        "주차 별", "쇼티지", "Shortage", "글로벌 쇼티지",
        "주 10대", "출하목표", "출하 목표",
    ],

    # ===== 메이저 모듈 전용 =====
    "챔버 진행현황": [
        "챔버", "Chamber", "PO 미접수", "PO미접수",
    ],
    "엔클로저 진행현황": [
        "엔클로저", "Enclosure", "EA",
    ],
    "하바플레이트 진행현황": [
        "하바플레이트", "하바", "하바 플레이트",
    ],
    "플레이팅 셀": [
        "플레이팅", "플레이팅 셀", "Plating",
    ],

    # ===== 반도체 장비 전용 =====
    "EFEM": [
        "EFEM", "efem", "이에프이엠",
    ],
    "VTM": [
        "VTM", "vtm", "Vacuum Transfer",
    ],
    "CEFEM": [
        "CEFEM", "cefem",
    ],
    "QUAROS": [
        "QUAROS", "Quaros", "쿠아로스",
    ],

    # ===== 내재화 전용 =====
    "CUP": ["CUP", "컵"],
    "자일란 코팅": ["자일란", "Xylan", "자일란코팅"],
    "톨론": ["톨론", "Torlon"],
    "Space X": ["Space X", "SpaceX", "스페이스X"],
    "KLA": ["KLA", "케이엘에이"],
    "세정": ["세정", "Cleaning", "클리닝"],
    "EOS 챔버": ["EOS", "EOS 챔버"],
    "램 캐스팅": ["램 캐스팅", "Ram Casting", "캐스팅"],
}


def normalize(text: str) -> str:
    """공백, 대소문자 정규화"""
    return re.sub(r"\s+", " ", text or "").strip().lower()


def extract_slide_text(slide) -> str:
    """슬라이드 안의 모든 텍스트를 하나로 합침"""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text:
                        texts.append(run.text)
        # 표 안의 글자도 추출
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    texts.append(cell.text)
    return " ".join(texts)


def match_section(slide_text: str) -> str | None:
    """슬라이드 텍스트를 보고 어느 섹션인지 자동 판단"""
    text_norm = normalize(slide_text)
    if not text_norm:
        return None

    # 키워드 매칭 점수 계산
    scores = {}
    for section, keywords in SECTION_KEYWORDS.items():
        score = 0
        for kw in keywords:
            kw_norm = normalize(kw)
            if kw_norm and kw_norm in text_norm:
                # 키워드 길이가 길수록 가중치 ↑ (특이한 단어 우선)
                score += len(kw_norm)
        if score > 0:
            scores[section] = score

    if not scores:
        return None

    # 가장 높은 점수 섹션 반환
    return max(scores, key=scores.get)


def emu_to_px(emu_value: int, slide_width_emu: int, image_width_px: int) -> int:
    """PPT 좌표(EMU) → 이미지 좌표(px) 변환"""
    return int(emu_value / slide_width_emu * image_width_px)


def find_chart_shapes(slide):
    """슬라이드에서 차트/표/이미지 도형 좌표 리스트 반환"""
    targets = []
    for shape in slide.shapes:
        is_chart = getattr(shape, "has_chart", False)
        is_table = getattr(shape, "has_table", False)
        # shape_type 13 = Picture
        is_picture = getattr(shape, "shape_type", None) == 13

        if is_chart or is_table or is_picture:
            try:
                targets.append({
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                    "kind": "chart" if is_chart else ("table" if is_table else "picture"),
                })
            except Exception:
                continue
    return targets


def crop_from_slide_image(
    slide_image_path: Path,
    shape_box: dict,
    slide_width_emu: int,
    slide_height_emu: int,
    output_path: Path,
    padding_ratio: float = 0.02,
) -> bool:
    """슬라이드 이미지에서 해당 도형 영역만 잘라서 저장"""
    try:
        img = Image.open(slide_image_path)
        iw, ih = img.size

        left = emu_to_px(shape_box["left"], slide_width_emu, iw)
        top = emu_to_px(shape_box["top"], slide_height_emu, ih)
        right = left + emu_to_px(shape_box["width"], slide_width_emu, iw)
        bottom = top + emu_to_px(shape_box["height"], slide_height_emu, ih)

        # 약간의 여백 추가
        pad_x = int(iw * padding_ratio)
        pad_y = int(ih * padding_ratio)
        left = max(0, left - pad_x)
        top = max(0, top - pad_y)
        right = min(iw, right + pad_x)
        bottom = min(ih, bottom + pad_y)

        cropped = img.crop((left, top, right, bottom))
        output_path.parent.mkdir(parents=True, exist_ok=True)
        cropped.save(output_path, "PNG")
        return True
    except Exception as e:
        print(f"❌ crop 실패: {e}")
        return False


def extract_charts_from_pptx(
    pptx_path: Path,
    slides_image_dir: Path,
    output_dir: Path,
    doc_id: str,
) -> dict:
    """
    PPT를 읽어서 섹션별 차트를 자동 crop하고,
    {섹션명: [crop된 이미지의 상대경로 리스트]} 반환
    """
    prs = Presentation(str(pptx_path))
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    section_images: dict[str, list[str]] = {}

    for idx, slide in enumerate(prs.slides, start=1):
        text = extract_slide_text(slide)
        section = match_section(text)
        if not section:
            continue

        # 해당 슬라이드에 변환된 PNG가 있어야 함
        slide_png = slides_image_dir / f"slide_{idx:02d}.png"
        if not slide_png.exists():
            # 파일명 다른 패턴도 시도
            candidates = list(slides_image_dir.glob(f"*{idx}*.png"))
            if not candidates:
                continue
            slide_png = candidates[0]

        shapes = find_chart_shapes(slide)
        if not shapes:
            continue

        # 가장 큰 도형 1개만 (보통 그게 메인 차트)
        shapes.sort(key=lambda s: s["width"] * s["height"], reverse=True)
        main_shape = shapes[0]

        safe_section = re.sub(r"[^\w가-힣]+", "_", section)
        out_name = f"{safe_section}_slide{idx:02d}.png"
        out_path = output_dir / doc_id / out_name

        if crop_from_slide_image(slide_png, main_shape, slide_w, slide_h, out_path):
            rel_url = f"/cropped/{doc_id}/{out_name}"
            section_images.setdefault(section, []).append(rel_url)
            print(f"✅ {section} ← 슬라이드 {idx} 에서 차트 추출")

    return section_images
